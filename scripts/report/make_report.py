#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
make_report.py — 一鍵：由底層數據(feed) 做齊報告數字表 → 一份 pptx（報告只係 ref，全部 data 生成）。

用法（簡單）：
    python scripts\\report\\make_report.py            # 預設 mgm
    python scripts\\report\\make_report.py wynn       # 換 entity

自動揾：feed = tableau_combined_25.csv（root）、清單 = data\\投資項目清單\\*{ENTITY}*、
       template = data\\reports\\*{ENTITY}*.pptx（跟 slide 尺寸；冇就用 10.83x7.5 報告標準）。
出 {entity}_report_llm.pptx：概述 + 期後 + 主要發現 + 金額匯總 + 設施vs活動 + 單項審查(25/24/23)
+ 附件走訪，全部 native table／文字，版式跟 layout.py（KPMG house style，對齊報告 scan）。
體檢：python scripts\\report\\inspect_pptx.py {entity}_report_llm.pptx [--preview]
"""
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
try:
    import pandas as pd
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("✗ pip install pandas python-pptx openpyxl"); sys.exit(1)

import layout as L                     # 版式引擎（KPMG house style；顏色/字體/表格/分頁全部喺嗰度）

# 報告配色 alias（真身喺 layout.py，跟 KPMG Visual identity overview）
HDR = L.NAVY                           # KPMG Blue 00338D：表頭白字 / 標題 / 導語
SEC = L.SECFILL                        # 範疇 section 行
SUB = L.SUBTOT                         # 小計
TOT = L.TOTAL                          # 總計
DARK = L.DARK                          # 封面 / 章節分隔深底
LIGHT = L.WHITE
GREY = L.GREY

# entity → 承批公司全名（封面用；公開上市公司法定名，非機密）
ENTITY_FULL = {
    "mgm": "美高梅金殿超濠股份有限公司",
    "galaxy": "銀河娛樂場股份有限公司",
    "sjm": "澳門博彩股份有限公司",
    "wynn": "永利渡假村（澳門）股份有限公司",
    "vml": "威尼斯人澳門股份有限公司",
    "melco": "新濠博亞博彩（澳門）股份有限公司",
}

# 公司 KPMG template（空 pptx，得 master + layouts）→ 開佢做 base，formatting 全部嚟自 master。
# 放 repo root / data/ 任一（gitignored，confidential）。搵唔到就 fallback 手砌（fresh Presentation）。
TEMPLATE_NAMES = ["template.pptx", "report_template.pptx", "kpmg_template.pptx"]
TEMPLATE_DIRS = [".", "data", "data/template", "conf/local"]
# 報告版 → template layout 名（跨 master 搵第一個 match；搵唔到就 fallback）
LAYOUTS = {
    "cover": ["TITLE SLIDE - Right vertical dark image", "TITLE SLIDE", "Title Slide"],
    "section": ["Section Divider"],
    "subsection": ["Subsection Divider"],
    "appendix": ["Appendix Divider"],
    "table_text": ["1_ANALYSIS narrow table", "ANALYSIS narrow table"],
    "two_col": ["Two Column Text", "ANALYSIS 2 col text"],
    "one_col": ["One Column Text", "RPOE One Column Text"],
    "title_only": ["Title Only"],
    "findings2": ["KEY FINDINGS 2_text only"],
    "toc": ["Table of content Storage Layout", "Table of content"],
}


def _find_template():
    for d in TEMPLATE_DIRS:
        for n in TEMPLATE_NAMES:
            p = Path(d) / n
            if p.exists():
                return p
    return None


USE_TEMPLATE = False    # 只有 --use-template 真係開咗 KPMG template 先可以用 template layout


def _layout(prs, key):
    """跨所有 master 搵第一個名 match LAYOUTS[key] 嘅 slide layout；搵唔到回 None。
    ⚠ 冇開 template 時一定回 None —— 否則會撞正 python-pptx 預設包嗰個 'Title Slide'
    layout，封面變一版白底光板（冇深底/冇 KPMG 字標/冇日期）。"""
    if not USE_TEMPLATE:
        return None
    wanted = LAYOUTS.get(key, [key])
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if lay.name in wanted:
                return lay
    return None


def _strip_slides(prs):
    """正確刪走 template 原有 content slides（drop rel + sldId）→ 避免 duplicate part 名 corruption。"""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        sldIdLst.remove(sldId)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


def _renumber_slides(prs, base=9000):
    """save 前把生成嘅 slide part 重編做高號（9000+）→ 就算 template 有殘留 orphan part（清唔切）
    都唔會撞名 → 徹底避開 duplicate-name corruption / 要 repair。"""
    try:
        from pptx.opc.packuri import PackURI
    except Exception:
        return
    for i, slide in enumerate(prs.slides, base):
        try:
            slide.part.partname = PackURI(f"/ppt/slides/slide{i}.xml")
        except Exception:
            pass


def _blank_layout(prs):
    """乾淨 layout（手砌 data 版用）：template 揀 'Title Only'/'Blank'（有 master title bar/footer），否則 index。"""
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if lay.name in ("Blank", "Title Only"):
                return lay
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]


def _ph(slide, idx):
    """由 placeholder idx 攞 placeholder；冇就 None。"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


import build_project_review_table as B
import feed_schema as FS
import build_summary_tables as S
import build_overview_tables as O
import build_narrative as N
import render_review_table_pptx as R
import biao2 as B2
from build_llm_narrative import (generate_llm_narrative, Workbench, tbl_key,   # bundler 會 inline
                                 proj_key, bkt_key)

BUILD_STAMP = "dev"          # bundler 會換做 git sha + 日期（睇 output 就知跑緊邊版）
FEED = "tableau_combined_25.csv"


def _is_num(v):
    try:
        float(v); return True
    except (TypeError, ValueError):
        return False


def _load_llm(entity):
    """有 {entity}_llm_narrative.json（build_llm_narrative.py 出）就用 LLM summary。"""
    import json
    p = Path(f"{entity}_llm_narrative.json")
    if p.exists():
        try:
            return json.loads(p.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


def _coverage_probe(df, qingdan):
    """一次性 coverage 探測（❓頁：主體/KPI/藝術品欄喺唔喺 inputs）→ print 供定 coverage。"""
    def pr(cols, *kw):
        return [c for c in cols if any(k in str(c) for k in kw)] or "無"
    print("  [coverage] feed 主體/執行公司欄:", pr(df.columns, "主體", "執行公司", "控股", "子公司"),
          "｜KPI欄:", pr(df.columns, "KPI", "收入", "留宿", "晚數", "毛收入", "國際"))
    if qingdan:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(qingdan, data_only=True, read_only=True)
            ws = next((wb[s] for s in wb.sheetnames if s.lower().startswith("database")), wb[wb.sheetnames[0]])
            rows = [r for i, r in enumerate(ws.iter_rows(values_only=True)) if i < 6]
            hr = next((i for i, r in enumerate(rows) if r and any("項目類型" in str(c or "") for c in r)), 2)
            hdr = [str(c or "").replace("\n", "") for c in rows[hr]]
            print("  [coverage] 清單 股權/主體欄:", pr(hdr, "股權", "主體", "執行公司", "控股", "子公司"))
            print("  [coverage] 清單 KPI欄:", pr(hdr, "KPI", "留宿", "晚數", "毛收入", "增長率", "國際住客"))
            print("  [coverage] 清單 藝術品欄:", pr(hdr, "藝術品", "展出", "館藏", "拍賣"))
        except Exception as e:
            print("  [coverage] 清單 probe skip:", e)


def _dump_pptx_text(prs, entity, with_tables=False):
    """逐版文字 dump → user paste 返，我就可以【逐句對返原報告 scan】（長度／風格／用字）。
    預設【只 dump 敘述文字，唔 dump 表格 cell】：表格係數字、已經另外驗過，而且會令個檔大到 paste 唔到。
    要連表格：build_report.py mgm --dump --dump-tables"""
    lines = []
    for i, sl in enumerate(prs.slides, 1):
        parts = []
        for sh in sl.shapes:
            if sh.has_table:
                if not with_tables:
                    t = sh.table
                    parts.append(f"〔表 {len(t.rows)}x{len(t.columns)}〕")
                    continue
                for row in sh.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                parts.append(sh.text_frame.text.strip())
        body = "\n".join(parts)
        lines.append(f"\n===== slide {i}（{len(body)} 字）=====\n" + body)
    out = Path(f"{entity}_報告_dump.txt")
    out.write_text("\n".join(lines), encoding="utf-8")
    return out


ENT_UP = "MGM"      # 由 main() 設；breadcrumb 右上角用（章節清單見 layout.SECTIONS）


def _page(prs, section_idx=0, crumb=None, headline=None):
    """開一版內容頁 = breadcrumb + footer(+頁碼) + 灰標題 + navy 導語。
    回 (slide, W, H, top_y)：top_y = 內容可以由邊開始。"""
    slide = L.blank(prs)
    W, H = L.size_of(prs)
    L.breadcrumb(slide, W, section_idx, ENT_UP)
    L.footer(slide, W, H, len(prs.slides._sldIdLst))
    top = L.page_head(slide, W, crumb, headline) if crumb else 0.5
    return slide, W, H, top


def _furniture(prs, slide, section_idx=0):
    """（保留舊 API）頂 nav tabs + 底 KPMG copyright + 初稿 + 頁碼。"""
    W, H = L.size_of(prs)
    L.breadcrumb(slide, W, section_idx, ENT_UP)
    L.footer(slide, W, H, len(prs.slides._sldIdLst))


def _dark_slide(prs):
    """新增一版深底（封面/分隔共用），回 (slide, w, h)。"""
    return L.dark_slide(prs)


def render_cover(prs, entity, date="2026年6月30日"):
    """封面（報告 p1）。template：用 TITLE SLIDE layout 填 placeholder（顏色/圖嚟自 master）；否則手砌深底。"""
    full = ENTITY_FULL.get(entity, entity.upper())
    lay = _layout(prs, "cover")
    if lay is not None:
        slide = prs.slides.add_slide(lay)
        t = _ph(slide, 0)
        if t is not None:
            t.text = f"{full}\n2025年年度投資計劃執行情況審查\n專項工作報告"
        b = _ph(slide, 11)
        if b is not None:
            b.text = f"初稿\n畢馬威會計師事務所\n{date}"
        return
    slide, w, h = _dark_slide(prs)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(7.2), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate([full, "2025年年度投資計劃執行情況審查", "專項工作報告"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(25); r.font.bold = True; r.font.color.rgb = LIGHT
        r.font.name = "微软雅黑"
    db = slide.shapes.add_textbox(Inches(0.6), Inches(4.7), Inches(4), Inches(0.5))
    dr = db.text_frame.paragraphs[0].add_run(); dr.text = "初稿"
    dr.font.size = Pt(16); dr.font.bold = True; dr.font.color.rgb = LIGHT; dr.font.name = "微软雅黑"
    fb = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(5), Inches(0.8))
    ftf = fb.text_frame
    for i, line in enumerate(["畢馬威會計師事務所", date]):
        p = ftf.paragraphs[0] if i == 0 else ftf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(11); r.font.color.rgb = LIGHT; r.font.name = "微软雅黑"


def divider(prs, title, number="", subitems=None):
    """章節分隔。template：用 Section/Appendix Divider layout（深底嚟自 master）+ 加標題 textbox；否則手砌。"""
    lay = _layout(prs, "appendix" if number == "6" else "section")
    if lay is not None:
        slide = prs.slides.add_slide(lay)
        w = prs.slide_width / 914400.0; h = prs.slide_height / 914400.0
        ny = h * 0.34; tx = 0.9
        if number:
            nb = slide.shapes.add_textbox(Inches(0.9), Inches(ny - 0.2), Inches(1.5), Inches(1.2))
            nr = nb.text_frame.paragraphs[0].add_run(); nr.text = f"{number}."
            nr.font.size = Pt(40); nr.font.bold = True; nr.font.color.rgb = LIGHT; nr.font.name = "Arial"
            tx = 2.2
        tb = slide.shapes.add_textbox(Inches(tx), Inches(ny), Inches(w - tx - 0.9), Inches(1.2))
        tr = tb.text_frame.paragraphs[0].add_run(); tr.text = title
        tr.font.size = Pt(24); tr.font.bold = True; tr.font.color.rgb = LIGHT; tr.font.name = "微软雅黑"
        if subitems:
            y = ny + 1.5
            for it in subitems:
                label, _pg = (it if isinstance(it, (tuple, list)) else (it, ""))
                rb = slide.shapes.add_textbox(Inches(tx), Inches(y), Inches(w - tx - 1.2), Inches(0.3))
                rr = rb.text_frame.paragraphs[0].add_run(); rr.text = label
                rr.font.size = Pt(12); rr.font.color.rgb = LIGHT; rr.font.name = "微软雅黑"
                y += 0.4
        return
    slide, w, h = _dark_slide(prs)
    ny = h * 0.30
    tx = 0.7
    if number:
        nb = slide.shapes.add_textbox(Inches(0.65), Inches(ny - 0.15), Inches(1.3), Inches(1.3))
        nr = nb.text_frame.paragraphs[0].add_run(); nr.text = f"{number}."
        nr.font.size = Pt(44); nr.font.bold = True; nr.font.color.rgb = LIGHT; nr.font.name = "Arial"
        tx = 2.05
    tb = slide.shapes.add_textbox(Inches(tx), Inches(ny), Inches(w - tx - 0.7), Inches(1.5))
    ttf = tb.text_frame; ttf.word_wrap = True
    tr = ttf.paragraphs[0].add_run(); tr.text = title
    tr.font.size = Pt(26); tr.font.bold = True; tr.font.color.rgb = LIGHT; tr.font.name = "微软雅黑"
    if subitems:
        y = ny + 1.7
        for it in subitems:
            label, page = (it if isinstance(it, (tuple, list)) else (it, ""))
            rb = slide.shapes.add_textbox(Inches(tx), Inches(y), Inches(w - tx - 1.4), Inches(0.32))
            rr = rb.text_frame.paragraphs[0].add_run(); rr.text = label
            rr.font.size = Pt(12); rr.font.color.rgb = RGBColor(0xC8, 0xC8, 0xD0); rr.font.name = "微软雅黑"
            if page != "":
                pb = slide.shapes.add_textbox(Inches(w - 1.3), Inches(y), Inches(0.8), Inches(0.32))
                pr = pb.text_frame.paragraphs[0].add_run(); pr.text = str(page)
                pr.font.size = Pt(12); pr.font.color.rgb = RGBColor(0xC8, 0xC8, 0xD0); pr.font.name = "微软雅黑"
            y += 0.42


def _find(dirp, entity, ext, prefer=None):
    d = Path(dirp)
    if not d.exists():
        return None
    cands = [p for p in sorted(d.rglob("*"))
             if p.suffix.lower() == ext and entity.lower() in p.name.lower()
             and not p.name.startswith("~$")]
    if not cands:
        return None
    for pk in (prefer or []):                    # 例：template 優先揀 2025 果份（唔好揀到 2023 舊報告）
        for p in cands:
            if pk in p.name:
                return p
    return cands[0]


def _fmt_ratio(v):
    """比例欄：scan 寫『(89.4%)』—— 調減（負數）用括號。"""
    try:
        f = float(v)
    except (TypeError, ValueError):
        return "" if v is None or str(v).strip() == "" else str(v)
    return f"({abs(f)*100:.1f}%)" if f < 0 else f"{f*100:.1f}%"


def _cell_txt(c, v):
    """跟欄名格式化：率／比例→%，數字→千分位（負數括號），其餘原文。"""
    if "比例" in str(c):
        return _fmt_ratio(v)
    if "率" in str(c):
        return R.fmt_pct(v)
    if _is_num(v):
        return R.fmt_money(v)
    return "" if v is None else str(v)


def _df_table(df, first_label=None):
    """DataFrame → (subs, rows, widths, supers)：認 `大組·細名` 做兩層表頭，
    第一欄空 = 範疇 section 行，尾『小計/合計/總計』= shaded。
    第一欄表頭：範疇表用「（萬澳門元）」（跟 scan 角位放單位），其餘用返欄名。"""
    cols = list(df.columns)
    if first_label is None:
        first_label = ("" if cols[0] == "序號" else       # 報告個角位空白，單位喺隔離格
                       ("萬澳門元" if cols[0] == "範疇" else cols[0]))
    grouped = any("·" in c for c in cols)
    widths = [0.4 if c == "序號" else
              (2.0 if c in ("範疇", "項目名稱", "潛在調整事項") else
               (2.8 if c == "主要涉及項目" else 0.95)) for c in cols]
    subs = [first_label] + [(c.split("·")[1] if "·" in c else c) for c in cols[1:]]
    if cols[0] == "序號" and len(cols) > 1 and cols[1] == "範疇":
        subs[1] = "萬澳門元"                          # 單位放範疇欄（跟 scan 角位）
    supers = None
    if grouped:
        groups = [""] + [(c.split("·")[0] if "·" in c else "") for c in cols[1:]]
        supers, ci = [], 1
        while ci < len(cols):
            cj = ci
            while cj + 1 < len(cols) and groups[cj + 1] == groups[ci]:
                cj += 1
            supers.append((groups[ci], ci, cj + 1)); ci = cj + 1
        supers.insert(0, ("", 0, 1))
    rows = []
    # ⚠ 加咗「序號」欄之後，範疇名喺第 2 欄 → 分類（sec/小計/總計）要睇標籤欄，唔可以睇 cols[0]，
    #   否則 小計/總計 全部當 data（冇粗體、冇橫線）。
    li = 1 if (cols[0] == "序號" and len(cols) > 1) else 0
    for _, row in df.iterrows():
        cells = [str(row[cols[0]]).strip()] + [_cell_txt(c, row[c]) for c in cols[1:]]
        lab = (cells[li] or cells[0]).strip()
        if all(str(row[c]).strip() == "" for c in cols[li + 1:]):
            rows.append(("sec", cells)); continue
        kind = ("tot" if lab.endswith("總計") else
                "subtot" if lab.endswith(("小計", "合計")) else "data")
        rows.append((kind, cells))
    return subs, rows, widths, supers


def _draw_table(slide, df, x, y, max_w, font=6.5):
    """喺 slide (x,y) 畫 navy 表（單 chunk，caller 自行分頁）。max_w=可用闊(吋)。"""
    subs, rows, widths, supers = _df_table(df)
    return L.draw_table(slide, x, y, max_w, subs, rows, widths, supers=supers,
                        font=font, hfont=font - 0.5)


def _bullets_into(box, bullets, size=8):
    """（保留舊 API）scan 敘述格式：navy 粗體小標題 + body 段落。"""
    L.prose(box, bullets, head_size=size - 1, body_size=size - 1.5)


# 概覽表顯示版式（對 scan slide 11/19）：加序號欄 + 兩層表頭。
# ⚠ 只喺【render 時】改，唔郁 overview_by_bucket 出嘅欄名 —— 所有算數/導語都靠原欄名。
_OV_GROUP = {
    "報告投資金額": "報告投資金額·金額",
    "投資計劃完成率": "報告投資金額·完成率",
    "潛在調整金額": "潛在調整後投資金額·潛在調整金額",
    "潛在調整後投資金額": "潛在調整後投資金額·金額",
    "潛在調整後投資計劃完成率": "潛在調整後投資金額·完成率",
    "設施建設/資本性支出": "潛在調整後投資金額·設施建設/資本性支出",
    "活動舉辦/營運性支出": "潛在調整後投資金額·活動舉辦/營運性支出",
}


def _hdr_cols(subs, supers):
    """表頭欄組色（項目組 2026-08-17 指定）：預設全部 HDR #1E49E2，只有【重點欄】用綠 #098E7E。
      · 1.2／1.3 概覽表 → 「獲批的計劃投資金額」
      · 4.1 金額匯總  → 最右邊「合計·潛在調整後投資金額」
    其餘表（4.2 設施/活動、期後概覽）暫時全藍——項目組未指定綠欄。"""
    out = {c: L.HDR_KEY for c, v in enumerate(subs) if str(v).strip() == "獲批的計劃投資金額"}
    last = len(subs) - 1
    if last >= 0 and str(subs[last]).strip() == "潛在調整後投資金額":
        for lab, c0, c1 in (supers or []):
            if c0 <= last < c1 and "合計" in str(lab):
                out[last] = L.HDR_KEY
    return out


def _overview_display(ov):
    """概覽表 → 顯示用：範疇前加【序號】（逐 scope 由 1 數起，section/小計/總計留空）+ 兩層表頭。"""
    if ov is None or ov.empty or "範疇" not in ov.columns:
        return ov
    d = ov.copy()
    seq, n = [], 0
    for _, r in d.iterrows():
        lab = str(r["範疇"]).strip()
        others = [c for c in d.columns if c != "範疇"]
        if all(str(r[c]).strip() == "" for c in others):      # section 行（博彩項目/非博彩項目）
            n = 0; seq.append("")
        elif lab.endswith(("小計", "總計", "合計")) or lab.startswith(("原計劃", "投資執行報告",
                                                                   "承諾的", "2025年投資支出")):
            seq.append("")
        else:
            n += 1; seq.append(str(n))
    d.insert(0, "序號", seq)
    return d.rename(columns={k: v for k, v in _OV_GROUP.items() if k in d.columns})


def render_overview_page(prs, crumb, headline, table_df, bullets, *, sec=0, table_name=None,
                         note=None):
    """報告概述式 2 欄版（對 scan slide 10/15）：crumb + navy 導語，左 表，右 敘述。"""
    slide, W, H, top = _page(prs, sec, crumb, headline)
    left_w = W * 0.60
    tbl_bot = top
    if table_df is not None and not table_df.empty:
        if table_name:
            top = L.caption_bar(slide, L.MARGIN, top, left_w, table_name)
        subs, rows, widths, supers = _df_table(_overview_display(table_df))
        wid = [w * left_w / sum(widths) for w in widths]
        avail = L.CONTENT_BOTTOM - top - 0.30          # 留位俾表下面個「註」
        hh = L.header_h(supers, subs, wid, L.SZ_TBL_HDR)
        font = L.SZ_TBL
        while font > 4.5 and sum(L.row_h(c, wid, font) for _, c in rows) > avail - hh:
            font -= 0.25
        tbl_bot, _ = L.draw_table(slide, L.MARGIN, top, left_w, subs, rows, widths,
                                  supers=supers, font=font, hfont=max(4.5, font - 0.5),
                                  fill_h=avail, left_cols=2,   # 序號 + 範疇 都左對齊（對報告）
                                  hdr_cols=_hdr_cols(subs, supers))
    if note:      # 「註」貼喺表底下，唔可以同底部嘅資料來源疊字
        L.put(slide, L.MARGIN, min(tbl_bot + 0.06, L.CONTENT_BOTTOM - 0.30), left_w, 0.3,
              note, size=L.SZ_NOTE - 1, italic=True, color=L.GREY)
    rx = L.MARGIN + left_w + 0.22
    L.prose_box(slide, rx, top - 0.02, W - rx - L.MARGIN, L.CONTENT_BOTTOM - top, bullets)
    L.source_note(slide, W)


def render_overview_pages(prs, crumb, headline, table_df, bullets, *, sec=0, table_name=None,
                          note=None, grouped=False):
    """同 render_overview_page，但右邊敘述長就自動分版，【左邊同一個表逐版重複】。
    對 scan slide 11-14：1.3 四版全部都係左邊 1.2 嗰個整體概況表 + 右邊唔同段落。
    grouped=True 時 bullets = [(右欄小標題, [(head, body)…])]，每組至少一版（報告 1/4…4/4）。"""
    if not bullets:
        return
    W, _H = L.size_of(prs)
    left_w = W * 0.60
    rx = L.MARGIN + left_w + 0.22
    colw = W - rx - L.MARGIN
    top = L.HEAD_Y + L.head_h(f"{headline}（1/9）", W)[0] + 0.10 + (0.20 if table_name else 0)
    avail = L.CONTENT_BOTTOM - top
    pages = []
    for grp in (bullets if grouped else [(None, bullets)]):
        head, items = grp if grouped else grp
        if not items:
            continue
        chunks = L.fit_prose(items, colw, avail - (0.24 if head else 0),
                             head_size=L.SZ_BODY_HEAD, body_size=L.SZ_BODY)
        for ci, ch in enumerate(chunks):
            pages.append(([(head + ("（續）" if ci else ""), "")] if head else []) + ch)
    for pi, page in enumerate(pages):
        render_overview_page(prs, crumb, headline + _pg(pi + 1, len(pages)), table_df, page,
                             sec=sec, table_name=table_name, note=note)


def _total_line(df):
    """由表自己嘅『總計』行砌一句機械導語（避免「淨得個表冇文字」）。"""
    cols = list(df.columns)
    tot = df[df[cols[0]].astype(str).str.strip().str.endswith("總計")]
    if tot.empty:
        return ""
    r = tot.iloc[0]
    parts = []
    for c in cols[1:]:
        v = r[c]
        if not _is_num(v) or float(v) == 0:
            continue
        lab = c.replace("·", "－")
        parts.append(f"{lab} {_cell_txt(c, v)}")
    return "；".join(parts[:6]) + "（單位：萬澳門元，除完成率外）。" if parts else ""


def _table_bullets(df):
    """冇 LLM 時嘅機械表旁敘述（由表自己嘅小計/總計行計，自洽）→ [(小標題, 內容)]。"""
    cols = list(df.columns)
    first = df[cols[0]].astype(str).str.strip()
    num = [c for c in cols[1:] if pd.to_numeric(df[c], errors="coerce").notna().any()]
    if not num:
        return []
    out = []
    tot = df[first.str.endswith("總計")]
    if not tot.empty:
        out.append(("整體情況", "全部範疇合計 " +
                    "；".join(f"{c.replace('·', '－')} {_cell_txt(c, tot.iloc[0][c])}"
                              for c in num[:5]) + "（單位：萬澳門元，除完成率外）。"))
    subs = df[first.str.endswith("小計")]
    if not subs.empty:
        out.append(("博彩／非博彩分佈", "；".join(
            f"{r[cols[0]]} " + "、".join(f"{c.replace('·', '－')} {_cell_txt(c, r[c])}"
                                         for c in num[:3]) for _, r in subs.iterrows()) + "。"))
    val = next((c for c in num if "報告" in c or "合計" in c), num[0])
    data = df[~first.str.endswith(("小計", "總計")) &
              df[cols[1:]].astype(str).apply(lambda r: any(x.strip() for x in r), axis=1)].copy()
    data["_v"] = pd.to_numeric(data[val], errors="coerce")
    top = data.sort_values("_v", ascending=False).head(4)
    if not top.empty:
        out.append((f"金額最大的範疇（按{val.replace('·', '－')}）", "、".join(
            f"{r[cols[0]]}（{_cell_txt(val, r[val])}）" for _, r in top.iterrows()) + "。"))
    return out


# 10年投資預算：唔喺 feed／清單／表2（2026-08-12 全 cell 搜尋確認），嚟自承批合同 → 外部 config。
# 放 data/10year_budget.yml（gitignored；public repo 唔可以有客戶數）：
#   mgm: {總計: 1970000, 博彩: 996800, 非博彩: 973200}
BUDGET_FILES = ["data/10year_budget.yml", "data/10year_budget.json", "conf/local/10year_budget.yml"]


def _load_budget(entity):
    for fn in BUDGET_FILES:
        p = Path(fn)
        if not p.exists():
            continue
        try:
            if p.suffix == ".json":
                import json
                d = json.loads(p.read_text(encoding="utf-8"))
            else:
                import yaml
                d = yaml.safe_load(p.read_text(encoding="utf-8"))
            b = (d or {}).get(entity) or (d or {}).get(entity.upper()) or {}
            if b:  # noqa: SIM102
                print(f"    10年投資預算：{fn} → 總計 {b.get('總計', 0):,.0f} 萬澳門元")
                return {k: float(v) for k, v in b.items()}
        except Exception as e:
            print(f"    ⚠ 讀唔到 {fn}: {e}")
    return {}


def _proj_counts(sdf, plan, ov=None):
    """(n_plan, n_impl, n_zero) —— 表尾數量行同 1.2 敘述共用，保證三個數同一母體、自洽。
    n_plan 優先攞【概況表自己個「總計」項目數量】，咁「總計 = 已實施 + 未實施」一定成立
    （項目組 2026-08-15：總計 84 但 79+10 唔等於 84）。冇表就用清單 2025 計劃金額 > 0 嘅碼
    （0 行唔算「獲批開展」，否則出 256 個）；再冇就退化成 feed 有支出嘅碼（n_zero=0，句子會略去）。"""
    d = sdf[sdf["_bucket"] == S.BUCKET_ORDER[0]]
    plan25 = {k for k, v in ((plan or {}).get(25, {}) or {}).items()
              if isinstance(v, (int, float)) and v > 0}
    spent = {(str(r["ng_scope"]) == "gaming", B._norm(r["dicj code"]))
             for _, r in d[pd.to_numeric(d["調整前_萬"], errors="coerce").fillna(0) != 0]
             .drop_duplicates(["ng_scope", "dicj code"]).iterrows()}
    n_impl = len(plan25 & spent) if plan25 else len(spent)
    n_plan = _tot_projects(ov) or (len(plan25) if plan25 else len(spent))
    n_plan = max(n_plan, n_impl)
    return n_plan, n_impl, n_plan - n_impl


def _tot_projects(ov):
    """概況表「總計」行嘅項目數量（0 = 攞唔到）。"""
    if ov is None or getattr(ov, "empty", True) or "項目數量" not in getattr(ov, "columns", []):
        return 0
    t = ov[ov["範疇"].astype(str).str.strip() == "總計"]
    if t.empty:
        return 0
    v = pd.to_numeric(pd.Series([t.iloc[0]["項目數量"]]), errors="coerce").iloc[0]
    return int(v) if pd.notna(v) and v > 0 else 0


def _overview_extra(ov, plan, sdf, budget, ent_up):
    """1.2 概況表尾段（scan slide 11）：原計劃未實施／已實施項目數量 + 承諾的10年投資預算 + 佔比。
    項目數量計得到；10年預算要 config，冇就唔出嗰兩行。"""
    cols = list(ov.columns)
    n_plan, n_impl, n_zero = _proj_counts(sdf, plan, ov)
    # 報告字眼（IMG_0441）：未實施個數寫括號，表示係計劃總數入面嗰部分
    rows = [{cols[0]: "原計劃中未實施的投資項目數", cols[1]: f"({n_zero})" if n_zero else "-"},
            {cols[0]: "投資執行報告中申報已實施的投資項目數", cols[1]: n_impl}]
    # 10年投資預算：全 cell 搜過清單 + 表2 都冇（2026-08-12 確認），嚟自承批合同。
    # 冇 config 就【照出行、留空】—— 保持報告結構，一眼睇到係待填而唔係漏咗（user 2026-08-12）。
    tot = ov[ov["範疇"].astype(str).str.strip() == "總計"]
    b_all = budget.get("總計") if budget else None
    # ⚠ 冇 budget 都要填「-」，唔可以留空 —— 成行全空會俾 _df_table 當做【範疇 section 行】
    #   （→ 變粗體）。「-」亦係報告表示「冇數」嘅寫法。
    _bcol = "獲批的計劃投資金額" if "獲批的計劃投資金額" in cols else "報告投資金額"
    rows.append({cols[0]: "承諾的10年投資預算", _bcol: b_all if b_all else "-"})
    r = {cols[0]: "2025年投資支出佔10年投資預算的完成率"}
    for c in ("報告投資金額", "潛在調整後投資金額"):
        if c in cols and len(tot) and b_all:
            v = pd.to_numeric(pd.Series([tot.iloc[0][c]]), errors="coerce").iloc[0]
            r[c] = B._rate(float(v or 0), b_all)
        else:
            r[c] = "-"
    rows.append(r)
    return pd.concat([ov, pd.DataFrame(rows)], ignore_index=True) if rows else ov


def _bucket_adj_table(ov):
    """期後調整事項匯總嘅左表（對 scan p-11）：範疇 × 報告(a)｜潛在調整(b)｜調整後(c=a+b)｜b/a。"""
    keep = ["範疇", "項目數量", "報告投資金額", "潛在調整金額", "潛在調整後投資金額"]
    d = ov[[c for c in keep if c in ov.columns]].copy()
    rep = pd.to_numeric(d.get("報告投資金額"), errors="coerce")
    adj = pd.to_numeric(d.get("潛在調整金額"), errors="coerce")
    ratio = (adj / rep).where(rep.abs() > 0).astype(object)   # object：section 行要填 ""（避 pandas FutureWarning）
    ratio[d["範疇"].astype(str).str.strip().eq("") | rep.isna()] = ""
    # ⚠ 分母 0（例如博彩全 0 嘅小計行）→ 出「-」，唔可以出 nan%
    ratio[ratio.isna()] = "-"
    d["潛在調整金額佔報告投資金額比例"] = ratio
    return d


def render_bucket_adjustment(prs, ent_up, bk, sdf, ov, narr, llm=None):
    """② 期後【調整事項匯總】（scan p-11 / p-13）：左 = 範疇 × 調整表，
    右 = navy 小標題 + 編號清單『N. 類型（金額）：說明』，編號跟七大類 canonical 序（會跳號）。"""
    yr = bk[:4]
    d = sdf[sdf["_bucket"] == bk].copy()
    if d.empty:
        return
    d["_adj"] = d["調整一級"].map(B.CANON).fillna(d["調整一級"])
    llm_bkt = (llm or {}).get("bkt", {})
    items = []
    for i, t in enumerate(B.ADJ7, start=1):      # 編號＝七大類 canonical 序
        sub = d[d["_adj"] == t]
        amt = pd.to_numeric(sub["調整_萬"], errors="coerce").sum()
        if abs(amt) < 0.5:
            continue
        body = llm_bkt.get(bkt_key(bk, t)) or (
            f"在{yr}年度投資計劃期後投資金額中，{ent_up}同樣申報了{t}。")
        body += (f"與前述的「2025年度投資計劃報告投資金額的潛在調整事項匯總」一致，"
                 f"我們就{yr}年度投資計劃期後投資金額同樣建議了該項調整。詳見後續主要發現「{t}」。")
        items.append((i, f"{t}（{_amt(amt)}）：", body))
    if not items:
        return
    tot = pd.to_numeric(d["調整_萬"], errors="coerce").sum()
    npj = int(d[pd.to_numeric(d["調整_萬"], errors="coerce") != 0]["dicj code"].nunique())
    aft = pd.to_numeric(d["調整後_萬"], errors="coerce").sum()
    head = (f"基於我們的各項審查程序，我們認為{ent_up}報告的{yr}年度投資計劃在2025年繼續執行的支出中，"
            f"存在{_cn(len(items))}大類的調整事項，潛在調減投資金額約{_amt(tot)}"
            f"（涉及{npj}個投資項目），經潛在調減後的{yr}年度計劃投資項目在2025年的投資金額約{_amt(aft)}。")
    S2 = "過往年度投資計劃在2025年繼續執行的審查跟進"
    tname = f"{ent_up} {yr}年度投資計劃於2025年申報的期後投資金額的潛在調整"
    W, H = L.size_of(prs)
    lw = W * 0.55
    tbl = _bucket_adj_table(ov)
    subs, rows, widths, supers = _df_table(tbl)
    wid = [w * lw / sum(widths) for w in widths]
    cw = W - (L.MARGIN + lw + 0.22) - L.MARGIN
    pages = L.fit_prose([(h, b) for _n, h, b in items], cw, L.CONTENT_BOTTOM - 1.9)
    idx = 0
    for pi, page in enumerate(pages):
        suffix = _pg(pi + 1, len(pages))
        slide, W, H, top = _page(prs, 1, f"{S2}  |  {yr}年度投資計劃報告投資金額的潛在調整事項匯總",
                                 head + suffix)
        if pi == 0:
            t2 = L.caption_bar(slide, L.MARGIN, top, lw, tname)
            L.draw_table(slide, L.MARGIN, t2, lw, subs, rows, widths, supers=supers,
                         font=L.SZ_TBL, hfont=L.SZ_TBL_HDR, fill_h=L.CONTENT_BOTTOM - t2 - 0.28)
            L.put(slide, L.MARGIN, L.CONTENT_BOTTOM - 0.26, lw, 0.3,
                  "註：金額單位為萬澳門元；括號表示調減。", size=L.SZ_NOTE - 1, italic=True, color=L.GREY)
        box = L._tb(slide, L.MARGIN + lw + 0.22, top - 0.02, cw, L.CONTENT_BOTTOM - top)
        L.prose_numbered(box, items[idx:idx + len(page)], size=L.SZ_BODY,
                         title=(tname if pi == 0 else tname + "（續）"))
        idx += len(page)
        L.source_note(slide, W, more=(pi < len(pages) - 1))


def _cum_table(df, plan, cat=None):
    """三年累計表（scan slide 26）→ DataFrame（`·` = 兩層表頭）。
    每個計劃年 Y：獲批(a)=清單計劃｜2025年前已獲認可(b)=報告年<25 調整後｜2025年期後(c)=報告年25 調整後
    ｜合計(d=b+c)｜完成率(d/a)。2025計劃冇 b。尾段＝三年累計 Σa｜Σd｜Σd/Σa。"""
    d = df[df["dicj code"].astype(str).str.match(r"^項目\s*\d")].copy()
    if d.empty or not plan:
        return pd.DataFrame()
    d["_sub"] = FS.sub_of(d)
    d["_g"] = (d["ng_scope"] == "gaming")
    d["_ry"] = pd.to_numeric(d["報告年"], errors="coerce")
    d["_af"] = pd.to_numeric(d["調整後_萬"], errors="coerce").fillna(0)
    d["_ngn"] = d["ng_code"].map(S._ngn)
    d["_go"] = d["_sub"].map(lambda x: S.GORDER.get(x, 5))
    order = (d.drop_duplicates(["_g", "_sub"]).sort_values(["_g", "_go", "_ngn", "_sub"],
             ascending=[False, True, True, True])[["_g", "_sub"]].values.tolist())
    code_sub = {(bool(r["_g"]), B._norm(r["dicj code"])): str(r["_sub"])
                for _, r in d.drop_duplicates(["_g", "dicj code"]).iterrows()}
    # feed 有嘅項目 → 直接知範疇；feed 冇（零投資／往年項目）→ 用清單『項目性質』學返
    d2sub = {}
    for (gm, code), sub in code_sub.items():
        dv = (cat or {}).get((gm, code))
        if dv:
            d2sub.setdefault(str(dv), set()).add(sub)
    d2sub1 = {k: next(iter(v)) for k, v in d2sub.items() if len(v) == 1}
    A, miss = {}, 0.0
    for yr in (23, 24, 25):
        for (gm, code), v in (plan.get(yr, {}) or {}).items():
            sub = code_sub.get((bool(gm), code)) or d2sub1.get(str((cat or {}).get((bool(gm), code), "")))
            if sub:
                A[(yr, bool(gm), sub)] = A.get((yr, bool(gm), sub), 0.0) + float(v or 0)
            else:
                miss += float(v or 0)
    if miss > 0.5:
        print(f"    ⚠ 2.5 三年累計：{miss:,.0f} 萬計劃金額派唔到範疇（清單項目性質對唔到）")

    def _af(yr, gm, sub, pre):
        m = ((d["_plan_year"] == yr) & (d["_g"] == gm) & (d["_sub"] == sub) &
             ((d["_ry"] < 25) if pre else (d["_ry"] >= 25)))
        return round(float(d.loc[m, "_af"].sum()), 1)

    Y3, Y5 = "2023年度投資計劃", "2025年度投資計劃"
    Y4, CUM = "2024年度投資計劃", "截至2025年末三年累計"
    cols = ["範疇"]
    for y in (Y3, Y4):
        cols += [f"{y}·獲批的計劃投資金額", f"{y}·2025年前已獲認可", f"{y}·2025年期後",
                 f"{y}·合計", f"{y}·完成率"]
    cols += [f"{Y5}·獲批的計劃投資金額", f"{Y5}·潛在調整後投資金額", f"{Y5}·完成率"]
    cols += [f"{CUM}·獲批的計劃投資金額", f"{CUM}·獲認可／潛在調整後投資金額", f"{CUM}·完成率"]

    def line(name, cells):
        r = {"範疇": name}
        r.update(dict(zip(cols[1:], cells)))
        return r

    def calc(items):
        """items = [(gm, sub)] → 12 個數（同一套公式，逐個範疇同小計/總計共用）。"""
        out, ta, td = [], 0.0, 0.0
        for yr, y in ((23, Y3), (24, Y4), (25, Y5)):
            a = round(sum(A.get((yr, g, sb), 0.0) for g, sb in items), 1)
            b = round(sum(_af(yr, g, sb, True) for g, sb in items), 1)
            c = round(sum(_af(yr, g, sb, False) for g, sb in items), 1)
            out += ([a, b, c, round(b + c, 1), B._rate(b + c, a)] if yr != 25
                    else [a, c, B._rate(c, a)])
            ta += a; td += b + c
        return out + [round(ta, 1), round(td, 1), B._rate(td, ta)]

    rows, all_items = [], []
    for gm, label in ((True, "博彩項目"), (False, "非博彩項目")):
        items = [(g, sb) for g, sb in order if g == gm]
        if not items:
            continue
        rows.append({"範疇": label})
        for g, sb in items:
            rows.append(line(sb, calc([(g, sb)])))
        rows.append(line(f"{label}小計", calc(items)))
        all_items += items
    rows.append(line("合計", calc(all_items)))
    return pd.DataFrame(rows, columns=cols)


def _collect_toc(prs, W, H):
    """由【已起好嘅版】反推目錄：深色分隔頁 = 章節；內容頁 crumb「章節 | 子題」= 子項。
    頁碼 = 插入目錄之後嘅位置（+1）。"""
    ent, last = [], None
    for i, sl in enumerate(prs.slides, 1):
        full = any(sh.width / 914400.0 > W - 0.1 and sh.height / 914400.0 > H - 0.1
                   for sh in sl.shapes)
        pg = i + 1                                  # 目錄會插喺第 2 版
        if full:
            texts = [sh.text_frame.text.strip() for sh in sl.shapes
                     if sh.has_text_frame and sh.text_frame.text.strip()]
            no = next((t for t in texts if re.fullmatch(r"\d+\.", t)), "")
            ttl = next((t for t in texts if len(t) >= 2 and t != "KPMG"
                        and not re.fullmatch(r"\d+\.", t)), "")
            if no and ttl:
                ent.append((no, ttl.split("\n")[0], False, pg)); last = None
            continue
        for sh in sl.shapes:
            if sh.has_text_frame and 0.28 < sh.top / 914400.0 < 0.46 and "  |  " in sh.text_frame.text:
                sub = sh.text_frame.text.split("  |  ")[-1].strip()
                sub = re.sub(r"（\d+/\d+）$", "", sub).strip()
                if sub and sub != last:
                    ent.append(("", sub, True, pg)); last = sub
                break
    return ent


def _renumber_footers(prs, W, H):
    """插咗目錄之後，全份『初稿 N』重編（頁碼喺建版時已寫死）。"""
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if not sh.has_text_frame or sh.top / 914400.0 < H - 0.45:
                continue
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.text.startswith("初稿"):
                        r.text = f"初稿　{i}"


def _sec_slides(prs, W, H):
    """{章 index(0-based): slide index} —— 深色分隔頁上嘅「N.」認章號，供 breadcrumb 頁籤跳頁。"""
    out = {}
    for i, sl in enumerate(prs.slides):
        if not any(sh.width / 914400.0 > W - 0.1 and sh.height / 914400.0 > H - 0.1
                   for sh in sl.shapes):
            continue
        for sh in sl.shapes:
            m = sh.has_text_frame and re.fullmatch(r"(\d+)\.", sh.text_frame.text.strip())
            if m:
                out.setdefault(int(m.group(1)) - 1, i)
                break
    return out


def _move_slide(prs, frm, to):
    """把第 frm 版（0-based）搬去 to。"""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[frm]); lst.insert(to, ids[frm])


def render_toc(prs, ent_up, entries):
    """報告 slide 7 目錄：六大章節 + 子項 + 頁碼（頁碼喺 build 完先知 → 由 caller 傳）。
    子項多過一版就自動分版。"""
    avail = L.CONTENT_BOTTOM - (L.HEAD_Y + 0.06) - 0.55
    pages, cur, used = [], [], 0.0
    for e in entries:
        h = 0.30 if e[2] else 0.34
        if cur and used + h > avail:
            pages.append(cur); cur, used = [], 0.0
        cur.append(e); used += h
    if cur:
        pages.append(cur)
    for pi, page in enumerate(pages):
        _render_toc_page(prs, ent_up, page, _pg(pi + 1, len(pages)))


def _render_toc_page(prs, ent_up, entries, suffix=""):
    slide, W, H, top = _page(prs, 0, f"{ent_up} 2025年年度投資計劃執行情況審查專項工作報告  |  目錄",
                             None)
    L.put(slide, L.MARGIN, top, W - 2 * L.MARGIN, 0.35, "目錄" + suffix, size=18, bold=True, color=HDR)
    y = top + 0.55
    for no, title, sub, pg in entries:
        L.put(slide, L.MARGIN, y, 0.6, 0.26, no, size=L.SZ_BODY_HEAD, bold=True, color=HDR)
        L.put(slide, L.MARGIN + 0.62, y, W - 2 * L.MARGIN - 1.5, 0.26, title,
              size=L.SZ_BODY_HEAD, bold=not sub, color=HDR if not sub else L.INK)
        if pg:
            L.put(slide, W - L.MARGIN - 0.7, y, 0.7, 0.26, str(pg), size=L.SZ_BODY,
                  color=L.GREY, align=PP_ALIGN.RIGHT)
        y += 0.30 if sub else 0.34
    L.source_note(slide, W, note="")


def render_visit_summary(prs, ent_up, df, threshold=2000):
    """報告 slide 71『設施建設項目現場走訪情況匯總』：樣本選取標準 + 樣本量表。
    全部由 feed 計（capex 項目母體 vs 走訪樣本），冇外部資料都做到。"""
    cap = df[df["final_capex_opex"] == "Capex"].copy()
    cap = cap[cap["dicj code"].astype(str).str.match(r"^項目\s*\d")]
    if cap.empty:
        return
    g = cap.groupby(["ng_scope", "dicj code"]).agg(報告=("調整前_萬", "sum")).reset_index()
    n_all, amt_all = len(g), float(g["報告"].sum())
    sel = g[g["報告"] >= threshold]
    n_sel, amt_sel = len(sel), float(sel["報告"].sum())
    rows = pd.DataFrame([
        {"項目": "設施建設（資本性支出）項目母體", "項目數量": n_all, "涉及金額": round(amt_all, 1),
         "佔母體金額比例": 1.0},
        {"項目": f"現場走訪樣本（單一項目資本性支出 ≥ {threshold:,.0f} 萬澳門元）",
         "項目數量": n_sel, "涉及金額": round(amt_sel, 1),
         "佔母體金額比例": (amt_sel / amt_all if amt_all else None)},
        {"項目": "未列入走訪樣本", "項目數量": n_all - n_sel,
         "涉及金額": round(amt_all - amt_sel, 1),
         "佔母體金額比例": ((amt_all - amt_sel) / amt_all if amt_all else None)},
    ])
    head = (f"我們在制定本次審查工作範圍時，計劃就{ent_up}報告投資金額中重大的設施建設項目開展現場走訪。"
            f"我們根據管理層提供的投資項目底層財務明細，就設施建設（資本性支出）項目共{n_all}個"
            f"（涉及{_amt(amt_all)}）作為母體，選取單一項目資本性支出達{threshold:,.0f}萬澳門元或以上之"
            f"{n_sel}個項目進行現場走訪，涉及金額{_amt(amt_sel)}，佔母體金額"
            f"{_pct(amt_sel / amt_all) if amt_all else '—'}。")
    render_generic(prs, f"{ent_up} 設施建設項目現場走訪情況匯總", rows, sec=3,
                   crumb="其他信息  |  本次審查工作執行的程序匯總", headline=head, side=False,
                   note="註：金額單位為萬澳門元；母體為報告投資金額中歸類為設施建設（資本性支出）之項目。")


def render_artwork(prs, ent_up, biao2_dir="data/表2", entity="mgm"):
    """報告 slide 101『藝術品展出情況清單』：表2 附件『藝術品』sheet 逐件列示。"""
    cols, body = B2.load_artwork(biao2_dir, entity, log=print)
    if not cols or not body:
        return
    KEEP = ["名稱", "類別", "Artist", "購入", "當前位置", "當前狀態", "展出紀錄"]
    idx = [i for i, c in enumerate(cols) if any(k in c for k in KEEP)]
    if not idx:
        idx = list(range(min(7, len(cols))))
    hdr = ["序號"] + [cols[i][:14] for i in idx]
    rows = pd.DataFrame([[str(n)] + [body[n - 1][i][:60] for i in idx]
                         for n in range(1, len(body) + 1)], columns=hdr)
    st_i = next((i for i, c in enumerate(cols) if "當前狀態" in c), None)
    shown = sum(1 for r in body if st_i is not None and "展出" in r[st_i]) if st_i is not None else 0
    head = (f"下表列示{ent_up}已購入之藝術品共{len(body)}件之展出情況"
            + (f"，其中{shown}件現正展出，{len(body) - shown}件未在展出中" if st_i is not None else "")
            + "。藝術品之社會價值須透過面向公眾持續展出方能體現，故其展出情況為本次審查"
              "「未完全實現投資目的的投資支出」之判斷依據。")
    render_generic(prs, f"{ent_up} 藝術品展出情況清單", rows, sec=5,
                   crumb="附件  |  藝術品展出情況清單", headline=head, side=False,
                   note="資料來源：承批公司提供之藝術品清單（審查底稿表2 附件），畢馬威分析")


def render_cumulative(prs, ent_up, df, plan, cat=None):
    """2.5 截至2025年末投資金額概覽（scan slide 26）。"""
    tbl = _cum_table(df, plan, cat)
    if tbl.empty:
        return
    t = tbl[tbl["範疇"].astype(str).str.strip() == "合計"]
    r = t.iloc[0] if len(t) else None

    def v(c):
        x = r.get(c) if r is not None else None
        return float(x) if isinstance(x, (int, float)) and not pd.isna(x) else 0.0
    CUM = "截至2025年末三年累計"
    a, dd = v(f"{CUM}·獲批的計劃投資金額"), v(f"{CUM}·獲認可／潛在調整後投資金額")
    g = tbl[tbl["範疇"].astype(str).str.strip() == "博彩項目小計"]
    ng = tbl[tbl["範疇"].astype(str).str.strip() == "非博彩項目小計"]

    def sv(x, c):
        return (float(x.iloc[0][c]) if len(x) and isinstance(x.iloc[0][c], (int, float))
                and not pd.isna(x.iloc[0][c]) else 0.0)
    head = (f"截至2025年末，{ent_up}於2023至2025年三年的年度投資計劃累計獲批的計劃投資金額為{_amt(a)}"
            f"（博彩項目{_amt(sv(g, f'{CUM}·獲批的計劃投資金額'))}和"
            f"非博彩項目{_amt(sv(ng, f'{CUM}·獲批的計劃投資金額'))}），"
            f"累計獲認可／潛在調整後投資金額為{_amt(dd)}，"
            f"累計投資計劃金額完成率為{_pct(B._rate(dd, a))}"
            f"（博彩項目完成率為{_pct(sv(g, f'{CUM}·完成率'))}，"
            f"非博彩項目完成率為{_pct(sv(ng, f'{CUM}·完成率'))}）。")
    render_generic(prs, f"{ent_up} 截至2025年末投資金額概覽", tbl.fillna(""), sec=1,
                   crumb="過往年度投資計劃在2025年繼續執行的審查跟進  |  截至2025年末投資金額概覽",
                   headline=head, side=False,
                   note="註：金額單位為萬澳門元。「2025年前已獲認可」＝該年度計劃於2025年之前"
                        "（即當年及往年審查）已認可之投資金額；「2025年期後」＝於2025年發生之期後投資。")


def render_generic(prs, title, df, *, sec=3, crumb=None, headline=None, note=None,
                   llm=None, tbl_id=None, side=None):
    """單張表（範疇/項目 + 數字欄；·=2-row group header）。逐頁：crumb + 導語 + caption bar
    + 表 + 資料來源，按【累積高度】分頁（唔會超出版面）。

    side=True → 報告式 2 欄（表左 + 敘述右，對 scan p-10~p-13）。敘述優先用 LLM 寫嘅
    `llm['tbl'][tbl_id]`，冇就用 _table_bullets 機械生成。欄數少（≤8）而且一版放得落先會用。"""
    if df is None or df.empty:
        return
    raw = (llm or {}).get("tbl", {}).get(tbl_id) or []
    if isinstance(raw, dict):                      # 新 shape：{"導語": …, "段落": [[h,b],…]}
        headline = raw.get("導語") or headline
        raw = raw.get("段落") or []
    bullets = [tuple(x) for x in raw] or _table_bullets(df)
    if side is None:
        side = len(df.columns) <= 8 and bool(bullets)
    if side and bullets:
        W, _H = L.size_of(prs)
        lw = W * 0.60
        s2, r2, w2, sp2 = _df_table(df)
        wid2 = [w * lw / sum(w2) for w in w2]
        need = L.header_h(sp2, s2, wid2, 5.0) + sum(L.row_h(c, wid2, 5.0) for _, c in r2)
        if need <= L.CONTENT_BOTTOM - 1.9:              # 一版放得落先用 2 欄，否則落返全闊分頁
            render_overview_page(prs, (crumb or title), headline or _total_line(df), df,
                                 bullets, sec=sec, table_name=title, note=note)
            return
    df = _overview_display(df) if "範疇" in df.columns and "項目數量" in df.columns else df
    subs, rows, widths, supers = _df_table(df)
    W, H = L.size_of(prs)
    tw = W - 2 * L.MARGIN
    wid = [w * tw / sum(widths) for w in widths]
    head = headline or _total_line(df)
    crumb = crumb or title
    # 先用一版試高度（導語行數會食掉可用高）
    probe_top = L.HEAD_Y + L.head_h(head, W)[0] + 0.10
    avail = L.CONTENT_BOTTOM - probe_top - 0.24
    hh = L.header_h(supers, subs, wid, 5.5)
    hcols = _hdr_cols(subs, supers)      # 全闊表一樣要派重點欄色（4.1 最右「潛在調整後投資金額」）
    pages = L.fit_rows(rows, wid, 6.5, avail, hh)
    for pi, chunk in enumerate(pages):
        suffix = f"（{pi+1}/{len(pages)}）" if len(pages) > 1 else ""
        slide, W, H, top = _page(prs, sec, crumb, (head or "") + suffix)
        top = L.caption_bar(slide, L.MARGIN, top, tw, title + suffix)
        L.draw_table(slide, L.MARGIN, top, tw, subs, chunk, widths, supers=supers,
                     font=L.SZ_TBL, hfont=L.SZ_TBL_HDR, fill_h=L.CONTENT_BOTTOM - top - 0.28,
                     hdr_cols=hcols)
        L.source_note(slide, W, note=note, more=(pi < len(pages) - 1))


def _cards(prs, sec, crumb, headline, recs, *, note=None):
    """逐個項目一張 card（navy 標題條 + 敘述段），按【累積高度】排版分頁 → 填滿版面唔留大白位。
    recs = [(bar_text, [(label, body)])]。"""
    W, H = L.size_of(prs)
    cw = W - 2 * L.MARGIN
    probe = L.HEAD_Y + L.head_h(headline, W)[0] + 0.10
    avail = L.CONTENT_BOTTOM - probe

    def card_h(items):
        return 0.24 + L.est_prose_h(items, cw - 0.12, head_size=L.SZ_BODY, body_size=L.SZ_BODY, gap=3) + 0.14
    pages, cur, used = [], [], 0.0
    for rec in recs:
        h = card_h(rec[1])
        if cur and used + h > avail:
            pages.append(cur); cur, used = [], 0.0
        cur.append(rec); used += min(h, avail)
    if cur:
        pages.append(cur)
    for pi, page in enumerate(pages):
        suffix = f"（{pi+1}/{len(pages)}）" if len(pages) > 1 else ""
        slide, W, H, y = _page(prs, sec, crumb, (headline or "") + suffix)
        for bar_text, items in page:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(L.MARGIN), Inches(y),
                                         Inches(cw), Inches(0.22))
            bar.fill.solid(); bar.fill.fore_color.rgb = HDR      # #00338D（card 條，唔跟表頭色）
            bar.line.fill.background(); bar.shadow.inherit = False
            btf = bar.text_frame; btf.word_wrap = True
            btf.margin_left = Emu(54000); btf.margin_top = btf.margin_bottom = Emu(0)
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            br = btf.paragraphs[0].add_run(); br.text = bar_text
            L.setfont(br, 8, bold=True, color=LIGHT)   # Arial latin + Microsoft YaHei ea（同全份一致）
            bh = min(L.est_prose_h(items, cw - 0.12, head_size=L.SZ_BODY, body_size=L.SZ_BODY, gap=3),
                     L.CONTENT_BOTTOM - y - 0.26)
            L.prose_box(slide, L.MARGIN + 0.06, y + 0.26, cw - 0.12, bh, items,
                        head_size=L.SZ_BODY, body_size=L.SZ_BODY, gap=3)
            y += 0.24 + bh + 0.14
        L.source_note(slide, W, note=note, more=(pi < len(pages) - 1))


def _finding_body(box, find, mgmt, grey=None):
    """（保留舊 API）KPMG分析發現 / 管理層解釋 兩段。"""
    L.prose(box, [(l + "：", t) for l, t in
                  [("KPMG分析發現", find), ("管理層解釋", mgmt)] if t],
            head_size=L.SZ_BODY, body_size=L.SZ_BODY, gap=3)


def render_findings(prs, ent_up, df, narr, llm=None, b2=None):
    """③ 主要發現：每 canonical 調整類型 → 受影響項目 card = navy 標題條(項目+金額) + body。
    body 優先用 LLM 寫嘅『事項描述』（ground 住表2＋清單），管理層原話照樣保留；
    冇 LLM 就 fallback 清單抄字（KPMG分析發現／管理層解釋）。"""
    llm_proj = (llm or {}).get("proj", {})
    d = df.copy()
    d["_adj"] = d["調整一級"].map(B.CANON).fillna(d["調整一級"])
    for adj in B.ADJ7:
        sub = d[(d["_adj"] == adj) & (pd.to_numeric(d["調整_萬"], errors="coerce") != 0)]
        if sub.empty:
            continue
        projs = sub.groupby(["ng_scope", "dicj code"]).agg(名稱=("project", "first"),
                             報告=("調整前_萬", "sum"), 調整=("調整_萬", "sum")).reset_index()
        projs = projs.reindex(projs["調整"].abs().sort_values(ascending=False).index)
        recs = []
        for _, p in projs.iterrows():
            nr = N.nlook(narr, p["ng_scope"], p["dicj code"])
            desc = llm_proj.get(proj_key(adj, p["ng_scope"], p["dicj code"]), "")
            if desc:      # LLM 寫嘅事項描述（ground 表2＋清單）＋ 管理層原話
                items = [("事項描述：", desc)]
                if nr.get("管理層解釋"):
                    items.append(("管理層解釋：", nr["管理層解釋"]))
            else:
                items = [(l + "：", t) for l, t in
                         [("KPMG分析發現", nr.get("KPMG分析發現", "")),
                          ("管理層解釋", nr.get("管理層解釋", "")),
                          ("跨司工作組／KPMG意見", nr.get("跨司回覆", "") or nr.get("KPMG回覆", ""))] if t]
                if not items and b2:      # 清單冇 → 用表2 抽到嘅原文頂住
                    t2 = B2.b2text(b2, p["ng_scope"], p["dicj code"])
                    if t2:
                        items = [("事項描述：", t2[:600])]
            if not items:
                items = [("", "清單未提供本項目之分析發現，待項目組補充。")]
            recs.append((f"{p['dicj code']}　{str(p['名稱'])[:34]}　│　報告 "
                         f"{R.fmt_money(p['報告'])}／潛在調整 {R.fmt_money(p['調整'])} 萬澳門元", items))
        tot = sub["調整_萬"].sum()
        head = (f"{ent_up} 報告的投資金額中，屬「{adj}」之潛在調減金額合計約{abs(tot):,.0f}萬澳門元，"
                f"涉及{len(recs)}個投資項目，逐項說明如下：")
        _cards(prs, 2, f"本年度審查工作的主要發現  |  {adj}", head, recs)


def render_site_visits(prs, ent_up, df, narr, threshold=2000):
    """附件二 現場走訪（slide 93-100）：capex ≥ 2,000萬 設施項目（報告走訪準則），
    配清單 實施地點 + 實際投資內容 做走訪概述。每頁 2 個項目 card。"""
    cap = df[df["final_capex_opex"] == "Capex"].copy()
    cap = cap[cap["dicj code"].astype(str).str.match(r"^項目\s*\d")]
    g = cap.groupby(["ng_scope", "dicj code"]).agg(
        名稱=("project", "first"), 報告=("調整前_萬", "sum")).reset_index()
    g = g[g["報告"] >= threshold]
    if g.empty:
        return
    g["_s"] = (g["ng_scope"] != "gaming").astype(int)
    g = g.sort_values(["_s", "報告"], ascending=[True, False])
    recs = []
    for _, p in g.iterrows():
        nr = N.nlook(narr, p["ng_scope"], p["dicj code"])
        items = [(l + "：", t) for l, t in
                 [("實施地點", nr.get("實施地點", "")),
                  ("現場走訪概述", nr.get("實際投資內容", "")),
                  ("現場走訪圖片", "〔待插入〕")] if t]
        recs.append((f"{p['dicj code']}　{str(p['名稱'])[:32]}　│　設施建設（資本性支出）"
                     f"{R.fmt_money(p['報告'])} 萬澳門元", items))
    head = (f"我們就{ent_up}報告投資金額中設施建設（資本性支出）達{threshold:,.0f}萬澳門元或以上之"
            f"{len(recs)}個投資項目進行了現場走訪，走訪情況如下：")
    _cards(prs, 5, "附件  |  部分項目的現場走訪情況", head, recs,
           note="資料來源：現場走訪記錄、管理層提供之項目資料，畢馬威分析")


def _prose_slide(prs, title, bullets, headline=None, *, sec=0):
    """一版敘述（crumb + navy 導語 + 段落），按估算高度自動分頁。"""
    W, H = L.size_of(prs)
    cw = W - 2 * L.MARGIN
    probe = L.HEAD_Y + L.head_h(headline, W)[0] + 0.10
    pages = L.fit_prose(bullets, cw, L.CONTENT_BOTTOM - probe, head_size=8, body_size=8)
    for pi, page in enumerate(pages):
        suffix = f"（{pi+1}/{len(pages)}）" if len(pages) > 1 else ""
        slide, W, H, top = _page(prs, sec, title, (headline or "") + suffix)
        L.prose_box(slide, L.MARGIN, top, cw, L.CONTENT_BOTTOM - top, page,
                    head_size=8, body_size=8, gap=7)
        L.source_note(slide, W, more=(pi < len(pages) - 1))


def _headline(ent_up, ov, df, plan):
    """slide 10 整體投資支出概況 → 回 (headline 句, bullets)。全 mechanical 由數字生成。"""
    tot = ov[ov["範疇"] == "總計"]
    if not len(tot):
        return "", []
    r = tot.iloc[0]

    def num(x):
        return float(x) if isinstance(x, (int, float)) and not pd.isna(x) else 0.0
    plan_amt = num(r.get("獲批的計劃投資金額"))
    report_amt = num(r.get("報告投資金額"))
    after_amt = num(r.get("潛在調整後投資金額"))
    adj_amt = report_amt - after_amt
    rate = r.get("投資計劃完成率")
    after_rate = r.get("潛在調整後投資計劃完成率")

    d = df.copy()
    d["_adj"] = pd.to_numeric(d["調整_萬"], errors="coerce").fillna(0)
    # ⚠ 三個數要自洽：計劃(母體) ⊇ 已實施；之前 n_impl 數 feed 全部有支出嘅碼，
    #   同 n_plan 唔同母體 → 出現「實施 112 > 計劃 89、0 個未發生」（項目組 2026-08-13 指出）
    n_plan, n_impl, n_zero = _proj_counts(df, plan, ov)
    n_adj = d[d["_adj"] != 0]["dicj code"].nunique()
    headline = (f"{ent_up} 2025年度原獲批計劃開展{n_plan}個投資項目，涉及計劃投資金額約{_amt(plan_amt)}；"
                f"{ent_up}提交的投資執行報告顯示2025年度投資金額約{_amt(report_amt)}"
                f"（投資計劃金額完成率{_pct(rate)}）。本次審查工作識別潛在調減金額約{_amt(adj_amt)}，"
                f"經潛在調整後的2025年度投資金額約{_amt(after_amt)}（經調整後投資計劃金額完成率{_pct(after_rate)}）。")
    bullets = [
        ("2025年度獲批的計劃投資金額與報告的投資金額：",
         f"根據{ent_up}提交的2025年度投資計劃方案與投資執行報告，{ent_up}獲批計劃開展{n_plan}個投資項目，"
         f"計劃投資金額約{_amt(plan_amt)}。投資執行報告顯示實際開展其中{n_impl}個投資項目"
         # 0 個未發生就唔好寫「計劃中有0個項目未發生投資金額」（項目組 2026-08-13 指出係錯）
         + (f"（計劃中有{n_zero}個項目未發生投資金額）" if n_zero else "")
         + f"，報告投資金額約{_amt(report_amt)}，報告投資計劃金額完成率為{_pct(rate)}。"),
        ("2025年度投資支出金額的潛在調整事項：",
         f"我們在本次審查工作中發現，{ent_up}報告投資金額中存在部分投資支出可能不應確認為2025年度計劃的投資支出"
         f"（涉及{n_adj}個投資項目，合計約{_amt(adj_amt)}）。考慮潛在調減事項後，{ent_up} 2025年度投資支出金額"
         f"約{_amt(after_amt)}，投資計劃金額完成率應為{_pct(after_rate)}。"),
    ]
    return headline, bullets


def _pct(v):
    return f"{v*100:.1f}%" if isinstance(v, (int, float)) and not pd.isna(v) else "—"


def _trim(s):
    """剝尾標點 —— 清單／表2 原文多數自帶「。」，接落我哋句式會變「。。」。"""
    return str(s or "").strip().rstrip("。．.；;，,、 ")


# ── 報告用字 helper（逐句對返 scan slide 10/11/15/19/23）────────────────────
_CN_NUM = "零一二三四五六七八九十"


def _cn(n):
    """1→一、7→七、12→十二（報告寫『存在七大類的調整事項』）。"""
    n = int(n)
    if n <= 10:
        return _CN_NUM[n]
    if n < 20:
        return "十" + (_CN_NUM[n - 10] if n > 10 else "")
    return f"{n}"


def _amt(wan):
    """萬 → 報告用字。scan：≥1億寫『6.4億澳門元』，<1億寫『5,527萬澳門元』。"""
    try:
        w = abs(float(wan or 0))
    except (TypeError, ValueError):
        return "—"
    return f"{w/10000:.1f}億澳門元" if w >= 10000 else f"{w:,.0f}萬澳門元"


def _cats_of(ov, col="報告投資金額", n=3, scope=None):
    """由 overview 表攞金額最大嘅幾個範疇名（scan：『主要涉及會議展覽、文化藝術、社區旅遊等…』）。"""
    if ov is None or ov.empty or col not in ov.columns:
        return ""
    d = ov[~ov["範疇"].astype(str).str.endswith(("小計", "總計", "項目"))].copy()
    d["_v"] = pd.to_numeric(d[col], errors="coerce").fillna(0)
    d = d[d["_v"] > 0].sort_values("_v", ascending=False)
    named = d[~d["範疇"].astype(str).isin(["其他", "其它"])]      # 「其他」唔好排喺點名範疇最前
    d = (named if len(named) >= min(n, 2) else d).head(n)
    return "、".join(str(r["範疇"]) for _, r in d.iterrows())


def _pg(i, n):
    """scan 導語尾有頁碼標記『（1/2）』。"""
    return f"（{i}/{n}）" if n > 1 else ""


def _zero_intro(ent_up, zi):
    """報告 1.2 頁尾段：未發生投資項目嘅大致介紹（跨年／內部研究／取消），並指去 1.3。"""
    if not zi:
        return None
    n, tot, groups = zi
    seg = []
    for kind, txt in (("跨年", "為跨年項目，{e}仍在就以前年度計劃項目進行持續投資，於2025年將投資額"
                              "作為2023年度計劃或2024年度計劃期後投資金額進行申報"),
                      ("內部研究", "個項目由於處於內部研究、計劃階段或未收到詳細指引未發生實際支出"),
                      ("取消", "個項目已取消")):
        g = groups.get(kind) or []
        if not g:
            continue
        seg.append((f"其中{len(g)}個" if kind == "跨年" else f"{len(g)}") + txt.format(e=ent_up))
    return ("未發生投資項目的大致介紹",
            f"{ent_up}於2025年度的投資計劃中有{n}個非博彩項目未產生投資金額，"
            + "；".join(seg) + "。請見後續「2025年度投資項目的整體執行概況」。")


def _bucket_headline(ent_up, bucket, ov):
    """②期後概覽導語（由表自己嘅總計行計，自洽）。"""
    tot = ov[ov["範疇"].astype(str).str.strip() == "總計"]
    if tot.empty:
        return ""
    r = tot.iloc[0]

    def n(c):
        v = r.get(c)
        return float(v) if isinstance(v, (int, float)) and not pd.isna(v) else 0.0
    rep, a, aft = n("報告投資金額"), n("潛在調整金額"), n("潛在調整後投資金額")
    yr, npj = bucket[:4], int(n("項目數量"))
    cats = _cats_of(ov, "報告投資金額", 3)
    tail = (f"，主要涉及{cats}等非博彩投資範疇的{npj}個項目" if cats else f"，涉及{npj}個投資項目")
    return (f"{ent_up}在2025年度執行報告中申報的「因發生期後事項需作後續調整之{yr}年度博彩／非博彩項目」"
            f"投資金額為{_amt(rep)}。本次審查工作識別潛在調減金額約{_amt(a)}，"
            f"經潛在調減後的{yr}年度計劃投資項目在2025年的投資金額約{_amt(aft)}{tail}。")


def _rate_of(df, name, col):
    r = df[df["範疇"] == name]
    if not len(r) or col not in df.columns:
        return None
    v = r.iloc[0][col]
    return v if isinstance(v, (int, float)) and not pd.isna(v) else None


def _prose_paginated(prs, title, bullets, per):
    if not bullets:
        return
    pages = [bullets[i:i + per] for i in range(0, len(bullets), per)]
    for pi, page in enumerate(pages):
        t = title + (f"（{pi+1}/{len(pages)}）" if len(pages) > 1 else "")
        _prose_slide(prs, t, page)


def _prose_2col(prs, title, bullets, per=12, subtitle=None, *, sec=0, headline=None):
    """報告式 2 欄敘述（對 scan slide 16-17：左右兩欄，每欄 navy 小標題 + body）。
    每頁裝幾多由【估算高度】決定，唔會爆版。"""
    if not bullets:
        return
    numbered = bool(bullets) and len(bullets[0]) == 3      # (no, head, body) = scan 編號清單
    W, H = L.size_of(prs)
    colw = (W - 2 * L.MARGIN - L.COL_GAP) / 2
    probe = L.HEAD_Y + L.head_h(headline, W)[0] + 0.10
    avail = L.CONTENT_BOTTOM - probe - (0.2 if subtitle else 0)
    if numbered:
        hs_all = [L.est_numbered_h([b], colw, size=L.SZ_BODY) for b in bullets]
        half_pages, cur, used = [], [], 0.0
        for b, hh in zip(bullets, hs_all):
            if cur and used + hh > avail * 2:
                half_pages.append(cur); cur, used = [], 0.0
            cur.append(b); used += hh
        if cur:
            half_pages.append(cur)
    else:
        half_pages = L.fit_prose(bullets, colw, avail * 2, head_size=L.SZ_BODY_HEAD, body_size=L.SZ_BODY)
    for pi, page in enumerate(half_pages):
        suffix = f"（{pi+1}/{len(half_pages)}）" if len(half_pages) > 1 else ""
        slide, W, H, top = _page(prs, sec, title, (headline or "") + suffix)
        if subtitle:
            L.put(slide, L.MARGIN, top, W - 2 * L.MARGIN, 0.18, subtitle, size=6.5,
                  italic=True, color=L.GREY)
            top += 0.20
        # 斷欄：以【總高一半】為目標令左右大致平均（對 scan），但唔可以超過一欄可用高
        lim = L.CONTENT_BOTTOM - top
        hs = [(L.est_numbered_h([it], colw, size=L.SZ_BODY) if numbered
               else L.est_prose_h([it], colw, head_size=L.SZ_BODY_HEAD, body_size=L.SZ_BODY)) for it in page]
        target = sum(hs) / 2.0
        cut, used = len(page), 0.0
        for i, ih in enumerate(hs):
            if i and (used >= target or used + ih > lim):
                cut = i; break
            used += ih
        cut = max(1, cut)
        if numbered:
            L.prose_numbered(L._tb(slide, L.MARGIN, top, colw, lim), page[:cut], size=L.SZ_BODY)
            if page[cut:]:
                L.prose_numbered(L._tb(slide, L.MARGIN + colw + L.COL_GAP, top, colw, lim),
                                 page[cut:], size=L.SZ_BODY)
        else:
            L.prose_box(slide, L.MARGIN, top, colw, lim, page[:cut], head_size=L.SZ_BODY_HEAD, body_size=L.SZ_BODY)
            if page[cut:]:
                L.prose_box(slide, L.MARGIN + colw + L.COL_GAP, top, colw, lim, page[cut:],
                            head_size=L.SZ_BODY_HEAD, body_size=L.SZ_BODY)
        L.source_note(slide, W, more=(pi < len(half_pages) - 1))


def render_category_overview(prs, ent_up, ov, df, narr, llm=None, ovx=None, note=None):
    """報告 slide 11-14（1.3 整體執行概況）：【左邊照舊擺 1.2 嗰個整體概況表】、右邊逐範疇敘述，
    敘述長就分版（scan 係 1/4 … 4/4，四版嘅表一模一樣）。LLM summary 優先，否則清單抄字。"""
    if not narr:
        return
    llm_cat = (llm or {}).get("cat", {})
    d = df.copy()
    d["_sub"] = FS.sub_of(d)
    proj = d.groupby(["_sub", "dicj code"])["調整前_萬"].sum().reset_index()
    cats = ov[~ov["範疇"].astype(str).str.endswith(("小計", "總計", "項目"))]
    g_bul, n_bul = [], []       # 按範疇概況：博彩 / 非博彩 各自一版（報告 3/4、4/4）
    for _, r in cats.iterrows():
        sub = str(r["範疇"]); rate = r.get("投資計劃完成率")
        if not isinstance(rate, (int, float)) or pd.isna(rate):
            continue
        if sub in llm_cat and llm_cat[sub]:               # LLM 寫嘅摘要優先
            txt = llm_cat[sub]
            txt = txt[len(sub) + 1:] if txt.startswith(sub + "：") else txt
            (g_bul if sub.startswith("博彩") else n_bul).append((f"{sub}：", txt)); continue
        scope = "gaming" if sub.startswith("博彩") else "non_gaming"
        pr = proj[proj["_sub"] == sub].sort_values("調整前_萬", ascending=False)
        content = reason = ""       # content=清單實際投資內容；reason=清單管理層解釋(變更原因)
        for _, pp in pr.iterrows():
            nr = N.nlook(narr, scope, pp["dicj code"])
            if not content:
                content = nr.get("實際投資內容", "")
            if not reason:
                reason = nr.get("管理層解釋", "")   # 業務原因；唔用 KPMG分析發現（審計腔）
            if content and reason:
                break
        content, reason = _trim(content), _trim(reason)   # 清單原文已有句號 → 唔剝就出「。。」
        summ = (content[:90] + "…") if len(content) > 90 else content
        rsn = ("，主要由於" + (reason[:80] + "…" if len(reason) > 80 else reason)) if reason else ""
        body = (f"主要包括{summ}。投資計劃金額完成率為{_pct(rate)}{rsn}。" if summ
                else f"投資計劃金額完成率為{_pct(rate)}{rsn}。")
        (g_bul if sub.startswith("博彩") else n_bul).append((f"{sub}：", body))
    # scan p-06 句式：著重於投入{博彩範疇}等博彩項目，以及{非博彩範疇}等非博彩投資項目 + 前後完成率
    gm = ov[ov["範疇"].astype(str).str.startswith("博彩") &
            ~ov["範疇"].astype(str).str.endswith(("小計", "項目"))]
    ng = ov[~ov["範疇"].astype(str).str.startswith("博彩") &
            ~ov["範疇"].astype(str).str.endswith(("小計", "總計", "項目"))]
    g_cats = _cats_of(gm, "報告投資金額", 2) or "博彩娛樂場場地的優化"
    n_cats = _cats_of(ng, "報告投資金額", 5)
    g_r, ng_r = _rate_of(ov, "博彩項目小計", "投資計劃完成率"), _rate_of(ov, "非博彩項目小計", "投資計劃完成率")
    g_a, ng_a = (_rate_of(ov, "博彩項目小計", "潛在調整後投資計劃完成率"),
                 _rate_of(ov, "非博彩項目小計", "潛在調整後投資計劃完成率"))
    # scan slide 11 尾句：「博彩項目各範疇…均超過100%，非博彩項目…平均完成率為44.2%，未有達到100%完成率的範疇」
    def _sub100(d):
        c = d[d["潛在調整後投資計劃完成率"].apply(
            lambda v: isinstance(v, (int, float)) and not pd.isna(v))] if \
            "潛在調整後投資計劃完成率" in d.columns else d.iloc[0:0]
        return c[c["潛在調整後投資計劃完成率"] < 1.0]
    g_all100 = "博彩項目各範疇的投資計劃金額完成率均超過100%" if len(_sub100(gm)) == 0 else \
               f"博彩項目的投資計劃金額完成率為{_pct(g_a)}"
    ng_lo = _sub100(ng)
    ng_tail = ("未有達到100%完成率的範疇" if len(ng_lo) == len(ng) and len(ng) else
               ("未達到100%完成率的範疇包括" +
                "、".join(f"{r['範疇']}（{_pct(r['潛在調整後投資計劃完成率'])}）"
                          for _, r in ng_lo.sort_values("潛在調整後投資計劃完成率").head(4).iterrows())
                if len(ng_lo) else "各範疇的投資計劃金額完成率均超過100%"))
    head = (f"{ent_up}的2025年度計劃投資項目著重於投入{g_cats}等博彩項目，以及{n_cats}等非博彩投資項目。"
            f"在報告投資金額中，博彩項目的投資計劃金額完成率為{_pct(g_r)}，"
            f"非博彩項目的投資計劃金額完成率為{_pct(ng_r)}。考慮投資金額的潛在調整後，"
            f"{g_all100}，非博彩項目的投資計劃金額平均完成率為{_pct(ng_a)}，{ng_tail}。")
    hi_s = "、".join(f"{r['範疇']}（{_pct(r['投資計劃完成率'])}）" for _, r in
                    cats[cats["投資計劃完成率"].apply(
                        lambda v: isinstance(v, (int, float)) and not pd.isna(v) and v >= 1.0)]
                    .sort_values("投資計劃完成率", ascending=False).iterrows())
    lo_s = "、".join(f"{r['範疇']}（{_pct(r['投資計劃完成率'])}）" for _, r in
                    cats[cats["投資計劃完成率"].apply(
                        lambda v: isinstance(v, (int, float)) and not pd.isna(v) and v < 0.5)]
                    .sort_values("投資計劃完成率").iterrows())
    exec_b = []               # 1/4：整體執行概況
    if hi_s:
        exec_b.append(("報告投資金額完成率較高的範疇", f"{hi_s}。"))
    if lo_s:
        exec_b.append(("報告投資金額完成率相對較低的範疇", f"{lo_s}。"))
    # ★ 報告 1.3 係【4 版、每版右邊唔同主題】（項目組 2026-08-17 指明），
    #   唔係逐範疇順住排落去分頁。左邊 4 版都係同一個 1.2 概況表。
    groups = [("2025年度投資項目的整體執行概況", exec_b),
              ("2025年度投資計劃區分設施建設/活動舉辦的投資金額", _fac_bullets(ent_up, ov)),
              ("按範疇的項目概況 — 博彩項目", g_bul),
              ("按範疇的項目概況 — 非博彩項目", n_bul)]
    render_overview_pages(prs, "2025年度投資計劃執行情況概述  |  2025年度投資項目的整體執行概況",
                          head, ovx if ovx is not None else ov, groups, sec=0, note=note,
                          table_name=f"{ent_up} 2025年度計劃的整體投資支出概況", grouped=True)


def _fac_bullets(ent_up, ov):
    """1.3 第 2 版：區分設施建設／活動舉辦嘅投資金額（由概況表自己嗰兩欄機械計）。"""
    F, A = "設施建設/資本性支出", "活動舉辦/營運性支出"
    if ov is None or ov.empty or F not in ov.columns:
        return []

    def num(row, c):
        v = row.get(c)
        return float(v) if isinstance(v, (int, float)) and not pd.isna(v) else 0.0

    def line(label):
        r = ov[ov["範疇"].astype(str).str.strip() == label]
        return (num(r.iloc[0], F), num(r.iloc[0], A)) if len(r) else (0.0, 0.0)
    tf, ta = line("總計")
    gf, ga = line("博彩項目小計")
    nf, na = line("非博彩項目小計")
    cat = ov[~ov["範疇"].astype(str).str.endswith(("小計", "總計", "項目"))].copy()
    if cat.empty or (tf + ta) == 0:
        return []
    cat["_f"] = cat.apply(lambda r: num(r, F), axis=1)
    cat["_a"] = cat.apply(lambda r: num(r, A), axis=1)
    top_f = "、".join(f"{r['範疇']}（{_amt(r['_f'])}）" for _, r in
                     cat.sort_values("_f", ascending=False).head(3).iterrows() if r["_f"] > 0)
    top_a = "、".join(f"{r['範疇']}（{_amt(r['_a'])}）" for _, r in
                     cat.sort_values("_a", ascending=False).head(3).iterrows() if r["_a"] > 0)
    return [
        ("設施建設／資本性支出",
         f"考慮潛在調整事項後，設施建設／資本性支出的投資金額約{_amt(tf)}"
         f"（佔{tf / (tf + ta) * 100:.1f}%），其中博彩項目約{_amt(gf)}、非博彩項目約{_amt(nf)}。"
         + (f"金額較大的範疇包括{top_f}。" if top_f else "")),
        ("活動舉辦／營運性支出",
         f"考慮潛在調整事項後，活動舉辦／營運性支出的投資金額約{_amt(ta)}"
         f"（佔{ta / (tf + ta) * 100:.1f}%），其中博彩項目約{_amt(ga)}、非博彩項目約{_amt(na)}。"
         + (f"金額較大的範疇包括{top_a}。" if top_a else "")),
    ]


def _adj_detail_bullets(ent_up, adj, df, narr, llm=None):
    """slide 16-17 潛在調整事項詳述 → 回 bullets：LLM summary 優先，否則清單分析發現。"""
    if not narr:
        return []
    llm_adj = (llm or {}).get("adj", {})
    pb = S.BUCKET_ORDER[0]      # 2025計劃 bucket（唔用合計；期後另計）
    d = df.copy()
    d["_adj"] = d["調整一級"].map(B.CANON).fillna(d["調整一級"])
    bullets = []
    for _, r in adj.iterrows():
        t = str(r["潛在調整事項"])
        if t in ("合計", "跨年及其他調整"):
            continue
        amt = r.get(pb, 0)
        if not isinstance(amt, (int, float)) or abs(amt) < 0.5:
            continue
        no = (B.ADJ7.index(t) + 1) if t in B.ADJ7 else len(B.ADJ7) + 1
        if t in llm_adj and llm_adj[t]:                   # LLM 寫嘅摘要優先
            bullets.append((no, f"{t}（{_amt(amt)}）：", llm_adj[t])); continue
        sub = d[(d["_adj"] == t) & (d["_bucket"] == pb) & (pd.to_numeric(d["調整_萬"], errors="coerce") != 0)]
        if sub.empty:
            continue
        names = "、".join(str(x) for x in sub.groupby("dicj code")["project"].first().tolist()[:3])
        reason = ruling = ""
        for _, pp in sub.drop_duplicates("dicj code").iterrows():
            nr = N.nlook(narr, pp["ng_scope"], pp["dicj code"])
            if not reason:
                reason = nr.get("KPMG分析發現", "") or nr.get("調整事項備註", "")
            if not ruling:
                ruling = nr.get("跨司回覆", "") or nr.get("KPMG回覆", "")
            if reason and ruling:
                break
        reason, ruling = _trim(reason), _trim(ruling)
        r2 = (reason[:150] + "…" if len(reason) > 150 else reason + "。") if reason else ""
        rl = ("跨司工作組／KPMG意見：" + (ruling[:90] + "…" if len(ruling) > 90 else ruling + "。")) if ruling else ""
        body = f"主要涉及{names}等項目。{r2}{rl}" if (r2 or rl) else f"主要涉及{names}等項目。"
        bullets.append((no, f"{t}（{_amt(amt)}）：", body))
    return bullets


def _exec_bullets(ent_up, ov):
    """整體執行概況敘述（報告 slide 11）：完成率高/低範疇，全部由 overview 數字生成 → 回 bullets。"""
    g = _rate_of(ov, "博彩項目小計", "投資計劃完成率")
    ng = _rate_of(ov, "非博彩項目小計", "投資計劃完成率")
    tot = _rate_of(ov, "總計", "投資計劃完成率")
    ga = _rate_of(ov, "博彩項目小計", "潛在調整後投資計劃完成率")
    nga = _rate_of(ov, "非博彩項目小計", "潛在調整後投資計劃完成率")
    cat = ov[~ov["範疇"].astype(str).str.endswith(("小計", "總計", "項目"))].copy()
    cat = cat[cat["投資計劃完成率"].apply(lambda v: isinstance(v, (int, float)) and not pd.isna(v))]
    high = cat[cat["投資計劃完成率"] >= 1.0].sort_values("投資計劃完成率", ascending=False)
    low = cat[cat["投資計劃完成率"] < 0.5].sort_values("投資計劃完成率")
    high_s = "、".join(f"{r['範疇']}（{_pct(r['投資計劃完成率'])}）" for _, r in high.iterrows())
    low_s = "、".join(f"{r['範疇']}（{_pct(r['投資計劃完成率'])}）" for _, r in low.iterrows())
    bullets = [
        ("", f"{ent_up}於2025年度計劃投資項目涵蓋博彩及非博彩範疇。在報告投資金額中，"
             f"博彩項目的投資計劃完成率為{_pct(g)}，非博彩項目為{_pct(ng)}，整體為{_pct(tot)}。"),
        ("", f"考慮投資金額的潛在調整後，博彩項目的投資計劃完成率為{_pct(ga)}，"
             f"非博彩項目的平均完成率為{_pct(nga)}。"),
    ]
    # ⚠ 高/低完成率範疇喺報告係喺【1.3 整體執行概況】嗰版，唔喺 1.2（項目組 2026-08-13）
    return bullets


def _adj_summary(ent_up, adj, ov=None, sdf=None):
    """潛在調整事項匯總（報告 slide 15）→ 回 (headline, bullets)。逐類型金額。
    ⚠ 用 2025計劃 bucket（唔係合計）：報告調整詳述只計 2025年度計劃，期後另有匯總。"""
    pb = S.BUCKET_ORDER[0]      # "2025年度投資計劃"
    tot_row = adj[adj["潛在調整事項"] == "合計"]
    total = tot_row.iloc[0][pb] if len(tot_row) else 0
    n_type = sum(1 for _, r in adj.iterrows()
                 if r["潛在調整事項"] not in ("合計", "跨年及其他調整")
                 and isinstance(r.get(pb), (int, float)) and abs(r.get(pb, 0)) >= 0.5)
    n_proj = after = None
    if sdf is not None:
        d = sdf[(sdf["_bucket"] == pb) & (pd.to_numeric(sdf["調整_萬"], errors="coerce") != 0)]
        n_proj = int(d["dicj code"].nunique())
    if ov is not None and not ov.empty:
        t = ov[ov["範疇"].astype(str).str.strip() == "總計"]
        if len(t):
            after = t.iloc[0].get("潛在調整後投資金額")
    # scan p-08 句式：存在七大類的調整事項，潛在調減約6.9億澳門元（涉及22個投資項目），經潛在調減後約11.8億澳門元
    headline = (f"基於我們的各項審查程序，我們認為{ent_up}報告的2025年度投資金額中，"
                f"存在{_cn(n_type)}大類的調整事項，潛在調減投資金額約{_amt(total)}"
                + (f"（涉及{n_proj}個投資項目）" if n_proj else "")
                + (f"，經潛在調減後的2025年度計劃的投資金額約{_amt(after)}。" if after is not None else "。"))
    bullets = []
    for _, r in adj.iterrows():
        name = r["潛在調整事項"]
        if name in ("合計", "跨年及其他調整"):
            continue
        amt = r[pb] if pb in adj.columns else 0
        if not isinstance(amt, (int, float)) or abs(amt) < 0.5:
            continue
        bullets.append((f"{name}：", f"約{abs(amt):,.0f}萬澳門元。"))
    return headline, bullets


def main():
    entity = (sys.argv[1] if len(sys.argv) > 1 and not sys.argv[1].startswith("-") else "mgm").lower()
    feed = Path(FEED)
    if not feed.exists():
        print(f"✗ 揾唔到 feed {feed}（root 應有 tableau_combined_25.csv）"); return
    qingdan = _find("data/投資項目清單", entity, ".xlsx")
    template = _find("data/reports", entity, ".pptx", prefer=["2025"])
    global ENT_UP
    ent_up = ENT_UP = entity.upper()
    print(f"build {BUILD_STAMP}")
    print(f"entity={ent_up}  feed={feed.name}  清單={qingdan.name if qingdan else '(冇)'}  "
          f"template={template.name if template else '(冇→用 13.33x7.5)'}")

    df = pd.read_csv(feed, low_memory=False)
    df = df[df["entity"].astype(str).str.lower() == entity]
    df["報告年"] = pd.to_numeric(df["報告年"], errors="coerce")
    df["_plan_year"] = df["year_bucket"].map(B._plan_year)
    plan = B.load_plan(qingdan) if qingdan else None
    cat = B.load_category(qingdan) if qingdan else None     # 項目性質(D)→派零投資項目計劃返範疇
    narr = N.load_narrative(qingdan) if qingdan else {}     # 清單 by-project narrative（抄字）
    if narr:
        print(f"    清單 narrative: {sum(1 for r in narr.values() if r.get('KPMG分析發現'))} 個項目有發現")
    _coverage_probe(df, str(qingdan) if qingdan else None)   # 探 ❓頁（主體/KPI/藝術品）coverage
    # 有 workbench creds 就自動即場生成 LLM 敘述（毋須任何 flag）；冇 creds 靜靜跳過用清單 fallback；--no-llm 強制跳過
    av = sys.argv
    if "--no-llm" not in av:
        model = av[av.index("--model") + 1] if "--model" in av else None
        try:
            has_creds = bool(Workbench(model=model).config_masked().get("key_ok"))
        except Exception:
            has_creds = False
        if has_creds:
            # default 4：8 條並行俾公司網關 WAF 全部擋（2026-08-15，60/60「request is blocked」）
            workers = int(av[av.index("--workers") + 1]) if "--workers" in av else 4
            biao2_dir = av[av.index("--biao2") + 1] if "--biao2" in av else "data/表2"
            print("  由 feed+清單+表2 即場生成 LLM 敘述…")
            try:
                generate_llm_narrative(str(feed), entity, str(qingdan) if qingdan else None,
                                       biao2_dir=biao2_dir, model=model, workers=workers)
            except Exception as e:
                print(f"  ⚠ LLM 生成失敗（{type(e).__name__}: {e}）→ 用現有 json / 清單 fallback")
    llm = _load_llm(entity)     # {entity}_llm_narrative.json 有就用 LLM 文字，否則清單 fallback
    if llm:
        print("    LLM narrative: " + "、".join(f"{k} {len(llm.get(k, {}))}"
                                              for k in ("adj", "cat", "tbl", "proj", "bkt")) + " 段")

    global USE_TEMPLATE
    tmpl = _find_template() if "--use-template" in sys.argv else None   # 預設 fresh 手砌（template 樣式已 hardcode）；--use-template 先開 template
    USE_TEMPLATE = bool(tmpl)
    if tmpl:
        prs = Presentation(str(tmpl))
        _strip_slides(prs)      # drop_rel 正確清走原有 content slides（唔會 duplicate 名 corrupt）
        print(f"    template（master-driven）：{tmpl}  layouts={sum(len(m.slide_layouts) for m in prs.slide_masters)}")
    else:
        prs = Presentation()
        if template:
            prs.slide_width, prs.slide_height = Presentation(str(template)).slide_width, Presentation(str(template)).slide_height
        else:
            prs.slide_width, prs.slide_height = Inches(L.SLIDE_W), Inches(L.SLIDE_H)
        print(f"    無 template → fallback 手砌 formatting（{prs.slide_width/914400:.2f}x"
              f"{prs.slide_height/914400:.2f}in）")

    sdf = S._load(feed, entity)     # 於2025發生 slice（概述 + 金額匯總 共用）

    render_cover(prs, entity)       # 封面（報告 p1）

    # ① 2025年度投資計劃執行情況概述（報告 slide 8-18）
    S1 = "2025年度投資計劃執行情況概述"
    divider(prs, S1, "1", [
        ("1.1  股權架構簡圖及發生投資支出的主體公司", ""),
        ("1.2  2025年度計劃的整體投資支出概況", ""),
        ("1.3  2025年度投資項目的整體執行概況", ""),
        ("1.4  2025年度投資計劃報告投資金額的潛在調整事項匯總", ""),
    ])
    budget = _load_budget(entity)
    ov = O.overview_by_bucket(sdf, "2025年度投資計劃", plan, cat)
    adj = O.adjustment_bridge(sdf)
    NOTE_RATE = ("註：投資計劃完成率 ＝ 報告投資金額 ／ 獲批的計劃投資金額；潛在調整後完成率 ＝ "
                 "潛在調整後投資金額 ／ 獲批的計劃投資金額。金額單位為萬澳門元。")
    if not ov.empty:      # slide 10-11：表左 + headline/執行敘述右（報告 2 欄式）
        zi = O.zero_investment_summary(sdf, plan, cat, narr, ent_up)
        hl, hlb = _headline(ent_up, ov, sdf, plan)
        exb = _exec_bullets(ent_up, ov)
        zintro = _zero_intro(ent_up, zi)
        ovx = _overview_extra(ov, plan, sdf, budget, ent_up).fillna("")
        render_overview_page(prs, f"{S1}  |  2025年度計劃的整體投資支出概況",
                             hl, ovx, hlb + exb + ([zintro] if zintro else []), sec=0,
                             table_name=f"{ent_up} 2025年度的整體投資支出概況", note=NOTE_RATE)
        # slide 11-14 逐範疇概況（LLM 優先）；表照 1.2 嗰個逐版重複，同 scan 一致
        render_category_overview(prs, ent_up, ov, sdf, narr, llm, ovx=ovx, note=NOTE_RATE)
        zit = O.zero_investment_text(zi, ent_up)
        if zit:      # 報告概述尾段：2025計劃申報投資為零嘅項目（跨年/內部研究/取消）
            _prose_slide(prs, f"{S1}  |  2025年度計劃申報投資支出為零的項目",
                         [("", x) for x in zit[1:]], headline=zit[0], sec=0)
    ahl, ab = _adj_summary(ent_up, adj, ov, sdf)   # slide 15：表左 + 匯總敘述右
    render_overview_page(prs, f"{S1}  |  2025年度投資計劃報告投資金額的潛在調整事項匯總",
                         ahl, adj.fillna(""), ab, sec=0,
                         table_name=f"{ent_up} 2025年度投資計劃報告投資金額的潛在調整事項匯總",
                         note="註：金額單位為萬澳門元；括號表示調減。")
    _prose_2col(prs, f"{S1}  |  2025年度投資計劃報告投資金額的潛在調整事項匯總（詳述）",
                _adj_detail_bullets(ent_up, adj, sdf, narr, llm), 6, sec=0,
                headline=ahl)   # slide 16-17 詳述（LLM 優先）

    # ② 過往年度投資計劃在2025年繼續執行的審查跟進（報告 slide 19-26）
    S2 = "過往年度投資計劃在2025年繼續執行的審查跟進"
    divider(prs, S2, "2", [      # 子項用返報告字眼（scan slide 18）
        ("2.1  2024年度投資計劃期後投資金額概覽", ""),
        ("2.2  2024年度投資計劃期後報告投資金額的潛在調整事項匯總", ""),
        ("2.3  2023年度投資計劃期後投資金額概覽", ""),
        ("2.4  2023年度投資計劃期後報告投資金額的潛在調整事項匯總", ""),
        ("2.5  截至2025年末投資金額概覽", ""),
    ])
    for bk in ["2024年度計劃期後投資", "2023年度計劃期後投資"]:
        ov = O.overview_by_bucket(sdf, bk, plan, cat)
        if not ov.empty:
            render_generic(prs, f"{ent_up} {bk}金額概覽", ov.fillna(""), sec=1,
                           crumb=f"{S2}  |  {bk}金額概覽",
                           headline=_bucket_headline(ent_up, bk, ov),
                           note="註：金額單位為萬澳門元；括號表示調減。",
                           llm=llm, tbl_id=tbl_key("期後概覽", bk))
            render_bucket_adjustment(prs, ent_up, bk, sdf, ov, narr, llm)   # 2.2 / 2.4
    render_cumulative(prs, ent_up, df, plan, cat)      # 2.5 截至2025年末投資金額概覽（scan slide 26）

    # ③ 本年度審查工作的主要發現（報告 slide 28-40）
    S3 = "本年度審查工作的主要發現"
    divider(prs, S3, "3")
    fs = O.finding_summary(sdf)
    if not fs.empty:
        render_generic(prs, f"{ent_up} 本年度審查工作的主要發現摘要", fs.fillna(""), sec=2,
                       crumb=f"{S3}  |  主要發現摘要",
                       headline=(f"本次審查工作就{ent_up}報告的投資金額識別出{len(fs)}類潛在調整事項，"
                                 f"合計潛在調減約{abs(pd.to_numeric(fs['調整額合計'], errors='coerce').sum()):,.0f}"
                                 f"萬澳門元，摘要如下；逐項說明見後頁。"),
                       note="註：金額單位為萬澳門元；括號表示調減。",
                       llm=llm, tbl_id=tbl_key("發現摘要"))
    if narr:      # 逐調整類型 × 項目：金額(feed) + 事項描述(LLM ground 表2＋清單) / 清單抄字
        b2 = {}
        try:            # 表2＝審查底稿，清單冇料時頂住（加密檔，開唔到就靜靜跳過）
            b2 = B2.load_biao2_struct(av[av.index("--biao2") + 1] if "--biao2" in av else "data/表2",
                               entity, log=lambda *a: None)
        except Exception:
            pass
        render_findings(prs, ent_up, sdf, narr, llm=llm, b2=b2)

    # ④ 其他信息（報告 slide 42-63）
    S4 = "其他信息"
    divider(prs, S4, "4")
    render_generic(prs, f"{ent_up} 2025年發生的投資金額匯總",
                   S.summary_amount(sdf).fillna(""), sec=3,
                   crumb=f"{S4}  |  2025年發生的投資金額匯總",
                   headline=(f"下表匯總{ent_up} 2025年度投資計劃及過往年度計劃期後投資"
                             f"於2025年發生的投資金額（報告投資金額及潛在調整後投資金額）。"),
                   note="註：金額單位為萬澳門元。", llm=llm, tbl_id=tbl_key("金額匯總"))
    for bk in S.BUCKET_ORDER:
        fa = S.facility_activity(sdf, bk)
        if not fa.empty:
            render_generic(prs, f"{ent_up} {bk}區分設施建設/活動舉辦的投資金額", fa.fillna(""), sec=3,
                           crumb=f"{S4}  |  2025年發生的投資金額區分設施建設/活動舉辦",
                           headline=(f"下表按範疇列示{ent_up} {bk}於2025年發生的投資金額，"
                                     f"區分設施建設（資本性支出）及活動舉辦（營運性支出）。"),
                           note="註：金額為潛在調整後金額，單位為萬澳門元。",
                           llm=llm, tbl_id=tbl_key("設施活動", bk))
    for yr in (25, 24, 23):     # 單個項目審查匯總（slide 46-63）
        tab, _ = B.build_year(df, yr, plan.get(yr) if plan else None)
        if tab is not None and not tab.empty:
            R.render_sheet(prs, f"報告年{yr}", tab.fillna(""), list(tab.columns),
                           ent_up=ent_up, sec=3, crumb=f"{S4}  |  單個項目審查結果匯總")

    render_visit_summary(prs, ent_up, sdf)      # 報告 slide 71 走訪情況匯總（樣本標準+樣本量）

    # ⑥ 附件（slide 93-105）
    divider(prs, "附件", "6")
    if narr:
        render_site_visits(prs, ent_up, sdf, narr)
    render_artwork(prs, ent_up, av[av.index("--biao2") + 1] if "--biao2" in av else "data/表2",
                   entity)                      # 報告 slide 101 藝術品展出情況清單

    # 目錄（報告 slide 7）：起完全部版先知頁碼 → 砌好插去第 2 版，再全份重編頁碼
    _W, _H = L.size_of(prs)
    toc = _collect_toc(prs, _W, _H)
    if toc:
        n0 = len(prs.slides._sldIdLst)
        render_toc(prs, ent_up, toc)
        for k in range(len(prs.slides._sldIdLst) - n0):        # 目錄可能多過一版
            _move_slide(prs, n0 + k, 1 + k)
        _renumber_footers(prs, _W, _H)
        print(f"    目錄：{sum(1 for e in toc if not e[2])} 章 / {sum(1 for e in toc if e[2])} 子項")
    L.wire_nav(prs, _sec_slides(prs, _W, _H), home=1 if toc else 0)   # ◀⌂▶ + 頁籤內部跳頁

    if tmpl:      # template mode：重編 slide 高號，徹底避開 template 殘留 orphan part 撞名 corruption
        _renumber_slides(prs)

    L.apply_theme_fonts(prs)      # deck theme 字體 → KPMG（唔明寫嘅地方唔會跌返 Calibri）
    out = Path(f"{entity}_report_llm.pptx")
    try:
        prs.save(out)
    except PermissionError:      # 舊檔喺 PowerPoint 開住鎖住 → 改名唔 crash
        import time
        out = Path(f"{entity}_report_llm_{time.strftime('%H%M%S')}.pptx")
        prs.save(out)
        print(f"⚠ 原檔開住(鎖住)，改存 → {out.name}（開之前記得閂舊 pptx）")
    print(f"✓ {out.resolve()}  共 {len(list(prs.slides))} 頁（概述 + 主要發現 + 金額匯總 + 設施 + 單項審查）")
    if "--dump" in sys.argv:      # 要 cross-check 先加 --dump（慳空間；ok 咗嘅唔使 dump）
        dump = _dump_pptx_text(prs, entity, with_tables="--dump-tables" in sys.argv)
        print(f"✓ text dump → {dump.name}（逐版文字，cross-check 用）")


if __name__ == "__main__":
    main()
