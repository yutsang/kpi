#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
inspect_pptx.py — 報告 pptx 版面體檢（唔使開 PowerPoint 逐版睇）。逐版報：

  OVERFLOW-X/Y   shape 超出 slide 邊界（右／底）——「超出 border」直接捉到
  TABLE-GROW     表格估算實際高度 > requested（PowerPoint 會自動長高 row）→ 會爆版
  TEXT-OVERFLOW  文字框估算 wrap 高度 > 框高
  NO-TEXT        淨得表 / 圖，冇任何敘述文字（報告唔應該有淨表冇字嘅版）
  EMPTY          乜都冇
  OFF-PALETTE    用咗 KPMG 品牌色以外嘅顏色（IMG_0420 Visual identity overview）
  NO-FURNITURE   冇 breadcrumb / footer / 頁碼

用法：
    python scripts\\report\\inspect_pptx.py mgm_report_llm.pptx            # 體檢
    python scripts\\report\\inspect_pptx.py mgm_report_llm.pptx --slide 12 # 淨睇某版 shape 清單
    python scripts\\report\\inspect_pptx.py mgm_report_llm.pptx --dump [--batch 12]
                                            # 逐版文字（唔使重 build）；出檔 + console 印方便 copy
    python scripts\\report\\inspect_pptx.py mgm_report_llm.pptx --render   # 用 PowerPoint 出 PDF/PNG（Mac）
    python scripts\\report\\inspect_pptx.py "data\\reports\\MGM…初稿.pptx" --spec  # 真報告嘅尺寸/字體/配色
    python scripts\\report\\inspect_pptx.py "data\\reports\\MGM…初稿.pptx" --fonts --range 10-63
                                                          # 真報告【逐個位置實際用幾多 pt】
"""
import re
import subprocess
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("✗ pip install python-pptx"); sys.exit(1)

import layout as L

EMU_IN = 914400.0
PALETTE = {  # 公司 template theme 實際配色（--spec 讀出，2026-08-12）＋ 派生 tint
    "00338D": "dk2/accent2 KPMG Blue", "0C233C": "accent3 深底", "1E49E2": "accent1",
    "00B8F5": "accent4", "7213EA": "accent5", "FD349C": "accent6", "E5E5E5": "lt2",
    "ACEAFF": "master3 accent1", "00C0AE": "master3 accent4",
    "005EB8": "Medium Blue", "0091DA": "Light Blue",
    "483698": "Violet", "470A68": "Purple", "6D2077": "Light Purple", "00A3A1": "Green",
    "FFFFFF": "White", "EEF1F8": "tint-section", "D9E1F2": "tint-subtotal",
    "BDD7EE": "tint-total", "BFBFBF": "格線", "000000": "小計/總計橫線", "1E5C46": "表頭綠（設施/活動欄組）", "808080": "欄組虛線", "F2F2F2": "breadcrumb banner（lt2 派生）", "222222": "內文", "333333": "內文",
    "595959": "註", "8C8C8C": "breadcrumb", "C8C8C8": "分隔", "0C233C": "封面深底",
    "C8C8D0": "分隔頁副題",
}


def _in(v):
    return (v or 0) / EMU_IN


def _shape_text(sh):
    try:
        return sh.text_frame.text if sh.has_text_frame else ""
    except Exception:
        return ""


def _colors(sh, out):
    """收集 shape 用到嘅 srgbClr（fill + font）。"""
    for el in sh._element.iter():
        if el.tag.endswith("}srgbClr"):
            v = el.get("val")
            if v:
                out.add(v.upper())


def _table_h(tbl):
    """(requested_h, 估算實際 h)：PowerPoint 會按內容長高 row。用返 layout.row_h 同一把尺，
    唔會同 build 時嘅分頁計算行開（否則成堆假警報）。"""
    req = sum(_in(r.height) for r in tbl.rows)
    widths = [_in(c.width) for c in tbl.columns]
    est = 0.0
    for r in tbl.rows:
        sizes = [run.font.size.pt for cell in r.cells for p in cell.text_frame.paragraphs
                 for run in p.runs if run.font.size]
        size = max(sizes) if sizes else 10.0
        txts, ws = [], []
        for ci, cell in enumerate(r.cells):
            if getattr(cell, "is_spanned", False):       # 被合併掉嘅格：唔重複計
                continue
            n = cell.span_width if getattr(cell, "is_merge_origin", False) else 1
            txts.append(cell.text)
            ws.append(sum(widths[ci:ci + n]))            # 合併格：加埋跨住嘅欄闊
        need = L.row_h(txts, ws, size)
        est += max(_in(r.height), need)
    return req, est


def _text_h(sh):
    """(框高, 估算 wrap 高)。"""
    box_h = _in(sh.height)
    w = _in(sh.width)
    need = 0.0
    for p in sh.text_frame.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            continue
        size = max((r.font.size.pt for r in runs if r.font.size), default=10.0)
        txt = "".join(r.text for r in runs)
        need += L.est_lines(txt, w, size) * size * 1.3 / 72.0
        need += (p.space_after.pt if p.space_after else 0) / 72.0
        need += (p.space_before.pt if p.space_before else 0) / 72.0
    return box_h, need


# 機械 fallback 敘述嘅固定小標題（用嚟分辨「LLM 寫」定「冇 LLM 頂上」）
_MECH_HEADS = ("整體情況", "博彩／非博彩分佈", "金額最大的範疇")


def _slide_kind(slide, W, H):
    full = any(_in(s.width) > W - 0.1 and _in(s.height) > H - 0.1 for s in slide.shapes)
    n_tbl = sum(1 for s in slide.shapes if s.has_table)
    bars = sum(1 for s in slide.shapes if not s.has_table and s.has_text_frame
               and _in(s.width) > W * 0.6 and _in(s.height) < 0.3
               and s.text_frame.text.strip() and _fill_hex(s) == "00338D")
    if full:
        return "封面/分隔"
    if n_tbl and _has_side_prose(slide, W):
        return "表+敘述"
    if n_tbl:
        return "表"
    if bars >= 1:
        return "card"
    return "敘述"


def _fill_hex(sh):
    try:
        c = sh.fill.fore_color.rgb
        return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"
    except Exception:
        return ""


def _has_side_prose(slide, W):
    """表右邊有冇敘述框（＝報告式 2 欄）。"""
    tbl_r = max((_in(s.left) + _in(s.width) for s in slide.shapes if s.has_table), default=0)
    for s in slide.shapes:
        if s.has_table or not s.has_text_frame:
            continue
        if _in(s.left) >= tbl_r - 0.1 and _in(s.top) < 3.0 and \
                len(s.text_frame.text.strip()) > 60:
            return True
    return False


def _crumb_of(slide):
    for s in slide.shapes:
        if s.has_text_frame and 0.28 < _in(s.top) < 0.46 and "  |  " in s.text_frame.text:
            return s.text_frame.text.strip()
    return ""


def _prose_stats(slide, W):
    """(段數, 字數, 係咪機械 fallback)。skip breadcrumb/footer/表。"""
    n = chars = 0
    mech = False
    for s in slide.shapes:
        if s.has_table or not s.has_text_frame:
            continue
        t = s.text_frame.text.strip()
        if not t or _in(s.top) < 0.28 or _in(s.top) > 7.0:
            continue
        if t.startswith("©") or t.startswith("初稿") or t.startswith("KPMG"):
            continue
        n += len([p for p in s.text_frame.paragraphs if p.text.strip()])
        chars += len(t)
        if any(t.startswith(h) or f"\n{h}" in t for h in _MECH_HEADS):
            mech = True
    return n, chars, mech


def _empty_cols(tbl):
    """整欄冇數（全部 '-' 或空）嘅欄名 → 報告出咗一堆得個殼嘅欄。"""
    out, ncol = [], len(tbl.columns)
    hdr_rows = min(2, len(tbl.rows))
    for ci in range(1, ncol):
        vals = [tbl.cell(ri, ci).text.strip() for ri in range(hdr_rows, len(tbl.rows))]
        if vals and all(v in ("", "-") for v in vals):
            name = " ".join(tbl.cell(ri, ci).text.strip() for ri in range(hdr_rows)).strip()
            out.append(re.sub(r"[\n\v\x0b]+", "", name) or f"col{ci}")
    return out


def audit(path, tol=0.02):
    prs = Presentation(str(path))
    W, H = _in(prs.slide_width), _in(prs.slide_height)
    n_slides = len(prs.slides._sldIdLst)
    print(f"── {Path(path).name}：{n_slides} 版，{W:.2f} x {H:.2f} in")
    if abs(W - L.SLIDE_W) > 0.05 or abs(H - L.SLIDE_H) > 0.05:
        print(f"   ⚠ slide 尺寸 ≠ 報告標準 {L.SLIDE_W} x {L.SLIDE_H} in")
    inv, warn, stat = [], [], {"tbl": 0, "side": 0, "mech": 0, "chars": 0, "para": 0,
                               "cells": 0, "todo": 0}
    seen_text = {}
    for i, slide in enumerate(prs.slides, 1):
        kind = _slide_kind(slide, W, H)
        crumb = _crumb_of(slide)
        npara, nchar, mech = _prose_stats(slide, W)
        dims, empties = [], []
        for sh in slide.shapes:
            if sh.has_table:
                t = sh.table
                dims.append(f"{len(t.rows)}x{len(t.columns)}")
                stat["tbl"] += 1
                stat["cells"] += len(t.rows) * len(t.columns)
                empties += _empty_cols(t)
        stat["para"] += npara; stat["chars"] += nchar
        if kind == "表+敘述":
            stat["side"] += 1
            if mech:
                stat["mech"] += 1
        inv.append((i, kind, crumb[:52], ",".join(dims), npara, nchar,
                    "機械" if mech else ("LLM/清單" if npara else "")))
        if dims and nchar < 60:
            warn.append(f"slide {i}：有表但敘述只有 {nchar} 字（報告唔應該淨表冇字）")
        if len(empties) >= 2:      # 1 欄空好平常（該頁啱啱冇嗰類調整）；≥2 先值得望一望
            warn.append(f"slide {i}：本頁有 {len(empties)} 個調整欄整欄無數 → "
                        f"{'、'.join(empties[:4])}（該頁項目冇呢幾類調整，正常；淨係提你留意）")
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text
                stat["todo"] += t.count("〔待插入〕")
                if _in(sh.top) < 1.75 or _in(sh.top) > 7.0:
                    continue                      # breadcrumb / 導語 / footer：跨版重複係設計（同 scan）
                for para in t.split("\n"):
                    q = para.strip()
                    if len(q) > 40:
                        seen_text.setdefault(q, []).append(i)

    print("\n── 內容清單")
    print(f"{'#':>3}  {'類型':<8} {'版面':<52} {'表':<10} {'段':>3} {'字':>6}  來源")
    for i, kind, crumb, dims, npara, nchar, src in inv:
        print(f"{i:>3}  {kind:<8} {crumb:<52} {dims:<10} {npara:>3} {nchar:>6}  {src}")

    print("\n── 覆蓋")
    print(f"  表格 {stat['tbl']} 張（{stat['cells']:,} 格）；敘述 {stat['para']} 段共 {stat['chars']:,} 字")
    print(f"  表旁敘述 {stat['side']} 版（其中 {stat['mech']} 版係機械 fallback"
          f"{'／建議跑 LLM' if stat['mech'] else '，全部 LLM/清單'}）")
    if stat["todo"]:
        print(f"  〔待插入〕placeholder {stat['todo']} 個（現場走訪相等）")

    dup = {t: v for t, v in seen_text.items() if len(v) >= 3}
    if dup:
        warn.append(f"有 {len(dup)} 段文字喺 ≥3 版重複出現（例：{list(dup)[0][:40]}…）")

    bad, palette_hits, geo = 0, {}, []
    for i, slide in enumerate(prs.slides, 1):
        issues, has_table, has_prose, has_foot = [], False, False, False
        for sh in slide.shapes:
            x, y, w, h = _in(sh.left), _in(sh.top), _in(sh.width), _in(sh.height)
            name = sh.shape_type
            if x + w > W + tol:
                issues.append(f"OVERFLOW-X  {name} 右邊到 {x+w:.2f}in > {W:.2f}")
            if y + h > H + tol:
                issues.append(f"OVERFLOW-Y  {name} 底到 {y+h:.2f}in > {H:.2f}")
            if sh.has_table:
                has_table = True
                req, est = _table_h(sh.table)
                if est > req + 0.05:
                    issues.append(f"TABLE-GROW  表 requested {req:.2f}in → 估實際 {est:.2f}in")
                if y + est > H + tol:
                    issues.append(f"TABLE-OVERFLOW 表底到 {y+est:.2f}in > {H:.2f}（會被切）")
            txt = _shape_text(sh)
            if txt.strip():
                if sh.has_text_frame and not sh.has_table:
                    bh, need = _text_h(sh)
                    if need > bh + 0.08 and y + need > H:
                        issues.append(f"TEXT-OVERFLOW 文字框 {bh:.2f}in → 需 {need:.2f}in（超版底）")
                if "初稿" in txt:
                    has_foot = True
                if len(txt.strip()) > 40 and "初稿" not in txt and "©" not in txt:
                    has_prose = True
            _colors(sh, palette_hits.setdefault(i, set()))
        n_sh = len(slide.shapes)
        # 封面 / 章節分隔＝滿版深色底 → 冇 footer 係正常（跟 scan）
        cover = any(_in(s.width) > W - 0.1 and _in(s.height) > H - 0.1 for s in slide.shapes)
        if n_sh == 0:
            issues.append("EMPTY  冇任何 shape")
        elif has_table and not has_prose:
            issues.append("NO-TEXT  淨得表冇敘述文字")
        if not has_foot and n_sh > 3 and not cover:
            issues.append("NO-FURNITURE  冇 footer/頁碼")
        off = {c for c in palette_hits.get(i, set()) if c not in PALETTE}
        if off:
            issues.append("OFF-PALETTE  " + ", ".join("#" + c for c in sorted(off)))
        if issues:
            bad += 1
            geo.append((i, n_sh, issues))
    print(f"\n── 版面體檢：{bad} / {n_slides} 版有問題")
    for i, n_sh, issues in geo:
        print(f"  [slide {i}] {n_sh} shapes")
        for m in issues:
            print(f"      ✗ {m}")

    print(f"\n── 內容提示：{len(warn)} 項")
    for m in warn:
        print(f"  ! {m}")
    return bad + len(warn)


def dump_slide(path, n):
    prs = Presentation(str(path))
    slide = list(prs.slides)[n - 1]
    W, H = _in(prs.slide_width), _in(prs.slide_height)
    print(f"── slide {n}（{W:.2f}x{H:.2f}in）")
    for sh in slide.shapes:
        x, y, w, h = _in(sh.left), _in(sh.top), _in(sh.width), _in(sh.height)
        kind = "TABLE" if sh.has_table else ("TEXT" if sh.has_text_frame else str(sh.shape_type))
        extra = ""
        if sh.has_table:
            req, est = _table_h(sh.table)
            extra = f"  {len(sh.table.rows)}x{len(sh.table.columns)} req {req:.2f} est {est:.2f}in"
        print(f"  {kind:7s} x={x:5.2f} y={y:5.2f} w={w:5.2f} h={h:5.2f}{extra}")
        t = _shape_text(sh).replace("\n", " ⏎ ")
        if t.strip():
            print(f"          「{t[:150]}」")


def _osa(script, timeout=600):
    return subprocess.run(["osascript", "-e", script], capture_output=True, text=True,
                          timeout=timeout)


def render(path, dpi=80):
    """用 Microsoft PowerPoint（Mac）出 PDF，再 pdftoppm 出 PNG → 睇真實版面（含 table auto-grow）。

    ⚠ PowerPoint 嘅 AppleScript `open` 只喺【冷啟動】先穩陣：app 已經開住時會 -9074。
      所以每次都先 quit 再 open；第一個 AppleEvent 會 -1712 timeout（其實開緊）→ 忽略，改為 poll。"""
    import time
    src = Path(path).resolve()
    out = src.parent / (src.stem + "_render")
    out.mkdir(exist_ok=True)
    pdf = out / (src.stem + ".pdf")
    if pdf.exists():
        pdf.unlink()
    _osa('tell application "Microsoft PowerPoint" to quit saving no', timeout=120)
    time.sleep(5)
    _osa(f'tell application "Microsoft PowerPoint" to open POSIX file "{src}"', timeout=600)
    for _ in range(60):          # poll：cold launch 慢，AppleEvent timeout 唔代表失敗
        r = _osa('tell application "Microsoft PowerPoint" to get count of presentations', 120)
        if r.stdout.strip().isdigit() and int(r.stdout.strip()) > 0:
            break
        time.sleep(2)
    else:
        print("✗ PowerPoint 開唔到檔"); return
    r = _osa(f'''with timeout of 540 seconds
tell application "Microsoft PowerPoint"
    save presentation 1 in POSIX file "{pdf}" as save as PDF
    close presentation 1 saving no
end tell
end timeout''', timeout=600)
    if not pdf.exists():
        print("✗ PowerPoint 出 PDF 失敗：", (r.stderr or "").strip()[:300]); return
    subprocess.run(["pdftoppm", "-png", "-r", str(dpi), str(pdf), str(out / "s")], check=False)
    print(f"✓ {pdf}\n✓ PNG → {out}/s-*.png")


# ── 自畫 preview（唔使 PowerPoint / LibreOffice）───────────────────────────
FONT_DIRS = ["/System/Library/Fonts/Supplemental", "/System/Library/Fonts",
             "C:/Windows/Fonts", str(Path.home() / "Library/Fonts")]
FONT_FILES = {  # (cjk?, bold?) → 候選檔名
    (False, False): ["Arial.ttf", "arial.ttf", "Helvetica.ttc"],
    (False, True): ["Arial Bold.ttf", "arialbd.ttf", "Arial.ttf"],
    (True, False): ["STHeiti Light.ttc", "msyh.ttc", "PingFang.ttc", "Arial Unicode.ttf"],
    (True, True): ["STHeiti Medium.ttc", "msyhbd.ttc", "PingFang.ttc", "Arial Unicode.ttf"],
}
_fc = {}


def _font(cjk, bold, px):
    from PIL import ImageFont
    key = (cjk, bold, px)
    if key in _fc:
        return _fc[key]
    for name in FONT_FILES[(cjk, bold)]:
        for d in FONT_DIRS:
            p = Path(d) / name
            if p.exists():
                try:
                    _fc[key] = ImageFont.truetype(str(p), px); return _fc[key]
                except Exception:
                    pass
    _fc[key] = ImageFont.load_default()
    return _fc[key]


def _runs(tf):
    """text_frame → [(para_align, [(text, size_pt, bold, rgb)])]。
    ⚠ 要行返 <a:br/>（硬換行）；python-pptx 嘅 p.runs 會 skip 咗佢哋。"""
    out = []
    for p in tf.paragraphs:
        rs, last = [], (12.0, False, (0, 0, 0))
        for el in p._p:
            tag = el.tag.split("}")[-1]
            if tag == "br":
                rs.append(("\n",) + last); continue
            if tag != "r":
                continue
            r = next((x for x in p.runs if x._r is el), None)
            if r is None or not r.text:
                continue
            try:
                col = r.font.color.rgb
                rgb = (col[0], col[1], col[2]) if col else (0, 0, 0)
            except Exception:
                rgb = (0, 0, 0)
            last = (r.font.size.pt if r.font.size else 12.0, bool(r.font.bold), rgb)
            rs.append((r.text,) + last)
        out.append((str(p.alignment), rs))
    return out


_ATOM = re.compile(r"[0-9A-Za-z]+(?:[.,%/\-][0-9A-Za-z]+)*%?|[\s\S]")
_NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _cell_edges(dr, cell, x0, y0, x1, y1):
    """照 tcPr 入面真係有嘅 a:lnT/B/L/R 畫線（報告嘅表冇逐格格線，唔可以照畫框）。"""
    tcPr = cell._tc.find(f"{_NS_A}tcPr")
    if tcPr is None:
        return
    for side, pts in (("T", (x0, y0, x1, y0)), ("B", (x0, y1, x1, y1)),
                      ("L", (x0, y0, x0, y1)), ("R", (x1, y0, x1, y1))):
        ln = tcPr.find(f"{_NS_A}ln{side}")
        if ln is None:
            continue
        clr = ln.find(f".//{_NS_A}srgbClr")
        v = (clr.get("val") if clr is not None else "000000") or "000000"
        rgb = tuple(int(v[i:i + 2], 16) for i in (0, 2, 4))
        if ln.find(f"{_NS_A}prstDash") is not None:      # 虛線：畫短 dash
            n = int(max(abs(pts[2] - pts[0]), abs(pts[3] - pts[1])) // 6) or 1
            for k in range(0, n, 2):
                f0, f1 = k / n, min((k + 1) / n, 1.0)
                dr.line([pts[0] + (pts[2] - pts[0]) * f0, pts[1] + (pts[3] - pts[1]) * f0,
                         pts[0] + (pts[2] - pts[0]) * f1, pts[1] + (pts[3] - pts[1]) * f1],
                        fill=rgb, width=1)
        else:
            dr.line(list(pts), fill=rgb, width=1)


def _draw_text(dr, runs, x, y, w, scale, *, align="LEFT", valign_h=None):
    """畫 wrap 文字（逐 run 接住排，中英文分開量度）。回結束 y。"""
    for _al, rs in runs:
        if not rs:
            y += 6 * scale / 72.0 * 72
            continue
        size = max(r[1] for r in rs)
        lh = size * 1.3 * scale
        line, lw = [], 0.0
        maxw = w * 72 * scale
        for text, sz, bold, rgb in rs:
            # ★ 逐「原子」排：中文一字一原子，但英數/百分比（144.5%）當一整嚿唔可以拆
            #   （之前逐字拆 → preview 見到「100\n%」、「6\n2.2%」，PowerPoint 其實唔會咁斷）
            for atom in _ATOM.findall(text):
                if atom == "\t":                     # hanging indent tab
                    tgt = 0.24 * 72 * scale
                    if lw < tgt:
                        line.append((" ", _font(False, False, max(6, int(round(sz * scale)))), rgb))
                        lw = tgt
                    continue
                if atom == "\n":                     # <a:br/> 硬換行
                    _flush(dr, line, x, y, maxw, scale, align)
                    y += lh; line, lw = [], 0.0
                    continue
                f = _font(L._is_cn(atom[0]), bold, max(6, int(round(sz * scale))))
                cw = f.getlength(atom)
                if lw + cw > maxw and line:
                    _flush(dr, line, x, y, maxw, scale, align)
                    y += lh; line, lw = [], 0.0
                for ch in atom:
                    line.append((ch, f, rgb)); lw += f.getlength(ch)
        if line:
            _flush(dr, line, x, y, maxw, scale, align)
            y += lh
    return y


def _flush(dr, line, x, y, maxw, scale, align):
    lw = sum(f.getlength(c) for c, f, _ in line)
    cx = x + (maxw - lw if align.startswith("RIGHT") else
              (maxw - lw) / 2 if align.startswith("CENTER") else 0)
    for c, f, rgb in line:
        dr.text((cx, y), c, font=f, fill=rgb)
        cx += f.getlength(c)


def preview(path, dpi=110, only=None):
    """自畫每版 PNG（唔使 Office）：rect fill、table（含估算 row 長高）、文字 run。
    ★ 表格用同 layout.row_h 一樣嘅估算 → 睇到嘅就係 audit 檢查嘅嘢。"""
    from PIL import Image, ImageDraw
    prs = Presentation(str(path))
    W, H = _in(prs.slide_width), _in(prs.slide_height)
    out = Path(path).resolve().parent / (Path(path).stem + "_preview")
    out.mkdir(exist_ok=True)
    sc = dpi / 72.0                                  # pt → px
    for i, slide in enumerate(prs.slides, 1):
        if only and i not in only:
            continue
        img = Image.new("RGB", (int(W * dpi), int(H * dpi)), "white")
        dr = ImageDraw.Draw(img)
        for sh in slide.shapes:
            x, y, w, h = _in(sh.left) * dpi, _in(sh.top) * dpi, _in(sh.width) * dpi, _in(sh.height) * dpi
            if sh.has_table:
                tbl = sh.table
                widths = [_in(c.width) for c in tbl.columns]
                cy = y
                for r in tbl.rows:
                    sizes = [run.font.size.pt for cell in r.cells
                             for p in cell.text_frame.paragraphs for run in p.runs if run.font.size]
                    size = max(sizes) if sizes else 10.0
                    txts, ws = [], []
                    for ci, cell in enumerate(r.cells):
                        if getattr(cell, "is_spanned", False):
                            continue
                        n = cell.span_width if getattr(cell, "is_merge_origin", False) else 1
                        txts.append(cell.text); ws.append(sum(widths[ci:ci + n]))
                    rh = max(_in(r.height), L.row_h(txts, ws, size)) * dpi
                    cx = x
                    for ci, cell in enumerate(r.cells):
                        if getattr(cell, "is_spanned", False):
                            continue          # merge origin 已經推進晒 cx，唔可以再加
                        n = cell.span_width if getattr(cell, "is_merge_origin", False) else 1
                        cw = sum(widths[ci:ci + n]) * dpi
                        try:
                            fc = cell.fill.fore_color.rgb
                            fill = (fc[0], fc[1], fc[2])
                        except Exception:
                            fill = (255, 255, 255)
                        dr.rectangle([cx, cy, cx + cw, cy + rh], fill=fill)
                        _cell_edges(dr, cell, cx, cy, cx + cw, cy + rh)   # 只畫真係有嘅邊
                        al = str(cell.text_frame.paragraphs[0].alignment or "LEFT")
                        _draw_text(dr, _runs(cell.text_frame), cx + 2, cy + 2,
                                   (cw - 6) / dpi, sc, align=al)
                        cx += cw
                    cy += rh
                if cy > H * dpi:                     # 表爆版 → 畫紅界
                    dr.line([0, H * dpi - 2, W * dpi, H * dpi - 2], fill=(220, 0, 0), width=3)
                continue
            if sh.shape_type is not None and not sh.has_text_frame:
                continue
            oval = str(getattr(sh, "shape_type", "")).startswith("OVAL")
            try:
                if sh.fill.type is not None and sh.fill.type == 1:
                    fc = sh.fill.fore_color.rgb
                    box, col = [x, y, x + w, y + h], (fc[0], fc[1], fc[2])
                    (dr.ellipse if oval else dr.rectangle)(box, fill=col,
                                                          outline=(0, 51, 141) if oval else None)
            except Exception:
                pass
            if sh.has_text_frame and sh.text_frame.text.strip():
                al = str(sh.text_frame.paragraphs[0].alignment or "LEFT")
                _draw_text(dr, _runs(sh.text_frame), x + 1, y + 1, (w - 2) / dpi, sc, align=al)
        dr.rectangle([0, 0, W * dpi - 1, H * dpi - 1], outline=(170, 170, 170))
        img.save(out / f"s-{i:02d}.png")
    print(f"✓ preview PNG → {out}/s-*.png")


def spec(path):
    """dump 一份 pptx 嘅【版式規格】：slide 尺寸、theme 字體／配色、master/layout 名。
    → 用喺項目組真報告上，就知我哋要對嘅確切數字（user 唔使自己揾）。"""
    prs = Presentation(str(path))
    W, H = _in(prs.slide_width), _in(prs.slide_height)
    print(f"── {Path(path).name}")
    print(f"slide 尺寸：{W:.4f} x {H:.4f} in  ({prs.slide_width} x {prs.slide_height} EMU)"
          f"  = {W*2.54:.2f} x {H*2.54:.2f} cm")
    print(f"（我哋而家用：{L.SLIDE_W} x {L.SLIDE_H} in —— "
          f"{'一致 ✓' if abs(W-L.SLIDE_W)<0.02 and abs(H-L.SLIDE_H)<0.02 else '★ 唔一致，要改 layout.SLIDE_W/H'}）")
    for mi, m in enumerate(prs.slide_masters, 1):
        print(f"\nmaster {mi}：{len(m.slide_layouts)} 個 layout")
        try:
            th = m.part.part_related_by(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
            root = th._element if hasattr(th, "_element") else None
            import re as _re
            xml = th.blob.decode("utf-8", "replace") if hasattr(th, "blob") else ""
            for tag, lab in (("majorFont", "標題字體"), ("minorFont", "內文字體")):
                m2 = _re.search(tag + r".*?</a:" + tag + ">", xml, _re.S)
                if m2:
                    lat = _re.search(r'<a:latin typeface="([^"]*)"', m2.group(0))
                    ea = _re.search(r'<a:ea typeface="([^"]*)"', m2.group(0))
                    print(f"  {lab}：latin={lat.group(1) if lat else '?'}"
                          f"  ea(中文)={ea.group(1) if ea else '?'}")
            cs = _re.search(r"<a:clrScheme.*?</a:clrScheme>", xml, _re.S)
            if cs:
                pairs = _re.findall(r'<a:(dk1|lt1|dk2|lt2|accent[1-6])>.*?'
                                    r'(?:srgbClr val="([0-9A-Fa-f]{6})"|sysClr[^/]*lastClr="([0-9A-Fa-f]{6})")',
                                    cs.group(0), _re.S)
                print("  theme 配色：" + "  ".join(f"{a}=#{(b or c).upper()}" for a, b, c in pairs))
        except Exception as e:
            print("  （theme 讀唔到：", e, "）")
        for lay in m.slide_layouts:
            print(f"    · {lay.name}")
    print(f"\n內容 slide：{len(prs.slides._sldIdLst)} 版")
    for i, sl in enumerate(prs.slides, 1):
        if i > 3:
            print("    …（只列頭 3 版）"); break
        print(f"  [slide {i}] layout={sl.slide_layout.name}")
        for sh in sl.shapes:
            sizes = sorted({r.font.size.pt for p in sh.text_frame.paragraphs for r in p.runs
                            if r.font.size} if sh.has_text_frame else [])
            names = sorted({r.font.name for p in sh.text_frame.paragraphs for r in p.runs
                            if r.font.name} if sh.has_text_frame else [])
            print(f"    {str(sh.shape_type):22s} x={_in(sh.left):5.2f} y={_in(sh.top):5.2f} "
                  f"w={_in(sh.width):5.2f} h={_in(sh.height):5.2f}  pt={sizes}  font={names}")


def fonts(path, rng=None):
    """dump 一份 pptx 【實際用緊】嘅字號：按版面位置分 role，出 pt 直方圖 + 樣本。
    → 用喺項目組真報告上，就知每個位應該用幾多 pt（--spec 只睇 theme，睇唔到內文）。"""
    prs = Presentation(str(path))
    W, H = _in(prs.slide_width), _in(prs.slide_height)
    lo, hi = (rng or (1, 10 ** 9))
    tbl_hdr, tbl_body, roles = {}, {}, {}

    def add(d, k, txt):
        e = d.setdefault(k, [0, ""])
        e[0] += 1
        if not e[1] and txt.strip():
            e[1] = txt.strip().replace("\n", " ")[:44]

    def role_of(y, h):
        if y < 0.30:
            return "① 頂 breadcrumb"
        if y < 0.50:
            return "② 章節｜子題"
        if y < 1.70:
            return "③ 導語 strapline"
        if y > H - 0.45:
            return "⑥ footer/頁碼"
        if y > H - 0.85:
            return "⑤ 資料來源/註"
        return "④ 內文 body"
    for i, sl in enumerate(prs.slides, 1):
        if not (lo <= i <= hi):
            continue
        for sh in sl.shapes:
            if sh.has_table:
                t = sh.table
                for ri, r in enumerate(t.rows):
                    for c in r.cells:
                        for p_ in c.text_frame.paragraphs:
                            for run in p_.runs:
                                if run.font.size:
                                    add(tbl_hdr if ri < 2 else tbl_body,
                                        round(run.font.size.pt, 1), run.text)
                continue
            if not sh.has_text_frame:
                continue
            y = _in(sh.top)
            for p_ in sh.text_frame.paragraphs:
                for run in p_.runs:
                    if run.font.size and run.text.strip():
                        roles.setdefault(role_of(y, _in(sh.height)), {})
                        add(roles[role_of(y, _in(sh.height))],
                            round(run.font.size.pt, 1), run.text)
    print(f"── {Path(path).name}：{len(prs.slides._sldIdLst)} 版"
          f"{f'（只計 slide {lo}-{hi}）' if rng else ''}\n")
    print("── 非表格文字：逐個版面位置嘅字號分佈")
    for role in sorted(roles):
        print(f"\n  {role}")
        for pt, (n, samp) in sorted(roles[role].items(), key=lambda kv: -kv[1][0])[:6]:
            print(f"      {pt:>5} pt  × {n:<5}  「{samp}」")
    print("\n── 表格字號")
    for lab, d in (("表頭（頭 2 行）", tbl_hdr), ("表身", tbl_body)):
        if not d:
            continue
        print(f"  {lab}")
        for pt, (n, samp) in sorted(d.items(), key=lambda kv: -kv[1][0])[:5]:
            print(f"      {pt:>5} pt  × {n:<6}  「{samp}」")
    print("\n── 我哋而家用緊（layout.py）")
    for k in sorted(v for v in dir(L) if v.startswith("SZ_")):
        print(f"      {k:<12} {getattr(L, k)} pt")


def dump(path, with_tables=False, batch=0):
    """由【現成 pptx】dump 逐版文字 → 唔使重新 build（user 2026-08-12）。
    預設唔 dump 表格 cell（表格係數字、另外驗；連表格會大到 paste 唔到）→ --dump-tables 先要。"""
    prs = Presentation(str(path))
    lines = []
    for i, sl in enumerate(prs.slides, 1):
        parts = []
        for sh in sl.shapes:
            if sh.has_table:
                t = sh.table
                if not with_tables:
                    parts.append(f"〔表 {len(t.rows)}x{len(t.columns)}〕"); continue
                for row in t.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                parts.append(sh.text_frame.text.strip())
        body = "\n".join(parts)
        lines.append(f"\n===== slide {i}（{len(body)} 字）=====\n" + body)
    out = Path(path).resolve().with_name(Path(path).stem + "_dump.txt")
    txt = "\n".join(lines)
    out.write_text(txt, encoding="utf-8")
    print(f"✓ {out}（{len(prs.slides._sldIdLst)} 版）")
    # 同時喺 console 印，user 直接 copy（--batch N = 每批 N 版；0 = 一次過）
    if batch:
        for i in range(0, len(lines), batch):
            print(f"\n########## BATCH {i//batch + 1}/{-(-len(lines)//batch)}"
                  f"（slide {i+1}-{min(i+batch, len(lines))}）##########")
            print("\n".join(lines[i:i + batch]))
    else:
        print(txt)


def main():
    args = [a for a in sys.argv[1:]]
    if not args:
        print(__doc__); return
    path = args[0]
    if "--dump" in args:
        b = int(args[args.index("--batch") + 1]) if "--batch" in args else 0
        dump(path, with_tables="--dump-tables" in args, batch=b); return
    if "--fonts" in args:
        rng = None
        if "--range" in args:
            a, b = args[args.index("--range") + 1].split("-")
            rng = (int(a), int(b))
        fonts(path, rng); return
    if "--spec" in args:
        spec(path); return
    if "--render" in args:
        render(path); return
    if "--preview" in args:
        only = None
        if "--slide" in args:
            only = {int(x) for x in args[args.index("--slide") + 1].split(",")}
        preview(path, only=only); return
    if "--slide" in args:
        dump_slide(path, int(args[args.index("--slide") + 1])); return
    audit(path)


if __name__ == "__main__":
    main()
