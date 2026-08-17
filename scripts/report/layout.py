#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
layout.py — 報告 pptx 版式引擎（KPMG house style）。全部版式常數集中喺呢度，
make_report / render_review_table_pptx 只負責「派數字」，唔再各自砌 formatting。

版式對齊 mgm_2025_report scan（10.83 x 7.5 in、每版：頂 breadcrumb → 灰色「章節 | 子題」
→ navy 粗體導語 → navy caption bar + 表／敘述 → 資料來源 → footer）。

顏色跟 KPMG Visual identity overview（品牌手冊）：
    Primary   KPMG Blue 00338D｜Medium Blue 005EB8｜Light Blue 0091DA
    Secondary Violet 483698｜Purple 470A68｜Light Purple 6D2077｜Green 00A3A1
表格 tint（由 KPMG Blue 派生）：section EEF1F8、小計 D9E1F2、總計 BDD7EE、格線 BFBFBF。
字體：PowerPoint 用 Arial（品牌手冊指定）行數字/英文，中文用微软雅黑。

★ 高度控制：PowerPoint 會自動長高 row 去就內容 → 純靠 row 數分頁一定爆版。
  呢度用 est_lines() 逐 cell 估 wrap 行數 → 逐 row 定實高度 → 按【累積高度】分頁，
  所以永遠唔會超出 slide 底（單項審查「超出 border」嘅正解）。
"""
import re

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── KPMG 品牌色 ──────────────────────────────────────────────────────────
NAVY = RGBColor(0x00, 0x33, 0x8D)          # KPMG Blue（表頭 / 標題 / 導語）
MBLUE = RGBColor(0x00, 0x5E, 0xB8)         # Medium Blue
LBLUE = RGBColor(0x00, 0x91, 0xDA)         # Light Blue
VIOLET = RGBColor(0x48, 0x36, 0x98)
PURPLE = RGBColor(0x47, 0x0A, 0x68)
LPURPLE = RGBColor(0x6D, 0x20, 0x77)
GREEN = RGBColor(0x00, 0xA3, 0xA1)
# 表格 tint（KPMG Blue 派生；報告通篇用呢 3 級）
SECFILL = RGBColor(0xEE, 0xF1, 0xF8)       # 範疇 section 行
SUBTOT = RGBColor(0xD9, 0xE1, 0xF2)        # 小計
TOTAL = RGBColor(0xBD, 0xD7, 0xEE)         # 總計
BORDER = "BFBFBF"                          # 格線（srgbClr hex）
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x22, 0x22, 0x22)           # 內文黑
GREY = RGBColor(0x59, 0x59, 0x59)          # 註 / 資料來源
LGREY = RGBColor(0x8C, 0x8C, 0x8C)         # breadcrumb 非當前
DARK = RGBColor(0x0C, 0x23, 0x3C)          # 封面 / 章節分隔深底

# 負數用括號表示（KPMG palette 冇紅色）→ 唔另外上色。想要紅色改呢個做 RGBColor(0xC0,0,0)。
NEG_COLOR = None

# 字體：由公司 template theme 實測（inspect_pptx --spec，2026-08-12）
#   majorFont latin=KPMG Bold  ea=Microsoft YaHei ｜ minorFont latin=Arial  ea=Microsoft YaHei
# ⚠ python-pptx 嘅 font.name 只寫 <a:latin>，中文字要寫 <a:ea>，否則 PowerPoint 會用 theme 預設。
FONT_CN = "Microsoft YaHei"                 # <a:ea>（中文）
FONT_NUM = "Arial"                          # <a:latin>（數字/英文）
FONT_HEAD = "KPMG Bold"                     # 標題 latin（中文照樣行 ea）

# ── 字號（集中一處；inspect_pptx --fonts 可印出嚟同真報告逐個位置對）────────
# ★ 由項目組真報告實測（inspect_pptx --fonts --range 10-63，2026-08-12）：
#   章節｜子題 12pt（202 runs 單一值）｜導語 13-14pt｜內文 body 9pt（575 runs）
#   資料來源/註 7pt｜native 表格 9pt（佢哋大表多數係 Tableau 截圖，所以表身抽到 0）
#   我哋原本細成 1.5 倍（8.5/8.5/8.0/6.0/6.5）→ 全部校準。
#   ⚠ 字大咗，同一版塞唔到咁多 → 自動分頁會出多幾版，呢個【正合】報告嘅版數
#     （報告 1.3 有 4 版、1.4 有 3 版，我哋之前一版塞晒）。
SZ_CRUMB = 7.0      # ① 頂 breadcrumb（scan 睇落細，唔跟 body）
SZ_TITLE = 12.0     # ② 章節｜子題
SZ_HEAD = 13.0      # ③ 導語 strapline
SZ_BODY = 9.0       # ④ 內文 body（prose 段落）
SZ_BODY_HEAD = 9.5  # ④ 內文小標題
SZ_TBL = 7.5        # 表身（一般表；真報告 native 表 9pt，但佢哋大表係圖）
SZ_TBL_HDR = 7.0    # 表頭
SZ_TBL_WIDE = 6.0   # 表身（>16 欄嘅大表，9pt 塞唔落 18 欄）
SZ_CAPTION = 7.5    # 表頂 navy caption bar
SZ_NOTE = 7.0       # ⑤ 資料來源 / 註
SZ_FOOT = 6.0       # ⑥ footer 版權
SZ_PAGE = 9.0       # ⑥ 頁碼

SLIDE_W = 10.83                             # 報告 slide 尺寸（scan 量度確認）
SLIDE_H = 7.5

# 版面錨點（吋）
MARGIN = 0.53          # template 實測：內容 x=0.53、闊 9.76（--spec）
COL_GAP = 0.21         # template 兩欄 gap 實測
CRUMB_Y = 0.13
SUBTITLE_Y = 0.34
HEAD_Y = 0.56
FOOT_Y = 7.16
CONTENT_BOTTOM = 6.98                       # 內容最底（資料來源之上）

SECTIONS = ["2025年度投資計劃執行情況概述", "過往年度投資計劃在2025年繼續執行的審查跟進",
            "本年度審查工作的主要發現", "其他信息", "投資計劃執行報告的六項KPI分析", "附件"]

_CN_RE = None


def set_ea(run_or_font, ea=None):
    """寫 <a:ea>（中文字體）—— python-pptx 只寫 <a:latin>，唔寫 ea 中文會跌返 theme 預設。
    OOXML 次序：… latin, ea, cs …，所以要 insert 喺 latin 之後。"""
    f = getattr(run_or_font, "font", run_or_font)
    try:
        rPr = f._rPr
    except AttributeError:
        return
    if rPr is None:
        return
    el = rPr.find(qn("a:ea"))
    if el is None:
        el = rPr.makeelement(qn("a:ea"), {})
        lat = rPr.find(qn("a:latin"))
        (lat.addnext(el) if lat is not None else rPr.append(el))
    el.set("typeface", ea or FONT_CN)


# openpyxl 讀 Excel 有啲 cell 帶住 _x0000_ 呢類轉義（原檔有控制字元），照抄落 pptx 會見到
#   「在泰國曼_xFFFF_」咁嘅怪字。喺【所有文字最後出口】清一次，唔使逐個 loader 補。
_ESC = re.compile(r"_x[0-9A-Fa-f]{4}_")


def scrub(t):
    return _ESC.sub("", str(t))


def setfont(run, size, *, bold=False, italic=False, color=None, heading=False, latin=None):
    """一次過設 size/bold/color + <a:latin> + <a:ea>（跟 template theme）。
    順手清走 Excel 轉義殘留（_xFFFF_ 之類）——每個 run 一定會行過呢度。"""
    if "_x" in run.text:
        run.text = scrub(run.text)
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    if color is not None:
        f.color.rgb = color
    f.name = latin or (FONT_HEAD if heading else FONT_NUM)
    set_ea(f)
    return run


def _is_cn(ch):
    return "⺀" <= ch <= "鿿" or "＀" <= ch <= "￯" or "　" <= ch <= "〿"


def has_cn(s):
    return any(_is_cn(c) for c in str(s))


def text_w(s, size):
    """估文字闊度（pt）：中文/全形 ≈ 1 em、英數 ≈ 0.52 em。"""
    w = 0.0
    for c in str(s):
        w += size * (1.0 if _is_cn(c) else 0.52)
    return w


def est_lines(s, col_w_in, size, margin_in=0.06):
    """估 wrap 行數（col_w_in = 欄闊吋）。認 \\n 明碼換行。"""
    avail = max((col_w_in - margin_in) * 72.0, 6.0)
    n = 0
    for seg in str(s).split("\n"):
        n += max(1, -(-text_w(seg, size) // avail))     # ceil
    return int(n)


def row_h(cells, widths, size, pad_in=0.045, min_h=0.155):
    """一行嘅需要高度（吋）＝ 最多 wrap 行數 × 行距 + 上下 padding。"""
    lines = 1
    for txt, w in zip(cells, widths):
        lines = max(lines, est_lines(txt, w, size))
    return max(min_h, lines * size * 1.24 / 72.0 + pad_in)


# ── 基本元件 ─────────────────────────────────────────────────────────────
def size_of(prs):
    return prs.slide_width / 914400.0, prs.slide_height / 914400.0


def blank(prs):
    lay = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(lay)


def _tb(slide, x, y, w, h, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return box


def put(slide, x, y, w, h, text, *, size=8, bold=False, color=INK, align=PP_ALIGN.LEFT,
        italic=False, font=None, wrap=True):
    """一行/一段文字框。wrap=False 用喺一定要一行嘅嘢（breadcrumb 頁籤）。"""
    box = _tb(slide, x, y, w, h, wrap)
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    p.font.size = Pt(size); p.font.name = FONT_NUM; set_ea(p.font)   # 空段落唔好跌返 theme 預設
    p._p.get_or_add_endParaRPr().set("sz", str(int(round(size * 100))))
    r = p.add_run(); r.text = str(text)
    setfont(r, size, bold=bold, italic=italic, color=color, latin=font)
    return box


BAND = RGBColor(0xF2, 0xF2, 0xF2)          # breadcrumb 淺灰底（scan 頂部 banner）


def _rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.5)
    sh.shadow.inherit = False
    return sh


def _name(shape, nm):
    try:
        shape.name = nm
    except Exception:      # noqa: BLE001 — 改唔到名只係少咗 hyperlink，唔好炸咗成個 build
        pass


def breadcrumb(slide, W, active=0, entity="MGM"):
    """頂 nav（對 scan p-23 放大）：白底、頁籤用「｜」分隔，當前頁籤 navy 粗體、其餘淺灰，
    右邊 entity + ◀ ⌂ ▶ 三粒圓掣。shape 改名做 nav:* ，wire_nav() 事後駁內部 hyperlink。"""
    x0 = MARGIN - 0.23
    d, gapc = 0.185, 0.05                                  # 圓掣直徑 / 間距
    right = W - x0
    for i, (nm, ch) in enumerate((("next", "▶"), ("home", "⌂"), ("prev", "◀"))):
        cx = right - d - i * (d + gapc)
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(CRUMB_Y - 0.035),
                                    Inches(d), Inches(d))
        sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
        sh.line.color.rgb = NAVY; sh.line.width = Pt(0.75); sh.shadow.inherit = False
        _name(sh, f"nav:{nm}")
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ch
        setfont(r, 6.0, bold=True, color=NAVY)
    ex = right - 3 * (d + gapc)                            # entity 靠圓掣左邊
    put(slide, ex - 0.85, CRUMB_Y, 0.8, 0.18, entity, size=SZ_CRUMB, bold=True,
        color=INK, align=PP_ALIGN.RIGHT)
    sep, avail = " ｜ ", (ex - 0.92) - x0
    # ×1.08：text_w 對粗體中文估細咗少少，唔留鬆位頁籤會撞埋一齊
    widths = [text_w(t, SZ_CRUMB) * 1.08 / 72.0 for t in SECTIONS]
    sw = text_w(sep, SZ_CRUMB) / 72.0
    scale = min(1.0, avail / (sum(widths) + sw * (len(SECTIONS) - 1)))
    x = x0
    for i, t in enumerate(SECTIONS):
        if i:
            put(slide, x, CRUMB_Y, sw * scale + 0.03, 0.18, sep, size=SZ_CRUMB * scale,
                color=LGREY, wrap=False)
            x += sw * scale
        w = widths[i] * scale
        _name(put(slide, x, CRUMB_Y, w + 0.05, 0.18, t, size=SZ_CRUMB * scale, wrap=False,
                  bold=(i == active), color=NAVY if i == active else LGREY), f"nav:sec{i}")
        x += w


def _hlink(shape, rid):
    """畀 shape 內所有 run 加內部跳頁 hyperlink（a:hlinkClick + ppaction://hlinksldjump）。
    hlinkClick 喺 CT_TextCharacterProperties 排 latin/ea 之後 → append 就啱序。"""
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            rPr = r._r.get_or_add_rPr()
            h = rPr.makeelement(qn("a:hlinkClick"),
                                {qn("r:id"): rid, "action": "ppaction://hlinksldjump"})
            rPr.append(h)


def wire_nav(prs, sec_slide=None, home=0):
    """全部 slide 砌完（連目錄插咗、重排咗）之後至駁：◀/▶ = 上/下頁、⌂ = 目錄、
    頁籤 = 該章分隔頁。sec_slide = {章 index: slide index}。"""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    slides = list(prs.slides)
    n = len(slides)
    for i, s in enumerate(slides):
        tgt = {"nav:prev": max(i - 1, 0), "nav:next": min(i + 1, n - 1), "nav:home": home}
        for k, v in (sec_slide or {}).items():
            tgt[f"nav:sec{k}"] = v
        for sh in s.shapes:
            j = tgt.get(sh.name)
            if j is None or j == i or not sh.has_text_frame:
                continue
            _hlink(sh, s.part.relate_to(slides[j].part, RT.SLIDE))


def footer(slide, W, H, page):
    """底：KPMG 字標 + 版權 + 初稿/頁碼（對 scan）。"""
    kb = _tb(slide, MARGIN - 0.23, H - 0.34, 0.7, 0.22)
    kr = kb.text_frame.paragraphs[0].add_run(); kr.text = "KPMG"
    setfont(kr, 11, bold=True, italic=True, color=NAVY)
    put(slide, MARGIN + 0.5, H - 0.30, W - 2.2, 0.2,
        "© 2026畢馬威會計師事務所 — 澳門特別行政區合夥制事務所。版權所有，不得轉載。",
        size=SZ_FOOT, color=LGREY)
    if page is not None:
        put(slide, W - 1.15, H - 0.32, 0.95, 0.2, f"初稿　{page}", size=SZ_PAGE, bold=True,
            color=NAVY, align=PP_ALIGN.RIGHT)


MAX_HEAD_H = 1.35      # 導語最多食呢咁多高（scan 一般 2-4 行）；再長就縮字，唔可以食晒成版


def head_h(headline, W, hsize=SZ_HEAD):
    """導語需要嘅高度 + 實際字號（長就自動縮到 MAX_HEAD_H 為止）→ (h, size)。"""
    if not headline:
        return 0.06, hsize
    while hsize > 6.0:
        h = est_lines(headline, W - 2 * MARGIN, hsize) * hsize * 1.35 / 72.0
        if h <= MAX_HEAD_H:
            return h, hsize
        hsize -= 0.5
    return MAX_HEAD_H, hsize


def page_head(slide, W, crumb, headline=None, *, hsize=SZ_HEAD):
    """灰色「章節 | 子題」+ navy 粗體導語 → 回內容起始 y。"""
    put(slide, MARGIN, SUBTITLE_Y, W - 2 * MARGIN, 0.2, crumb, size=SZ_TITLE, bold=True, color=NAVY)
    if not headline:
        return HEAD_Y + 0.06
    h, hsize = head_h(headline, W, hsize)
    box = _tb(slide, MARGIN, HEAD_Y, W - 2 * MARGIN, h)
    p = box.text_frame.paragraphs[0]
    r = p.add_run(); r.text = str(headline)
    setfont(r, hsize, bold=True, color=NAVY, heading=True)
    return HEAD_Y + h + 0.10


def caption_bar(slide, x, y, w, text, *, size=SZ_CAPTION):
    """表頂 caption bar（重覆表名，對 scan 每張表都有）。
    ⚠ 用深色 HDR2 —— IMG_0441 量到 caption 條比表頭嗰排藍【深啲】，唔係同一隻色。"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.17))
    bar.fill.solid(); bar.fill.fore_color.rgb = CAPTION_FILL
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.margin_left = Emu(36000); tf.margin_right = Emu(18000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = str(text)
    setfont(r, size, bold=True, color=WHITE)
    return y + 0.17


def source_note(slide, W, y=None, *, note=None, more=False):
    """表下：資料來源（左）+（下頁待續）（右）。"""
    y = CONTENT_BOTTOM if y is None else y
    put(slide, MARGIN, y, W - 2.0, 0.16,
        note or "資料來源：管理層提供之項目投資計劃及執行報告資料，畢馬威分析",
        size=SZ_NOTE, color=NOTE_FG)
    if more:
        put(slide, W - MARGIN - 1.2, y, 1.2, 0.16, "（下頁待續）", size=SZ_NOTE, color=NOTE_FG,
            align=PP_ALIGN.RIGHT)


# ── 表格 ─────────────────────────────────────────────────────────────────
# ★ 真報告（IMG_0441 彩色版）嘅表：【冇逐格格線、冇 row 底色】。
#   得返：表頭 navy（設施建設/活動舉辦 嗰組 teal）＋ 小計/總計行上下幼橫線
#   ＋ 欄組之間虛線直線。之前全格線 + sec/小計/總計 3 級藍底 = 自己作，同報告唔同。
# ★ 表格配色（項目組 2026-08-17 逐項指定 hex，唔再靠影相估）：
RULE = "00338D"                             # 表格線（橫線 + 欄組虛線）＝ KPMG Blue
HDR_FILL = RGBColor(0x1E, 0x49, 0xE2)       # 表頭預設（accent1 亮藍）
#   ⚠ 唔可以叫 HDR —— make_report 已經有個 HDR(=NAVY)，bundle dedup 會靜靜丟咗佢
HDR_KEY = RGBColor(0x09, 0x8E, 0x7E)        # 綠：重點欄（獲批的計劃投資金額／潛在調整後／三年累計／比例）
HDR_SKY = RGBColor(0x00, 0xB8, 0xF5)        # 天藍：調整事項欄組（1-7+合計）／2025年度／潛在調整金額
HDR_PUR = RGBColor(0x48, 0x36, 0x98)        # 紫：2024年度（KPMG Violet；項目組未畀 hex，暫用品牌紫）
CAPTION_FILL = NAVY                         # caption 條 #00338D
SEC_FG = NAVY                               # 「博彩項目 / 非博彩項目」字色
NOTE_FG = NAVY                              # 註 / 資料來源
HDR1, HDR2, HDR3 = HDR_FILL, HDR_FILL, HDR_KEY        # 舊名保留（唔好散落 import error）
TEAL = HDR_KEY


def _edge(cell, side, *, w=9525, color=RULE, dash=None):
    """畫單一條邊（side ∈ T/B/L/R）。ln* 要插喺 tcPr 最前，否則 PowerPoint 會叫修復。
    同一邊重覆設就換走舊嗰條。"""
    tcPr = cell._tc.get_or_add_tcPr()
    tag = qn(f"a:ln{side}")
    old = tcPr.find(tag)
    if old is not None:
        tcPr.remove(old)
    ln = tcPr.makeelement(tag, {"w": str(w), "cap": "flat"})
    fill = ln.makeelement(qn("a:solidFill"), {})
    clr = fill.makeelement(qn("a:srgbClr"), {"val": color})
    fill.append(clr); ln.append(fill)
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    tcPr.insert(0, ln)


def set_cell(cell, text, *, size=SZ_TBL, bold=False, fill=None, align=PP_ALIGN.RIGHT,
             color=None, wrap=True, anchor=MSO_ANCHOR.MIDDLE, italic=False):
    cell.margin_left = cell.margin_right = Emu(18000)
    cell.margin_top = cell.margin_bottom = Emu(9000)
    cell.vertical_anchor = anchor
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill if fill is not None else WHITE
    tf = cell.text_frame; tf.word_wrap = wrap
    tf.clear()          # merge 會把被合併格嘅字搬入 origin → 唔清就會出「萬澳門元萬澳門元」
    p = tf.paragraphs[0]; p.alignment = align
    txt = "" if text is None else str(text)
    if color is None:
        color = NEG_COLOR if (NEG_COLOR is not None and txt.startswith("(")) else INK
    # ★ 空格一定要定死字號：PowerPoint 見到「冇 run／空 run」會跌返去 endParaRPr／預設 text
    #   style（python-pptx fresh deck ＝ Calibri 18pt）→ row 被撐到 ~0.3in，成張表爆版。
    #   所以：paragraph 層 defRPr + endParaRPr 都寫死，而且空字串索性唔加 run。
    p.font.size = Pt(size); p.font.bold = bold
    p.font.name = FONT_NUM; set_ea(p.font)
    epr = p._p.get_or_add_endParaRPr()
    epr.set("sz", str(int(round(size * 100))))
    if not txt:
        return
    # ⚠ DrawingML 入面 "\n" 唔係換行（會當空白）→ 一定要用 <a:br/>，否則表頭喺 PowerPoint 會擠成一行
    for i, seg in enumerate(txt.split("\n")):
        if i:
            p.add_line_break()
        if not seg:
            continue
        r = p.add_run(); r.text = seg
        setfont(r, size, bold=bold, italic=italic, color=color)


ROW_FILL = {"sec": None, "subtot": None, "tot": None, "data": None, "formula": None}   # 報告：body 全白，靠橫線分層


def header_h(supers, subs, widths, hfont):
    """表頭需要高度（吋）。"""
    h = row_h(subs, widths, hfont, pad_in=0.05, min_h=0.20)
    return (0.17 + h) if supers else h


def fit_rows(rows, widths, font, avail_h, hh):
    """按【累積高度】切頁 → 保證唔會超出可用高度。rows = [(kind, cells)]。
    keep=True 嘅 row（範疇 block）盡量唔拆：見 fit_blocks。"""
    out, cur, used = [], [], 0.0
    cap = max(avail_h - hh, 0.6)          # guard：導語太長時唔好變 0/負數（會無限開版）
    for kind, cells in rows:
        h = row_h(cells, widths, font)
        if cur and used + h > cap:
            out.append(cur); cur, used = [], 0.0
        cur.append((kind, cells)); used += h
    if cur:
        out.append(cur)
    return out or [[]]


def fit_blocks(blocks, widths, font, avail_h, hh):
    """block = 一個範疇（section + data + 小計）。整個 block 唔拆頁（對 scan：全報告冇「續」）；
    單一 block 大過一版先逼住切。"""
    cap = max(avail_h - hh, 0.6)          # guard：同上
    pages, cur, used = [], [], 0.0
    for blk in blocks:
        bh = sum(row_h(c, widths, font) for _, c in blk)
        if cur and used + bh > cap:
            pages.append(cur); cur, used = [], 0.0
        if bh > cap:                                    # 單一範疇爆版 → 逐行切（安全網）
            for kind, cells in blk:
                h = row_h(cells, widths, font)
                if cur and used + h > cap:
                    pages.append(cur); cur, used = [], 0.0
                cur.append((kind, cells)); used += h
            continue
        cur.extend(blk); used += bh
    if cur:
        pages.append(cur)
    return pages or [[]]


def draw_table(slide, x, y, w, subs, rows, widths, *, supers=None, font=SZ_TBL, hfont=SZ_TBL_HDR,
               left_cols=1, fill_h=None, max_row_h=0.26, hdr_cols=None):
    """畫 navy 表。subs=欄名（可含 \\n）；rows=[(kind, cells)]；widths=相對闊度（會 scale 到 w）。
    supers=[(label, c0, c1_exclusive)] 兩層表頭。fill_h=想填滿嘅高度（行數少時撐開行高，
    唔好剩一大橛白位；每行最多 max_row_h）。回 (bottom_y, 實際高度)。"""
    ncol = len(subs)
    scale = w / sum(widths)
    wid = [v * scale for v in widths]
    nhdr = 2 if supers else 1
    heights = [row_h(cells, wid, font) for _, cells in rows]
    hsub = row_h(subs, wid, hfont, pad_in=0.05, min_h=0.20)
    if fill_h and heights:
        slack = fill_h - ((0.17 if supers else 0) + hsub + sum(heights))
        if slack > 0.05:
            add = min(slack / len(heights), max(0.0, max_row_h - max(heights)))
            heights = [h + add for h in heights]
    total_h = (0.17 if supers else 0) + hsub + sum(heights)
    tbl = slide.shapes.add_table(nhdr + len(rows), ncol, Inches(x), Inches(y),
                                 Inches(w), Inches(total_h)).table
    tbl.first_row = False; tbl.horz_banding = False
    for i, v in enumerate(wid):
        tbl.columns[i].width = Inches(v)
    # 三色欄組——【要 caller 明示】：4.2 表都有「設施建設/活動舉辦」欄但成排 navy（scan p.24），
    #   所以唔可以淨靠欄名估。hdr_cols = {欄 index: 顏色}，冇指定就 HDR1。
    hc = dict(hdr_cols or {})
    if supers:
        for c in range(ncol):
            set_cell(tbl.cell(0, c), "", size=hfont, fill=hc.get(c, HDR1), color=WHITE)
        # 一個欄組跨住兩隻表頭色（報告：潛在調整後 = 深藍嗰兩欄 + 綠嗰兩欄）→ 拆開兩格，
        #   個 label 兩邊都寫（同 IMG_0441 一樣，「潛在調整後投資金額」出現兩次）。
        for label, c0, c1 in supers:
            a = c0
            while a < c1:
                b = a + 1
                while b < c1 and hc.get(b, HDR1) == hc.get(a, HDR1):
                    b += 1
                if b - a > 1:
                    tbl.cell(0, a).merge(tbl.cell(0, b - 1))
                set_cell(tbl.cell(0, a), label or "", size=hfont + 0.5, bold=True,
                         fill=hc.get(a, HDR1), color=WHITE, align=PP_ALIGN.CENTER)
                a = b
        tbl.rows[0].height = Emu(int(0.17 * 914400))
    for c, s in enumerate(subs):
        # 報告嘅欄名喺表頭【貼底】（wrap 做兩行時尤其明顯）
        set_cell(tbl.cell(nhdr - 1, c), s, size=hfont, bold=True, anchor=MSO_ANCHOR.BOTTOM,
                 fill=hc.get(c, HDR1), color=WHITE,
                 align=PP_ALIGN.LEFT if c < left_cols else PP_ALIGN.CENTER)
    # 角位：報告係一格（序號欄冇字），單位「萬澳門元」貼住最左
    if left_cols >= 2 and not str(subs[0]).strip() and str(subs[1]).strip():
        tbl.cell(nhdr - 1, 0).merge(tbl.cell(nhdr - 1, 1))
        set_cell(tbl.cell(nhdr - 1, 0), subs[1], size=hfont, bold=True, anchor=MSO_ANCHOR.BOTTOM,
                 fill=hc.get(0, HDR1), color=WHITE, align=PP_ALIGN.LEFT)
    tbl.rows[nhdr - 1].height = Emu(int(hsub * 914400))
    for ri, (kind, cells) in enumerate(rows, start=nhdr):
        bold = kind in ("sec", "subtot", "tot")
        if kind == "formula":      # 報告表頭下面嗰行斜體公式（a｜1..7｜b｜c=a+b｜d=b/a）
            for c, v in enumerate(cells):
                set_cell(tbl.cell(ri, c), v, size=max(4.5, font - 1.0), italic=True,
                         color=GREY, align=PP_ALIGN.LEFT if c < left_cols else PP_ALIGN.RIGHT)
            tbl.rows[ri].height = Emu(int(max(0.14, (font - 1.0) * 1.24 / 72.0 + 0.03) * 914400))
            continue
        # 標籤（範疇/小計/總計/表尾說明行）喺報告係【由最左邊起】，唔係縮喺名稱欄：
        #   序號欄空 + 名稱欄有字 → merge 埋，個 label 先有位唔會 wrap
        k = 0
        if left_cols >= 2 and not str(cells[0]).strip() and str(cells[1]).strip():
            k = 1
            tbl.cell(ri, 0).merge(tbl.cell(ri, 1))
        for c, v in enumerate(cells):
            if 0 < c <= k:
                continue                      # 已 merge 入 col 0
            al = PP_ALIGN.LEFT if c < left_cols else PP_ALIGN.RIGHT
            set_cell(tbl.cell(ri, c), cells[k] if c == 0 and k else v,
                     size=font, bold=bold, fill=ROW_FILL.get(kind), align=al,
                     color=SEC_FG if kind == "sec" else None)
        tbl.rows[ri].height = Emu(int(heights[ri - nhdr] * 914400))
    # ── 線：只有小計/總計橫線 + 欄組虛線直線（報告冇逐格格線）────────────
    # 只喺【有名嘅欄組】邊界畫虛線；標籤欄自成一「組」（label=""）唔算
    gsep = {c0 for _l, c0, _c1 in (supers or []) if c0 > 0 and str(_l).strip()}
    last = nhdr + len(rows) - 1
    for ri, (kind, _c) in enumerate(rows, start=nhdr):
        if kind in ("subtot", "tot"):
            for c in range(ncol):
                _edge(tbl.cell(ri, c), "T")
                _edge(tbl.cell(ri, c), "B")      # 報告：小計同總計【上下都有】幼線
    for ri in range(nhdr, nhdr + len(rows)):
        for c in gsep:
            _edge(tbl.cell(ri, c), "L", w=6350, color=RULE, dash="sysDash")
    return y + total_h, total_h


# ── 敘述 ─────────────────────────────────────────────────────────────────
def prose(box, items, *, head_size=SZ_BODY_HEAD, body_size=SZ_BODY, gap=6):
    """scan 敘述格式：navy 粗體小標題一行 + 下面 body 段落（唔用 ■ bullet）。
    items = [(head, body)]；head 可為空。"""
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for head, body in items:
        if head:
            head = str(head).rstrip("：:")        # 項目組：小標題唔應該有冒號
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_before = Pt(0 if first else gap); p.space_after = Pt(1)
            p.font.size = Pt(head_size); p.font.name = FONT_NUM; set_ea(p.font)
            p._p.get_or_add_endParaRPr().set("sz", str(int(round(head_size * 100))))
            r = p.add_run(); r.text = str(head)
            setfont(r, head_size, bold=True, color=NAVY)
            first = False
        if body:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_before = Pt(0 if first else (0 if head else gap)); p.space_after = Pt(1)
            p.font.size = Pt(body_size); p.font.name = FONT_NUM; set_ea(p.font)
            p._p.get_or_add_endParaRPr().set("sz", str(int(round(body_size * 100))))
            r = p.add_run(); r.text = str(body)
            setfont(r, body_size, color=RGBColor(0x33, 0x33, 0x33))
            first = False


def prose_numbered(box, items, *, size=SZ_BODY, gap=7, indent=0.24, title=None, tsize=SZ_BODY_HEAD):
    """scan 表旁格式（p-11/p-13 右欄）：navy 粗體小標題 + 編號清單
        1.  {粗體類型}（{金額}）：{內文…}       ← hanging indent，內文對齊類型名
    items = [(編號, 粗體引子, 內文)]；編號跟七大類 canonical 序（會跳號）。"""
    tf = box.text_frame; tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]; p.space_after = Pt(4)
        p.font.size = Pt(tsize)
        p._p.get_or_add_endParaRPr().set("sz", str(int(round(tsize * 100))))
        r = p.add_run(); r.text = str(title)
        setfont(r, tsize, bold=True, color=NAVY)
        first = False
    emu = int(indent * 914400)
    for no, head, body in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(gap); p.space_after = Pt(0)
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(emu)); pPr.set("indent", str(-emu))     # hanging indent
        p.font.size = Pt(size); p.font.name = FONT_NUM; set_ea(p.font)
        p._p.get_or_add_endParaRPr().set("sz", str(int(round(size * 100))))
        rn = p.add_run(); rn.text = f"{no}.\t"
        setfont(rn, size, bold=True, color=NAVY)
        rh = p.add_run(); rh.text = str(head)
        setfont(rh, size, bold=True, color=NAVY)
        rb = p.add_run(); rb.text = str(body)
        setfont(rb, size, color=RGBColor(0x33, 0x33, 0x33))


def est_numbered_h(items, w, size=SZ_BODY, gap=7, title=None, tsize=SZ_BODY_HEAD, indent=0.24):
    h = (est_lines(title, w, tsize) * tsize * 1.3 / 72.0 + 4 / 72.0) if title else 0.0
    for _no, head, body in items:
        h += est_lines(f"　{head}{body}", w - indent, size) * size * 1.35 / 72.0 + gap / 72.0
    return h


def prose_box(slide, x, y, w, h, items, **kw):
    box = _tb(slide, x, y, w, h)
    prose(box, items, **kw)
    return box


def est_prose_h(items, w, head_size=SZ_BODY_HEAD, body_size=SZ_BODY, gap=6):
    """估敘述高度（吋）→ 用嚟分頁，唔會爆版。"""
    h = 0.0
    for head, body in items:
        if head:
            h += est_lines(head, w, head_size) * head_size * 1.3 / 72.0 + gap / 72.0
        if body:
            h += est_lines(body, w, body_size) * body_size * 1.35 / 72.0 + 2 / 72.0
    return h


def fit_prose(items, w, avail_h, **kw):
    """按估算高度切頁 → [[items]]。"""
    pages, cur, used = [], [], 0.0
    avail_h = max(avail_h, 0.6)           # guard
    for it in items:
        ih = est_prose_h([it], w, **kw)
        if cur and used + ih > avail_h:
            pages.append(cur); cur, used = [], 0.0
        cur.append(it); used += ih
    if cur:
        pages.append(cur)
    return pages or [[]]


# ── 深色版（封面 / 章節分隔）─────────────────────────────────────────────
def apply_theme_fonts(prs):
    """把生成 deck 嘅 theme 字體改成公司 template 嗰套（major KPMG Bold / minor Arial，
    ea 兩者都 Microsoft YaHei）。python-pptx 開新檔用 Office 預設 theme（Calibri），
    凡係我哋冇明寫字體嘅地方（placeholder、空段落、表格預設）都會跌返 Calibri。"""
    import re as _re
    try:
        for m in prs.slide_masters:
            part = m.part.part_related_by(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
            xml = part.blob.decode("utf-8")
            for tag, latin in (("majorFont", FONT_HEAD), ("minorFont", FONT_NUM)):
                def _fix(mo, latin=latin):
                    seg = mo.group(0)
                    seg = _re.sub(r'<a:latin typeface="[^"]*"', f'<a:latin typeface="{latin}"', seg, count=1)
                    seg = _re.sub(r'<a:ea typeface="[^"]*"', f'<a:ea typeface="{FONT_CN}"', seg, count=1)
                    return seg
                xml = _re.sub(r"<a:" + tag + r">.*?</a:" + tag + r">", _fix, xml, flags=_re.S)
            part._blob = xml.encode("utf-8")
    except Exception:
        pass            # theme 改唔到唔應該搞冧成個 build（每個 run 都已經明寫咗字體）


def dark_slide(prs):
    slide = blank(prs)
    W, H = size_of(prs)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = DARK
    rect.line.fill.background(); rect.shadow.inherit = False
    kb = _tb(slide, 0.55, 0.35, 2.2, 0.4)
    kr = kb.text_frame.paragraphs[0].add_run(); kr.text = "KPMG"
    setfont(kr, 20, bold=True, italic=True, color=WHITE)
    return slide, W, H
