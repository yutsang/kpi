#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
inspect_template.py — dump pptx 嘅 slide master + 每個 slide layout + placeholders（idx/type/名/位置/字體）。
用嚟評估：可唔可以直接開間公司空 template 做 base、用佢 layouts 填內容（formatting 來自 master），
唔使 Claude 手砌顏色/furniture。

用法：python scripts\\report\\inspect_template.py "你的空template.pptx"
   （用【空 template】——得 master+layouts、冇內容；唔好用填好嘅初稿報告）
"""
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("✗ pip install python-pptx"); sys.exit(1)


def _in(emu):
    try:
        return round(Emu(emu).inches, 2)
    except Exception:
        return emu


def _ph_type(ph):
    try:
        return str(ph.placeholder_format.type).split()[0]
    except Exception:
        return "?"


def _font_of(ph):
    try:
        r = ph.text_frame.paragraphs[0].runs
        f = (r[0].font if r else ph.text_frame.paragraphs[0].font)
        col = ""
        try:
            col = f"#{f.color.rgb}" if f.color and f.color.type is not None else ""
        except Exception:
            col = ""
        return f"{f.name or ''} {f.size.pt if f.size else ''}pt {'B' if f.bold else ''} {col}".strip()
    except Exception:
        return ""


def main():
    if len(sys.argv) < 2:
        print("俾 template pptx 路徑"); return
    prs = Presentation(sys.argv[1])
    print(f"slide size: {_in(prs.slide_width)} x {_in(prs.slide_height)} in")
    print(f"slide_masters: {len(prs.slide_masters)}　| 直接可用 layouts: {len(prs.slide_layouts)}")
    print("=" * 90)
    for mi, master in enumerate(prs.slide_masters):
        print(f"\n■ MASTER {mi}: name='{master.name}'")
        # master 背景色
        try:
            fill = master.background.fill
            print(f"   background fill type: {fill.type}")
        except Exception:
            pass
        for li, layout in enumerate(master.slide_layouts):
            phs = list(layout.placeholders)
            print(f"\n   ▸ Layout {li}: '{layout.name}'  （{len(phs)} placeholders）")
            for ph in phs:
                print(f"       [idx {ph.placeholder_format.idx}] {_ph_type(ph):<10} "
                      f"name='{ph.name}'  pos=({_in(ph.left)},{_in(ph.top)},{_in(ph.width)},{_in(ph.height)})"
                      f"  font：{_font_of(ph)}")
    _dump_theme(prs)
    _dump_bg(prs)
    print("\n" + "=" * 90)
    print("→ paste 返（尤其【THEME 色盤/字體】+【背景色】兩段）：我 hardcode 精確 spec 入 code")


A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _clr_hex(el):
    """a:*Clr（srgbClr/sysClr）→ #hex。"""
    for ch in el:
        tag = ch.tag
        if tag == A + "srgbClr":
            return "#" + (ch.get("val") or "")
        if tag == A + "sysClr":
            return "#" + (ch.get("lastClr") or ch.get("val") or "")
        if tag == A + "schemeClr":
            return "scheme:" + (ch.get("val") or "")
    return "?"


def _dump_theme(prs):
    """由每個 master 嘅 theme part 抽 色盤(clrScheme) + 字體(fontScheme)。"""
    try:
        from lxml import etree
    except ImportError:
        print("（冇 lxml，跳過 theme）"); return
    print("\n" + "=" * 90)
    for mi, master in enumerate(prs.slide_masters):
        theme = None
        for rel in master.part.rels.values():
            if "theme" in getattr(rel, "reltype", ""):
                theme = rel.target_part
                break
        if theme is None:
            continue
        root = etree.fromstring(theme.blob)
        print(f"\n■ MASTER {mi} — THEME 色盤：")
        clr = root.find(f".//{A}clrScheme")
        if clr is not None:
            for c in clr:
                print(f"    {etree.QName(c).localname:<10}: {_clr_hex(c)}")
        print(f"■ MASTER {mi} — THEME 字體：")
        fs = root.find(f".//{A}fontScheme")
        if fs is not None:
            for grp in fs:
                nm = etree.QName(grp).localname
                if nm not in ("majorFont", "minorFont"):
                    continue
                lat = grp.find(f"{A}latin"); ea = grp.find(f"{A}ea"); cs = grp.find(f"{A}cs")
                print(f"    {nm}: latin={lat.get('typeface') if lat is not None else ''}"
                      f"  ea={ea.get('typeface') if ea is not None else ''}"
                      f"  cs={cs.get('typeface') if cs is not None else ''}")


def _dump_bg(prs):
    """master + 關鍵 layout（封面/分隔）背景填色。"""
    print("\n" + "=" * 90 + "\n■ 背景色（master + 封面/分隔 layout）：")
    def _fillinfo(obj, label):
        try:
            f = obj.background.fill
            t = str(f.type)
            hexv = ""
            try:
                hexv = f"#{f.fore_color.rgb}"
            except Exception:
                pass
            print(f"    {label}: type={t} {hexv}")
        except Exception as e:
            print(f"    {label}: (讀唔到 {e})")
    for mi, master in enumerate(prs.slide_masters):
        _fillinfo(master, f"master{mi}")
        for lay in master.slide_layouts:
            if any(k in lay.name for k in ("TITLE SLIDE", "Section Divider", "Divider", "Back Cover")):
                _fillinfo(lay, f"  layout '{lay.name}'")


if __name__ == "__main__":
    main()
