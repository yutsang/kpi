#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
demo_page.py — 自足 demo：淨用 code 砌「概述 整體執行概況」1 版，模仿 scan
（2 層表頭、navy、小計/總計 shaded、數字右對齊、2 欄：表左+敘述右、breadcrumb、footer）。
數據讀 gitignored results\\project_dump.tsv（inspect_project.py 出，報告年25 aggregate 到範疇）；
冇檔就用明顯假數（本檔【零真實金額】，可 commit）。出 demo_overview.pptx。

用法：python scripts\\report\\inspect_project.py ... 先出 tsv，再 python scripts\\report\\demo_page.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# 由 template theme hardcode（精確 spec，唔再估）
NAVY = RGBColor(0x00, 0x33, 0x8D)      # dk2/accent2 主色
DARK = RGBColor(0x0C, 0x23, 0x3C)      # accent3 封面/分隔深底
LT2 = RGBColor(0xE5, 0xE5, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUBTOT = RGBColor(0xD9, 0xE1, 0xF2)
TOT = RGBColor(0xBD, 0xD7, 0xEE)
SECHDR = RGBColor(0xEE, 0xF1, 0xF8)
GREY = RGBColor(0x7F, 0x7F, 0x7F)
FONT_CN = "微软雅黑"      # 中文 body（theme minorFont ea）
FONT_NUM = "Arial"        # 數字/英文（theme minorFont latin）★表格數字用呢個
FONT_HEAD = "KPMG Bold"   # 標題（theme majorFont；Win 有裝，Mac render 會 fallback）
FONT = FONT_CN            # 舊 alias

SUB = ["（萬澳門元）", "項目\n數量", "獲批的計劃\n投資金額", "報告\n投資金額", "投資計劃\n完成率",
       "潛在調整後\n投資金額", "潛在調整後\n完成率", "設施建設/\n資本性支出", "活動舉辦/\n營運性支出"]
# 博彩範疇擺前（其餘當非博彩）；顯示順序
GAMING_SUBS = ["博彩娛樂場優化", "博彩設施設備優化"]
NONG_ORDER = ["吸引外國客源", "會議展覽", "娛樂表演", "體育盛事", "文化藝術", "健康養生",
              "主題遊樂", "美食之都", "社區旅遊", "海上旅遊", "其他"]

# fallback 假數（零真實金額；只為展示版式）
FAKE = [
    ("gaming", "博彩娛樂場優化", 2, 1000, 2000, 2000, 1900, 100),
    ("gaming", "博彩設施設備優化", 2, 1000, 1500, 1500, 1400, 100),
    ("non_gaming", "吸引外國客源", 3, 2000, 1800, 900, 100, 1700),
    ("non_gaming", "會議展覽", 2, 1500, 600, 590, 300, 290),
    ("non_gaming", "娛樂表演", 2, 3000, 4000, 1000, 100, 900),
    ("non_gaming", "其他", 2, 500, 300, 280, 0, 280),
]


def _fmt(v):
    try:
        n = float(v)
    except (TypeError, ValueError):
        return str(v)
    if n == 0:
        return "-"
    return f"({abs(n):,.0f})" if n < 0 else f"{n:,.0f}"


def _rate(rep, pl):
    return f"{rep / pl * 100:.1f}%" if pl else "-"


import re as _re
_NUMRE = _re.compile(r"^\(?-?[\d,]+(\.\d+)?\)?%?$")


def _tonum(s):
    s = str(s).replace(",", "").replace("%", "")
    neg = s.startswith("(") and s.endswith(")")
    s = s.strip("()")
    try:
        v = float(s)
        return -v if neg else v
    except ValueError:
        return 0.0


def load_agg():
    """讀 results/project_dump.tsv（報告年25）aggregate 到 (scope,範疇)；冇就用假數。回 dict。
    ⚠ tab 常被 copy/paste 變 space + 項目名含空格 → 唔逐 column split，改『由尾認數字』：
    前 3 token = 報告年/scope/範疇；尾段連續數字 = [計劃,調整前,調整,調整後,(完成率?),設施,活動]。"""
    p = Path("results/project_dump.tsv")
    recs = []
    if p.exists():
        lines = [ln for ln in p.read_text(encoding="utf-8").splitlines() if ln.strip()]
        for ln in lines[1:]:      # skip header
            tok = ln.split()
            if len(tok) < 9:
                continue
            nums = []
            while tok and _NUMRE.match(tok[-1]):
                nums.insert(0, tok.pop())
            if len(nums) < 6 or len(tok) < 3:
                continue
            yr, scope, sub = tok[0], tok[1], tok[2]
            if yr != "25":
                continue
            計劃, 調整前, 調整後 = _tonum(nums[0]), _tonum(nums[1]), _tonum(nums[3])
            設施, 活動 = _tonum(nums[-2]), _tonum(nums[-1])
            recs.append((scope, sub, 1, 計劃, 調整前, 調整後, 設施, 活動))
        src = f"results/project_dump.tsv（真數，{len(recs)} 個 25年項目）"
    else:
        recs = FAKE
        src = "假數（fallback；run inspect_project.py 出真數）"
    agg = {}
    for scope, sub, n, pl, rep, aft, fac, act in recs:
        a = agg.setdefault((scope, sub), [0, 0, 0, 0, 0, 0])
        a[0] += n; a[1] += pl; a[2] += rep; a[3] += aft; a[4] += fac; a[5] += act
    return agg, src


def _rows(agg):
    """agg → 顯示 ROWS：[(label, kind, values[8] or None)]。"""
    def line(scope, subs, label_tot):
        out = []
        tot = [0, 0, 0, 0, 0, 0]
        for sub in subs:
            a = agg.get((scope, sub))
            if not a:
                continue
            n, pl, rep, aft, fac, act = a
            out.append((sub, "data", [str(n), _fmt(pl), _fmt(rep), _rate(rep, pl),
                                      _fmt(aft), _rate(aft, pl), _fmt(fac), _fmt(act)]))
            for i in range(6):
                tot[i] += a[i]
        n, pl, rep, aft, fac, act = tot
        sub_row = (label_tot, "subtot", [str(n), _fmt(pl), _fmt(rep), _rate(rep, pl),
                                         _fmt(aft), _rate(aft, pl), _fmt(fac), _fmt(act)])
        return out, sub_row, tot

    rows = [("博彩項目", "sec", None)]
    g_rows, g_sub, g_tot = line("gaming", GAMING_SUBS, "博彩項目小計")
    rows += g_rows + [g_sub]
    rows.append(("非博彩項目", "sec", None))
    ng_rows, ng_sub, ng_tot = line("non_gaming", NONG_ORDER, "非博彩項目小計")
    rows += ng_rows + [ng_sub]
    grand = [g_tot[i] + ng_tot[i] for i in range(6)]
    n, pl, rep, aft, fac, act = grand
    rows.append(("總計", "tot", [str(n), _fmt(pl), _fmt(rep), _rate(rep, pl),
                                 _fmt(aft), _rate(aft, pl), _fmt(fac), _fmt(act)]))
    # 簡單敘述（由數計，非機密）
    cats = [(s, agg[(sc, s)]) for sc in ("gaming", "non_gaming") for s in (GAMING_SUBS if sc == "gaming" else NONG_ORDER) if (sc, s) in agg]
    hi = [f"{s}（{a[2]/a[1]*100:.1f}%）" for s, a in cats if a[1] and a[2] / a[1] >= 1.0]
    lo = [f"{s}（{a[2]/a[1]*100:.1f}%）" for s, a in cats if a[1] and a[2] / a[1] < 0.5]
    bullets = [
        ("整體概況：", f"博彩項目完成率 {_rate(g_tot[2], g_tot[1])}，非博彩項目 {_rate(ng_tot[2], ng_tot[1])}，"
                     f"整體 {_rate(grand[2], grand[1])}。"),
        ("潛在調整後：", f"博彩項目 {_rate(g_tot[3], g_tot[1])}，非博彩項目平均 {_rate(ng_tot[3], ng_tot[1])}。"),
        ("完成率較高的範疇：", "、".join(hi) + "。" if hi else "—"),
        ("完成率相對較低的範疇：", "、".join(lo) + "。" if lo else "—"),
    ]
    return rows, bullets


def load_projects(year="25"):
    """讀 project_dump 逐項目（year）：(scope,範疇,序號,名稱,計劃,報告,調整,調整後,設施,活動)。"""
    p = Path("results/project_dump.tsv")
    if not p.exists():
        return []
    out = []
    for ln in p.read_text(encoding="utf-8").splitlines()[1:]:
        tok = ln.split()
        if len(tok) < 9:
            continue
        nums = []
        while tok and _NUMRE.match(tok[-1]):
            nums.insert(0, tok.pop())
        if len(nums) < 6 or len(tok) < 4 or tok[0] != year:
            continue
        out.append((tok[1], tok[2], tok[3], " ".join(tok[4:]),
                    _tonum(nums[0]), _tonum(nums[1]), _tonum(nums[2]), _tonum(nums[3]),
                    _tonum(nums[-2]), _tonum(nums[-1])))
    return out


# 單項審查 18 欄
RV_SUPER = [("項目基本信息", 0, 5), ("投資金額的潛在調整事項", 5, 14), ("潛在調整後投資金額", 14, 18)]
RV_SUB = ["項目\n序號", "項目名稱", "計劃\n投資金額", "報告\n投資金額", "投資計劃\n完成率",
          "一般支持\n人工成本", "其他日常\n營運調整", "超出可計入\n內部資源", "酒店客房\n改造", "不符吸引\n外國客源",
          "未完全\n實現目的", "獲批前\n未認可", "〔跨年〕", "潛在調整\n合計", "調整後\n投資金額",
          "潛在調整後\n完成率", "設施建設/\n資本性", "活動舉辦/\n營運性"]
RV_W = [0.42, 2.0, 0.6, 0.55, 0.5] + [0.5] * 9 + [0.55, 0.5, 0.55, 0.55]


def _footer(slide, page):
    W = 10.83
    ft = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(7), Inches(0.25))
    fr = ft.text_frame.paragraphs[0].add_run()
    fr.text = "KPMG　© 2026畢馬威會計師事務所 — 澳門特別行政區合夥制事務所。版權所有，不得轉載。"
    fr.font.size = Pt(6); fr.font.color.rgb = GREY; fr.font.name = FONT_CN
    pg = slide.shapes.add_textbox(Inches(W - 1.1), Inches(7.15), Inches(0.9), Inches(0.25))
    pr = pg.text_frame.paragraphs[0].add_run(); pr.text = f"初稿　{page}"
    pr.font.size = Pt(7); pr.font.bold = True; pr.font.color.rgb = NAVY; pr.font.name = FONT_CN


def _breadcrumb(slide, active=0):
    W = 10.83
    tabs = ["2025年度投資計劃執行情況概述", "過往年度…", "主要發現", "其他信息", "六項KPI分析", "附件"]
    bc = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(W - 1.5), Inches(0.2))
    bp = bc.text_frame.paragraphs[0]
    for i, t in enumerate(tabs):
        if i:
            sep = bp.add_run(); sep.text = "   |   "
            sep.font.size = Pt(6); sep.font.color.rgb = RGBColor(0xC8, 0xC8, 0xC8); sep.font.name = FONT_CN
        r = bp.add_run(); r.text = t
        r.font.size = Pt(6); r.font.name = FONT_CN
        r.font.bold = (i == active); r.font.color.rgb = NAVY if i == active else GREY
    mg = slide.shapes.add_textbox(Inches(W - 1.0), Inches(0.12), Inches(0.8), Inches(0.2))
    mr = mg.text_frame.paragraphs[0].add_run(); mr.text = "MGM  ◀ ⌂ ▶"
    mr.font.size = Pt(6); mr.font.color.rgb = GREY; mr.font.name = FONT_CN


def _dark(prs):
    from pptx.enum.shapes import MSO_SHAPE
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = DARK; rect.line.fill.background()
    kb = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(2.2), Inches(0.4))
    kr = kb.text_frame.paragraphs[0].add_run(); kr.text = "KPMG"
    kr.font.size = Pt(20); kr.font.bold = True; kr.font.color.rgb = WHITE; kr.font.name = FONT_HEAD
    return slide


def _cover(prs):
    slide = _dark(prs)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(7.2), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(["美高梅金殿超濠股份有限公司", "2025年年度投資計劃執行情況審查", "專項工作報告"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(25); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT_HEAD
    for y, txt, sz, bd in [(4.7, "初稿", 16, True), (5.5, "畢馬威會計師事務所", 11, False), (5.85, "2026年6月30日", 11, False)]:
        db = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(5), Inches(0.4))
        dr = db.text_frame.paragraphs[0].add_run(); dr.text = txt
        dr.font.size = Pt(sz); dr.font.bold = bd; dr.font.color.rgb = WHITE; dr.font.name = FONT_CN


def _divider(prs, num, title, subs):
    slide = _dark(prs)
    W, H, ny = 10.83, 7.5, 7.5 * 0.30
    nb = slide.shapes.add_textbox(Inches(0.7), Inches(ny - 0.15), Inches(1.3), Inches(1.3))
    nr = nb.text_frame.paragraphs[0].add_run(); nr.text = f"{num}."
    nr.font.size = Pt(44); nr.font.bold = True; nr.font.color.rgb = WHITE; nr.font.name = FONT_HEAD
    tb = slide.shapes.add_textbox(Inches(2.05), Inches(ny), Inches(W - 2.75), Inches(1.4))
    tb.text_frame.word_wrap = True
    tr = tb.text_frame.paragraphs[0].add_run(); tr.text = title
    tr.font.size = Pt(26); tr.font.bold = True; tr.font.color.rgb = WHITE; tr.font.name = FONT_HEAD
    y = ny + 1.7
    for s in subs:
        rb = slide.shapes.add_textbox(Inches(2.05), Inches(y), Inches(W - 3.0), Inches(0.32))
        rr = rb.text_frame.paragraphs[0].add_run(); rr.text = s
        rr.font.size = Pt(12); rr.font.color.rgb = RGBColor(0xC8, 0xC8, 0xD0); rr.font.name = FONT_CN
        y += 0.42


def _review_pages(prs, year, page_start):
    """單項審查 18 欄，全部範疇，自動分頁（每頁 ~26 行），模仿 scan slide 35-39。回 next page。"""
    projs = load_projects(year)
    W = 10.83
    order = [("gaming", s) for s in GAMING_SUBS] + [("non_gaming", s) for s in NONG_ORDER]
    # 逐 NG（範疇）砌 block（sec header + data 行 + 小計）
    blocks, gtot = [], [0.0] * 6
    for scope, sub in order:
        ps = [p for p in projs if p[0] == scope and p[1] == sub]
        if not ps:
            continue
        blk = [("sec", ("博彩項目" if scope == "gaming" else "非博彩項目") + "—" + sub, None)]
        st_ = [0.0] * 6
        for p in sorted(ps, key=lambda x: x[2]):
            _, _, code, name, pl, rep, adj, aft, fac, act = p
            blk.append(("data", (code, name), [_fmt(pl), _fmt(rep), _rate(rep, pl)] + ["-"] * 8 +
                        [_fmt(adj), _fmt(aft), _rate(aft, pl), _fmt(fac), _fmt(act)]))
            for i, x in enumerate([pl, rep, adj, aft, fac, act]):
                st_[i] += x
        pl, rep, adj, aft, fac, act = st_
        blk.append(("subtot", sub + " 小計", [_fmt(pl), _fmt(rep), _rate(rep, pl)] + ["-"] * 8 +
                    [_fmt(adj), _fmt(aft), _rate(aft, pl), _fmt(fac), _fmt(act)]))
        blocks.append(blk)
        for i in range(6):
            gtot[i] += st_[i]
    pl, rep, adj, aft, fac, act = gtot
    blocks.append([("tot", "總計", [_fmt(pl), _fmt(rep), _rate(rep, pl)] + ["-"] * 8 +
                   [_fmt(adj), _fmt(aft), _rate(aft, pl), _fmt(fac), _fmt(act)])])

    # 分頁 = 完全 by NG（範疇）：每個範疇【完整】放同一頁、絕不拆（對 scan —— 每個 block 都有自己小計、
    # 全份報告冇「續」）。一頁盡量塞多幾個完整範疇；下個範疇放唔落就開新頁。
    # 只有單一範疇本身大過一頁先逼住切（MGM 各範疇 ≤21 行 < CAP，實際唔會 hit）。
    CAP = 25
    chunks, cur = [], []
    for blk in blocks:
        if cur and len(cur) + len(blk) > CAP:            # 下個範疇放唔落 → 開新頁（唔拆）
            chunks.append(cur); cur = []
        cur.extend(blk)
        while len(cur) > CAP:                            # 安全網：單一範疇 > 一頁先切（MGM 唔會 hit）
            chunks.append(cur[:CAP]); cur = list(cur[CAP:])
    if cur:
        chunks.append(cur)
    if not chunks:
        chunks = [[]]
    tbl_name = f"MGM {year}年度投資計劃單個項目截至2025年末的審查結果匯總表"
    page = page_start
    for ci, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _breadcrumb(slide, 3)
        # ① 小 subtitle（灰）② 大 navy 標題 —— 對 scan slide 51 頂部
        sub = slide.shapes.add_textbox(Inches(0.3), Inches(0.30), Inches(W - 0.6), Inches(0.18))
        subr = sub.text_frame.paragraphs[0].add_run()
        subr.text = "其他信息  |  單個項目審查結果匯總"
        subr.font.size = Pt(8.5); subr.font.color.rgb = RGBColor(0x59, 0x59, 0x59); subr.font.name = FONT_CN
        tt = slide.shapes.add_textbox(Inches(0.3), Inches(0.46), Inches(W - 0.6), Inches(0.30))
        ttr = tt.text_frame.paragraphs[0].add_run()
        ttr.text = f"{tbl_name}（{ci + 1}/{len(chunks)}）"
        ttr.font.size = Pt(13); ttr.font.bold = True; ttr.font.color.rgb = NAVY; ttr.font.name = FONT_HEAD
        # ③ 表頂 navy caption bar（重覆表名，貼 scan）
        from pptx.enum.shapes import MSO_SHAPE
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(0.82),
                                     Inches(W - 0.4), Inches(0.17))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
        bar.shadow.inherit = False
        btf = bar.text_frame; btf.margin_top = btf.margin_bottom = Emu(0)
        btf.margin_left = Emu(36000); btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bpr = btf.paragraphs[0]; bpr.alignment = PP_ALIGN.LEFT
        brun = bpr.add_run(); brun.text = tbl_name
        brun.font.size = Pt(6); brun.font.bold = True; brun.font.color.rgb = WHITE; brun.font.name = FONT_CN
        ncol = 18
        nr = 2 + len(chunk)
        gt = slide.shapes.add_table(nr, ncol, Inches(0.2), Inches(0.99),
                                    Inches(W - 0.4), Inches(min(5.5, nr * 0.2))).table
        gt.first_row = False; gt.horz_banding = False
        tw = sum(RV_W)
        for i, wd in enumerate(RV_W):
            gt.columns[i].width = Inches(wd * (W - 0.4) / tw)
        for c in range(ncol):
            _set_cell(gt.cell(0, c), "", size=4.5, fill=NAVY, white=True)
        for name, a, b in RV_SUPER:
            gt.cell(0, a).merge(gt.cell(0, b - 1))
            _set_cell(gt.cell(0, a), name, size=5, bold=True, fill=NAVY, white=True, align=PP_ALIGN.CENTER)
        for c in range(ncol):
            _set_cell(gt.cell(1, c), RV_SUB[c], size=4.3, bold=True, fill=NAVY, white=True,
                      align=PP_ALIGN.LEFT if c == 1 else PP_ALIGN.CENTER)
        for ri, (kind, label, vals) in enumerate(chunk, start=2):
            fill = {"sec": SECHDR, "subtot": SUBTOT, "tot": TOT}.get(kind)
            bold = kind in ("sec", "subtot", "tot")
            if kind == "data":
                code, name = label
                _set_cell(gt.cell(ri, 0), code, size=4.3, align=PP_ALIGN.LEFT)
                _set_cell(gt.cell(ri, 1), name, size=4.3, align=PP_ALIGN.LEFT)
                for c, v in enumerate(vals, start=2):
                    _set_cell(gt.cell(ri, c), v, size=4.3)
            else:
                _set_cell(gt.cell(ri, 0), "", fill=fill)
                _set_cell(gt.cell(ri, 1), label, size=4.5, bold=bold, fill=fill, align=PP_ALIGN.LEFT)
                for c in range(2, ncol):
                    _set_cell(gt.cell(ri, c), (vals[c - 2] if vals else ""), size=4.3, bold=bold, fill=fill)
        _thin_borders(gt)
        for rr in gt.rows:
            rr.height = Emu(150000)      # ~0.164in compact 行高（貼 scan 密度）
        # ④ 表下：資料來源（左）+（下頁待續）（右，最後一頁除外）—— 對 scan
        # note 固定放頁底（footer 7.15 之上）—— 唔跟表估算高度（項目名 wrap 會令表比 requested 高）
        ty = 6.92
        src = slide.shapes.add_textbox(Inches(0.2), Inches(ty), Inches(W - 2.0), Inches(0.16))
        srcr = src.text_frame.paragraphs[0].add_run()
        srcr.text = "資料來源：管理層提供之項目投資計劃及執行報告資料，畢馬威分析"
        srcr.font.size = Pt(6.5); srcr.font.color.rgb = RGBColor(0x59, 0x59, 0x59); srcr.font.name = FONT_CN
        if ci < len(chunks) - 1:
            con = slide.shapes.add_textbox(Inches(W - 1.7), Inches(ty), Inches(1.5), Inches(0.16))
            conp = con.text_frame.paragraphs[0]; conp.alignment = PP_ALIGN.RIGHT
            conr = conp.add_run(); conr.text = "（下頁待續）"
            conr.font.size = Pt(6.5); conr.font.color.rgb = RGBColor(0x59, 0x59, 0x59); conr.font.name = FONT_CN
        _footer(slide, page); page += 1
    return page


def _set_cell(cell, text, *, size=5.5, bold=False, fill=None, align=PP_ALIGN.RIGHT, white=False):
    cell.margin_left = cell.margin_right = Emu(18000)
    cell.margin_top = cell.margin_bottom = Emu(2500)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid(); cell.fill.fore_color.rgb = fill if fill is not None else WHITE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold
    r.font.name = FONT_CN if any("一" <= c <= "鿿" for c in str(text)) else FONT_NUM
    r.font.color.rgb = WHITE if white else RGBColor(0x22, 0x22, 0x22)


def _thin_borders(table):
    for row in table.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
                ln = tcPr.makeelement(qn(tag), {"w": "3175", "cap": "flat"})
                fl = ln.makeelement(qn("a:solidFill"), {})
                cl = fl.makeelement(qn("a:srgbClr"), {"val": "BFBFBF"})
                fl.append(cl); ln.append(fl); tcPr.append(ln)


def _overview_slide(prs, page):
    agg, src = load_agg()
    ROWS, BULLETS = _rows(agg)
    W = 10.83
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _breadcrumb(slide, 0)

    # 真敘述（gitignored results/demo_narrative.txt）
    headline_txt, bullets_real = None, BULLETS
    _np = Path("results/demo_narrative.txt")
    if _np.exists():
        _nl = [x for x in _np.read_text(encoding="utf-8").splitlines() if x.strip()]
        if _nl:
            headline_txt = _nl[0]
            bullets_real = [(x.split("|||", 1) if "|||" in x else ("", x)) for x in _nl[1:]]
    if headline_txt is None:
        headline_txt = "MGM 2025年度計劃投資項目涵蓋博彩及非博彩範疇（示範；真 headline 由 build 填）。"

    # title（subtitle line）0.53,0.48 —— KPMG Bold navy
    st = slide.shapes.add_textbox(Inches(0.53), Inches(0.44), Inches(9.76), Inches(0.4))
    sr = st.text_frame.paragraphs[0].add_run()
    sr.text = "2025年度投資計劃執行情況概述  |  2025年度投資項目的整體執行概況"
    sr.font.size = Pt(12); sr.font.bold = True; sr.font.color.rgb = NAVY; sr.font.name = FONT_HEAD

    # headline（strapline）navy 段
    hl = slide.shapes.add_textbox(Inches(0.53), Inches(1.02), Inches(9.76), Inches(1.0))
    hl.text_frame.word_wrap = True
    hr = hl.text_frame.paragraphs[0].add_run(); hr.text = headline_txt
    hr.font.size = Pt(7.5); hr.font.color.rgb = NAVY; hr.font.name = FONT_CN

    # 左 table title + 右 narrative subtitle（收緊：貼近 headline 下）
    tt = slide.shapes.add_textbox(Inches(0.53), Inches(1.72), Inches(6.46), Inches(0.2))
    ttr = tt.text_frame.paragraphs[0].add_run(); ttr.text = "MGM 2025年度計劃的整體投資執行概況"
    ttr.font.size = Pt(7.5); ttr.font.bold = True; ttr.font.color.rgb = NAVY; ttr.font.name = FONT_CN
    ns = slide.shapes.add_textbox(Inches(7.1), Inches(1.72), Inches(3.19), Inches(0.2))
    nsr = ns.text_frame.paragraphs[0].add_run(); nsr.text = "2025年度投資項目的整體執行概況"
    nsr.font.size = Pt(7.5); nsr.font.bold = True; nsr.font.color.rgb = NAVY; nsr.font.name = FONT_CN

    # 表（收緊 y=1.94）
    ncol = 9
    gt = slide.shapes.add_table(2 + len(ROWS), ncol, Inches(0.53), Inches(1.94), Inches(6.46), Inches(4.5)).table
    gt.first_row = False; gt.horz_banding = False
    widths = [1.35, 0.5, 0.72, 0.7, 0.62, 0.72, 0.62, 0.6, 0.6]
    tot_w = sum(widths)
    for i, wd in enumerate(widths):
        gt.columns[i].width = Inches(wd * 6.46 / tot_w)
    super_hdr = ["", "", "", "報告投資金額", "", "潛在調整後投資金額", "", "", ""]
    for c in range(ncol):
        _set_cell(gt.cell(0, c), super_hdr[c], size=5.5, bold=True, fill=NAVY, white=True, align=PP_ALIGN.CENTER)
    gt.cell(0, 3).merge(gt.cell(0, 4))
    gt.cell(0, 5).merge(gt.cell(0, 6))
    for c in range(ncol):
        _set_cell(gt.cell(1, c), SUB[c], size=5.5, bold=True, fill=NAVY, white=True,
                  align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    for ri, (label, kind, vals) in enumerate(ROWS, start=2):
        fill = {"sec": SECHDR, "subtot": SUBTOT, "tot": TOT}.get(kind)
        bold = kind in ("subtot", "tot", "sec")
        _set_cell(gt.cell(ri, 0), label, size=5.5, bold=bold, fill=fill, align=PP_ALIGN.LEFT)
        for c in range(1, ncol):
            _set_cell(gt.cell(ri, c), (vals[c - 1] if vals else ""), size=5.5, bold=bold, fill=fill)
    _thin_borders(gt)

    nb2 = slide.shapes.add_textbox(Inches(0.53), Inches(6.62), Inches(6.46), Inches(0.35))
    b = nb2.text_frame.paragraphs[0].add_run()
    b.text = "註：投資計劃完成率 ＝ 報告投資金額 ／ 獲批的計劃投資金額；潛在調整後完成率 ＝ 潛在調整後投資金額 ／ 獲批的計劃投資金額。"
    b.font.size = Pt(5); b.font.italic = True; b.font.color.rgb = GREY; b.font.name = FONT_CN

    # 右敘述（精確位 7.1,3.19）密實
    nb = slide.shapes.add_textbox(Inches(7.1), Inches(1.94), Inches(3.19), Inches(4.9))
    ntf = nb.text_frame; ntf.word_wrap = True
    for j, (lead, body) in enumerate(bullets_real):
        p = ntf.paragraphs[0] if j == 0 else ntf.add_paragraph()
        p.space_after = Pt(3)
        rb = p.add_run(); rb.text = "■ "
        rb.font.size = Pt(6.5); rb.font.color.rgb = NAVY; rb.font.name = FONT_CN
        if lead:
            rl = p.add_run(); rl.text = lead
            rl.font.size = Pt(6.5); rl.font.bold = True; rl.font.color.rgb = NAVY; rl.font.name = FONT_CN
        rt = p.add_run(); rt.text = body
        rt.font.size = Pt(6.5); rt.font.color.rgb = RGBColor(0x33, 0x33, 0x33); rt.font.name = FONT_CN

    _footer(slide, page)


def build():
    print(f"  數據來源：{load_agg()[1]}")
    prs = Presentation()
    prs.slide_width = Inches(10.83); prs.slide_height = Inches(7.5)
    p = 1
    _cover(prs); p += 1                                          # 封面
    _divider(prs, 1, "2025年度投資計劃執行情況概述", [           # 章節分隔
        "1.1  股權架構簡圖及發生投資支出的主體公司",
        "1.2  2025年度計劃的整體投資支出概況",
        "1.3  2025年度投資項目的整體執行概況",
        "1.4  2025年度投資計劃報告投資金額的潛在調整事項匯總"]); p += 1
    _overview_slide(prs, p); p += 1                             # 概述 整體執行概況
    _divider(prs, 4, "其他信息", [
        "4.1  2025年度投資計劃及過往年度期後投資於2025年發生的投資金額匯總",
        "4.2  單個項目審查結果匯總"]); p += 1
    for yr in ("25", "24", "23"):                               # 單項審查 25/24/23（多頁）
        p = _review_pages(prs, yr, p)
    prs.save("demo_overview.pptx")
    print(f"✓ demo_overview.pptx（{len(list(prs.slides))} 版：封面 + 分隔 + 概述 + 單項審查25/24/23；PowerPoint 開睇）")


if __name__ == "__main__":
    build()
