#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
check_numbers.py — 數字 tie 唔 tie 檢查：掃 folder 入面全部 pptx，逐張表自己驗自己。

點用（同 check_text 一樣，掉檔入 folder 就得）：
    python scripts\\report\\check_numbers.py            # 掃 file_check\\ → 出 md + docx
    python scripts\\report\\check_numbers.py --dir 某資料夾
    python scripts\\report\\check_numbers.py --no-docx  # 淨係出 md

出：file_check\\_數字檢查\\{檔名}.md ＋ .docx ＋ console 總結表

★ 全部係【機械計算】，冇 LLM —— 數字唔可以靠 LLM 加。

三種檢查：
  ① 橫向（跟表頭下面嗰行公式）：報告啲表本身印住 a｜1..8｜b｜c=a+b｜d=b/a，
     直接攞嚟逐行驗。4.2 嗰啲 a¹ b¹ c¹ d¹=b¹+c¹ 一樣識。
  ② 直向：逐個範疇加埋 =唔=小計；各小計加埋 =唔= 合計／總計。
  ③ 跨頁：全份文件所有「合計／總計」行抽出嚟，同一個欄名喺唔同版出唔同數就標出嚟
     （例如 1.2 總計報告投資金額 應該 = 4.1 2025年度報告投資金額）。

容差：每格已經 round 到整數萬，所以 N 項加埋最多可以爭 N×0.5 →
      tol = max(1, ceil(0.5×N))，另外加 0.1% 相對容差（大數用）。
"""
from __future__ import annotations

import argparse
import hashlib
import json
import math
import re
import sys
from collections import defaultdict
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

_HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(_HERE))
for _p in (_HERE.parent.parent, _HERE.parent.parent / "src"):      # repo root / src
    sys.path.insert(0, str(_p))
from check_text import _Ticker                                     # tqdm 每秒 refresh

try:
    from tqdm import tqdm
except ImportError:
    def tqdm(it=None, **kw):
        return it
    tqdm.write = print

DEFAULT_DIR = "file_check"
OUT_SUB = "_數字檢查"
CACHE = ".num_cache.json"

# ★ 分工死線：LLM 淨係做【閱讀同配對】—— 邊句敘述講緊表入面邊個數。
#   加減乘除、比大細、換算單位，全部 Python 做。LLM 永遠唔會計數。
SYS_CLAIM = (
    "你係財務報告校對員。俾你一版嘅【敘述文字】同埋嗰版表格入面嘅【事實清單】（已編號）。"
    "任務：揾出敘述文字入面每一個【引用緊表格數字】嘅講法，指出佢對應邊個事實編號，"
    "同埋佢喺文字入面聲稱嘅數值（照抄，唔好自己計）。\n"
    "★ 唔好做算術、唔好推算、唔好改寫數字 —— 你只係負責配對。\n"
    "★ 配唔到任何事實嘅講法唔好報。純粹描述性、冇數字嘅句子唔好報。\n"
    "★ 單位照抄原文（億／萬／個／%）。\n"
    '回覆純 JSON：{"claims":[{"quote":"<原文片段,20字內>","fact":<事實編號int>,'
    '"value":<數字>,"unit":"億|萬|個|%"}]}　冇就 {"claims":[]}'
)
SYS_PAIR = (
    "你係財務報告校對員。俾你一份報告入面全部【合計／總計】嘅事實清單（已編號，"
    "每項有頁數、表名、欄名、數值）。任務：憑報告結構判斷邊幾對【應該完全相等】。\n"
    "例：同一年度嘅『概覽表合計』同『調整事項匯總表合計』嘅報告投資金額應該相等；"
    "『整體投資支出概況』嘅總計應該等於『金額匯總』入面同一年度嘅欄組。\n"
    "★ 唔好比較數值、唔好計差額 —— 你只係指出邊兩項【定義上應該相等】。\n"
    "★ 唔同年度、唔同口徑（報告 vs 調整後 vs 計劃）唔好配埋一齊。冇把握就唔好報。\n"
    '回覆純 JSON：{"pairs":[{"a":<編號>,"b":<編號>,"why":"<15字內理由>"}]}'
)
_UNIT = {"億": 10000.0, "萬": 1.0, "個": 1.0, "%": 1.0}

_NUM = re.compile(r"^\(?\s*-?[\d,]+(?:\.\d+)?\s*\)?%?$")
_FORMULA = re.compile(r"^[a-zA-Z][¹²³]?(?:\s*=\s*[a-zA-Z¹²³\d+/*\-\s]+)?$")
_TOTAL_EXACT = {"合計", "總計"}          # 淨係「合計」／「總計」先算全表總數
_SUBTOT_SUF = ("小計", "合計")            # 「非博彩項目合計」呢類 = 小計
_PAGED = re.compile(r"（\s*(\d+)\s*/\s*(\d+)\s*）")   # caption 有（3/3）＝ 拆咗頁嘅表
_SKIP_ROW = ("涉及項目數量", "承諾的10年投資預算")     # 唔係加總關係嘅行


def parse_num(s):
    """'(1,234)'→-1234｜'1,234'→1234｜'12.3%'→0.123｜'-'／空→None。"""
    t = str(s or "").strip().replace(" ", "").replace(" ", "")
    if not t or t in {"-", "—", "–", "n/a", "N/A"}:
        return None
    neg = t.startswith("(") and t.endswith(")")
    t = t.strip("()")
    pct = t.endswith("%")
    t = t.rstrip("%").replace(",", "")
    try:
        v = float(t)
    except ValueError:
        return None
    if pct:
        v /= 100.0
    return -v if neg else v


def tol_of(vals):
    n = max(1, len(vals))
    return max(1.0, math.ceil(0.5 * n), abs(sum(v for v in vals if v is not None)) * 0.001)


def grid_of(tbl):
    """pptx table → (rows[list[str]], 合併咗嘅 cell 唔會重複)。"""
    out = []
    for r in tbl.rows:
        out.append([(c.text or "").strip().replace("\n", " ") for c in r.cells])
    return out


def split_head(rows):
    """→ (表頭行數, 公式行 index 或 None, letter→col)。"""
    hdr, fml, letters = 0, None, {}
    for ri, row in enumerate(rows[:5]):
        cells = [c for c in row if c]
        if not cells:
            continue
        nums = sum(1 for c in cells if _NUM.match(c))
        # ⚠ 公式行入面「1」「2」…「7」（調整類編號）同時 match 得 _NUM，
        #   所以唔可以用 nums 嚟否決 —— 要睇有冇【字母】公式（a／b／c=a+b／d¹=b¹+c¹）。
        alpha = [c for c in cells if re.match(r"^[a-zA-Z][¹²³]?(\s*=.+)?$", c)]
        if len(alpha) >= 2:
            fml = ri
            for ci, c in enumerate(row):
                c = c.strip()
                if not c:
                    continue
                key = c.split("=")[0].strip()
                if key and key not in letters:
                    letters[key] = ci
            hdr = ri + 1
            break               # 公式行永遠係表頭最後一行 → 之後全部係 body
        if nums == 0:
            hdr = ri + 1
        else:
            break
    return hdr, fml, letters


def _eval(expr, letters, row):
    """'b¹+c¹' / 'a+b' / 'b/a' → 數值（用該行嘅格）；有一項冇數就回 None。"""
    toks = re.split(r"([+\-/])", expr.replace(" ", ""))
    vals, ops = [], []
    for i, t in enumerate(toks):
        if i % 2:
            ops.append(t); continue
        if t in letters:
            v = parse_num(row[letters[t]]) if letters[t] < len(row) else None
        elif re.fullmatch(r"\d+(?:\.\d+)?", t):
            v = float(t)
        else:
            return None
        if v is None:
            return None
        vals.append(v)
    r = vals[0]
    for op, v in zip(ops, vals[1:]):
        if op == "+":
            r += v
        elif op == "-":
            r -= v
        elif op == "/":
            if abs(v) < 1e-9:
                return None
            r /= v
    return r


def check_across(rows, hdr, fml, letters, subs):
    """① 橫向：跟公式行逐行驗。"""
    out = []
    if fml is None or not letters:
        return out
    exprs = [(ci, rows[fml][ci].split("=", 1)[0].strip(), rows[fml][ci].split("=", 1)[1].strip())
             for ci in range(len(rows[fml])) if "=" in rows[fml][ci]]
    for ri in range(hdr, len(rows)):
        row = rows[ri]
        lab = row_label(row)
        if not lab or lab in _SKIP_ROW:
            continue
        for ci, lhs, expr in exprs:
            if ci >= len(row):
                continue
            got = parse_num(row[ci])
            want = _eval(expr, letters, row)
            if got is None or want is None:
                continue
            pct = "/" in expr
            t = 0.0015 if pct else tol_of([got, want])      # 率：容 0.15 個百分點
            if abs(got - want) > t:
                out.append({"row": lab, "col": subs[ci] if ci < len(subs) else f"c{ci}",
                            "rule": f"{lhs}={expr}", "got": got, "want": want,
                            "kind": "橫向", "pct": pct})
    return out


def row_label(row):
    for c in row[:3]:
        c = (c or "").strip()
        if c and not _NUM.match(c):
            return c
    return ""


def check_down(rows, hdr, fml, subs, paged=False):
    """② 直向：範疇加埋 = 小計；小計加埋 = 合計／總計。
    paged=True（表拆咗幾版，例如單項審查「（3/3）」）→ 唔驗全表合計 ——
    嗰行喺最後一版，但明細行散喺前面幾版，喺呢一版加梗唔夠數（會全部誤報）。"""
    out = []
    body = [(ri, rows[ri]) for ri in range(hdr, len(rows))]
    ncol = max((len(r) for _i, r in body), default=0)
    bucket, subtots = [], []
    for ri, row in body:
        lab = row_label(row)
        if not lab or lab in _SKIP_ROW:
            continue
        if ri == fml:
            continue
        vals = [parse_num(row[c]) if c < len(row) else None for c in range(ncol)]
        if all(v is None for v in vals):          # section 標題行（博彩項目／非博彩項目）
            bucket = []; continue
        if lab.strip() in _TOTAL_EXACT:
            if not paged:
                out += _cmp(subtots or bucket, vals, lab, subs, ncol)
            continue
        if lab.endswith(_SUBTOT_SUF):
            out += _cmp(bucket, vals, lab, subs, ncol)
            subtots.append(vals); bucket = []; continue
        bucket.append(vals)
    return out


def _cmp(parts, tot, lab, subs, ncol):
    out = []
    if not parts:
        return out
    for c in range(ncol):
        col = subs[c] if c < len(subs) else f"c{c}"
        if "率" in col or "比例" in col:            # 率唔可以直接加
            continue
        terms = [p[c] for p in parts if c < len(p) and p[c] is not None]
        got = tot[c] if c < len(tot) else None
        if got is None or len(terms) < 2:
            continue
        want = sum(terms)
        if abs(got - want) > tol_of(terms):
            out.append({"row": lab, "col": col, "rule": f"{len(terms)} 項加總",
                        "got": got, "want": want, "kind": "直向"})
    return out


def subs_of(rows, hdr, fml=None):
    """欄名 = 表頭最後一行（有 super 就同上面接埋）。
    ⚠ 公式行（a｜b｜c=a+b）都算表頭，但佢唔係欄名 → 要跳過佢再向上攞。"""
    end = fml if fml is not None else hdr      # 有公式行 → 欄名喺公式行【上面】嗰行
    if end == 0:
        return []
    last = rows[end - 1]
    if end >= 2:
        sup = rows[end - 2]
        return [(f"{(sup[i] or '').strip()}·{(last[i] or '').strip()}".strip("·")
                 if i < len(sup) else (last[i] or "").strip()) for i in range(len(last))]
    return [(c or "").strip() for c in last]


def scan_pptx(path):
    """→ (issues, totals, pages)；
    totals = [(slide, caption, 欄名, 值)]；pages = {版號: (敘述文字, [事實])}（餵 LLM 用）。"""
    from pptx import Presentation
    prs = Presentation(str(path))
    issues, totals, pages = [], [], {}
    for si, sl in enumerate(prs.slides, 1):
        cap = ""
        for sh in sl.shapes:
            if sh.has_text_frame and 0.5 < sh.top / 914400.0 < 2.2:
                t = (sh.text_frame.text or "").strip()
                if t and len(t) < 60 and not cap:
                    cap = t.replace("\n", " ")
        narr = " ".join(
            (sh.text_frame.text or "").strip().replace("\n", " ")
            for sh in sl.shapes
            if sh.has_text_frame and len((sh.text_frame.text or "").strip()) > 30)
        facts = []
        for sh in sl.shapes:
            if not getattr(sh, "has_table", False):
                continue
            rows = grid_of(sh.table)
            if len(rows) < 3:
                continue
            hdr, fml, letters = split_head(rows)
            subs = subs_of(rows, hdr, fml)
            paged = bool(_PAGED.search(cap)) and _PAGED.search(cap).group(2) != "1"
            for x in (check_across(rows, hdr, fml, letters, subs)
                      + check_down(rows, hdr, fml, subs, paged)):
                issues.append({**x, "page": si, "caption": cap})
            for ri in range(hdr, len(rows)):
                lab = row_label(rows[ri])
                if lab.strip() in _TOTAL_EXACT or lab.endswith(_SUBTOT_SUF):
                    grand = lab.strip() in _TOTAL_EXACT
                    for c in range(len(rows[ri])):
                        v = parse_num(rows[ri][c])
                        col = subs[c] if c < len(subs) else ""
                        if v is None or not col:
                            continue
                        rate = "率" in col or "比例" in col
                        if grand and not rate:
                            totals.append((si, cap, col, v))
                        facts.append({"page": si, "cap": cap, "row": lab, "col": col,
                                      "val": v, "rate": rate})
        if narr and facts:
            pages[si] = (narr[:1800], facts)
    return issues, totals, pages


# ── LLM 層（只做閱讀／配對，計數一律 Python）────────────────────────────
def _fact_lines(facts, with_page=False):
    return "\n".join(
        f"[{i}] " + (f"p{f['page']} " if with_page else "")
        + f"{f['cap'][:26]}｜{f['row']}｜{f['col']}＝"
        + (f"{f['val']*100:.1f}%" if f["rate"] else f"{f['val']:,.0f}")
        for i, f in enumerate(facts, 1))


def llm_claims(wb, doc, page, narr, facts, model=None):
    """敘述文字 vs 表：LLM 配對，Python 驗數。→ issues"""
    r = wb.chat_json(f"《{doc}》第 {page} 頁。\n\n【敘述文字】\n{narr}\n\n"
                     f"【事實清單】\n{_fact_lines(facts)}", system=SYS_CLAIM,
                     model=model, reasoning_effort="high")
    out = []
    for c in (r or {}).get("claims", []) or []:
        try:
            f = facts[int(c["fact"]) - 1]
            said = float(c["value"])
        except Exception:
            continue
        unit = str(c.get("unit", "萬"))
        if f["rate"]:
            said_w, real, tol, pct = said / 100.0, f["val"], 0.0015, True
        else:
            said_w = said * _UNIT.get(unit, 1.0)
            real, pct = f["val"], False
            # 億 只印到 1 位小數 → 容差要跟返嗰個精度（0.05億 = 500萬）
            tol = max(tol_of([real]), 500.0 if unit == "億" else 1.0)
        if abs(said_w - real) > tol:
            out.append({"page": page, "caption": f["cap"], "kind": "文表",
                        "row": f"「{str(c.get('quote', ''))[:22]}」", "col": f"{f['row']}·{f['col']}",
                        "rule": "敘述 vs 表", "got": said_w, "want": real, "pct": pct})
    return out


def llm_pairs(wb, doc, facts, model=None):
    """LLM 提出邊幾對合計【應該相等】→ Python 逐對比。→ checked"""
    r = wb.chat_json(f"《{doc}》全份報告嘅合計清單：\n{_fact_lines(facts, True)}",
                     system=SYS_PAIR, model=model, reasoning_effort="high")
    out = []
    for pr in (r or {}).get("pairs", []) or []:
        try:
            a_, b_ = facts[int(pr["a"]) - 1], facts[int(pr["b"]) - 1]
        except Exception:
            continue
        d = a_["val"] - b_["val"]
        out.append({"name": f"（LLM）{str(pr.get('why', ''))[:20]}",
                    "left": (a_["page"], f"{a_['cap'][:16]}·{a_['col']}", a_["val"]),
                    "right": (b_["page"], f"{b_['cap'][:16]}·{b_['col']}", b_["val"]),
                    "diff": d, "ok": abs(d) <= max(1.0, abs(a_["val"]) * 0.001)})
    return out


# ③ 跨表【應該相等】嘅關係（報告結構決定，唔係估）。每條 = (名, 左邊, 右邊)，
#    左右 = (caption 關鍵字, 欄名關鍵字)；喺全份文件搵到兩邊就比較。
CROSS_RULES = [
    ("1.2 總計 = 1.4 合計（報告投資金額）",
     ("整體投資支出概況", "報告投資金額"), ("報告投資金額潛在調整", "報告投資金額")),
    ("1.2 總計 = 1.4 合計（潛在調整後）",
     ("整體投資支出概況", "潛在調整後投資金額"), ("報告投資金額潛在調整", "潛在調整後投資金額")),
    ("2.1 合計 = 2.2 合計（報告投資金額）",
     ("2024年度計劃期後投資金額概覽", "報告投資金額"),
     ("2024年度投資計劃於2025年申報", "報告投資金額")),
    ("2.3 合計 = 2.4 合計（報告投資金額）",
     ("2023年度計劃期後投資金額概覽", "報告投資金額"),
     ("2023年度投資計劃於2025年申報", "報告投資金額")),
    ("1.2 總計 = 4.1「2025年度投資計劃」欄組（報告投資金額）",
     ("整體投資支出概況", "報告投資金額"), ("發生的投資金額匯總", "2025年度投資計劃·報告投資金額")),
]


def _pick(totals, cap_kw, col_kw):
    hit = [(s, c, v) for s, c, col, v in totals if cap_kw in c and col_kw == col]
    if not hit:
        hit = [(s, c, v) for s, c, col, v in totals if cap_kw in c and col_kw in col]
    return hit[0] if hit else None


def cross_page(totals):
    """→ (明確關係嘅結果, 參考附錄)。附錄唔算「錯」，淨係俾人 eyeball。"""
    checked = []
    for name, (lc, lk), (rc, rk) in CROSS_RULES:
        a_, b_ = _pick(totals, lc, lk), _pick(totals, rc, rk)
        if not a_ or not b_:
            continue
        d = a_[2] - b_[2]
        checked.append({"name": name, "left": a_, "right": b_, "diff": d,
                        "ok": abs(d) <= max(1.0, abs(a_[2]) * 0.001)})
    ref = defaultdict(list)
    for si, cap, col, v in totals:
        ref[col].append((si, cap, v))
    return checked, {k: v for k, v in sorted(ref.items()) if len(v) > 1}


def write_md(path, src, issues, checked, ref, npages):
    bad_x = [c for c in checked if not c["ok"]]
    L = [f"# {src.name} 數字 tie 檢查", "",
         f"- 來源：`{src.resolve()}`", f"- 頁數：{npages}",
         f"- 表內對唔上：**{len(issues)}** 處", 
         f"- 跨表關係：驗咗 {len(checked)} 條，**{len(bad_x)}** 條唔 tie", ""]
    if not issues:
        L += ["✅ 表內橫向／直向加總全部 tie。", ""]
    else:
        L += ["> 類型：**橫向**＝跟表頭公式行逐行驗｜**直向**＝小計／合計加總｜"
              "**文表**＝敘述文字講嘅數 vs 表上嘅數（LLM 負責配對，數值一律 Python 比）。", ""]
    cur = None
    for x in sorted(issues, key=lambda x: (x["page"], x["kind"], x["row"])):
        if x["page"] != cur:
            cur = x["page"]
            L += ["", f"## 第 {cur} 頁　{x['caption']}", "",
                  "| 類型 | 行 | 欄 | 規則 | 表上寫 | 應該係 | 差額 |",
                  "|---|---|---|---|---|---|---|"]
        d = x["got"] - x["want"]
        if x.get("pct"):
            g, w, dd = f"{x['got']*100:.1f}%", f"{x['want']*100:.1f}%", f"{d*100:+.1f}pp"
        else:
            g, w, dd = f"{x['got']:,.0f}", f"{x['want']:,.0f}", f"{d:+,.0f}"
        L.append(f"| {x['kind']} | {x['row']} | {x['col']} | `{x['rule']}` | "
                 f"{g} | {w} | **{dd}** |")
    if checked:
        L += ["", "## 跨表關係（報告結構決定，應該相等）", "",
              "| 關係 | 左 | 右 | 差額 | 結果 |", "|---|---|---|---|---|"]
        for c in checked:
            (ls, lc, lv), (rs, rc, rv) = c["left"], c["right"]
            L.append(f"| {c['name']} | p{ls} = {lv:,.0f} | p{rs} = {rv:,.0f} | "
                     f"{c['diff']:+,.0f} | {'✅ tie' if c['ok'] else '❌ 唔 tie'} |")
    if ref:
        L += ["", "## 附錄：全份文件嘅「合計／總計」一覽（參考，唔算錯）", "",
              "> 同名唔同數好正常（唔同年度／唔同口徑），要人 eyeball 有冇應該相等但唔等。", "",
              "| 欄名 | 各版數值 |", "|---|---|"]
        for col, hits in ref.items():
            L.append(f"| {col} | " + "；".join(f"p{s} {c[:16]} = {v:,.0f}"
                                               for s, c, v in hits) + " |")
    path.write_text("\n".join(L) + "\n", encoding="utf-8")


def main():
    ap = argparse.ArgumentParser(description="掃 folder 入面嘅 pptx，驗表內／跨頁數字 tie 唔 tie")
    ap.add_argument("--dir", default=DEFAULT_DIR, help=f"要掃嘅資料夾（預設 {DEFAULT_DIR}）")
    ap.add_argument("--no-docx", action="store_true", help="淨係出 md")
    ap.add_argument("--no-llm", action="store_true", help="淨係機械檢查，唔叫 API")
    ap.add_argument("--workers", type=int, default=4, help="LLM 並行數（同 build_report 一致）")
    ap.add_argument("--model", default=None)
    ap.add_argument("--fresh", action="store_true", help="唔用 cache")
    a = ap.parse_args()

    root = Path(a.dir)
    if not root.exists():
        root.mkdir(parents=True, exist_ok=True)
        print(f"✓ 開咗 {root.resolve()} —— 掉 pptx 入去再跑一次"); return
    files = sorted(p for p in root.rglob("*")
                   if p.suffix.lower() in {".pptx", ".potx", ".pptm"}
                   and not p.name.startswith("~$") and OUT_SUB not in p.parts)
    if not files:
        print(f"✗ {root.resolve()} 入面搵唔到 pptx"); return
    outdir = root / OUT_SUB
    outdir.mkdir(exist_ok=True)
    cachef = outdir / CACHE
    cache = {}
    if not a.fresh and cachef.exists():
        try:
            cache = json.loads(cachef.read_text(encoding="utf-8"))
        except Exception:
            cache = {}
    wb = None
    if not a.no_llm:
        try:
            from kpi.lib.workbench import Workbench
            wb = Workbench(model=a.model)
            if not wb.config_masked().get("key_ok"):
                print("  ⚠ 冇 workbench api key → 淨係跑機械檢查"); wb = None
            else:
                wb.ping()
        except Exception as e:
            print(f"  ⚠ LLM 連唔到（{str(e)[:110]}）→ 淨係跑機械檢查"); wb = None
    print(f"── 掃 {root.resolve()}：{len(files)} 個檔"
          f"｜機械檢查一定行｜LLM {'開' if wb else '關'}（只做閱讀配對，計數 Python 做）")

    # ── 第一浸：全部檔跑機械層 ──
    docs = []
    for f in tqdm(files, desc="機械對數", unit="檔", ncols=76):
        try:
            from pptx import Presentation
            npages = len(Presentation(str(f)).slides)
            issues, totals, pages = scan_pptx(f)
        except Exception as e:
            tqdm.write(f"  ✗ {f.name}：{str(e)[:140]}"); continue
        checked, ref = cross_page(totals)
        docs.append({"file": f, "npages": npages, "issues": issues,
                     "checked": checked, "ref": ref, "pages": pages,
                     "tot_facts": [{"page": s_, "cap": c_, "row": "合計", "col": col,
                                    "val": v, "rate": False} for s_, c_, col, v in totals]})

    # ── 第二浸：LLM（逐版「文表對照」+ 每份「應該相等」）—— 一條 task list、一條 bar ──
    if wb:
        tasks = []
        for di, d in enumerate(docs):
            doc = d["file"].stem
            for pi, (narr, facts) in d["pages"].items():
                k = f"claim|{doc}|{pi}|" + hashlib.sha1(
                    (narr + _fact_lines(facts)).encode("utf-8")).hexdigest()[:16]
                tasks.append(("claim", di, doc, pi, narr, facts, k))
            if d["tot_facts"]:
                k = f"pair|{doc}|" + hashlib.sha1(
                    _fact_lines(d["tot_facts"], True).encode("utf-8")).hexdigest()[:16]
                tasks.append(("pair", di, doc, 0, "", d["tot_facts"], k))
        todo = [t for t in tasks if t[-1] not in cache]
        for kind, di, _doc, pi, _n, _f, k in tasks:
            if k in cache:
                (docs[di]["issues"] if kind == "claim" else docs[di]["checked"]).extend(cache[k])
        if todo:
            bar = _Ticker(tqdm(total=len(todo), desc="LLM 對照", unit="項", ncols=80,
                               mininterval=0.5))
            with ThreadPoolExecutor(max_workers=a.workers) as ex:
                fut = {}
                for kind, di, doc, pi, narr, facts, k in todo:
                    fn = (llm_claims if kind == "claim" else llm_pairs)
                    args = ((wb, doc, pi, narr, facts, a.model) if kind == "claim"
                            else (wb, doc, facts, a.model))
                    fut[ex.submit(fn, *args)] = (kind, di, doc, pi, k)
                for fu in as_completed(fut):
                    kind, di, doc, pi, k = fut[fu]
                    try:
                        res = fu.result()
                        cache[k] = res
                        (docs[di]["issues"] if kind == "claim"
                         else docs[di]["checked"]).extend(res)
                    except Exception as e:
                        tqdm.write(f"    ⚠ {doc} p{pi} LLM 失敗：{str(e)[:100]}")
                    bar.update(1)
            bar.close()
            cachef.write_text(json.dumps(cache, ensure_ascii=False), encoding="utf-8")

    summary = []
    for d in docs:
        f, npages = d["file"], d["npages"]
        issues, checked, ref = d["issues"], d["checked"], d["ref"]
        md = outdir / f"{f.stem}.md"
        write_md(md, f, issues, checked, ref, npages)
        if not a.no_docx:
            try:
                import md2doc
                md2doc.to_docx(md2doc.parse_md(md.read_text(encoding="utf-8")),
                               outdir / f"{f.stem}.docx", f.stem)
            except Exception as e:
                tqdm.write(f"  ⚠ {f.name} 出唔到 docx：{str(e)[:100]}")
        summary.append((f.name, npages, len(issues),
                        sum(1 for c in checked if not c["ok"]), len(checked)))

    print(f"\n{'檔名':<46}{'頁':>4}{'表內對唔上':>11}{'跨表唔tie':>11}")
    print("-" * 92)
    for name, npg, ni, nx, nt in summary:
        print(f"{name[:44]:<46}{npg:>4}{ni:>11}{f'{nx}/{nt}':>11}")
    print(f"\n✓ 報告寫咗落 {outdir.resolve()}")


if __name__ == "__main__":
    main()
