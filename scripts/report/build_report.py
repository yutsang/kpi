#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from __future__ import annotations
"""
build_report.py — 單一自足檔：由底層數據（feed + 清單）生成表二審查報告 pptx。
毋須任何 prerequisite / 其他模組 / 前置 command：
    python build_report.py [entity]            # 預設 mgm；用現有 {entity}_llm_narrative.json 或清單 fallback
    python build_report.py [entity] --llm      # 即場生成 LLM 敘述（需 KPMG 網 + workbench creds）再出報告
（此檔由各 build/LLM 模組自動合併；LLM 相關 heavy import [openai/msoffcrypto] 全 lazy；報告只作 ref。）
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import re
import sys
from pathlib import Path
import json
import os
import time
from typing import Any
import io
from concurrent.futures import ThreadPoolExecutor, as_completed


# ── from layout ──
NAVY = RGBColor(0x00, 0x33, 0x8D)


# ── from layout ──
MBLUE = RGBColor(0x00, 0x5E, 0xB8)


# ── from layout ──
LBLUE = RGBColor(0x00, 0x91, 0xDA)


# ── from layout ──
VIOLET = RGBColor(0x48, 0x36, 0x98)


# ── from layout ──
PURPLE = RGBColor(0x47, 0x0A, 0x68)


# ── from layout ──
LPURPLE = RGBColor(0x6D, 0x20, 0x77)


# ── from layout ──
GREEN = RGBColor(0x00, 0xA3, 0xA1)


# ── from layout ──
SECFILL = RGBColor(0xEE, 0xF1, 0xF8)


# ── from layout ──
SUBTOT = RGBColor(0xD9, 0xE1, 0xF2)


# ── from layout ──
TOTAL = RGBColor(0xBD, 0xD7, 0xEE)


# ── from layout ──
BORDER = "BFBFBF"


# ── from layout ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ── from layout ──
INK = RGBColor(0x22, 0x22, 0x22)


# ── from layout ──
GREY = RGBColor(0x59, 0x59, 0x59)


# ── from layout ──
LGREY = RGBColor(0x8C, 0x8C, 0x8C)


# ── from layout ──
DARK = RGBColor(0x0C, 0x23, 0x3C)


# ── from layout ──
NEG_COLOR = None


# ── from layout ──
FONT_CN = "Microsoft YaHei"


# ── from layout ──
FONT_NUM = "Arial"


# ── from layout ──
FONT_HEAD = "KPMG Bold"


# ── from layout ──
SZ_CRUMB = 7.0


# ── from layout ──
SZ_TITLE = 12.0


# ── from layout ──
SZ_HEAD = 13.0


# ── from layout ──
SZ_BODY = 9.0


# ── from layout ──
SZ_BODY_HEAD = 9.5


# ── from layout ──
SZ_TBL = 7.5


# ── from layout ──
SZ_TBL_HDR = 7.0


# ── from layout ──
SZ_TBL_WIDE = 6.0


# ── from layout ──
SZ_CAPTION = 7.5


# ── from layout ──
SZ_NOTE = 7.0


# ── from layout ──
SZ_FOOT = 6.0


# ── from layout ──
SZ_PAGE = 9.0


# ── from layout ──
SLIDE_W = 10.83


# ── from layout ──
SLIDE_H = 7.5


# ── from layout ──
MARGIN = 0.53


# ── from layout ──
COL_GAP = 0.21


# ── from layout ──
CRUMB_Y = 0.13


# ── from layout ──
SUBTITLE_Y = 0.34


# ── from layout ──
HEAD_Y = 0.56


# ── from layout ──
FOOT_Y = 7.16


# ── from layout ──
CONTENT_BOTTOM = 6.98


# ── from layout ──
SECTIONS = ["2025年度投資計劃執行情況概述", "過往年度投資計劃在2025年繼續執行的審查跟進",
            "本年度審查工作的主要發現", "其他信息", "投資計劃執行報告的六項KPI分析", "附件"]


# ── from layout ──
_CN_RE = None


# ── from layout ──
def set_ea(run_or_font, ea=None):
    """寫 <a:ea>（中文字體）—— python-pptx 只寫 <a:latin>，唔寫 ea 中文會跌返 theme 預設。
    OOXML 次序：… latin, ea, cs …，所以要 insert 喺 latin 之後。"""
    f = getattr(run_or_font, "font", run_or_font)
    try:
        rPr = f._rPr
    except AttributeError:
        return
    if rPr is None:
        return
    el = rPr.find(qn("a:ea"))
    if el is None:
        el = rPr.makeelement(qn("a:ea"), {})
        lat = rPr.find(qn("a:latin"))
        (lat.addnext(el) if lat is not None else rPr.append(el))
    el.set("typeface", ea or FONT_CN)


# ── from layout ──
def setfont(run, size, *, bold=False, italic=False, color=None, heading=False, latin=None):
    """一次過設 size/bold/color + <a:latin> + <a:ea>（跟 template theme）。"""
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    if color is not None:
        f.color.rgb = color
    f.name = latin or (FONT_HEAD if heading else FONT_NUM)
    set_ea(f)
    return run


# ── from layout ──
def _is_cn(ch):
    return "⺀" <= ch <= "鿿" or "＀" <= ch <= "￯" or "　" <= ch <= "〿"


# ── from layout ──
def has_cn(s):
    return any(_is_cn(c) for c in str(s))


# ── from layout ──
def text_w(s, size):
    """估文字闊度（pt）：中文/全形 ≈ 1 em、英數 ≈ 0.52 em。"""
    w = 0.0
    for c in str(s):
        w += size * (1.0 if _is_cn(c) else 0.52)
    return w


# ── from layout ──
def est_lines(s, col_w_in, size, margin_in=0.06):
    """估 wrap 行數（col_w_in = 欄闊吋）。認 \\n 明碼換行。"""
    avail = max((col_w_in - margin_in) * 72.0, 6.0)
    n = 0
    for seg in str(s).split("\n"):
        n += max(1, -(-text_w(seg, size) // avail))     # ceil
    return int(n)


# ── from layout ──
def row_h(cells, widths, size, pad_in=0.045, min_h=0.155):
    """一行嘅需要高度（吋）＝ 最多 wrap 行數 × 行距 + 上下 padding。"""
    lines = 1
    for txt, w in zip(cells, widths):
        lines = max(lines, est_lines(txt, w, size))
    return max(min_h, lines * size * 1.24 / 72.0 + pad_in)


# ── from layout ──
def size_of(prs):
    return prs.slide_width / 914400.0, prs.slide_height / 914400.0


# ── from layout ──
def blank(prs):
    lay = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(lay)


# ── from layout ──
def _tb(slide, x, y, w, h, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return box


# ── from layout ──
def put(slide, x, y, w, h, text, *, size=8, bold=False, color=INK, align=PP_ALIGN.LEFT,
        italic=False, font=None, wrap=True):
    """一行/一段文字框。wrap=False 用喺一定要一行嘅嘢（breadcrumb 頁籤）。"""
    box = _tb(slide, x, y, w, h, wrap)
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    p.font.size = Pt(size); p.font.name = FONT_NUM; set_ea(p.font)   # 空段落唔好跌返 theme 預設
    p._p.get_or_add_endParaRPr().set("sz", str(int(round(size * 100))))
    r = p.add_run(); r.text = str(text)
    setfont(r, size, bold=bold, italic=italic, color=color, latin=font)
    return box


# ── from layout ──
BAND = RGBColor(0xF2, 0xF2, 0xF2)


# ── from layout ──
def _rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.5)
    sh.shadow.inherit = False
    return sh


# ── from layout ──
def _name(shape, nm):
    try:
        shape.name = nm
    except Exception:      # noqa: BLE001 — 改唔到名只係少咗 hyperlink，唔好炸咗成個 build
        pass


# ── from layout ──
def breadcrumb(slide, W, active=0, entity="MGM"):
    """頂 nav（對 scan p-23 放大）：白底、頁籤用「｜」分隔，當前頁籤 navy 粗體、其餘淺灰，
    右邊 entity + ◀ ⌂ ▶ 三粒圓掣。shape 改名做 nav:* ，wire_nav() 事後駁內部 hyperlink。"""
    x0 = MARGIN - 0.23
    d, gapc = 0.185, 0.05                                  # 圓掣直徑 / 間距
    right = W - x0
    for i, (nm, ch) in enumerate((("next", "▶"), ("home", "⌂"), ("prev", "◀"))):
        cx = right - d - i * (d + gapc)
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(CRUMB_Y - 0.035),
                                    Inches(d), Inches(d))
        sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
        sh.line.color.rgb = NAVY; sh.line.width = Pt(0.75); sh.shadow.inherit = False
        _name(sh, f"nav:{nm}")
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ch
        setfont(r, 6.0, bold=True, color=NAVY)
    ex = right - 3 * (d + gapc)                            # entity 靠圓掣左邊
    put(slide, ex - 0.85, CRUMB_Y, 0.8, 0.18, entity, size=SZ_CRUMB, bold=True,
        color=INK, align=PP_ALIGN.RIGHT)
    sep, avail = " ｜ ", (ex - 0.92) - x0
    # ×1.08：text_w 對粗體中文估細咗少少，唔留鬆位頁籤會撞埋一齊
    widths = [text_w(t, SZ_CRUMB) * 1.08 / 72.0 for t in SECTIONS]
    sw = text_w(sep, SZ_CRUMB) / 72.0
    scale = min(1.0, avail / (sum(widths) + sw * (len(SECTIONS) - 1)))
    x = x0
    for i, t in enumerate(SECTIONS):
        if i:
            put(slide, x, CRUMB_Y, sw * scale + 0.03, 0.18, sep, size=SZ_CRUMB * scale,
                color=LGREY, wrap=False)
            x += sw * scale
        w = widths[i] * scale
        _name(put(slide, x, CRUMB_Y, w + 0.05, 0.18, t, size=SZ_CRUMB * scale, wrap=False,
                  bold=(i == active), color=NAVY if i == active else LGREY), f"nav:sec{i}")
        x += w


# ── from layout ──
def _hlink(shape, rid):
    """畀 shape 內所有 run 加內部跳頁 hyperlink（a:hlinkClick + ppaction://hlinksldjump）。
    hlinkClick 喺 CT_TextCharacterProperties 排 latin/ea 之後 → append 就啱序。"""
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            rPr = r._r.get_or_add_rPr()
            h = rPr.makeelement(qn("a:hlinkClick"),
                                {qn("r:id"): rid, "action": "ppaction://hlinksldjump"})
            rPr.append(h)


# ── from layout ──
def wire_nav(prs, sec_slide=None, home=0):
    """全部 slide 砌完（連目錄插咗、重排咗）之後至駁：◀/▶ = 上/下頁、⌂ = 目錄、
    頁籤 = 該章分隔頁。sec_slide = {章 index: slide index}。"""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    slides = list(prs.slides)
    n = len(slides)
    for i, s in enumerate(slides):
        tgt = {"nav:prev": max(i - 1, 0), "nav:next": min(i + 1, n - 1), "nav:home": home}
        for k, v in (sec_slide or {}).items():
            tgt[f"nav:sec{k}"] = v
        for sh in s.shapes:
            j = tgt.get(sh.name)
            if j is None or j == i or not sh.has_text_frame:
                continue
            _hlink(sh, s.part.relate_to(slides[j].part, RT.SLIDE))


# ── from layout ──
def footer(slide, W, H, page):
    """底：KPMG 字標 + 版權 + 初稿/頁碼（對 scan）。"""
    kb = _tb(slide, MARGIN - 0.23, H - 0.34, 0.7, 0.22)
    kr = kb.text_frame.paragraphs[0].add_run(); kr.text = "KPMG"
    setfont(kr, 11, bold=True, italic=True, color=NAVY)
    put(slide, MARGIN + 0.5, H - 0.30, W - 2.2, 0.2,
        "© 2026畢馬威會計師事務所 — 澳門特別行政區合夥制事務所。版權所有，不得轉載。",
        size=SZ_FOOT, color=LGREY)
    if page is not None:
        put(slide, W - 1.15, H - 0.32, 0.95, 0.2, f"初稿　{page}", size=SZ_PAGE, bold=True,
            color=NAVY, align=PP_ALIGN.RIGHT)


# ── from layout ──
MAX_HEAD_H = 1.35


# ── from layout ──
def head_h(headline, W, hsize=SZ_HEAD):
    """導語需要嘅高度 + 實際字號（長就自動縮到 MAX_HEAD_H 為止）→ (h, size)。"""
    if not headline:
        return 0.06, hsize
    while hsize > 6.0:
        h = est_lines(headline, W - 2 * MARGIN, hsize) * hsize * 1.35 / 72.0
        if h <= MAX_HEAD_H:
            return h, hsize
        hsize -= 0.5
    return MAX_HEAD_H, hsize


# ── from layout ──
def page_head(slide, W, crumb, headline=None, *, hsize=SZ_HEAD):
    """灰色「章節 | 子題」+ navy 粗體導語 → 回內容起始 y。"""
    put(slide, MARGIN, SUBTITLE_Y, W - 2 * MARGIN, 0.2, crumb, size=SZ_TITLE, bold=True, color=NAVY)
    if not headline:
        return HEAD_Y + 0.06
    h, hsize = head_h(headline, W, hsize)
    box = _tb(slide, MARGIN, HEAD_Y, W - 2 * MARGIN, h)
    p = box.text_frame.paragraphs[0]
    r = p.add_run(); r.text = str(headline)
    setfont(r, hsize, bold=True, color=NAVY, heading=True)
    return HEAD_Y + h + 0.10


# ── from layout ──
def caption_bar(slide, x, y, w, text, *, size=SZ_CAPTION):
    """表頂 caption bar（重覆表名，對 scan 每張表都有）。
    ⚠ 用深色 HDR2 —— IMG_0441 量到 caption 條比表頭嗰排藍【深啲】，唔係同一隻色。"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.17))
    bar.fill.solid(); bar.fill.fore_color.rgb = CAPTION_FILL
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.margin_left = Emu(36000); tf.margin_right = Emu(18000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = str(text)
    setfont(r, size, bold=True, color=WHITE)
    return y + 0.17


# ── from layout ──
def source_note(slide, W, y=None, *, note=None, more=False):
    """表下：資料來源（左）+（下頁待續）（右）。"""
    y = CONTENT_BOTTOM if y is None else y
    put(slide, MARGIN, y, W - 2.0, 0.16,
        note or "資料來源：管理層提供之項目投資計劃及執行報告資料，畢馬威分析",
        size=SZ_NOTE, color=NOTE_FG)
    if more:
        put(slide, W - MARGIN - 1.2, y, 1.2, 0.16, "（下頁待續）", size=SZ_NOTE, color=NOTE_FG,
            align=PP_ALIGN.RIGHT)


# ── from layout ──
RULE = "00338D"


# ── from layout ──
HDR_FILL = RGBColor(0x1E, 0x49, 0xE2)


# ── from layout ──
HDR_KEY = RGBColor(0x09, 0x8E, 0x7E)


# ── from layout ──
HDR_SKY = RGBColor(0x00, 0xB8, 0xF5)


# ── from layout ──
HDR_PUR = RGBColor(0x48, 0x36, 0x98)


# ── from layout ──
CAPTION_FILL = NAVY


# ── from layout ──
SEC_FG = NAVY


# ── from layout ──
NOTE_FG = NAVY


# ── from layout ──
HDR1, HDR2, HDR3 = HDR_FILL, HDR_FILL, HDR_KEY


# ── from layout ──
TEAL = HDR_KEY


# ── from layout ──
def _edge(cell, side, *, w=9525, color=RULE, dash=None):
    """畫單一條邊（side ∈ T/B/L/R）。ln* 要插喺 tcPr 最前，否則 PowerPoint 會叫修復。
    同一邊重覆設就換走舊嗰條。"""
    tcPr = cell._tc.get_or_add_tcPr()
    tag = qn(f"a:ln{side}")
    old = tcPr.find(tag)
    if old is not None:
        tcPr.remove(old)
    ln = tcPr.makeelement(tag, {"w": str(w), "cap": "flat"})
    fill = ln.makeelement(qn("a:solidFill"), {})
    clr = fill.makeelement(qn("a:srgbClr"), {"val": color})
    fill.append(clr); ln.append(fill)
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    tcPr.insert(0, ln)


# ── from layout ──
def set_cell(cell, text, *, size=SZ_TBL, bold=False, fill=None, align=PP_ALIGN.RIGHT,
             color=None, wrap=True, anchor=MSO_ANCHOR.MIDDLE, italic=False):
    cell.margin_left = cell.margin_right = Emu(18000)
    cell.margin_top = cell.margin_bottom = Emu(9000)
    cell.vertical_anchor = anchor
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill if fill is not None else WHITE
    tf = cell.text_frame; tf.word_wrap = wrap
    tf.clear()          # merge 會把被合併格嘅字搬入 origin → 唔清就會出「萬澳門元萬澳門元」
    p = tf.paragraphs[0]; p.alignment = align
    txt = "" if text is None else str(text)
    if color is None:
        color = NEG_COLOR if (NEG_COLOR is not None and txt.startswith("(")) else INK
    # ★ 空格一定要定死字號：PowerPoint 見到「冇 run／空 run」會跌返去 endParaRPr／預設 text
    #   style（python-pptx fresh deck ＝ Calibri 18pt）→ row 被撐到 ~0.3in，成張表爆版。
    #   所以：paragraph 層 defRPr + endParaRPr 都寫死，而且空字串索性唔加 run。
    p.font.size = Pt(size); p.font.bold = bold
    p.font.name = FONT_NUM; set_ea(p.font)
    epr = p._p.get_or_add_endParaRPr()
    epr.set("sz", str(int(round(size * 100))))
    if not txt:
        return
    # ⚠ DrawingML 入面 "\n" 唔係換行（會當空白）→ 一定要用 <a:br/>，否則表頭喺 PowerPoint 會擠成一行
    for i, seg in enumerate(txt.split("\n")):
        if i:
            p.add_line_break()
        if not seg:
            continue
        r = p.add_run(); r.text = seg
        setfont(r, size, bold=bold, italic=italic, color=color)


# ── from layout ──
ROW_FILL = {"sec": None, "subtot": None, "tot": None, "data": None, "formula": None}


# ── from layout ──
def header_h(supers, subs, widths, hfont):
    """表頭需要高度（吋）。"""
    h = row_h(subs, widths, hfont, pad_in=0.05, min_h=0.20)
    return (0.17 + h) if supers else h


# ── from layout ──
def fit_rows(rows, widths, font, avail_h, hh):
    """按【累積高度】切頁 → 保證唔會超出可用高度。rows = [(kind, cells)]。
    keep=True 嘅 row（範疇 block）盡量唔拆：見 fit_blocks。"""
    out, cur, used = [], [], 0.0
    cap = max(avail_h - hh, 0.6)          # guard：導語太長時唔好變 0/負數（會無限開版）
    for kind, cells in rows:
        h = row_h(cells, widths, font)
        if cur and used + h > cap:
            out.append(cur); cur, used = [], 0.0
        cur.append((kind, cells)); used += h
    if cur:
        out.append(cur)
    return out or [[]]


# ── from layout ──
def fit_blocks(blocks, widths, font, avail_h, hh):
    """block = 一個範疇（section + data + 小計）。整個 block 唔拆頁（對 scan：全報告冇「續」）；
    單一 block 大過一版先逼住切。"""
    cap = max(avail_h - hh, 0.6)          # guard：同上
    pages, cur, used = [], [], 0.0
    for blk in blocks:
        bh = sum(row_h(c, widths, font) for _, c in blk)
        if cur and used + bh > cap:
            pages.append(cur); cur, used = [], 0.0
        if bh > cap:                                    # 單一範疇爆版 → 逐行切（安全網）
            for kind, cells in blk:
                h = row_h(cells, widths, font)
                if cur and used + h > cap:
                    pages.append(cur); cur, used = [], 0.0
                cur.append((kind, cells)); used += h
            continue
        cur.extend(blk); used += bh
    if cur:
        pages.append(cur)
    return pages or [[]]


# ── from layout ──
def draw_table(slide, x, y, w, subs, rows, widths, *, supers=None, font=SZ_TBL, hfont=SZ_TBL_HDR,
               left_cols=1, fill_h=None, max_row_h=0.26, hdr_cols=None):
    """畫 navy 表。subs=欄名（可含 \\n）；rows=[(kind, cells)]；widths=相對闊度（會 scale 到 w）。
    supers=[(label, c0, c1_exclusive)] 兩層表頭。fill_h=想填滿嘅高度（行數少時撐開行高，
    唔好剩一大橛白位；每行最多 max_row_h）。回 (bottom_y, 實際高度)。"""
    ncol = len(subs)
    scale = w / sum(widths)
    wid = [v * scale for v in widths]
    nhdr = 2 if supers else 1
    heights = [row_h(cells, wid, font) for _, cells in rows]
    hsub = row_h(subs, wid, hfont, pad_in=0.05, min_h=0.20)
    if fill_h and heights:
        slack = fill_h - ((0.17 if supers else 0) + hsub + sum(heights))
        if slack > 0.05:
            add = min(slack / len(heights), max(0.0, max_row_h - max(heights)))
            heights = [h + add for h in heights]
    total_h = (0.17 if supers else 0) + hsub + sum(heights)
    tbl = slide.shapes.add_table(nhdr + len(rows), ncol, Inches(x), Inches(y),
                                 Inches(w), Inches(total_h)).table
    tbl.first_row = False; tbl.horz_banding = False
    for i, v in enumerate(wid):
        tbl.columns[i].width = Inches(v)
    # 三色欄組——【要 caller 明示】：4.2 表都有「設施建設/活動舉辦」欄但成排 navy（scan p.24），
    #   所以唔可以淨靠欄名估。hdr_cols = {欄 index: 顏色}，冇指定就 HDR1。
    hc = dict(hdr_cols or {})
    if supers:
        for c in range(ncol):
            set_cell(tbl.cell(0, c), "", size=hfont, fill=hc.get(c, HDR1), color=WHITE)
        # 一個欄組跨住兩隻表頭色（報告：潛在調整後 = 深藍嗰兩欄 + 綠嗰兩欄）→ 拆開兩格，
        #   個 label 兩邊都寫（同 IMG_0441 一樣，「潛在調整後投資金額」出現兩次）。
        for label, c0, c1 in supers:
            a = c0
            while a < c1:
                b = a + 1
                while b < c1 and hc.get(b, HDR1) == hc.get(a, HDR1):
                    b += 1
                if b - a > 1:
                    tbl.cell(0, a).merge(tbl.cell(0, b - 1))
                set_cell(tbl.cell(0, a), label or "", size=hfont + 0.5, bold=True,
                         fill=hc.get(a, HDR1), color=WHITE, align=PP_ALIGN.CENTER)
                a = b
        tbl.rows[0].height = Emu(int(0.17 * 914400))
    for c, s in enumerate(subs):
        # 報告嘅欄名喺表頭【貼底】（wrap 做兩行時尤其明顯）
        set_cell(tbl.cell(nhdr - 1, c), s, size=hfont, bold=True, anchor=MSO_ANCHOR.BOTTOM,
                 fill=hc.get(c, HDR1), color=WHITE,
                 align=PP_ALIGN.LEFT if c < left_cols else PP_ALIGN.CENTER)
    # 角位：報告係一格（序號欄冇字），單位「萬澳門元」貼住最左
    if left_cols >= 2 and not str(subs[0]).strip() and str(subs[1]).strip():
        tbl.cell(nhdr - 1, 0).merge(tbl.cell(nhdr - 1, 1))
        set_cell(tbl.cell(nhdr - 1, 0), subs[1], size=hfont, bold=True, anchor=MSO_ANCHOR.BOTTOM,
                 fill=hc.get(0, HDR1), color=WHITE, align=PP_ALIGN.LEFT)
    tbl.rows[nhdr - 1].height = Emu(int(hsub * 914400))
    for ri, (kind, cells) in enumerate(rows, start=nhdr):
        bold = kind in ("sec", "subtot", "tot")
        if kind == "formula":      # 報告表頭下面嗰行斜體公式（a｜1..7｜b｜c=a+b｜d=b/a）
            for c, v in enumerate(cells):
                set_cell(tbl.cell(ri, c), v, size=max(4.5, font - 1.0), italic=True,
                         color=GREY, align=PP_ALIGN.LEFT if c < left_cols else PP_ALIGN.RIGHT)
            tbl.rows[ri].height = Emu(int(max(0.14, (font - 1.0) * 1.24 / 72.0 + 0.03) * 914400))
            continue
        # 標籤（範疇/小計/總計/表尾說明行）喺報告係【由最左邊起】，唔係縮喺名稱欄：
        #   序號欄空 + 名稱欄有字 → merge 埋，個 label 先有位唔會 wrap
        k = 0
        if left_cols >= 2 and not str(cells[0]).strip() and str(cells[1]).strip():
            k = 1
            tbl.cell(ri, 0).merge(tbl.cell(ri, 1))
        for c, v in enumerate(cells):
            if 0 < c <= k:
                continue                      # 已 merge 入 col 0
            al = PP_ALIGN.LEFT if c < left_cols else PP_ALIGN.RIGHT
            set_cell(tbl.cell(ri, c), cells[k] if c == 0 and k else v,
                     size=font, bold=bold, fill=ROW_FILL.get(kind), align=al,
                     color=SEC_FG if kind == "sec" else None)
        tbl.rows[ri].height = Emu(int(heights[ri - nhdr] * 914400))
    # ── 線：只有小計/總計橫線 + 欄組虛線直線（報告冇逐格格線）────────────
    # 只喺【有名嘅欄組】邊界畫虛線；標籤欄自成一「組」（label=""）唔算
    gsep = {c0 for _l, c0, _c1 in (supers or []) if c0 > 0 and str(_l).strip()}
    last = nhdr + len(rows) - 1
    for ri, (kind, _c) in enumerate(rows, start=nhdr):
        if kind in ("subtot", "tot"):
            for c in range(ncol):
                _edge(tbl.cell(ri, c), "T")
                _edge(tbl.cell(ri, c), "B")      # 報告：小計同總計【上下都有】幼線
    for ri in range(nhdr, nhdr + len(rows)):
        for c in gsep:
            _edge(tbl.cell(ri, c), "L", w=6350, color=RULE, dash="sysDash")
    return y + total_h, total_h


# ── from layout ──
def prose(box, items, *, head_size=SZ_BODY_HEAD, body_size=SZ_BODY, gap=6):
    """scan 敘述格式：navy 粗體小標題一行 + 下面 body 段落（唔用 ■ bullet）。
    items = [(head, body)]；head 可為空。"""
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for head, body in items:
        if head:
            head = str(head).rstrip("：:")        # 項目組：小標題唔應該有冒號
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_before = Pt(0 if first else gap); p.space_after = Pt(1)
            p.font.size = Pt(head_size); p.font.name = FONT_NUM; set_ea(p.font)
            p._p.get_or_add_endParaRPr().set("sz", str(int(round(head_size * 100))))
            r = p.add_run(); r.text = str(head)
            setfont(r, head_size, bold=True, color=NAVY)
            first = False
        if body:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_before = Pt(0 if first else (0 if head else gap)); p.space_after = Pt(1)
            p.font.size = Pt(body_size); p.font.name = FONT_NUM; set_ea(p.font)
            p._p.get_or_add_endParaRPr().set("sz", str(int(round(body_size * 100))))
            r = p.add_run(); r.text = str(body)
            setfont(r, body_size, color=RGBColor(0x33, 0x33, 0x33))
            first = False


# ── from layout ──
def prose_numbered(box, items, *, size=SZ_BODY, gap=7, indent=0.24, title=None, tsize=SZ_BODY_HEAD):
    """scan 表旁格式（p-11/p-13 右欄）：navy 粗體小標題 + 編號清單
        1.  {粗體類型}（{金額}）：{內文…}       ← hanging indent，內文對齊類型名
    items = [(編號, 粗體引子, 內文)]；編號跟七大類 canonical 序（會跳號）。"""
    tf = box.text_frame; tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]; p.space_after = Pt(4)
        p.font.size = Pt(tsize)
        p._p.get_or_add_endParaRPr().set("sz", str(int(round(tsize * 100))))
        r = p.add_run(); r.text = str(title)
        setfont(r, tsize, bold=True, color=NAVY)
        first = False
    emu = int(indent * 914400)
    for no, head, body in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(gap); p.space_after = Pt(0)
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(emu)); pPr.set("indent", str(-emu))     # hanging indent
        p.font.size = Pt(size); p.font.name = FONT_NUM; set_ea(p.font)
        p._p.get_or_add_endParaRPr().set("sz", str(int(round(size * 100))))
        rn = p.add_run(); rn.text = f"{no}.\t"
        setfont(rn, size, bold=True, color=NAVY)
        rh = p.add_run(); rh.text = str(head)
        setfont(rh, size, bold=True, color=NAVY)
        rb = p.add_run(); rb.text = str(body)
        setfont(rb, size, color=RGBColor(0x33, 0x33, 0x33))


# ── from layout ──
def est_numbered_h(items, w, size=SZ_BODY, gap=7, title=None, tsize=SZ_BODY_HEAD, indent=0.24):
    h = (est_lines(title, w, tsize) * tsize * 1.3 / 72.0 + 4 / 72.0) if title else 0.0
    for _no, head, body in items:
        h += est_lines(f"　{head}{body}", w - indent, size) * size * 1.35 / 72.0 + gap / 72.0
    return h


# ── from layout ──
def prose_box(slide, x, y, w, h, items, **kw):
    box = _tb(slide, x, y, w, h)
    prose(box, items, **kw)
    return box


# ── from layout ──
def est_prose_h(items, w, head_size=SZ_BODY_HEAD, body_size=SZ_BODY, gap=6):
    """估敘述高度（吋）→ 用嚟分頁，唔會爆版。"""
    h = 0.0
    for head, body in items:
        if head:
            h += est_lines(head, w, head_size) * head_size * 1.3 / 72.0 + gap / 72.0
        if body:
            h += est_lines(body, w, body_size) * body_size * 1.35 / 72.0 + 2 / 72.0
    return h


# ── from layout ──
def fit_prose(items, w, avail_h, **kw):
    """按估算高度切頁 → [[items]]。"""
    pages, cur, used = [], [], 0.0
    avail_h = max(avail_h, 0.6)           # guard
    for it in items:
        ih = est_prose_h([it], w, **kw)
        if cur and used + ih > avail_h:
            pages.append(cur); cur, used = [], 0.0
        cur.append(it); used += ih
    if cur:
        pages.append(cur)
    return pages or [[]]


# ── from layout ──
def apply_theme_fonts(prs):
    """把生成 deck 嘅 theme 字體改成公司 template 嗰套（major KPMG Bold / minor Arial，
    ea 兩者都 Microsoft YaHei）。python-pptx 開新檔用 Office 預設 theme（Calibri），
    凡係我哋冇明寫字體嘅地方（placeholder、空段落、表格預設）都會跌返 Calibri。"""
    import re as _re
    try:
        for m in prs.slide_masters:
            part = m.part.part_related_by(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
            xml = part.blob.decode("utf-8")
            for tag, latin in (("majorFont", FONT_HEAD), ("minorFont", FONT_NUM)):
                def _fix(mo, latin=latin):
                    seg = mo.group(0)
                    seg = _re.sub(r'<a:latin typeface="[^"]*"', f'<a:latin typeface="{latin}"', seg, count=1)
                    seg = _re.sub(r'<a:ea typeface="[^"]*"', f'<a:ea typeface="{FONT_CN}"', seg, count=1)
                    return seg
                xml = _re.sub(r"<a:" + tag + r">.*?</a:" + tag + r">", _fix, xml, flags=_re.S)
            part._blob = xml.encode("utf-8")
    except Exception:
        pass


# ── from layout ──
def dark_slide(prs):
    slide = blank(prs)
    W, H = size_of(prs)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = DARK
    rect.line.fill.background(); rect.shadow.inherit = False
    kb = _tb(slide, 0.55, 0.35, 2.2, 0.4)
    kr = kb.text_frame.paragraphs[0].add_run(); kr.text = "KPMG"
    setfont(kr, 20, bold=True, italic=True, color=WHITE)
    return slide, W, H


# ── from feed_schema ──
_YB = re.compile(r"^(\d{2})(?:_(\d{2})SY)?$")


# ── from feed_schema ──
MEASURES = {
    "報告投資金額": "調整前_萬",
    "潛在調整金額": "調整_萬",
    "潛在調整後投資金額": "調整後_萬",
}


# ── from feed_schema ──
def split_year(yb):
    """year_bucket → (plan_year, spend_year)，兩位數 int；認唔到回 (None, None)。

        "25"      → (25, 25)   2025年計劃、2025年發生
        "25_24SY" → (24, 25)   2024年計劃、2025年發生（＝期後）
        "24_23SY" → (23, 24)
        "23"      → (23, 23)
    ⚠ 前面嗰個數字係【發生年】，_NNSY 嗰個先係【計劃年】。
    """
    m = _YB.match(str(yb).strip())
    if not m:
        return None, None
    spend = int(m.group(1))
    return (int(m.group(2)) if m.group(2) else spend), spend


# ── from feed_schema ──
_SUB_DISPLAY = {
    "博彩娛樂場優化": "博彩娛樂場場地的優化",
    "博彩設施設備優化": "博彩設施及設備的優化",
}


# ── from feed_schema ──
def sub_display(s):
    """單一名 → 報告寫法；「博彩項目—博彩娛樂場優化」呢類 section 標籤都認。"""
    t = str(s).strip()
    if t in _SUB_DISPLAY:
        return _SUB_DISPLAY[t]
    for sep in ("—", "－", "-"):
        if sep in t:
            a, _, b = t.partition(sep)
            if b.strip() in _SUB_DISPLAY:
                return f"{a}{sep}{_SUB_DISPLAY[b.strip()]}"
    return s


# ── from feed_schema ──
def sub_of(df):
    """報告 row label「範疇」：博彩用 vertical_label、非博彩用 ng_label。
    feed 已經有物化嘅「範疇」欄就直接用（prep_tableau 出）。"""
    if "範疇" in df.columns:
        s = df["範疇"].astype(str).str.strip()
        if not s.isin(["", "nan", "None"]).all():
            return s
    v = df["vertical_label"] if "vertical_label" in df.columns else ""
    n = df["ng_label"] if "ng_label" in df.columns else ""
    gm = df["ng_scope"].astype(str).str.strip().eq("gaming")
    import pandas as pd
    return pd.Series(v, index=df.index).where(gm, pd.Series(n, index=df.index))


# ── from feed_schema ──
def add_dims(df):
    """就地加 plan_year / spend_year / 範疇（已經有就唔郁）→ 回 df。"""
    if "year_bucket" in df.columns and ("plan_year" not in df.columns
                                        or "spend_year" not in df.columns):
        yb = df["year_bucket"].astype(str).str.strip()
        pairs = {v: split_year(v) for v in yb.unique()}
        if "plan_year" not in df.columns:
            df["plan_year"] = yb.map(lambda v: pairs[v][0])
        if "spend_year" not in df.columns:
            df["spend_year"] = yb.map(lambda v: pairs[v][1])
    if "範疇" not in df.columns:
        df["範疇"] = sub_of(df)
    return df


# ── from render_review_table_pptx ──
try:
    import pandas as pd
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("✗ 需要 pandas + python-pptx → pip install pandas python-pptx openpyxl"); sys.exit(1)


# ── from render_review_table_pptx ──
RATE_COLS = {"投資計劃完成率", "潛在調整後投資計劃完成率"}


# ── from render_review_table_pptx ──
TEXT_COLS = {"項目序號", "項目名稱"}


# ── from render_review_table_pptx ──
G1 = ["項目序號", "項目名稱", "計劃投資金額", "報告投資金額", "投資計劃完成率"]


# ── from render_review_table_pptx ──
G3 = ["調整後投資金額", "潛在調整後投資計劃完成率", "設施建設/資本性支出", "活動舉辦/營運性支出"]


# ── from render_review_table_pptx ──
GROUP_LABEL = {"G1": "項目基本信息", "G2": "投資金額的潛在調整事項", "G3": "潛在調整後投資金額"}


# ── from render_review_table_pptx ──
YEAR_TITLE = {
    "報告年25": "{e} 2025年度投資計劃單個項目審查結果匯總表",
    "報告年24": "{e} 2024年度投資計劃單個項目截至2025年末的審查結果匯總表",
    "報告年23": "{e} 2023年度投資計劃單個項目截至2025年末的審查結果匯總表",
}


# ── from render_review_table_pptx ──
SHORT = {
    "項目序號": "項目\n序號", "項目名稱": "項目名稱",
    "計劃投資金額": "計劃\n投資金額", "報告投資金額": "報告\n投資金額", "投資計劃完成率": "投資計劃\n完成率",
    "一般支持性部門的人工成本": "一般支持\n人工成本", "其他日常營運支出調整": "其他日常\n營運調整",
    "超出可計入範圍的內部資源支出": "超出可計入\n內部資源", "酒店客房改造支出": "酒店客房\n改造",
    "不符合“吸引外國客源”定義的相關投資支出": "不符吸引\n外國客源",
    "未完全實現投資目的的投資支出": "未完全\n實現目的",
    "投資計劃獲批前發生且未被認可的投資支出": "獲批前\n未認可",
    "潛在調整合計": "潛在調整\n合計", "調整後投資金額": "調整後\n投資金額",
    "潛在調整後投資計劃完成率": "潛在調整後\n完成率",
    "設施建設/資本性支出": "設施建設/\n資本性", "活動舉辦/營運性支出": "活動舉辦/\n營運性",
}


# ── from render_review_table_pptx ──
def fmt_money(v):
    if v is None or v == "" or (isinstance(v, float) and pd.isna(v)):
        return ""
    try:
        f = float(v)
    except (TypeError, ValueError):
        return str(v)
    n = int(round(f))
    if n == 0:
        return "-"
    return f"({abs(n):,})" if n < 0 else f"{n:,}"


# ── from render_review_table_pptx ──
def fmt_pct(v):
    # 分母 0／算唔到 → 「-」（同表內其餘「冇數」寫法一致；之前留空同 - 混用）
    if v is None or (isinstance(v, float) and pd.isna(v)):
        return "-"
    if v == "":
        return ""
    try:
        return f"{float(v)*100:.1f}%"
    except (TypeError, ValueError):
        return str(v)


# ── from render_review_table_pptx ──
def row_kind(seq: str) -> str:
    s = str(seq)
    if s.endswith("小計") or s.endswith("合計"):
        return "subtotal"
    if "—" in s or "－" in s:
        return "section"
    return "data"


# ── from render_review_table_pptx ──
def col_group(c):
    if c in G1:
        return "G1"
    if c in G3:
        return "G3"
    return "G2"


# ── from render_review_table_pptx ──
def _blocks(df, cols):
    """df → 逐【範疇】block：[(kind, cells)]。kind = sec / data / subtot / tot。
    scan 每個範疇有自己嘅小計、全份報告冇「續」→ 分頁時整個 block 唔拆。
    ⚠ 範疇／小計標籤要擺【項目名稱】欄（闊），唔可以留喺「項目序號」窄欄 —— 否則 wrap 5 行撐爆表。"""
    lab_c = cols.index("項目名稱") if "項目名稱" in cols else 0
    out, cur = [], []
    for _, row in df.iterrows():
        seq = str(row["項目序號"])
        kind = row_kind(seq)
        if kind == "section":
            if cur:
                out.append(cur)
            cells = [""] * len(cols); cells[lab_c] = sub_display(seq)
            cur = [("sec", cells)]
            continue
        cells = [("" if row.get(c, "") is None else str(row.get(c, ""))) if c in TEXT_COLS
                 else (fmt_pct(row.get(c, "")) if c in RATE_COLS else fmt_money(row.get(c, "")))
                 for c in cols]
        k = "tot" if seq.endswith("合計") else ("subtot" if kind == "subtotal" else "data")
        if k != "data":                       # 小計／合計：標籤搬去項目名稱欄
            cells[lab_c] = sub_display(seq); cells[0] = ""
        cur.append((k, cells))
    if cur:
        out.append(cur)
    return out


# ── from render_review_table_pptx ──
def render_sheet(prs, sheet_name, df, cols, *, ent_up="MGM", sec=3, crumb=None, page_cb=None):
    """單個項目審查結果匯總表（對 scan slide 46-63）：navy 2 層表頭、逐範疇 block 唔拆頁、
    按【累積高度】分頁（唔會超出版面）、表頂 caption bar、表底 資料來源／（下頁待續）。"""
    ncol = len(cols)
    W, H = size_of(prs)
    tw = W - 2 * MARGIN
    title = YEAR_TITLE.get(sheet_name, sheet_name).format(e=ent_up)
    subs = [SHORT.get(c, c) for c in cols]
    # 欄寬（相對）：序號窄、名闊、其餘等闊
    widths = [0.45 if c == "項目序號" else (2.0 if c == "項目名稱" else 0.62) for c in cols]
    wid = [w * tw / sum(widths) for w in widths]
    # super header（3 大組）
    supers, ci = [], 0
    while ci < ncol:
        g = col_group(cols[ci]); cj = ci
        while cj + 1 < ncol and col_group(cols[cj + 1]) == g:
            cj += 1
        supers.append((GROUP_LABEL[g], ci, cj + 1)); ci = cj + 1
    font = SZ_TBL_WIDE if ncol > 16 else SZ_TBL
    yr = "20" + (sheet_name[-2:] if sheet_name[-2:].isdigit() else "25")
    head = (f"下表匯總了我們在審查{ent_up} {yr}年度投資計劃各項目投資執行情況時，識別出的各項目"
            f"投資支出涉及的潛在調整事項，以及相關的影響金額。")
    probe = HEAD_Y + head_h(head, W)[0] + 0.10
    avail = CONTENT_BOTTOM - probe - 0.24 - 0.17          # 減 caption bar
    hh = header_h(supers, subs, wid, font - 0.5)
    pages = fit_blocks(_blocks(df, cols), wid, font, avail, hh)
    for pi, chunk in enumerate(pages):
        suffix = f"（{pi+1}/{len(pages)}）"
        slide = blank(prs)
        breadcrumb(slide, W, sec, ent_up)
        footer(slide, W, H, len(prs.slides._sldIdLst))
        top = page_head(slide, W, crumb or "其他信息  |  單個項目審查結果匯總", head + suffix)
        top = caption_bar(slide, MARGIN, top, tw, title + suffix)
        draw_table(slide, MARGIN, top, tw, subs, chunk, widths, supers=supers,
                     font=font, hfont=font - 0.5, left_cols=2,
                     fill_h=CONTENT_BOTTOM - top - 0.24)
        source_note(slide, W, note="註：金額單位為萬澳門元；括號表示調減。",
                      more=(pi < len(pages) - 1))
        print(f"    {sheet_name} 第 {pi+1}/{len(pages)} 頁：{len(chunk)} 行 × {ncol} 欄")


# ── from build_narrative ──
try:
    import openpyxl
    from openpyxl.utils import get_column_letter
except ImportError:
    print("✗ pip install openpyxl"); sys.exit(1)


# ── from build_narrative ──
CONCEPTS = [
    ("項目狀況", ["項目狀況"]),
    ("實施地點", ["實際實施地點", "實際實施地點/空間"]),
    ("計劃投資內容", ["計劃投資內容"]),
    ("實際投資內容", ["實際投資內容"]),
    ("KPMG分析發現", ["分析發現", "投資偏離"]),
    ("管理層解釋", ["管理層解釋", "變更原因"]),
    ("期後調整內容", ["期後調整内容", "期後調整內容"]),
    ("調整事項備註", ["調整事項備註", "調整備註"]),
    # 跨司工作組審計軌跡（報告調整詳述會引）：政府裁決 + KPMG 最終立場
    ("跨司回覆", ["跨司工作組的回", "跨司工作组的回", "跨司工作組回", "跨司工作组回", "第二輪意見", "第二輪意见"]),
    ("KPMG回覆", ["KPMG回覆", "KPMG回复"]),
]


# ── from build_narrative ──
def _norm_code(v):
    s = re.sub(r"\s+", "", str(v if v is not None else ""))
    s = re.sub(r"^項目", "", s)
    m = re.match(r"^0*(\d+)$", s)
    return m.group(1) if m else s


# ── from build_narrative ──
_PLACEHOLDER_RE = re.compile(r"參閱附件|請參閱|見附件|同附件|詳見|組織架構進展|^n/?a$|^-+$")


# ── from build_narrative ──
def _is_placeholder(s):
    s = str(s).strip()
    return len(s) < 8 or bool(_PLACEHOLDER_RE.search(s))


# ── from build_narrative ──
def nlook(narr, ng_scope, code):
    """由 (ng_scope, dicj code) 揾清單 rec —— 處理博彩/非博彩共用項目N 撞號（先試 exact）。"""
    g = (ng_scope == "gaming")
    c = _norm_code(code)
    return narr.get((g, c)) or narr.get((not g, c)) or {}


# ── from build_narrative ──
OUT_COLS = ["項目序號", "項目名稱", "項目類型"] + [n for n, _ in CONCEPTS]


# ── from build_narrative ──
def load_narrative(path, log=lambda *a: None) -> dict:
    """讀清單 → {正規化項目編號: rec}，rec 有 項目序號/名稱/類型 + 每概念（取最新非空）。"""
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    ws = None
    for sn in wb.sheetnames:
        if sn.lower().startswith("database") or sn == "Database":
            ws = wb[sn]; break
    ws = ws or wb[wb.sheetnames[0]]
    rows = []
    for i, r in enumerate(ws.iter_rows(values_only=True)):
        rows.append(r)
        if i > 2500:
            break
    hdr_r = code_c = name_c = type_c = None
    for ri in range(min(14, len(rows))):
        for ci, v in enumerate(rows[ri] or []):
            sv = "" if v is None else str(v)
            if "承批公司項目序號" in sv:
                hdr_r, code_c = ri, ci
            if name_c is None and sv.strip() == "項目名稱":
                name_c = ci
            if type_c is None and "項目類型" in sv:
                type_c = ci
        if hdr_r is not None:
            break
    if hdr_r is None:
        log("✗ 揾唔到『承批公司項目序號』表頭"); return {}
    hdr = [("" if v is None else str(v).replace("\n", "")) for v in rows[hdr_r]]
    concept_cols = {name: [ci for ci, h in enumerate(hdr) if any(k in h for k in keys)]
                    for name, keys in CONCEPTS}
    log(f"清單 header row {hdr_r+1}；概念欄命中：" +
        ", ".join(f"{n}:{len(concept_cols[n])}" for n, _ in CONCEPTS))
    out = {}
    for ri in range(hdr_r + 1, len(rows)):
        row = rows[ri]
        code = _norm_code(row[code_c]) if code_c < len(row) else ""
        if not code:
            continue
        ptype = (row[type_c] if type_c is not None and type_c < len(row) else "") or ""
        gaming = str(ptype).strip().startswith("博彩")     # 博彩/非博彩 共用項目N → key 帶 gaming
        rec = {"項目序號": f"項目{code}",
               "項目名稱": (row[name_c] if name_c is not None and name_c < len(row) else "") or "",
               "項目類型": ptype}
        for name, _ in CONCEPTS:
            vals = []
            for ci in concept_cols[name]:
                if ci < len(row) and row[ci] not in (None, "", 0, "0"):
                    sv = str(row[ci]).strip()
                    if not sv or sv.lower() in ("n/a", "nan"):
                        continue
                    if name in ("實際投資內容", "計劃投資內容") and _is_placeholder(sv):
                        continue                            # 跳 「參閱附件/請參閱/見附件」等 placeholder
                    if sv not in vals:
                        vals.append(sv)
            rec[name] = vals[-1] if vals else ""
        out.setdefault((gaming, code), rec)                 # key = (博彩?, 正規化碼)
    return out


# ── from build_project_review_table ──
try:
    import pandas as pd
except ImportError:
    print("✗ 未裝 pandas → pip install pandas openpyxl"); sys.exit(1)


# ── from build_project_review_table ──
pd.set_option("display.max_columns", 40)


# ── from build_project_review_table ──
pd.set_option("display.width", 260)


# ── from build_project_review_table ──
CANON = {
    "計入報告投資金額的高管及一般支持人員人工成本": "一般支持性部門的人工成本",
    "一般性支持部門的人工成本": "一般支持性部門的人工成本",
    "一般支持性部門的人工成本": "一般支持性部門的人工成本",
    "其他日常營運支出調整": "其他日常營運支出調整",
    "其他日常營運支出調整調整": "其他日常營運支出調整",
    "超過可計入範圍的内部資源支出": "超出可計入範圍的內部資源支出",
    "超出可計入範圍的內部資源支出": "超出可計入範圍的內部資源支出",
    "酒店客房改造支出": "酒店客房改造支出",
    "不符合“吸引外國客源”定義的相關投資支出": "不符合“吸引外國客源”定義的相關投資支出",
    "不符合吸引外國客源": "不符合“吸引外國客源”定義的相關投資支出",
    "未完全實現投資目的的投資支出": "未完全實現投資目的的投資支出",
    "未能實現投資目的的投資支出": "未完全實現投資目的的投資支出",
    "計劃獲批前的投資": "投資計劃獲批前發生且未被認可的投資支出",
    # 以下係跨年機制，報告 7 欄冇（user defer 跨年）→ 落 trailing other，之後 filter 當年就消失
    "將2024年的支出計入2025年報告投資金額": "〔跨年〕將往年支出計入本年報告投資金額",
    "2024年度計劃與2023年度計劃期後調整之間的報告投資金額跨期調整": "〔跨年〕年度計劃期後跨期調整",
}


# ── from build_project_review_table ──
ADJ7 = [
    "一般支持性部門的人工成本", "其他日常營運支出調整", "超出可計入範圍的內部資源支出",
    "酒店客房改造支出", "不符合“吸引外國客源”定義的相關投資支出",
    "未完全實現投資目的的投資支出", "投資計劃獲批前發生且未被認可的投資支出",
]


# ── from build_project_review_table ──
ADJ_POST = "不符合期後事項定義的投資支出"


# ── from build_project_review_table ──
ADJ_ALL = ADJ7 + [ADJ_POST]


# ── from build_project_review_table ──
def adj_no(t):
    """調整類型 → 報告嘅【固定編號】（1-8）。報告會 skip 冇金額嗰幾類，編號唔會重排。"""
    return (ADJ_ALL.index(t) + 1) if t in ADJ_ALL else len(ADJ_ALL) + 1


# ── from build_project_review_table ──
BASE_L = ["計劃投資金額", "報告投資金額", "投資計劃完成率"]


# ── from build_project_review_table ──
TAIL_L = ["潛在調整合計", "調整後投資金額", "潛在調整後投資計劃完成率", "設施建設/資本性支出", "活動舉辦/營運性支出"]


# ── from build_project_review_table ──
_PLAN_RE = {
    25: re.compile(r"2025.*預計投資金額.*合計|2025.*預計投資金額（萬"),
    24: re.compile(r"2024.*預計投資金額.*合計|2024.*預計投資金額（萬"),
    23: re.compile(r"2023.*預計投資金額"),
}


# ── from build_project_review_table ──
def _norm(c) -> str:
    s = re.sub(r"\s+", "", str(c if c is not None else ""))
    s = re.sub(r"^項目", "", s)
    m = re.match(r"^0*(\d+(?:\.\d+)?)$", s)
    return (m.group(1) if m else s).lower()


# ── from build_project_review_table ──
def _plan_year(yb) -> int:
    """year_bucket → 計劃年份（報告匯總表按計劃年分，唔係支出年）：
    25→25(2025計劃), 25_24SY→24(2024計劃嘅2025期後), 25_23SY→23, 24→24, 24_23SY→23, 23→23。"""
    s = str(yb).strip()
    m = re.search(r"_(\d+)SY", s)          # 有 _NNSY = 該計劃年嘅期後
    if m:
        return int(m.group(1))
    m2 = re.match(r"^(\d{2})", s)
    return int(m2.group(1)) if m2 else -1


# ── from build_project_review_table ──
def load_plan(path: Path, log=print) -> dict:
    """{報告年: {(is_gaming, 正規化項目編號): 計劃_萬}}（清單 database；博彩/非博彩共用項目N → key 帶 scope）。"""
    import openpyxl
    out = {25: {}, 24: {}, 23: {}}
    try:
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    except Exception as e:
        log(f"  ⚠ 清單開唔到 {path}: {e}"); return out
    for sn in wb.sheetnames:
        ws = wb[sn]
        rows = []
        for i, r in enumerate(ws.iter_rows(values_only=True)):
            rows.append(r)
            if i > 2500:
                break
        hdr_r = code_c = type_c = None
        for ri in range(min(14, len(rows))):
            for ci, v in enumerate(rows[ri] or []):
                sv = "" if v is None else str(v)
                if "承批公司項目序號" in sv:
                    hdr_r, code_c = ri, ci
                if type_c is None and "項目類型" in sv:
                    type_c = ci
            if hdr_r is not None:
                break
        if hdr_r is None:
            continue
        hdr = [("" if v is None else str(v).replace("\n", "")) for v in rows[hdr_r]]
        plan_c = {}
        for yr, rgx in _PLAN_RE.items():
            for ci, h in enumerate(hdr):
                if rgx.search(h):
                    plan_c[yr] = ci; break
        if not plan_c:
            continue
        log("  清單 " + sn + ": 計劃欄 " + ", ".join(f"{yr}→{hdr[ci][:18]!r}" for yr, ci in plan_c.items())
            + (f"；項目類型欄 col{type_c}" if type_c is not None else "；⚠冇項目類型欄(scope 分唔到)"))
        for ri in range(hdr_r + 1, len(rows)):
            row = rows[ri]
            code = _norm(row[code_c]) if code_c < len(row) else ""
            if not code:
                continue
            gaming = False
            if type_c is not None and type_c < len(row) and row[type_c] is not None:
                gaming = str(row[type_c]).strip().startswith("博彩")
            for yr, ci in plan_c.items():
                if ci < len(row):
                    try:
                        out[yr].setdefault((gaming, code), float(row[ci]))
                    except (TypeError, ValueError):
                        pass
        break
    return out


# ── from build_project_review_table ──
def load_category(path: Path, log=lambda *a: None) -> dict:
    """{(is_gaming, 正規化項目編號): 項目性質(D)}（清單；用嚟將零投資項目計劃 attribute 返範疇）。
    含 feed 冇嘅零投資項目（清單有齊全部項目）。"""
    import openpyxl
    out = {}
    try:
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    except Exception as e:
        log(f"  ⚠ 清單開唔到 {path}: {e}"); return out
    for sn in wb.sheetnames:
        ws = wb[sn]
        rows = []
        for i, r in enumerate(ws.iter_rows(values_only=True)):
            rows.append(r)
            if i > 2500:
                break
        hdr_r = code_c = type_c = sub_c = None
        for ri in range(min(14, len(rows))):
            for ci, v in enumerate(rows[ri] or []):
                sv = "" if v is None else str(v)
                if "承批公司項目序號" in sv:
                    hdr_r, code_c = ri, ci
                if type_c is None and "項目類型" in sv:
                    type_c = ci
                if sub_c is None and sv.strip() == "項目性質":
                    sub_c = ci
            if hdr_r is not None:
                break
        if hdr_r is None or sub_c is None:
            continue
        for ri in range(hdr_r + 1, len(rows)):
            row = rows[ri]
            code = _norm(row[code_c]) if code_c < len(row) else ""
            if not code:
                continue
            gaming = (type_c is not None and type_c < len(row) and row[type_c] is not None
                      and str(row[type_c]).strip().startswith("博彩"))
            sub = row[sub_c] if sub_c < len(row) else None
            if sub is not None and str(sub).strip():
                out.setdefault((gaming, code), str(sub).strip())
        log(f"  清單 項目性質(D)：{len(out)} 個項目")
        break
    return out


# ── from build_project_review_table ──
def _rate(rep, plan):
    try:
        return round(rep / plan, 4) if plan else None
    except Exception:
        return None


# ── from build_project_review_table ──
def build_year(df: pd.DataFrame, year: int, plan: dict | None = None):
    # 按計劃年份(plan year)分表：跨年期後(25_24SY 等)落返啱嘅計劃年，唔會爆 2025 表小計
    d = df[df["_plan_year"] == year].copy()
    # 只留乾淨「項目N」碼（丟 項目CAPEX-5 等 pseudo/分攤碼）
    d = d[d["dicj code"].astype(str).str.match(r"^項目\s*\d")]
    if d.empty:
        return None, []
    d["_adj"] = d["調整一級"].map(CANON).fillna(d["調整一級"])
    d["_sub"] = sub_of(d)
    key = ["ng_scope", "dicj code"]

    def _mode(s):
        m = s.dropna()
        return m.mode().iat[0] if len(m.mode()) else (m.iloc[0] if len(m) else "")
    base = d.groupby(key, dropna=False).agg(
        項目名稱=("project", "first"),
        _sub=("_sub", _mode),
        _ngcode=("ng_code", _mode),
        報告投資金額=("調整前_萬", "sum"),
        潛在調整合計=("調整_萬", "sum"),
        調整後投資金額=("調整後_萬", "sum"),
    )
    cap = d[d["final_capex_opex"] == "Capex"].groupby(key)["調整後_萬"].sum()
    ope = d[d["final_capex_opex"] == "Opex"].groupby(key)["調整後_萬"].sum()
    base["設施建設/資本性支出"] = cap.reindex(base.index).fillna(0)
    base["活動舉辦/營運性支出"] = ope.reindex(base.index).fillna(0)
    adj = d[d["調整一級"].notna()]
    pv = adj.pivot_table(index=key, columns="_adj", values="調整_萬", aggfunc="sum", fill_value=0)
    tab = base.join(pv).fillna(0).reset_index()
    tab = tab.rename(columns={"dicj code": "項目序號"})
    # 計劃 + 完成率（(scope,code) join；博彩/非博彩 collision → code-only fallback）
    if plan:
        codeonly = {}
        for (g, c), v in plan.items():
            codeonly.setdefault(c, v)
        tab["計劃投資金額"] = pd.to_numeric(tab.apply(
            lambda r: plan.get((r["ng_scope"] == "gaming", _norm(r["項目序號"])),
                               codeonly.get(_norm(r["項目序號"]))), axis=1), errors="coerce")
    else:
        tab["計劃投資金額"] = pd.NA
    tab["投資計劃完成率"] = tab.apply(lambda r: _rate(r["報告投資金額"], r["計劃投資金額"]), axis=1)
    tab["潛在調整後投資計劃完成率"] = tab.apply(lambda r: _rate(r["調整後投資金額"], r["計劃投資金額"]), axis=1)
    adj_cols = [c for c in ADJ7 if c in tab.columns]
    other = [c for c in pv.columns if c not in ADJ7]     # 跨年/未預期
    num = ["計劃投資金額", "報告投資金額"] + adj_cols + other + ["潛在調整合計",
          "調整後投資金額", "設施建設/資本性支出", "活動舉辦/營運性支出"]
    for c in num:
        tab[c] = pd.to_numeric(tab[c], errors="coerce").round(1)

    # 排序：博彩先，博彩內 vertical_label（娛樂場→設施設備），非博彩 ng_code(數字)；內部項目N
    def _pn(s):
        m = re.search(r"(\d+(?:\.\d+)?)", str(s)); return float(m.group(1)) if m else 9e9
    def _ngn(s):
        m = re.search(r"(\d+)", str(s)); return int(m.group(1)) if m else 99
    gorder = {"博彩娛樂場優化": 0, "博彩娛樂場場地的優化": 0, "博彩設施設備優化": 1, "博彩設施及設備的優化": 1}
    tab["_scope"] = (tab["ng_scope"] != "gaming").astype(int)          # 博彩=0 先
    tab["_go"] = tab["_sub"].map(lambda s: gorder.get(s, 5))
    tab["_ngn"] = tab["_ngcode"].map(_ngn)
    tab["_pn"] = tab["項目序號"].map(_pn)
    tab = tab.sort_values(["_scope", "_go", "_ngn", "_pn"]).reset_index(drop=True)

    # 砌 rows：section header + data + 範疇小計 + 類型合計
    ALL = ["項目序號", "項目名稱"] + BASE_L + adj_cols + other + TAIL_L
    numcols = ["計劃投資金額", "報告投資金額"] + adj_cols + other + ["潛在調整合計",
              "調整後投資金額", "設施建設/資本性支出", "活動舉辦/營運性支出"]
    rows = []

    def _agg_row(sub, label):
        r = {c: "" for c in ALL}
        r["項目序號"] = label
        for c in numcols:
            r[c] = round(pd.to_numeric(sub[c], errors="coerce").sum(), 1)
        r["投資計劃完成率"] = _rate(r["報告投資金額"], r["計劃投資金額"])
        r["潛在調整後投資計劃完成率"] = _rate(r["調整後投資金額"], r["計劃投資金額"])
        return r

    for scope in [s for s in ["gaming", "non_gaming"] if (tab["ng_scope"] == s).any()]:
        sc = tab[tab["ng_scope"] == scope]
        scope_name = "博彩項目" if scope == "gaming" else "非博彩項目"
        for sub_name, sub in sc.groupby(sc["_sub"], sort=False):
            rows.append({**{c: "" for c in ALL}, "項目序號": f"{scope_name}—{sub_name}"})  # section
            for _, row in sub.iterrows():
                rows.append({c: row.get(c, "") for c in ALL})
            rows.append(_agg_row(sub, f"{scope_name}—{sub_name} 小計"))
        rows.append(_agg_row(sc, f"{scope_name}合計"))
    rows.append(_agg_row(tab, "總計"))          # 全表總計（scan 每年表最後一行都有）

    out_df = pd.DataFrame(rows, columns=ALL)
    return out_df, other


# ── from build_summary_tables ──
try:
    import pandas as pd
except ImportError:
    print("✗ pip install pandas openpyxl"); sys.exit(1)


# ── from build_summary_tables ──
pd.set_option("display.max_columns", 30)


# ── from build_summary_tables ──
pd.set_option("display.width", 200)


# ── from build_summary_tables ──
BUCKET = {"25": "2025年度投資計劃", "25_24SY": "2024年度計劃期後投資", "25_23SY": "2023年度計劃期後投資"}


# ── from build_summary_tables ──
BUCKET_ORDER = ["2025年度投資計劃", "2024年度計劃期後投資", "2023年度計劃期後投資"]


# ── from build_summary_tables ──
GORDER = {"博彩娛樂場優化": 0, "博彩娛樂場場地的優化": 0, "博彩設施設備優化": 1, "博彩設施及設備的優化": 1}


# ── from build_summary_tables ──
def _ngn(s):
    m = re.search(r"(\d+)", str(s)); return int(m.group(1)) if m else 99


# ── from build_summary_tables ──
def _load(feed: Path, entity: str) -> pd.DataFrame:
    df = pd.read_csv(feed, low_memory=False)
    if entity and "entity" in df.columns:
        df = df[df["entity"].astype(str).str.lower() == entity]
    df = df[df["dicj code"].astype(str).str.match(r"^項目\s*\d")]     # 丟 pseudo 碼
    df["_yb"] = df["year_bucket"].astype(str).str.strip()
    df["_bucket"] = df["_yb"].map(BUCKET)
    df = df[df["_bucket"].notna()].copy()                             # 只留「於2025發生」3 bucket
    add_dims(df)                    # plan_year / spend_year / 範疇（一處派生）
    df["_sub"] = df["範疇"]
    df["_scope"] = (df["ng_scope"] != "gaming").astype(int)           # 博彩=0 先
    df["_go"] = df["_sub"].map(lambda s: GORDER.get(s, 5))
    df["_ngn"] = df["ng_code"].map(_ngn)
    return df


# ── from build_summary_tables ──
def _order(agg):
    return agg.sort_values(["_scope", "_go", "_ngn", "_sub"])


# ── from build_summary_tables ──
def _emit(agg, valcols):
    """agg：每 (ng_scope,_sub) 一行 + valcols。→ 加 博彩/非博彩 小計 + 總計。"""
    ALL = ["範疇"] + valcols
    rows = []

    def agg_row(sub, label):
        r = {c: "" for c in ALL}; r["範疇"] = label
        for c in valcols:
            r[c] = round(pd.to_numeric(sub[c], errors="coerce").sum(), 1)
        return r
    for scope in [0, 1]:
        sc = _order(agg[agg["_scope"] == scope])
        if sc.empty:
            continue
        nm = "博彩項目" if scope == 0 else "非博彩項目"
        rows.append({c: "" for c in ALL} | {"範疇": nm})   # section 標題行（對 scan p42）
        for _, row in sc.iterrows():
            r = {"範疇": row["_sub"]}
            for c in valcols:
                r[c] = round(float(row[c]), 1)
            rows.append(r)
        rows.append(agg_row(sc, f"{nm}小計"))
    rows.append(agg_row(_order(agg), "合計"))
    return pd.DataFrame(rows, columns=ALL)


# ── from build_summary_tables ──
def summary_amount(df) -> pd.DataFrame:
    """4.1 金額匯總：範疇 × bucket → 報告投資金額 / 潛在調整後投資金額 + 合計。"""
    # ★ 對 scan p42（項目組 2026-08-17）：每個年度出【報告 / 投資金額的潛在調整事項 / 潛在調整後】
    #   三欄；唔要「合計」欄組；最尾多一組「潛在調整後投資金額」拆設施建設／活動舉辦。
    g = df.groupby(["_scope", "_go", "_ngn", "_sub", "_bucket"], dropna=False).agg(
        報告=("調整前_萬", "sum"), 調整=("調整_萬", "sum"), 調整後=("調整後_萬", "sum")).reset_index()
    cap = df[df["final_capex_opex"] == "Capex"].groupby(
        ["_scope", "_go", "_ngn", "_sub"], dropna=False)["調整後_萬"].sum()
    ope = df[df["final_capex_opex"] == "Opex"].groupby(
        ["_scope", "_go", "_ngn", "_sub"], dropna=False)["調整後_萬"].sum()
    # pivot bucket → 兩個 measure
    base = g.groupby(["_scope", "_go", "_ngn", "_sub"], dropna=False)
    idx = base.size().reset_index()[["_scope", "_go", "_ngn", "_sub"]]
    out = idx.copy()
    valcols = []
    for bk in BUCKET_ORDER:
        sub = g[g["_bucket"] == bk].set_index(["_scope", "_go", "_ngn", "_sub"])
        for meas, lab in [("報告", "報告投資金額"), ("調整", "投資金額的潛在調整事項"),
                          ("調整後", "潛在調整後投資金額")]:
            col = f"{bk}·{lab}"
            out[col] = out.set_index(["_scope", "_go", "_ngn", "_sub"]).index.map(
                lambda k: sub[meas].get(k, 0.0)).astype(float).round(1).values
            valcols.append(col)
    _k = out.set_index(["_scope", "_go", "_ngn", "_sub"]).index
    for lab, ser in (("設施建設/資本性支出", cap), ("活動舉辦/營運性支出", ope)):
        col = f"潛在調整後投資金額·{lab}"
        out[col] = [round(float(ser.get(k, 0.0)), 1) for k in _k]
        valcols.append(col)
    return _emit(out, valcols)


# ── from build_summary_tables ──
def facility_activity(df, bucket_label) -> pd.DataFrame:
    """4.2 設施vs活動（一個 bucket）：範疇 × 設施建設(capex調整後)/活動舉辦(opex調整後)/合計。"""
    d = df[df["_bucket"] == bucket_label]
    if d.empty:
        return pd.DataFrame()
    cap = d[d["final_capex_opex"] == "Capex"].groupby(["_scope", "_go", "_ngn", "_sub"])["調整後_萬"].sum()
    ope = d[d["final_capex_opex"] == "Opex"].groupby(["_scope", "_go", "_ngn", "_sub"])["調整後_萬"].sum()
    idx = d.groupby(["_scope", "_go", "_ngn", "_sub"]).size().reset_index()[["_scope", "_go", "_ngn", "_sub"]]
    out = idx.copy()
    out["設施建設/資本性支出"] = out.set_index(["_scope", "_go", "_ngn", "_sub"]).index.map(
        lambda k: cap.get(k, 0.0)).astype(float).round(1).values
    out["活動舉辦/營運性支出"] = out.set_index(["_scope", "_go", "_ngn", "_sub"]).index.map(
        lambda k: ope.get(k, 0.0)).astype(float).round(1).values
    out["合計"] = (out["設施建設/資本性支出"] + out["活動舉辦/營運性支出"]).round(1)
    return _emit(out, ["設施建設/資本性支出", "活動舉辦/營運性支出", "合計"])


# ── from build_overview_tables ──
try:
    import pandas as pd
except ImportError:
    print("✗ pip install pandas openpyxl"); sys.exit(1)


# ── from build_overview_tables ──
BUCKET_PLANYR = {"2025年度投資計劃": 25, "2024年度計劃期後投資": 24, "2023年度計劃期後投資": 23}


# ── from build_overview_tables ──
def _plan_tot(plan, yr, gaming=None):
    d = (plan or {}).get(yr, {})
    return round(sum(v for (g, c), v in d.items() if gaming is None or g == gaming), 1)


# ── from build_overview_tables ──
def overview_by_bucket(df, bucket, plan, category=None):
    """整體投資概況（S10-14 = 2025計劃有計劃/完成率；S19-26 = 期後冇計劃、多潛在調整金額欄）。
    逐範疇 + 博彩/非博彩小計 + 總計。欄跟報告 IMG_0104/0105：
      2025計劃：項目數量 | 獲批的計劃投資金額 | 報告投資金額 | 完成率 | 潛在調整後投資金額 | 完成率 | 設施建設 | 活動舉辦
      期後    ：項目數量 | 報告投資金額 | 潛在調整金額 | 潛在調整後投資金額 | 設施建設 | 活動舉辦
    ⚠ 項目數量 = 執行報告披露嘅項目數，【含申報投資支出為零嘅項目】（scan p10 註釋2）→ 清單全部行都數。"""
    d = df[df["_bucket"] == bucket]
    if d.empty:
        return pd.DataFrame()
    yr = BUCKET_PLANYR[bucket]
    is_py = (bucket == "2025年度投資計劃")
    idx = ["_scope", "_go", "_ngn", "_sub"]
    g = d.groupby(idx, dropna=False).agg(
        項目數量=("dicj code", "nunique"), 報告=("調整前_萬", "sum"),
        調整=("調整_萬", "sum"), 後=("調整後_萬", "sum")).reset_index()
    cap = d[d["final_capex_opex"] == "Capex"].groupby(idx)["調整後_萬"].sum().rename("設施")
    ope = d[d["final_capex_opex"] == "Opex"].groupby(idx)["調整後_萬"].sum().rename("活動")
    g = g.merge(cap.reset_index(), on=idx, how="left").merge(ope.reset_index(), on=idx, how="left")
    g[["設施", "活動"]] = g[["設施", "活動"]].fillna(0.0)
    g = g.sort_values(idx)

    # 項目數量：2025計劃表要數【計劃項目】（清單），唔係 feed 出現嘅碼 —— 否則同「已實施/未實施」
    #   兩行對唔上（項目組 2026-08-15：總計 84 但 79+10≠84）。期後表冇計劃概念，照用 feed。
    plan_by_sub, n_by_sub, n_by_scope = {}, {}, {}
    if is_py and plan:
        cat = category or {}
        sub_of = {}       # (gaming,碼)→範疇；博彩/非博彩共用項目N 都要留（唔可以 drop 淨 code）
        d2sub = {}        # 項目性質(D)→set(範疇)：由 feed 有嘅項目學，用嚟派零投資項目計劃
        for _, r in d.drop_duplicates(["ng_scope", "dicj code"]).iterrows():
            key = (r["ng_scope"] == "gaming", _norm(r["dicj code"]))
            sub_of[key] = r["_sub"]
            dv = cat.get(key)
            if dv:
                d2sub.setdefault(str(dv), set()).add(r["_sub"])
        d2sub1 = {k: next(iter(v)) for k, v in d2sub.items() if len(v) == 1}  # 唯一對應先用
        miss = 0
        for (gm, code), v in (plan.get(yr, {}) or {}).items():
            sub = sub_of.get((gm, code))
            if sub is None:      # 零投資項目：feed 冇行 → 用 D→範疇 學到嘅 map 派返
                sub = d2sub1.get(str(cat.get((gm, code), "")))
            if sub is not None:
                plan_by_sub[sub] = plan_by_sub.get(sub, 0.0) + v
                # ★ scan p10 註釋2：「項目數量」＝執行報告披露嘅項目數，【包含申報投資支出為零嘅部分】
                #   → 唔可以再用 v > 0 過濾（之前 89，報告 95；未實施 10 vs 16 都係同一原因）
                n_by_sub[sub] = n_by_sub.get(sub, 0) + 1
                n_by_scope[0 if gm else 1] = n_by_scope.get(0 if gm else 1, 0) + 1
            elif v:
                miss += 1
        if miss:
            print(f"    ⚠ overview {bucket}：{miss} 個有計劃項目派唔到範疇（缺項目性質對應）")

    def mk(name, cnt, pl, rep, adj, aft, fac, act):
        r = {"範疇": name, "項目數量": int(cnt), "報告投資金額": round(rep, 1),
             "潛在調整金額": round(adj, 1), "潛在調整後投資金額": round(aft, 1),
             "設施建設/資本性支出": round(fac, 1), "活動舉辦/營運性支出": round(act, 1)}
        if is_py:
            r["獲批的計劃投資金額"] = round(pl, 1)
            r["投資金額的潛在調整事項"] = round(adj, 1)
            r["投資計劃完成率"] = _rate(rep, pl)
            r["潛在調整後投資計劃完成率"] = _rate(aft, pl)
        return r

    rows, n_tot = [], 0
    for scope in [0, 1]:
        sc = g[g["_scope"] == scope]
        name = "博彩項目" if scope == 0 else "非博彩項目"
        if sc.empty:
            # ⚠ 報告就算該 scope 全 0 都會【逐個範疇出行】＋小計（全部「-」）——
            #   scan p20 博彩項目下面照樣有「博彩娛樂場場地的優化 / 博彩設施及設備的優化」。
            rows.append({"範疇": name})
            if scope == 0:
                for gsub in ("博彩娛樂場優化", "博彩設施設備優化"):
                    rows.append(mk(gsub, 0, 0, 0, 0, 0, 0, 0))
            rows.append(mk(f"{name}小計", 0, _plan_tot(plan, yr, scope == 0), 0, 0, 0, 0, 0))
            continue
        rows.append({"範疇": name})     # section 標題行（跟報告 IMG_0105：博彩項目 / 非博彩項目）
        for _, row in sc.iterrows():
            rows.append(mk(row["_sub"], n_by_sub.get(row["_sub"], row["項目數量"]),
                           plan_by_sub.get(row["_sub"], 0.0),
                           row["報告"], row["調整"], row["後"], row["設施"], row["活動"]))
        # ⚠ 項目數量：小計/總計要用【去重】distinct，唔可以逐範疇加總 ——
        #   一個項目跨兩個範疇會被計兩次（user 2026-08-12 對數揭到）
        n_sc = n_by_scope.get(scope) or d[d["_scope"] == scope]["dicj code"].nunique()
        n_tot += n_sc                # ⚠ 總計唔可以全表 nunique：博彩「項目19」同非博彩「項目19」
                                     #   係兩個唔同項目（撞號），全表去重會少計 → 總計 = 兩個 scope 之和
        rows.append(mk(f"{name}小計", n_sc, _plan_tot(plan, yr, scope == 0),
                       sc["報告"].sum(), sc["調整"].sum(), sc["後"].sum(), sc["設施"].sum(), sc["活動"].sum()))
    # 尾行字眼跟報告：1.2 用「總計」（scan p10）、期後 2.1／2.3 用「合計」（scan p19-20）
    rows.append(mk("總計" if is_py else "合計", n_tot, _plan_tot(plan, yr, None),
                   g["報告"].sum(), g["調整"].sum(), g["後"].sum(), g["設施"].sum(), g["活動"].sum()))
    if is_py:
        # ★ scan p10：1.2 只有 5 條數字欄，冇完成率欄（完成率係表尾嘅【行】）、亦冇設施/活動欄。
        #   兩條完成率欄留喺【最後】俾下游文字邏輯用，_overview_extra 出表前會 drop 走。
        cols = ["範疇", "項目數量", "獲批的計劃投資金額", "報告投資金額",
                "投資金額的潛在調整事項", "潛在調整後投資金額",
                "投資計劃完成率", "潛在調整後投資計劃完成率"]
    else:
        cols = ["範疇", "項目數量", "報告投資金額", "潛在調整金額", "潛在調整後投資金額",
                "設施建設/資本性支出", "活動舉辦/營運性支出"]
    return pd.DataFrame(rows)[cols]


# ── from build_overview_tables ──
ADJ_SUPER = "投資金額的潛在調整事項"


# ── from build_overview_tables ──
ADJ_RESID = ADJ_POST


# ── from build_overview_tables ──
def adjustment_by_sub(df, bucket):
    """範疇 × 七大類調整 → DataFrame（`·` = 兩層表頭）。
       欄：報告投資金額(a) | 七大類(1..7)[+其他] | 潛在調整合計(b) | 潛在調整後投資金額(c=a+b)
           | 潛在調整金額佔報告投資金額比例(d=b/a)
    ⚠ b 一定等於【實際調整合計】（含殘差），咁 c=a+b 先會 tie 返概況表個「潛在調整後」。"""
    d = df[df["_bucket"] == bucket].copy()
    if d.empty:
        return pd.DataFrame()
    d["_adj"] = d["調整一級"].map(CANON).fillna(d["調整一級"])
    d.loc[~d["_adj"].isin(ADJ7), "_adj"] = ADJ_RESID
    d["_rep"] = pd.to_numeric(d["調整前_萬"], errors="coerce").fillna(0.0)
    d["_chg"] = pd.to_numeric(d["調整_萬"], errors="coerce").fillna(0.0)
    # 只出【有金額】嘅調整類；全 0 嗰幾類唔出欄（scan p15 出 1-7、p21 只出 1,2,4,5,6,8）
    types = [t for t in ADJ_ALL if abs(d.loc[d["_adj"] == t, "_chg"].sum()) > 0.5]
    cols = (["範疇", "報告投資金額"] + [f"{ADJ_SUPER}·{t}" for t in types]
            + ["潛在調整合計", "潛在調整後投資金額", "潛在調整金額佔報告投資金額比例"])

    def line(name, sub):
        rep = sub["_rep"].sum()
        r = {"範疇": name, "報告投資金額": round(rep, 1)}
        tot = 0.0
        for t in types:
            v = sub.loc[sub["_adj"] == t, "_chg"].sum()
            r[f"{ADJ_SUPER}·{t}"] = round(v, 1); tot += v
        r["潛在調整合計"] = round(tot, 1)
        r["潛在調整後投資金額"] = round(rep + tot, 1)
        r["潛在調整金額佔報告投資金額比例"] = _rate(tot, rep) if abs(rep) > 0.05 else None
        return r

    rows = []
    for scope in [0, 1]:
        sc = d[d["_scope"] == scope]
        name = "博彩項目" if scope == 0 else "非博彩項目"
        rows.append({"範疇": name})
        if not sc.empty:
            for sub in sc.sort_values(["_go", "_ngn", "_sub"])["_sub"].unique():
                rows.append(line(sub, sc[sc["_sub"] == sub]))
        rows.append(line(f"{name}小計", sc))
    rows.append(line("合計", d))
    # 報告最後一行＝涉及項目數量（逐欄：該調整類型涉及幾多個 distinct 項目）
    n = {"範疇": "涉及項目數量",
         "報告投資金額": int(d["dicj code"].nunique())}
    for t in types:
        sub = d[(d["_adj"] == t) & (d["_chg"].abs() > 0.05)]
        n[f"{ADJ_SUPER}·{t}"] = int(sub["dicj code"].nunique())
    n["潛在調整合計"] = int(d[d["_chg"].abs() > 0.05]["dicj code"].nunique())
    n["潛在調整後投資金額"] = ""
    n["潛在調整金額佔報告投資金額比例"] = ""
    rows.append(n)
    return pd.DataFrame(rows)[cols]


# ── from build_overview_tables ──
def overview_formula_row(cols):
    """期後概覽表（2.1／2.3）表頭下面嗰行斜體公式：a | b | c=a+b（對 scan p20）。"""
    m = {"報告投資金額": "a", "潛在調整金額": "b", "潛在調整後投資金額": "c=a+b"}
    return [m.get(str(c).split("·")[-1], "") for c in cols]


# ── from build_overview_tables ──
def adj_formula_row(cols):
    """報告表頭下面嗰行斜體公式：a | 調整類【固定編號】 | b | c=a+b | d=b/a（對 scan p15／p21）。
    ⚠ 編號要跟 ADJ_ALL 嘅位置，唔可以 1,2,3… 順住數 —— 報告 skip 咗嘅類會斷號（p21 = 1,2,4,5,6,8）。"""
    out = []
    for c in cols:
        c = str(c)
        if c == "範疇":
            out.append("")
        elif c == "報告投資金額":
            out.append("a")
        elif c.startswith(ADJ_SUPER + "·"):
            out.append(str(adj_no(c.split("·", 1)[1])))
        elif c == "潛在調整合計":
            out.append("b")
        elif c == "潛在調整後投資金額":
            out.append("c=a+b")
        elif c.endswith("比例"):
            out.append("d=b/a")
        else:
            out.append("")
    return out


# ── from build_overview_tables ──
def adjustment_bridge(df):
    """S15-17：7 canonical 調整類型 × {2025計劃/2024期後/2023期後/合計}。"""
    d = df.copy()
    d["_adj"] = d["調整一級"].map(CANON).fillna(d["調整一級"])
    d.loc[~d["_adj"].isin(ADJ7), "_adj"] = ADJ_POST      # 殘差＝第 8 類
    rows = []
    for adj in ADJ_ALL:
        r = {"潛在調整事項": adj}
        for bk in BUCKET_ORDER:
            r[bk] = round(d[(d["_bucket"] == bk) & (d["_adj"] == adj)]["調整_萬"].sum(), 1)
        r["合計"] = round(sum(r[bk] for bk in BUCKET_ORDER), 1)
        rows.append(r)
    # 跨年及其他調整（唔喺報告 7 類，多數係期後嘅跨期/將往年計入本年）→ 令 bucket 合計對返概況潛在調整
    other = {"潛在調整事項": "跨年及其他調整"}
    for bk in BUCKET_ORDER:
        tot_bk = round(d[d["_bucket"] == bk]["調整_萬"].sum(), 1)
        other[bk] = round(tot_bk - sum(x[bk] for x in rows), 1)
    other["合計"] = round(sum(other[bk] for bk in BUCKET_ORDER), 1)
    if any(abs(other[bk]) > 0.05 for bk in BUCKET_ORDER):
        rows.append(other)
    tot = {"潛在調整事項": "合計"}
    for bk in BUCKET_ORDER + ["合計"]:
        tot[bk] = round(sum(x[bk] for x in rows), 1)
    rows.append(tot)
    return pd.DataFrame(rows, columns=["潛在調整事項"] + BUCKET_ORDER + ["合計"])


# ── from build_overview_tables ──
def zero_investment_summary(df, plan, cat, narr, ent_up="MGM"):
    """報告概述尾段：2025計劃申報投資為零嘅非博彩項目（原計劃有金額、2025 實際=0）分類。
    → (n, 總計劃_萬, groups)；groups={跨年/內部研究/取消: [範疇…]}。分類靠 期後有支出=跨年、
    項目狀況含取消=取消、其餘=內部研究（best-effort，項目組口徑為準）。"""
    if not plan or not plan.get(25):
        return None
    from collections import Counter

    def _key(r):
        return (str(r["ng_scope"]) == "gaming", _norm(r["dicj code"]))

    def _has(sub):
        return {_key(r) for _, r in sub.iterrows()
                if abs(pd.to_numeric(r.get("調整前_萬", 0), errors="coerce") or 0) > 0.05}

    spent25 = _has(df[df["_bucket"] == "2025年度投資計劃"])
    post = _has(df[df["_bucket"].isin(["2024年度計劃期後投資", "2023年度計劃期後投資"])])
    groups = {"跨年": [], "內部研究": [], "取消": []}
    tot = 0.0
    for (gm, code), ev in (plan.get(25) or {}).items():
        if gm or (ev or 0) <= 0.05 or (gm, code) in spent25:
            continue        # 只計非博彩、原計劃>0、2025計劃無實際支出
        rec = (narr or {}).get((gm, code)) or (narr or {}).get((not gm, code)) or {}
        status = str(rec.get("項目狀況", ""))
        sub = (cat or {}).get((gm, code), "") or "其他"
        kind = "跨年" if (gm, code) in post else ("取消" if "取消" in status else "內部研究")
        groups[kind].append(sub)
        tot += ev
    n = sum(len(v) for v in groups.values())
    if n == 0:
        return None
    return n, round(tot, 1), groups


# ── from build_overview_tables ──
def zero_investment_text(zi, ent_up="MGM"):
    """zero_investment_summary → 報告式段落文字。"""
    if not zi:
        return None
    from collections import Counter
    n, tot, groups = zi
    lines = [f"{ent_up}有{n}個非博彩項目，原計劃項目於2025年度投資執行報告中申報的投資支出為零"
             f"（原計劃金額約{tot:,.0f}萬澳門元），主要包括："]
    lab = {"跨年": "個項目的2025年支出被申報為2023／2024年度計劃在2025年的期後投資金額（跨年項目）",
           "內部研究": "個項目仍處於內部研究、重新規劃或選址階段，未發生實際支出",
           "取消": "個項目已取消"}
    for kind in ("跨年", "內部研究", "取消"):
        g = groups.get(kind) or []
        if not g:
            continue
        c = Counter(g)
        detail = "、".join(f"{k}{v}個" for k, v in c.items())
        lines.append(f"{len(g)}{lab[kind]}（涉及{detail}）")
    return lines


# ── from build_overview_tables ──
def finding_summary(df):
    """S28-40：每個 canonical 調整類型 → 調整額合計 / 涉及項目數 / 主要涉及項目(top3)。"""
    d = df.copy()
    d["_adj"] = d["調整一級"].map(CANON).fillna(d["調整一級"])
    d.loc[~d["_adj"].isin(ADJ7), "_adj"] = ADJ_POST      # 殘差＝第 8 類
    rows = []
    for adj in ADJ_ALL:
        sub = d[(d["_adj"] == adj) & (pd.to_numeric(d["調整_萬"], errors="coerce") != 0)]
        if sub.empty:
            continue
        projs = sub.groupby("project")["調整_萬"].sum().abs().sort_values(ascending=False)
        rows.append({"潛在調整事項": adj, "調整額合計": round(sub["調整_萬"].sum(), 1),
                     "涉及項目數": int(sub["dicj code"].nunique()),
                     "主要涉及項目": "、".join(str(p) for p in projs.index[:3])})
    return pd.DataFrame(rows, columns=["潛在調整事項", "調整額合計", "涉及項目數", "主要涉及項目"])


# ── from workbench ──
_RETRY_STATUS = {408, 409, 429, 500, 502, 503, 504}


# ── from workbench ──
_TRANSIENT = ("too many requests", "rate limit", "timeout", "timed out")


# ── from workbench ──
_BACKOFF = (3, 9, 25)


# ── from workbench ──
_RETRIES = len(_BACKOFF)


# ── from workbench ──
def _retryable(e):
    sc = getattr(e, "status_code", None) or getattr(getattr(e, "response", None), "status_code", None)
    if sc is not None:
        return sc in _RETRY_STATUS
    return any(k in str(e).lower() for k in _TRANSIENT)


# ── from workbench ──
DEFAULT_PROVIDER = "azure"


# ── from workbench ──
DEFAULT_BASE_URL = "https://api.workbench.kpmg/genai/azure/openai"


# ── from workbench ──
DEFAULT_API_VERSION = "2024-12-01-preview"


# ── from workbench ──
DEFAULT_MODEL = "5.5"


# ── from workbench ──
DEFAULT_MODELS = {
    "5.5": "gpt-5-5-2026-04-24-gs-sdc",
    "5.4": "gpt-5-4-2026-03-05-gs-sdc",
}


# ── from workbench ──
_CRED = Path("conf/local/credentials.yml")


# ── from workbench ──
def _cred_section() -> dict:
    if not _CRED.exists():
        return {}
    try:
        import yaml
        d = yaml.safe_load(_CRED.read_text(encoding="utf-8")) or {}
        return dict(d.get("workbench") or {})
    except Exception:
        return {}


# ── from workbench ──
def _resolve(name: str, env: str, default: str = "") -> str:
    """env > config > default（字串）。"""
    v = os.environ.get(env, "").strip()
    if v:
        return v
    v = str(_cred_section().get(name, "") or "").strip()
    return v or default


# ── from workbench ──
def _parse_json_lenient(text: str) -> dict:
    s = (text or "").strip()
    if s.startswith("```"):
        s = re.sub(r"^```[a-zA-Z]*\s*", "", s)
        s = re.sub(r"\s*```$", "", s).strip()
    try:
        return json.loads(s)
    except json.JSONDecodeError:
        m = re.search(r"\{[\s\S]*\}", s)
        if not m:
            raise
        return json.loads(m.group(0))


# ── from workbench ──
class Workbench:
    def __init__(self, model: str | None = None, *, provider: str | None = None,
                 charge_code: str | None = None, region: str | None = None,
                 base_url: str | None = None, api_version: str | None = None,
                 api_key: str | None = None, models: dict | None = None):
        sec = _cred_section()
        # alias → 實際名 對照表：內置 ← config.models 覆蓋 ← 呼叫者覆蓋
        self.models = {**DEFAULT_MODELS, **(sec.get("models") or {}), **(models or {})}
        self.provider = (provider or _resolve("provider", "WB_PROVIDER", DEFAULT_PROVIDER)).lower()
        self.model_alias = model or _resolve("model", "WB_MODEL", DEFAULT_MODEL)
        self.base_url = base_url or _resolve("base_url", "WB_BASE_URL", DEFAULT_BASE_URL)
        self.api_version = api_version or _resolve("api_version", "WB_API_VERSION", DEFAULT_API_VERSION)
        self.charge_code = charge_code or _resolve("charge_code", "WB_CHARGE_CODE", "0000")
        self.region = region or _resolve("region", "WB_REGION", "westeurope")
        self._api_key = (api_key or "").strip() or None
        self._client = None

    def resolve_model(self, model: str | None = None) -> str:
        """alias → 實際 deployment/model 名（唔喺對照表就當已經係實名）。"""
        m = model or self.model_alias
        return self.models.get(m, m)

    @property
    def model(self) -> str:
        return self.resolve_model()

    # ── key / headers ────────────────────────────────────────────────
    @property
    def api_key(self) -> str:
        if not self._api_key:
            k = _resolve("api_key", "WB_API_KEY")
            if not k:
                raise RuntimeError(
                    "冇 API key：set env WB_API_KEY，"
                    "或 conf/local/credentials.yml 加 workbench.api_key（gitignored）")
            self._api_key = k
        return self._api_key

    def _headers(self) -> dict:
        # KPMG Workbench(azure) 專用 header；openai-compatible（qwen）唔加，SDK 自己用 api_key。
        if self.provider == "azure":
            return {
                "Ocp-Apim-Subscription-Key": self.api_key,
                "x-kpmg-charge-code": self.charge_code,
                "x-kpmg-region-override": self.region,
            }
        return {}

    def _client_obj(self):
        if self._client is None:
            if self.provider == "azure":
                from openai import AzureOpenAI
                self._client = AzureOpenAI(
                    api_key=self.api_key, base_url=self.base_url,
                    api_version=self.api_version, default_headers=self._headers() or None)
            else:  # openai-compatible（qwen 等）
                from openai import OpenAI
                self._client = OpenAI(
                    api_key=self.api_key, base_url=self.base_url,
                    default_headers=self._headers() or None)
        return self._client

    # ── chat ─────────────────────────────────────────────────────────
    def chat(self, user: str, system: str = "You are a helpful assistant.", *,
             model: str | None = None, reasoning_effort: str | None = "high",
             temperature: float | None = 1, max_tokens: int | None = None,
             json_mode: bool = False) -> str:
        """一問一答，回 content。對唔支持嘅參數自動 drop 再試，唔會因 model/provider 差異炒。"""
        cli = self._client_obj()
        kw: dict[str, Any] = {
            "model": self.resolve_model(model),
            "messages": [
                {"role": "system", "content": system},
                {"role": "user", "content": user},
            ],
        }
        hdr = self._headers()
        if hdr:
            kw["extra_headers"] = {"Content-Type": "application/json", **hdr}
        if temperature is not None:
            kw["temperature"] = temperature
        if reasoning_effort:
            kw["reasoning_effort"] = reasoning_effort
        if max_tokens:
            kw["max_tokens"] = max_tokens
        if json_mode:
            kw["response_format"] = {"type": "json_object"}

        droppable = ["reasoning_effort", "temperature", "response_format", "max_tokens"]
        tries = 0
        for _ in range(len(droppable) + 1 + _RETRIES):
            try:
                resp = cli.chat.completions.create(**kw)
                return resp.choices[0].message.content or ""
            except Exception as e:
                msg = str(e).lower().replace("_", "")
                hit = next((p for p in droppable if p in kw and p.replace("_", "") in msg), None)
                if hit is not None:
                    kw.pop(hit, None); continue
                # 塞車／上游 5xx → 退避重試；401/403 即刻放棄（重試極都一樣）
                if tries < _RETRIES and _retryable(e):
                    time.sleep(_BACKOFF[tries]); tries += 1; continue
                raise
        raise RuntimeError("chat 失敗：所有可 drop 參數都試過")

    def chat_json(self, user: str, system: str = "You are a helpful assistant. Reply with JSON only.",
                  **kw) -> dict:
        kw.setdefault("json_mode", True)
        return _parse_json_lenient(self.chat(user, system, **kw))

    def ping(self, model: str | None = None) -> str:
        return self.chat("hi", model=model)

    # ── config 檢視（key 遮蔽）────────────────────────────────────────
    def config_masked(self) -> dict:
        try:
            k = self.api_key
            km = f"{k[:4]}…({len(k)} chars)" if k else "(缺)"
            key_ok = True
        except Exception as e:
            km, key_ok = f"⚠ {e}", False
        return {
            "provider": self.provider,
            "model_alias": self.model_alias, "model(resolved)": self.model,
            "models_map": self.models,
            "base_url": self.base_url, "api_version": self.api_version,
            "charge_code": self.charge_code, "region": self.region,
            "api_key": km, "key_ok": key_ok,
        }


# ── from workbench ──
def _cli() -> None:
    import argparse
    ap = argparse.ArgumentParser(description="LLM client（databook 用，config-driven）")
    ap.add_argument("--model", default=None, help="alias（5.5/5.4/qwen…）或完整名；預設由 config")
    ap.add_argument("--provider", default=None, help="azure / openai（覆蓋 config）")
    ap.add_argument("--ping", action="store_true", help="送 'hi' 驗連線（要喺 KPMG 網內跑）")
    ap.add_argument("--config", action="store_true", help="offline：只印解析到嘅 config（key 遮蔽）")
    ap.add_argument("--ask", help="送一句 user prompt 睇回覆")
    a = ap.parse_args()
    wb = Workbench(model=a.model, provider=a.provider)
    if a.config or not (a.ping or a.ask):
        print("LLM config（offline，唔出網）:")
        for k, v in wb.config_masked().items():
            print(f"  {k}: {v}")
        if not (a.ping or a.ask):
            print("\n→ 連線測試喺 KPMG 網內跑： python -m kpi.lib.workbench --ping")
            return
    if a.ping:
        print("\n--ping →", wb.ping())
    if a.ask:
        print("\n--ask →", wb.chat(a.ask))


# ── from inspect_biao2 ──
try:
    import openpyxl
    from openpyxl.utils import get_column_letter
except ImportError:
    print("✗ pip install openpyxl"); sys.exit(1)


# ── from inspect_biao2 ──
PASSWORD = "dicj_kpmg"


# ── from inspect_biao2 ──
def load_wb(path, password=PASSWORD):
    """開 xlsx；『not a zip file』＝加密 → msoffcrypto 用密碼解。回 openpyxl workbook。"""
    try:
        return openpyxl.load_workbook(path, data_only=True, read_only=True)
    except Exception:
        pass
    import msoffcrypto
    buf = io.BytesIO()
    with open(path, "rb") as f:
        off = msoffcrypto.OfficeFile(f)
        off.load_key(password=password)
        off.decrypt(buf)
    buf.seek(0)
    return openpyxl.load_workbook(buf, data_only=True, read_only=True)


# ── from inspect_biao2 ──
KEY_HINT = ["投資項目序號", "項目序號及名稱", "項目序號"]


# ── from inspect_biao2 ──
NARR_HINT = ["KPMG分析", "關注事項", "管理層解釋", "調整金額", "調整原因", "分析意見",
             "反饋意見", "調整後金額", "項目分類", "備註", "狀態"]


# ── from inspect_biao2 ──
def _find_header(rows):
    """揾 detail header row（含投資項目序號）；上一行＝group header。回 (grp_r, det_r)。"""
    for ri in range(min(12, len(rows))):
        joined = "".join(str(v) for v in (rows[ri] or []) if v)
        if any(h in joined for h in KEY_HINT):
            return (ri - 1 if ri > 0 else ri), ri
    return None, None


# ── from inspect_biao2 ──
def _ffill(seq):
    out, last = [], ""
    for v in seq:
        s = "" if v is None else str(v).replace("\n", "").strip()
        if s:
            last = s
        out.append(last)
    return out


# ── from inspect_biao2 ──
def inspect_one(path: Path):
    print(f"\n{'#'*84}\n# {path.name}")
    try:
        wb = load_wb(path)
    except ImportError:
        print("  ✗ 加密檔要 msoffcrypto → pip install msoffcrypto-tool"); return
    except Exception as e:
        print(f"  ✗ 開唔到: {type(e).__name__}: {e}"); return
    for sn in wb.sheetnames:
        ws = wb[sn]
        rows = []
        for i, r in enumerate(ws.iter_rows(values_only=True)):
            rows.append(r)
            if i > 800:
                break
        grp_r, det_r = _find_header(rows)
        print(f"\n{'='*72}\n## sheet {sn!r}  約 {len(rows)} 行"
              + (f"  header：組 r{grp_r+1} / 欄 r{det_r+1}" if det_r is not None else "  ⚠揾唔到表頭"))
        if det_r is None:
            for ri in range(min(6, len(rows))):
                cells = [("" if v is None else str(v).replace("\n", " ")[:20]) for v in (rows[ri] or [])]
                print(f"  r{ri}:", " | ".join(c for c in cells if c)[:200])
            continue
        grp = _ffill(rows[grp_r]) if grp_r is not None and grp_r >= 0 else [""] * len(rows[det_r])
        det = [("" if v is None else str(v).replace("\n", "").strip()) for v in rows[det_r]]
        ncol = max(len(grp), len(det))
        ndata = sum(1 for ri in range(det_r + 1, len(rows))
                    if rows[ri] and any(v not in (None, "") for v in rows[ri]))
        print(f"  資料行約 {ndata}；共 {ncol} 欄：")
        for ci in range(ncol):
            g = grp[ci] if ci < len(grp) else ""
            dcol = det[ci] if ci < len(det) else ""
            sample = ""
            for ri in range(det_r + 1, min(det_r + 60, len(rows))):
                row = rows[ri]
                if row and ci < len(row) and row[ci] not in (None, ""):
                    sample = str(row[ci]).replace("\n", " ").strip()[:50]; break
            if not (g or dcol or sample):
                continue
            mark = ""
            if any(h in (dcol) for h in KEY_HINT):
                mark = "  ⟸ KEY(項目)"
            elif any(h in (g + dcol) for h in NARR_HINT):
                mark = "  ⟸ finding/調整"
            print(f"    {get_column_letter(ci+1):>3} | {g[:16]:<16} | {dcol[:24]:<24} | {sample}{mark}")


# ── from biao2 ──
_CODE_RE = re.compile(r"^(項目\s*\d+|[A-Za-z]{1,5}\d+(?:\.\d+)?)$")


# ── from biao2 ──
_JUNK_RE = re.compile(r"^(無新增問題|無|是|否|已回覆|未回覆|不適用|n/?a|請參閱附件)")


# ── from biao2 ──
_ENT_ALIASES = {
    "mgm": ["mgm", "美高梅"],
    "galaxy": ["galaxy", "銀河"],
    "sjm": ["sjm", "澳博", "新葡京"],
    "wynn": ["wynn", "永利"],
    "vml": ["vml", "威尼斯", "金沙"],
    "melco": ["melco", "新濠"],
}


# ── from biao2 ──
def _match_entity(fname, entity):
    fl = fname.lower()
    for a in _ENT_ALIASES.get(entity, [entity]):
        if a.lower() in fl:
            return True
    return False


# ── from biao2 ──
def _txt(v):
    return re.sub(r"\s+", " ", str(v if v is not None else "")).strip()


# ── from biao2 ──
B2_FIELDS = {
    "關注事項": ["畢馬威關注事項"],                 # ★ 調整嘅正文
    "調整類型": ["需溝通關注事項"],                 # e.g.「一般支持性部門的人工成本」
    "關注事項金額": ["該關注事項涉及調整金額"],
    "建議調整金額": ["建議調整金額"],
    "調整原因": ["調整原因"],
    "建議調整後金額": ["建議調整後金額"],
    "KPMG分析": ["KPMG分析"],                     # 第一/二輪諮詢各一，兩個都收
    "管理層解釋": ["承批公司管理層解釋"],
    "承批公司反饋": ["承批公司的反饋意見"],
    "跨司回覆": ["跨司工作組的回覆", "跨司工作組最新反饋意見", "跨司工作組的反饋意見",
                 "跨司工作組主責部門"],
    # 跨司工作組審閱意見 block：實測 mgm 成欄係空（跨司未填，標題仲係 2026.07.XX）→ 抽到 0 係正常
    "項目分析意見": ["項目分析意見"],
    "建議接納調整後金額": ["建議接納之調整後金額"],
    # ⚠ 唔好加「擬投資金額／已投放金額」：嗰啲係【該性質範疇各項目之加總金額】，唔掛喺項目碼上，
    #   forward-fill 落去會亂咁派。per-project 投資內容一律由清單『實際投資內容』攞。
}


# ── from biao2 ──
B2_ORDER = ["調整類型", "關注事項", "調整原因", "關注事項金額", "建議調整金額",
            "建議調整後金額", "KPMG分析", "管理層解釋", "承批公司反饋", "跨司回覆",
            "項目分析意見", "建議接納調整後金額"]


# ── from biao2 ──
_NUMISH = {"關注事項金額", "建議調整金額", "建議調整後金額", "建議接納調整後金額"}


# ── from biao2 ──
_HDR_HINT = [k for ks in B2_FIELDS.values() for k in ks] + ["項目編號", "資料要求", "問題狀態"]


# ── from biao2 ──
def _detail_header_row(rows, band=10):
    """detail 表頭行 = 頭 band 行入面 match 到最多概念嗰行（group 行喺佢上面）。"""
    best, bn = 0, 0
    for ri in range(min(band, len(rows))):
        n = sum(1 for v in (rows[ri] or []) if any(k in _txt(v) for k in _HDR_HINT))
        if n > bn:
            bn, best = n, ri
    return best, bn


# ── from biao2 ──
def _field_cols(hdr):
    """detail 表頭 → {概念: [col_idx]}。"""
    cols = {}
    for ci, v in enumerate(hdr):
        s = _txt(v)
        if not s:
            continue
        for concept, keys in B2_FIELDS.items():
            if any(k in s for k in keys):
                cols.setdefault(concept, []).append(ci)
    return cols


# ── from biao2 ──
def load_biao2_struct(folder, entity, log=lambda *a: None):
    """{(gaming, 正規化碼): {概念: 文字}} —— 按【欄名】structured 抽（唔再盲抓）。
    表2 layout：group 行 + detail 表頭行（概念名）+ 逐項目 data 行；欄序浮動 → 全部按名認。
    同一個碼有多行（matrix）→ 每個概念收最長嗰個非空值，其餘唔同值就接落去。"""
    out = {}
    d = Path(folder)
    if not d.exists():
        log(f"（冇 {folder}）"); return out
    allx = [p for p in sorted(d.rglob("*.xls*")) if not p.name.startswith("~$")]
    files = [p for p in allx if _match_entity(p.name, entity.lower()) and "提供附件" not in p.name]
    log(f"表2 folder {folder}：共 {len(allx)} 個 xls*，match「{entity}」{len(files)} 檔")
    n_field = 0
    for p in files:
        gaming = ("博監局" in p.name)
        try:
            wb = load_wb(p)
        except Exception as e:
            log(f"  ⚠ 開唔到 {p.name}: {e}"); continue
        for sn in wb.sheetnames:
            try:
                rows = []
                for i, r in enumerate(wb[sn].iter_rows(values_only=True)):
                    rows.append(r)
                    if i > 700:
                        break
                ncol = max((len(r) for r in rows if r), default=0)
                if ncol == 0:
                    continue
                hr, nmatch = _detail_header_row(rows)
                if nmatch < 3:
                    continue                     # 唔似標準表2 sheet（附件/圖片頁）→ 跳
                fcols = _field_cols(rows[hr])
                if not fcols:
                    continue
                # code 欄：優先『項目編號』表頭（實測佢喺 group 行，唔喺 detail 行）→ 兩行都揾；
                # 都揾唔到就用 _CODE_RE 命中最多嗰欄
                hdr_band = list(rows[hr]) + list(rows[hr - 1] if hr else [])
                code_c = next((ci % max(ncol, 1) for ci, v in enumerate(hdr_band)
                               if "項目編號" in _txt(v)), None)
                if code_c is None:
                    best, bestn = None, 0
                    for ci in range(ncol):
                        n = sum(1 for r in rows if ci < len(r) and r[ci] is not None
                                and _CODE_RE.match(re.sub(r"\s+", "", str(r[ci]))))
                        if n > bestn:
                            bestn, best = n, ci
                    code_c = best
                if code_c is None:
                    continue
                log(f"  · {p.name}｜{sn}：detail表頭 r{hr+1}、code欄 col{code_c}、"
                    f"{len(fcols)} 個概念（{'博彩' if gaming else '非博彩'}）")
                for r in rows[hr + 1:]:
                    if code_c >= len(r) or r[code_c] is None:
                        continue
                    if not _CODE_RE.match(re.sub(r"\s+", "", str(r[code_c]))):
                        continue
                    rec = out.setdefault((gaming, _norm(r[code_c])), {})
                    for concept, cis in fcols.items():
                        for ci in cis:
                            if ci >= len(r):
                                continue
                            s = _txt(r[ci])
                            if len(s) < 3 or _JUNK_RE.match(s):
                                continue
                            cur = rec.get(concept, "")
                            if s in cur:
                                continue
                            sep = "／" if concept in _NUMISH else "　"
                            rec[concept] = (cur + sep + s).strip(sep) if cur else s
                            n_field += 1
            except Exception as e:
                log(f"  ⚠ {p.name}｜{sn}: {e}")
    n_all = len(out)
    out = {k: v for k, v in out.items() if v}      # 冇任何欄值 = 該項目冇 finding，唔留空 key
    log(f"表2（structured）：{len(files)} 檔 → {len(out)} 個項目有內容"
        f"（另 {n_all - len(out)} 個碼冇 finding）、{n_field} 個欄值")
    return out


# ── from biao2 ──
ART_FIELDS = ["購入時間", "Artwork", "名稱", "類別", "產權歸屬", "Artist",
              "購入", "添置原因", "展出紀錄", "當前位置", "當前狀態", "未來處置"]


# ── from biao2 ──
def load_artwork(folder, entity, log=lambda *a: None):
    """表2 附件『藝術品』sheet → (欄名 list, 行 list)。冇就回 (None, [])。"""
    d = Path(folder)
    if not d.exists():
        return None, []
    for p in sorted(d.rglob("*.xls*")):
        if p.name.startswith("~$") or not _match_entity(p.name, entity.lower()):
            continue
        try:
            wb = load_wb(p)
        except Exception:
            continue
        for sn in wb.sheetnames:
            if "藝術品" not in sn and "艺术品" not in sn:
                continue
            rows = [r for i, r in enumerate(wb[sn].iter_rows(values_only=True)) if i < 400]
            hr = next((i for i, r in enumerate(rows[:8])
                       if sum(1 for v in (r or []) if any(k in _txt(v) for k in ART_FIELDS)) >= 4), None)
            if hr is None:
                continue
            hdr = [_txt(v) for v in rows[hr]]
            keep = [i for i, h in enumerate(hdr) if h]
            body = []
            for r in rows[hr + 1:]:
                vals = [_txt(r[i]) if i < len(r) else "" for i in keep]
                if sum(1 for v in vals if v) >= 3:
                    body.append(vals)
            log(f"  · 藝術品清單：{p.name}｜{sn} → {len(body)} 件 × {len(keep)} 欄")
            return [hdr[i] for i in keep], body
    return None, []


# ── from biao2 ──
def b2rec(b2s, ng_scope, code):
    """由 (ng_scope, code) 攞 structured rec（撞號先試 exact）。"""
    g = (ng_scope == "gaming")
    c = _norm(code)
    return b2s.get((g, c)) or b2s.get((not g, c)) or {}


# ── from biao2 ──
def b2text(b2s, ng_scope, code, limit=1400):
    """structured rec → 有 label 嘅文字（餵 LLM 用；比盲抓清楚好多）。"""
    rec = b2rec(b2s, ng_scope, code)
    if not rec:
        return ""
    parts = [f"{k}：{rec[k]}" for k in B2_ORDER if rec.get(k)]
    s = "；".join(parts)
    return s[:limit]


# ── from biao2 ──
def load_biao2(folder, entity, log=lambda *a: None):
    """{(gaming, 正規化碼): [finding 文字…]}。best-effort，逐檔逐 sheet try。
    （舊版盲抓；新 code 請用 load_biao2_struct + b2text。）"""
    out = {}
    d = Path(folder)
    if not d.exists():
        log(f"（冇 {folder}）"); return out
    allx = [p for p in sorted(d.rglob("*.xls*")) if not p.name.startswith("~$")]
    files = [p for p in allx if _match_entity(p.name, entity.lower()) and "提供附件" not in p.name]
    log(f"表2 folder {folder}：共 {len(allx)} 個 xls*，match「{entity}」{len(files)} 檔")
    for p in files:
        log(f"  ✓ {p.name}")
    if not files and allx:
        log(f"  ⚠ 一個都 match 唔到！全部檔名：" + " | ".join(p.name for p in allx[:20]))
    for p in files:
        gaming = ("博監局" in p.name)     # 博監局檔 = 博彩項目；其餘範疇 = 非博彩
        try:
            wb = load_wb(p)
        except Exception as e:
            log(f"  ⚠ 開唔到 {p.name}: {e}"); continue
        for sn in wb.sheetnames:
            try:
                ws = wb[sn]
                rows = []
                for i, r in enumerate(ws.iter_rows(values_only=True)):
                    rows.append(r)
                    if i > 700:
                        break
                ncol = max((len(r) for r in rows if r), default=0)
                if ncol == 0:
                    continue
                # code 欄 = 匹配 _CODE_RE 最多嘅欄
                best, bestn = None, 0
                for ci in range(ncol):
                    n = sum(1 for r in rows if ci < len(r) and r[ci] is not None
                            and _CODE_RE.match(re.sub(r"\s+", "", str(r[ci]))))
                    if n > bestn:
                        bestn, best = n, ci
                if best is None or bestn < 2:
                    continue
                log(f"  · {p.name}｜{sn}：code欄 col{best}（{bestn}個碼，{'博彩' if gaming else '非博彩'}）")
                for r in rows:
                    if best < len(r) and r[best] is not None \
                            and _CODE_RE.match(re.sub(r"\s+", "", str(r[best]))):
                        code = _norm(r[best])
                        snips = []
                        for c in r:
                            if c is None:
                                continue
                            s = str(c).replace("\n", " ").strip()
                            if len(s) >= 30 and not _JUNK_RE.match(s):
                                snips.append(s[:400])
                        if snips:
                            out.setdefault((gaming, code), [])
                            for s in snips[:4]:
                                if s not in out[(gaming, code)]:
                                    out[(gaming, code)].append(s)
            except Exception as e:
                log(f"  ⚠ {p.name}｜{sn}: {e}")
    log(f"表2：{len(files)} 檔 → {len(out)} 個 (gaming,碼) 有 finding")
    return out


# ── from biao2 ──
def b2look(b2, ng_scope, code, joiner="　"):
    """由 (ng_scope, code) 攞表2 finding 文字（合併），撞號用 exact 先。"""
    g = (ng_scope == "gaming")
    c = _norm(code)
    snips = b2.get((g, c)) or b2.get((not g, c)) or []
    return joiner.join(snips)


# ── from build_llm_narrative ──
try:
    import pandas as pd
except ImportError:
    print("✗ pip install pandas openpyxl"); sys.exit(1)


# ── from build_llm_narrative ──
ADJ_TAIL = ("輸出淨係連貫文字（唔好 markdown 標題／項目符號／開場白／結語），可以分 2 段；"
            "內容要齊全，唔好為咗短而略去項目、金額或理據。")


# ── from build_llm_narrative ──
SYS_ADJ = ("你係畢馬威（KPMG）投資計劃執行情況審查報告嘅專業撰稿員。用【繁體中文】書面語，"
           "審查報告語氣：精簡、客觀、專業、第三人稱（用『我們』）。"
           "只可根據所提供嘅資料撰寫，嚴禁虛構、誇大或加入未提供嘅事實/數字。"
           "直接寫有嘅內容，切勿寫『未獲提供』『資料不足』等 meta/免責語句。" + ADJ_TAIL)


# ── from build_llm_narrative ──
SYS_CAT = ("你係為承批公司『2025年度投資執行報告』撰寫『按範疇的項目概況』嘅撰稿員。用【繁體中文】書面語，"
           "站喺投資執行角度：描述該範疇實際投資咗啲乜（可含具體項目、子項目、活動／賽事／音樂會場次），"
           "再講完成率點解係咁。語氣自然、貼近企業投資執行報告，唔好似審計底稿、唔好似機器砌 list。"
           "★嚴禁用審計／調整用語：『剔除』『調減』『申報口徑』『可計入範圍』『超出範圍』『再次申報』"
           "『偏離』『不符合定義』『未在計劃中列示』等 —— 呢啲屬另一節（調整事項），概況絕不出現。"
           "完成率原因只用管理層嘅業務解釋（例：施工進度較預期延遲、實際所需資金低於預算、"
           "進度高於預期、活動如期舉辦、發現結構性問題增加成本），唔好用審計理由。"
           "直接寫有嘅內容，切勿寫『未獲提供』『資料不足』等 meta 語。"
           "輸出淨係一段連貫文字（唔好項目符號／開場白／結語），語句要順，忌逐點堆砌。")


# ── from build_llm_narrative ──
SYS_TBL_ADJ = (SYS_ADJ.replace(ADJ_TAIL, "")
               + "你而家寫嘅係【一張報告表格旁邊嘅敘述】：解釋張表講緊乜、關鍵金額同背後原因。"
               "只可引用表格入面真係有嘅數字，唔可以自己計新數或估數。"
"★金額一律用返報告嘅正式名：『報告投資金額』『潛在調整後投資金額』『獲批的計劃投資金額』"
               "『潛在調減』。【嚴禁自創】報告冇嘅字眼，例如『經後續管理檢視後』『經覆核後』"
               "『管理層檢視』『調整後淨額』等 —— 一律用『潛在調整後投資金額』。"
               "輸出 JSON：{\"導語\":\"…\",\"段落\":[{\"小標題\":\"…\",\"內容\":\"…\"}]}。"
               "『導語』＝成版最頂嗰句總結（80-160 字，跟報告句式，見 prompt 內示範）；"
               "『段落』2 至 4 段，每段小標題 ≤14 字、內容 60-130 字。")


# ── from build_llm_narrative ──
SYS_TBL_DESC = (SYS_CAT.replace("輸出淨係一段連貫文字（唔好項目符號／開場白／結語），語句要順，忌逐點堆砌。", "")
                + "你而家寫嘅係【一張報告表格旁邊嘅敘述】：解釋張表講緊乜、邊啲範疇金額最大、"
                "設施建設同活動舉辦嘅比重點樣。只可引用表格入面真係有嘅數字，唔可以自己計新數或估數。"
                "★金額一律用返報告嘅正式名（報告投資金額／潛在調整後投資金額／獲批的計劃投資金額）；"
                "【嚴禁自創】『經後續管理檢視後』『經覆核後』等報告冇嘅字眼。"
 "★金額一律用返報告嘅正式名：『報告投資金額』『潛在調整後投資金額』『獲批的計劃投資金額』"
               "『潛在調減』。【嚴禁自創】報告冇嘅字眼，例如『經後續管理檢視後』『經覆核後』"
               "『管理層檢視』『調整後淨額』等 —— 一律用『潛在調整後投資金額』。"
               "輸出 JSON：{\"導語\":\"…\",\"段落\":[{\"小標題\":\"…\",\"內容\":\"…\"}]}。"
                "『導語』＝成版最頂嗰句總結（80-160 字）；『段落』2 至 4 段，"
                "每段小標題 ≤14 字、內容 60-130 字。")


# ── from build_llm_narrative ──
def tbl_key(kind, arg=""):
    """表旁 comment 嘅 key（generator 同 make_report 必須用同一個）。"""
    return f"{kind}|{arg}"


# ── from build_llm_narrative ──
def proj_key(adj_type, ng_scope, code):
    """主要發現 card 逐項目『事項描述』嘅 key（generator 同 make_report 必須用同一個）。"""
    return f"{adj_type}|{ng_scope}|{_norm(code)}"


# ── from build_llm_narrative ──
def bkt_key(bucket, adj_type):
    """期後調整事項匯總（scan p-11/p-13）逐類開場句嘅 key。"""
    return f"{bucket}|{adj_type}"


# ── from build_llm_narrative ──
def _bkt_prompt(yr, adj_type, amt_wan, projects):
    lines = [f"年度：{yr}年度投資計劃期後投資（於2025年發生）", f"調整類型：{adj_type}",
             f"該類潛在調減金額：約{abs(amt_wan):,.0f}萬澳門元", "涉及項目："]
    for nm, amt, b2, find in projects[:5]:
        seg = f"- {nm}（{abs(amt):,.0f}萬澳門元）"
        if b2:
            seg += f"；【審查底稿表2】{b2[:600]}"
        elif find:
            seg += f"；KPMG分析發現：{find[:300]}"
        lines.append(seg)
    return ("請寫【一句至兩句】開場描述，講清楚喺該年度期後投資金額中，承批公司申報咗啲乜"
            "而我哋認為要調整。示範句式（要用返下面嘅真數同項目，唔好照抄）：\n"
            "　『在2024年度投資計劃期後投資金額中，MGM申報了澳門美高梅國際旗艦級藝術珍寶博物館"
            "營運後的營運成本（827萬澳門元）。』\n"
            "　『在2024年度投資計劃期後投資金額中，MGM仍申報了酒店客房改造支出，主要包括："
            "1）非博彩項目111多功能娛樂體驗區塊（娛樂表演範疇）的相關支出420萬澳門元；"
            "2）非博彩項目21美獅美高梅高端康養醫療中心的相關支出3,532萬澳門元。』\n"
            "★只寫呢一兩句，唔好寫調整建議／跨司意見／結論（後面有固定句接落去）。\n\n"
            + "\n".join(lines))


# ── from build_llm_narrative ──
def _proj_prompt(adj_type, name, code, rep, adjv, find, mgmt, b2, ruling, content=""):
    ctx = [f"投資項目：{code}　{name}", f"潛在調整類型：{adj_type}",
           f"報告投資金額：{rep:,.0f}萬澳門元；本類潛在調整：{adjv:,.0f}萬澳門元"]
    if content:
        ctx.append(f"實際投資內容（項目清單）：{content[:400]}")
    if b2:      # 表2＝審查底稿，最權威，俾最多
        ctx.append(f"審查底稿表2（關注事項／調整原因／跨司意見，事實依據）：{b2[:1200]}")
    if find:
        ctx.append(f"KPMG分析發現（項目清單）：{find[:400]}")
    if mgmt:
        ctx.append(f"承批公司管理層解釋：{mgmt[:300]}")
    if ruling:
        ctx.append(f"跨司工作組／KPMG回覆（項目清單）：{ruling[:300]}")
    return ("請為報告『本年度審查工作的主要發現』其中一個投資項目，寫一段【事項描述】"
            "（約150-250字）：講清楚該項目投資咗啲乜、我們喺審查中發現咗咩、"
            "點解相關支出不應／只可部分計入報告投資金額、以及調整金額。"
            "★如有跨司工作組回覆，用『跨司工作組』集體稱呼（例：『根據我們向跨司工作組諮詢得到的回覆，"
            "跨司工作組認為…』），切勿逐個司局點名，亦切勿自創『KPMG最終立場』等標籤。"
            "只可用下面提供嘅事實同數字，唔可以虛構。\n\n" + "\n".join(ctx))


# ── from build_llm_narrative ──
def _tbl_text(df, max_rows=40):
    """DataFrame → 精簡 TSV 餵 LLM（數字原封不動）。"""
    cols = list(df.columns)
    lines = ["\t".join(str(c) for c in cols)]
    for _, r in df.head(max_rows).iterrows():
        lines.append("\t".join("" if pd.isna(r[c]) else str(r[c]) for c in cols))
    return "\n".join(lines)


# ── from build_llm_narrative ──
def _tbl_prompt(title, df, sources, unit="萬澳門元"):
    src = ("\n".join(f"- {s}" for s in sources[:6])) if sources else "（無額外資料，只根據表格數字撰寫）"
    return (f"以下係報告入面一張表，請寫佢【旁邊】嘅敘述，同埋成版最頂嗰句【導語】。\n\n"
            f"★導語要跟返呢份報告一貫句式（示範，唔好照抄字眼，要用返下面表格嘅真數）：\n"
            f"　『…在2025年度執行報告中申報的「因發生期後事項需作後續調整之2024年度博彩／非博彩項目」"
            f"投資金額為6.4億澳門元，主要包括…以及…。本次審查工作識別潛在調減金額約4.8億澳門元，"
            f"調減後金額為1.6億澳門元，主要涉及會議展覽、文化藝術、社區旅遊等非博彩投資範疇的37個項目。』\n"
            f"　金額單位跟報告習慣：≥1億寫『X.X億澳門元』（一位小數），唔夠1億寫『X,XXX萬澳門元』"
            f"（【整數、千分位、冇小數】—— 唔可以寫『5,528.9萬澳門元』，要寫『5,529萬澳門元』）。\n\n"
            f"表名：{title}\n金額單位：{unit}（括號 = 負數／調減，「-」= 零）\n\n"
            f"【表格內容】\n{_tbl_text(df)}\n\n"
            f"【其他來源（項目清單／審查底稿表2，用嚟解釋原因，唔好抄佢嘅措辭）】\n{src}")


# ── from build_llm_narrative ──
def _adj_prompt(adj_type, amt_wan, projects):
    lines = [f"潛在調整類型：{adj_type}", f"涉及潛在調減金額：約{abs(amt_wan):,.0f}萬澳門元",
             "涉及項目及審查發現（審查底稿表2 為最權威來源，優先採用其跨司裁決及具體內容）："]
    for name, find, mgmt, b2, ruling in projects[:6]:
        seg = f"- 項目「{name}」"
        if b2:      # 表2＝審查底稿，最可信，放最前、俾最多
            seg += f"；【審查底稿表2】{b2[:900]}"
        if find:
            seg += f"；KPMG分析發現：{find[:260]}"
        if mgmt:
            seg += f"；管理層解釋：{mgmt[:200]}"
        if ruling:
            seg += f"；跨司工作組／KPMG裁決（清單）：{ruling[:220]}"
        lines.append(seg)
    ctx = "\n".join(lines)
    # ★ 字數：報告 p16-17 逐類說明係【成段長文】（單一類最長 600+ 字，仲有 1)2)3)4) 分項），
    #   之前封頂 120-200 字寫得太薄（項目組 2026-08-17：「他們的文字明顯較多」）。
    #   以【反映事實 + 重點齊全】為準，寧長勿漏；分頁由版式自己處理（1/3、2/3、3/3）。
    return (f"以下係一項『潛在調整事項』嘅底層資料（審查底稿表2 內容最詳盡，可用作事實依據）。"
            f"請寫報告正文（250-550字，事實愈齊愈好，唔好為咗短而略去項目／金額／理據）："
            f"說明該調整類型、金額、【逐個】主要涉及嘅投資項目（點名 + 各自金額 + 具體投資內容）同調減原因。"
            f"★涉及多個項目／多筆支出時，用「1）…；2）…；3）…」逐項列出，跟原報告寫法。"
            f"★之後另起一段講審查過程同結論（我們就上述事項已取得…的明確回覆／結合跨司工作組意見及我們對上述項目支出的審查，我們認為…）。"
            f"★用字須跟原報告：如有向跨司工作組諮詢得到嘅回覆，用『跨司工作組』集體稱呼帶出其立場"
            f"（例如『根據我們向跨司工作組諮詢得到的回覆，跨司工作組認為／未同意…』），"
            f"【切勿】逐個司局點名（如社會文化司、旅遊局、文化局），亦【切勿】自創『KPMG最終立場』等標籤。"
            f"最後點出審查建議（通常為建議剔除／調減）。\n\n{ctx}")


# ── from build_llm_narrative ──
def _cat_prompt(sub, rate_pct, projects, reason, b2=""):
    """projects = [(序號, 名稱, 報告金額萬, 實際投資內容)]，按金額大到細。
    ⚠ 一定要逐個項目【點名】—— 之前只餵一個項目嘅內容去寫成個範疇，讀者分唔清邊句開始
      講新項目（項目組 2026-08-13 反映）。"""
    lines = []
    for code, name, amt, content in projects[:2]:      # 2 個夠：一版要放晒 11 個範疇
        lines.append(f"- {code}「{name}」（報告投資金額 {amt:,.0f} 萬澳門元）：{str(content)[:240]}")
    ctx = (f"投資範疇：{sub}\n投資計劃金額完成率：{rate_pct}\n"
           f"該範疇金額最大嘅投資項目（項目清單）：\n" + "\n".join(lines) + "\n"
           f"管理層變更原因／業務解釋：{str(reason)[:340]}\n"
           f"表2 補充（只可攞嚟豐富『投資內容』，例如子項目／活動場次／金額明細；"
           f"切勿抄佢嘅審計措辭或調整理由）：{str(b2)[:700]}")
    # 字數：報告 p19-20「按範疇的項目概況」每個範疇約 60-120 字（有項目名／內容／金額）。
    #   之前封到 40-55 字寫得太薄（項目組 2026-08-17）；上限 120 係為咗博彩／非博彩各放一版。
    return (f"請為承批公司投資執行報告寫『按範疇的項目概況』"
            f"（60-120 字，以反映事實同重點為主；唔好客套話、唔好重覆範疇名）。\n"
            f"★格式：「主要包括{{項目序號}}「{{項目名稱}}」……；{{項目序號}}「{{項目名稱}}」……。"
            f"完成率{rate_pct}，主要由於……（管理層業務原因，一句起兩句止）」\n"
            f"★【每個項目必須先寫返項目序號同項目名稱】先講佢做咗乜，項目與項目之間用「；」分開，"
            f"令讀者一眼睇到邊句係講邊個項目 —— 唔可以將幾個項目嘅內容混埋一齊寫。\n"
            f"完成率原因只用管理層業務解釋，唔好用審計／調整措辭。\n\n{ctx}")


# ── from build_llm_narrative ──
def _short_err(e):
    """gateway 錯誤成日回一版 HTML → 抽 <h2>/<title> 或者剝晒 tag，最多 120 字。"""
    t = str(e)
    m = re.search(r"<h2[^>]*>(.*?)</h2>|<title[^>]*>(.*?)</title>", t, re.S | re.I)
    if m:
        t = (m.group(1) or m.group(2) or "").strip()
    elif "<html" in t.lower():
        t = re.sub(r"<[^>]+>", " ", t)
    t = re.sub(r"\s+", " ", t).strip()
    return t[:120] + ("…" if len(t) > 120 else "")


# ── from build_llm_narrative ──
def _gen(wb, prompt, effort, sysp, want_json=False):
    if not want_json:
        return wb.chat(prompt, sysp, reasoning_effort=effort).strip()
    d = wb.chat_json(prompt, sysp, reasoning_effort=effort)
    segs = (d or {}).get("段落") or []
    out = [[str(s.get("小標題", "")).strip(), str(s.get("內容", "")).strip()]
           for s in segs if isinstance(s, dict) and s.get("內容")]
    if not out:
        raise ValueError("LLM 冇回到『段落』")
    return {"導語": str((d or {}).get("導語", "")).strip(), "段落": out}


# ── from build_llm_narrative ──
def _proj_sources(d, narr, b2, mask, kind="content", n=5):
    """由 feed 一段 slice 抽 top 項目嘅來源片段（清單 + 表2）→ [str]，餵表旁 comment。"""
    sub = d[mask]
    if sub.empty:
        return []
    key = "調整_萬" if kind == "finding" else "調整前_萬"
    top = (sub.groupby(["ng_scope", "dicj code"])
              .agg(nm=("project", "first"), v=(key, "sum")).reset_index())
    top = top.reindex(top["v"].abs().sort_values(ascending=False).index).head(n)
    out = []
    for _, p in top.iterrows():
        nr = nlook(narr, p["ng_scope"], p["dicj code"])
        b2t = b2text(b2, p["ng_scope"], p["dicj code"])
        if kind == "finding":
            txt = nr.get("KPMG分析發現", "") or b2t
            mg = nr.get("管理層解釋", "")
            seg = f"項目「{p['nm']}」（{p['v']:,.0f}萬）：{str(txt)[:300]}"
            if mg:
                seg += f"；管理層解釋：{mg[:160]}"
        else:
            txt = nr.get("實際投資內容", "") or b2t
            seg = f"項目「{p['nm']}」（{p['v']:,.0f}萬）：{str(txt)[:300]}"
        if str(txt).strip():
            out.append(seg)
    return out


# ── from build_llm_narrative ──
def generate_llm_narrative(feed_path, entity, qingdan, biao2_dir="data/表2",
                           model=None, workers=8, out_path=None, log=print):
    """由 feed + 清單 + 表2 用 Workbench 生成 {adj,cat} 敘述；寫 {entity}_llm_narrative.json，回 dict。
    可被 build_report.py --llm 直接調用（唔使另跑 command）。"""
    wb = Workbench(model=model)
    df = _load(Path(feed_path), entity)
    narr = load_narrative(Path(qingdan)) if qingdan else {}
    b2 = load_biao2_struct(biao2_dir, entity or "", log=log)
    plan = load_plan(Path(qingdan)) if qingdan else None
    cat = load_category(Path(qingdan)) if qingdan else None
    ov = overview_by_bucket(df, "2025年度投資計劃", plan, cat)
    adj = adjustment_bridge(df)
    d = df.copy()
    d["_adj"] = d["調整一級"].map(CANON).fillna(d["調整一級"])

    # 併裝 tasks（(kind, key, prompt, effort, sys)）
    pb = BUCKET_ORDER[0]      # 2025計劃 bucket：調整詳述只計 2025年度計劃（期後另計，對返報告）
    tasks = []
    for _, r in adj.iterrows():
        t = r["潛在調整事項"]
        if t in ("合計", "跨年及其他調整"):
            continue
        amt = r.get(pb, 0)
        if not isinstance(amt, (int, float)) or abs(amt) < 0.5:
            continue
        sub = d[(d["_adj"] == t) & (pd.to_numeric(d["調整_萬"], errors="coerce") != 0)]
        agg = sub.groupby(["ng_scope", "dicj code"]).agg(
            nm=("project", "first"), rep=("調整前_萬", "sum"), adjv=("調整_萬", "sum")).reset_index()
        agg = agg.reindex(agg["adjv"].abs().sort_values(ascending=False).index)
        projs = []
        for _, pp in agg.iterrows():
            nr = nlook(narr, pp["ng_scope"], pp["dicj code"])
            b2t = b2text(b2, pp["ng_scope"], pp["dicj code"])
            ruling = "；".join(x for x in (nr.get("跨司回覆", ""), nr.get("KPMG回覆", "")) if x)
            projs.append((str(pp["nm"]), nr.get("KPMG分析發現", ""),
                          nr.get("管理層解釋", ""), b2t, ruling))
            # 逐項目『事項描述』（主要發現 card）：用表2 + 清單寫返報告嗰種 narrative
            tasks.append(("proj", proj_key(t, pp["ng_scope"], pp["dicj code"]),
                          _proj_prompt(t, pp["nm"], pp["dicj code"], pp["rep"], pp["adjv"],
                                       nr.get("KPMG分析發現", ""), nr.get("管理層解釋", ""),
                                       b2t, ruling, nr.get("實際投資內容", "")),
                          "medium", SYS_ADJ))
        tasks.append(("adj", t, _adj_prompt(t, amt, projs), "medium", SYS_ADJ))

    proj = d.groupby(["_sub", "dicj code"])["調整前_萬"].sum().reset_index()
    for _, r in ov[~ov["範疇"].astype(str).str.endswith(("小計", "總計", "項目"))].iterrows():
        sub = str(r["範疇"]); rate = r.get("投資計劃完成率")
        if not isinstance(rate, (int, float)) or pd.isna(rate):
            continue
        scope = "gaming" if sub.startswith("博彩") else "non_gaming"
        projs, reason, b2t = [], "", ""
        for _, pp in proj[proj["_sub"] == sub].sort_values("調整前_萬", ascending=False).iterrows():
            nr = nlook(narr, scope, pp["dicj code"])
            if len(projs) < 3:
                projs.append((str(pp["dicj code"]), str(nr.get("項目名稱", "") or pp["dicj code"]),
                              float(pp["調整前_萬"] or 0), nr.get("實際投資內容", "")))
            reason = reason or nr.get("管理層解釋", "")   # 業務原因；唔用 KPMG分析發現（審計腔）
            b2t = b2t or b2text(b2, scope, pp["dicj code"])
            if len(projs) >= 3 and reason and b2t:
                break
        tasks.append(("cat", sub, _cat_prompt(sub, f"{rate*100:.1f}%", projs, reason, b2t),
                      "low", SYS_CAT))

    # ── 表旁 comment（scan p-10~p-13 表左＋敘述右）：由表格數字 + 清單/表2 來源寫 ──
    for bk in BUCKET_ORDER[1:]:                       # 2024 / 2023 期後概覽
        ovb = overview_by_bucket(df, bk, plan, cat)
        if ovb.empty:
            continue
        src = _proj_sources(d, narr, b2, (d["_bucket"] == bk) &
                            (pd.to_numeric(d["調整_萬"], errors="coerce") != 0), "finding")
        tasks.append(("tbl", tbl_key("期後概覽", bk),
                      _tbl_prompt(f"{(entity or '').upper()} {bk}金額概覽", ovb.fillna(""), src),
                      "medium", SYS_TBL_ADJ, True))
    for bk in BUCKET_ORDER:                           # 設施建設 vs 活動舉辦
        fa = facility_activity(df, bk)
        if fa.empty:
            continue
        src = _proj_sources(d, narr, b2, d["_bucket"] == bk, "content")
        tasks.append(("tbl", tbl_key("設施活動", bk),
                      _tbl_prompt(f"{(entity or '').upper()} {bk}區分設施建設／活動舉辦的投資金額",
                                  fa.fillna(""), src),
                      "low", SYS_TBL_DESC, True))
    amt = summary_amount(df)                          # 4.1 金額匯總
    if not amt.empty:
        tasks.append(("tbl", tbl_key("金額匯總"),
                      _tbl_prompt(f"{(entity or '').upper()} 2025年發生的投資金額匯總",
                                  amt.fillna(""), _proj_sources(d, narr, b2,
                                                                d["_bucket"] == pb, "content")),
                      "low", SYS_TBL_DESC, True))
    for bk in BUCKET_ORDER[1:]:                       # 期後調整事項匯總：逐類開場句
        dd = d[d["_bucket"] == bk]
        for t in ADJ7:
            sub = dd[dd["_adj"] == t]
            amt = pd.to_numeric(sub["調整_萬"], errors="coerce").sum()
            if abs(amt) < 0.5:
                continue
            g = (sub.groupby(["ng_scope", "dicj code"])
                    .agg(nm=("project", "first"), v=("調整_萬", "sum")).reset_index())
            g = g.reindex(g["v"].abs().sort_values(ascending=False).index)
            projs = [(str(r["nm"]), r["v"], b2text(b2, r["ng_scope"], r["dicj code"]),
                      nlook(narr, r["ng_scope"], r["dicj code"]).get("KPMG分析發現", ""))
                     for _, r in g.iterrows()]
            tasks.append(("bkt", bkt_key(bk, t), _bkt_prompt(bk[:4], t, amt, projs),
                          "low", SYS_ADJ))

    fs = finding_summary(df)                          # ③ 主要發現摘要
    if not fs.empty:
        tasks.append(("tbl", tbl_key("發現摘要"),
                      _tbl_prompt(f"{(entity or '').upper()} 本年度審查工作的主要發現摘要", fs.fillna(""),
                                  _proj_sources(d, narr, b2,
                                                pd.to_numeric(d["調整_萬"], errors="coerce") != 0,
                                                "finding", n=6)),
                      "medium", SYS_TBL_ADJ, True))

    # ★ preflight：先試一個最平嘅 call。網關擋／key 唔啱嘅話即刻知，唔使等 60 個 task
    #   逐個 fail（2026-08-15 白等 10 分鐘）。
    try:
        wb.chat("ok", "Reply with the single word: ok", reasoning_effort=None, max_tokens=5)
    except Exception as e:
        log(f"  ✗ LLM 連唔到（{type(e).__name__}: {_short_err(e)}）→ 跳過全部 {len(tasks)} 個 summary，"
            "今次用清單／表2 原文 fallback。")
        log("    403／blocked = 網關擋：check 係咪喺 KPMG 內網、key 有冇過期、charge code／region 啱唔啱。")
        out = {"adj": {}, "cat": {}, "tbl": {}, "proj": {}, "bkt": {}}
        outp = Path(out_path) if out_path else Path(f"{entity or 'all'}_llm_narrative.json")
        outp.write_text(json.dumps(out, ensure_ascii=False, indent=1), encoding="utf-8")
        return out
    log(f"（{entity}）批 {len(tasks)} 個 summary，workers={workers}…")
    out = {"adj": {}, "cat": {}, "tbl": {}, "proj": {}, "bkt": {}}
    try:                                    # tqdm 進度條（冇裝就照 log 逐個）
        from tqdm import tqdm
    except ImportError:
        tqdm = None
    tasks = [(t + (False,))[:6] for t in tasks]         # 補齊 want_json（adj/cat = 純文字）
    with ThreadPoolExecutor(max_workers=workers) as ex:
        fut = {ex.submit(_gen, wb, p, eff, sysp, js): (kind, key)
               for kind, key, p, eff, sysp, js in tasks}
        it = as_completed(fut)
        bar = tqdm(it, total=len(tasks), desc=f"LLM {entity}", unit="段", ncols=90) if tqdm else it
        nfail = 0
        for f in bar:
            kind, key = fut[f]
            try:
                out[kind][key] = f.result()
                msg = f"  ✓ {kind}｜{key[:22]}"
                nfail = 0
            except Exception as e:
                nfail += 1
                # ⚠ err 可能係成版 HTML（KPMG gateway 擋 request 會回錯誤頁）→ 一定要縮短，
                #   否則 console 會俾 60 版 HTML 洗晒版（2026-08-15 實際發生過）
                msg = f"  ⚠ {kind}｜{key[:22]}: {type(e).__name__}: {_short_err(e)}"
            tqdm.write(msg) if tqdm else log(msg)   # tqdm.write 唔會撞爛進度條
        if nfail and not any(out.values()):
            log("  ⚠ LLM 全部失敗 → 今次報告用清單／表2 原文 fallback（唔會空白，但用字唔會似報告）。"
                "\n    常見成因：公司網關擋（並發太多／唔喺 KPMG 網／key 過期）。"
                "\n    試：build_report.py mgm --workers 2；仲係唔得就唔喺內網 or 換 key。")

    outp = Path(out_path) if out_path else Path(f"{entity or 'all'}_llm_narrative.json")
    outp.write_text(json.dumps(out, ensure_ascii=False, indent=1), encoding="utf-8")
    log(f"✓ {outp.resolve()}（adj {len(out['adj'])}、cat {len(out['cat'])}、"
        f"tbl {len(out['tbl'])}、proj {len(out['proj'])}、bkt {len(out['bkt'])} 段）")
    return out


# ── from make_report ──
try:
    import pandas as pd
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("✗ pip install pandas python-pptx openpyxl"); sys.exit(1)


# ── from make_report ──
HDR = NAVY


# ── from make_report ──
SEC = SECFILL


# ── from make_report ──
SUB = SUBTOT


# ── from make_report ──
TOT = TOTAL


# ── from make_report ──
LIGHT = WHITE


# ── from make_report ──
ENTITY_FULL = {
    "mgm": "美高梅金殿超濠股份有限公司",
    "galaxy": "銀河娛樂場股份有限公司",
    "sjm": "澳門博彩股份有限公司",
    "wynn": "永利渡假村（澳門）股份有限公司",
    "vml": "威尼斯人澳門股份有限公司",
    "melco": "新濠博亞博彩（澳門）股份有限公司",
}


# ── from make_report ──
TEMPLATE_NAMES = ["template.pptx", "report_template.pptx", "kpmg_template.pptx"]


# ── from make_report ──
TEMPLATE_DIRS = [".", "data", "data/template", "conf/local"]


# ── from make_report ──
LAYOUTS = {
    "cover": ["TITLE SLIDE - Right vertical dark image", "TITLE SLIDE", "Title Slide"],
    "section": ["Section Divider"],
    "subsection": ["Subsection Divider"],
    "appendix": ["Appendix Divider"],
    "table_text": ["1_ANALYSIS narrow table", "ANALYSIS narrow table"],
    "two_col": ["Two Column Text", "ANALYSIS 2 col text"],
    "one_col": ["One Column Text", "RPOE One Column Text"],
    "title_only": ["Title Only"],
    "findings2": ["KEY FINDINGS 2_text only"],
    "toc": ["Table of content Storage Layout", "Table of content"],
}


# ── from make_report ──
def _find_template():
    for d in TEMPLATE_DIRS:
        for n in TEMPLATE_NAMES:
            p = Path(d) / n
            if p.exists():
                return p
    return None


# ── from make_report ──
USE_TEMPLATE = False


# ── from make_report ──
def _layout(prs, key):
    """跨所有 master 搵第一個名 match LAYOUTS[key] 嘅 slide layout；搵唔到回 None。
    ⚠ 冇開 template 時一定回 None —— 否則會撞正 python-pptx 預設包嗰個 'Title Slide'
    layout，封面變一版白底光板（冇深底/冇 KPMG 字標/冇日期）。"""
    if not USE_TEMPLATE:
        return None
    wanted = LAYOUTS.get(key, [key])
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if lay.name in wanted:
                return lay
    return None


# ── from make_report ──
def _strip_slides(prs):
    """正確刪走 template 原有 content slides（drop rel + sldId）→ 避免 duplicate part 名 corruption。"""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        sldIdLst.remove(sldId)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


# ── from make_report ──
def _renumber_slides(prs, base=9000):
    """save 前把生成嘅 slide part 重編做高號（9000+）→ 就算 template 有殘留 orphan part（清唔切）
    都唔會撞名 → 徹底避開 duplicate-name corruption / 要 repair。"""
    try:
        from pptx.opc.packuri import PackURI
    except Exception:
        return
    for i, slide in enumerate(prs.slides, base):
        try:
            slide.part.partname = PackURI(f"/ppt/slides/slide{i}.xml")
        except Exception:
            pass


# ── from make_report ──
def _blank_layout(prs):
    """乾淨 layout（手砌 data 版用）：template 揀 'Title Only'/'Blank'（有 master title bar/footer），否則 index。"""
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if lay.name in ("Blank", "Title Only"):
                return lay
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]


# ── from make_report ──
def _ph(slide, idx):
    """由 placeholder idx 攞 placeholder；冇就 None。"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


# ── from make_report ──
BUILD_STAMP = "base b92c52e · bundled 2026-08-17 15:47"


# ── from make_report ──
FEED = "tableau_combined_25.csv"


# ── from make_report ──
def _is_num(v):
    try:
        float(v); return True
    except (TypeError, ValueError):
        return False


# ── from make_report ──
def _load_llm(entity):
    """有 {entity}_llm_narrative.json（build_llm_narrative.py 出）就用 LLM summary。"""
    import json
    p = Path(f"{entity}_llm_narrative.json")
    if p.exists():
        try:
            return json.loads(p.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


# ── from make_report ──
def _coverage_probe(df, qingdan):
    """一次性 coverage 探測（❓頁：主體/KPI/藝術品欄喺唔喺 inputs）→ print 供定 coverage。"""
    def pr(cols, *kw):
        return [c for c in cols if any(k in str(c) for k in kw)] or "無"
    print("  [coverage] feed 主體/執行公司欄:", pr(df.columns, "主體", "執行公司", "控股", "子公司"),
          "｜KPI欄:", pr(df.columns, "KPI", "收入", "留宿", "晚數", "毛收入", "國際"))
    if qingdan:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(qingdan, data_only=True, read_only=True)
            ws = next((wb[s] for s in wb.sheetnames if s.lower().startswith("database")), wb[wb.sheetnames[0]])
            rows = [r for i, r in enumerate(ws.iter_rows(values_only=True)) if i < 6]
            hr = next((i for i, r in enumerate(rows) if r and any("項目類型" in str(c or "") for c in r)), 2)
            hdr = [str(c or "").replace("\n", "") for c in rows[hr]]
            print("  [coverage] 清單 股權/主體欄:", pr(hdr, "股權", "主體", "執行公司", "控股", "子公司"))
            print("  [coverage] 清單 KPI欄:", pr(hdr, "KPI", "留宿", "晚數", "毛收入", "增長率", "國際住客"))
            print("  [coverage] 清單 藝術品欄:", pr(hdr, "藝術品", "展出", "館藏", "拍賣"))
        except Exception as e:
            print("  [coverage] 清單 probe skip:", e)


# ── from make_report ──
def _dump_pptx_text(prs, entity, with_tables=False):
    """逐版文字 dump → user paste 返，我就可以【逐句對返原報告 scan】（長度／風格／用字）。
    預設【只 dump 敘述文字，唔 dump 表格 cell】：表格係數字、已經另外驗過，而且會令個檔大到 paste 唔到。
    要連表格：build_report.py mgm --dump --dump-tables"""
    lines = []
    for i, sl in enumerate(prs.slides, 1):
        parts = []
        for sh in sl.shapes:
            if sh.has_table:
                if not with_tables:
                    t = sh.table
                    parts.append(f"〔表 {len(t.rows)}x{len(t.columns)}〕")
                    continue
                for row in sh.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                parts.append(sh.text_frame.text.strip())
        body = "\n".join(parts)
        lines.append(f"\n===== slide {i}（{len(body)} 字）=====\n" + body)
    out = Path(f"{entity}_報告_dump.txt")
    out.write_text("\n".join(lines), encoding="utf-8")
    return out


# ── from make_report ──
ENT_UP = "MGM"


# ── from make_report ──
def _page(prs, section_idx=0, crumb=None, headline=None):
    """開一版內容頁 = breadcrumb + footer(+頁碼) + 灰標題 + navy 導語。
    回 (slide, W, H, top_y)：top_y = 內容可以由邊開始。"""
    slide = blank(prs)
    W, H = size_of(prs)
    breadcrumb(slide, W, section_idx, ENT_UP)
    footer(slide, W, H, len(prs.slides._sldIdLst))
    top = page_head(slide, W, crumb, headline) if crumb else 0.5
    return slide, W, H, top


# ── from make_report ──
def _furniture(prs, slide, section_idx=0):
    """（保留舊 API）頂 nav tabs + 底 KPMG copyright + 初稿 + 頁碼。"""
    W, H = size_of(prs)
    breadcrumb(slide, W, section_idx, ENT_UP)
    footer(slide, W, H, len(prs.slides._sldIdLst))


# ── from make_report ──
def _dark_slide(prs):
    """新增一版深底（封面/分隔共用），回 (slide, w, h)。"""
    return dark_slide(prs)


# ── from make_report ──
def render_cover(prs, entity, date="2026年6月30日"):
    """封面（報告 p1）。template：用 TITLE SLIDE layout 填 placeholder（顏色/圖嚟自 master）；否則手砌深底。"""
    full = ENTITY_FULL.get(entity, entity.upper())
    lay = _layout(prs, "cover")
    if lay is not None:
        slide = prs.slides.add_slide(lay)
        t = _ph(slide, 0)
        if t is not None:
            t.text = f"{full}\n2025年年度投資計劃執行情況審查\n專項工作報告"
        b = _ph(slide, 11)
        if b is not None:
            b.text = f"初稿\n畢馬威會計師事務所\n{date}"
        return
    slide, w, h = _dark_slide(prs)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(7.2), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate([full, "2025年年度投資計劃執行情況審查", "專項工作報告"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(25); r.font.bold = True; r.font.color.rgb = LIGHT
        r.font.name = "微软雅黑"
    db = slide.shapes.add_textbox(Inches(0.6), Inches(4.7), Inches(4), Inches(0.5))
    dr = db.text_frame.paragraphs[0].add_run(); dr.text = "初稿"
    dr.font.size = Pt(16); dr.font.bold = True; dr.font.color.rgb = LIGHT; dr.font.name = "微软雅黑"
    fb = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(5), Inches(0.8))
    ftf = fb.text_frame
    for i, line in enumerate(["畢馬威會計師事務所", date]):
        p = ftf.paragraphs[0] if i == 0 else ftf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(11); r.font.color.rgb = LIGHT; r.font.name = "微软雅黑"


# ── from make_report ──
def divider(prs, title, number="", subitems=None):
    """章節分隔。template：用 Section/Appendix Divider layout（深底嚟自 master）+ 加標題 textbox；否則手砌。"""
    lay = _layout(prs, "appendix" if number == "6" else "section")
    if lay is not None:
        slide = prs.slides.add_slide(lay)
        w = prs.slide_width / 914400.0; h = prs.slide_height / 914400.0
        ny = h * 0.34; tx = 0.9
        if number:
            nb = slide.shapes.add_textbox(Inches(0.9), Inches(ny - 0.2), Inches(1.5), Inches(1.2))
            nr = nb.text_frame.paragraphs[0].add_run(); nr.text = f"{number}."
            nr.font.size = Pt(40); nr.font.bold = True; nr.font.color.rgb = LIGHT; nr.font.name = "Arial"
            tx = 2.2
        tb = slide.shapes.add_textbox(Inches(tx), Inches(ny), Inches(w - tx - 0.9), Inches(1.2))
        tr = tb.text_frame.paragraphs[0].add_run(); tr.text = title
        tr.font.size = Pt(24); tr.font.bold = True; tr.font.color.rgb = LIGHT; tr.font.name = "微软雅黑"
        if subitems:
            y = ny + 1.5
            for it in subitems:
                label, _pg = (it if isinstance(it, (tuple, list)) else (it, ""))
                rb = slide.shapes.add_textbox(Inches(tx), Inches(y), Inches(w - tx - 1.2), Inches(0.3))
                rr = rb.text_frame.paragraphs[0].add_run(); rr.text = label
                rr.font.size = Pt(12); rr.font.color.rgb = LIGHT; rr.font.name = "微软雅黑"
                y += 0.4
        return
    slide, w, h = _dark_slide(prs)
    ny = h * 0.30
    tx = 0.7
    if number:
        nb = slide.shapes.add_textbox(Inches(0.65), Inches(ny - 0.15), Inches(1.3), Inches(1.3))
        nr = nb.text_frame.paragraphs[0].add_run(); nr.text = f"{number}."
        nr.font.size = Pt(44); nr.font.bold = True; nr.font.color.rgb = LIGHT; nr.font.name = "Arial"
        tx = 2.05
    tb = slide.shapes.add_textbox(Inches(tx), Inches(ny), Inches(w - tx - 0.7), Inches(1.5))
    ttf = tb.text_frame; ttf.word_wrap = True
    tr = ttf.paragraphs[0].add_run(); tr.text = title
    tr.font.size = Pt(26); tr.font.bold = True; tr.font.color.rgb = LIGHT; tr.font.name = "微软雅黑"
    if subitems:
        y = ny + 1.7
        for it in subitems:
            label, page = (it if isinstance(it, (tuple, list)) else (it, ""))
            rb = slide.shapes.add_textbox(Inches(tx), Inches(y), Inches(w - tx - 1.4), Inches(0.32))
            rr = rb.text_frame.paragraphs[0].add_run(); rr.text = label
            rr.font.size = Pt(12); rr.font.color.rgb = RGBColor(0xC8, 0xC8, 0xD0); rr.font.name = "微软雅黑"
            if page != "":
                pb = slide.shapes.add_textbox(Inches(w - 1.3), Inches(y), Inches(0.8), Inches(0.32))
                pr = pb.text_frame.paragraphs[0].add_run(); pr.text = str(page)
                pr.font.size = Pt(12); pr.font.color.rgb = RGBColor(0xC8, 0xC8, 0xD0); pr.font.name = "微软雅黑"
            y += 0.42


# ── from make_report ──
def _find(dirp, entity, ext, prefer=None):
    d = Path(dirp)
    if not d.exists():
        return None
    cands = [p for p in sorted(d.rglob("*"))
             if p.suffix.lower() == ext and entity.lower() in p.name.lower()
             and not p.name.startswith("~$")]
    if not cands:
        return None
    for pk in (prefer or []):                    # 例：template 優先揀 2025 果份（唔好揀到 2023 舊報告）
        for p in cands:
            if pk in p.name:
                return p
    return cands[0]


# ── from make_report ──
def _fmt_ratio(v):
    """比例欄：scan 寫『(89.4%)』—— 調減（負數）用括號。"""
    try:
        f = float(v)
    except (TypeError, ValueError):
        return "" if v is None or str(v).strip() == "" else str(v)
    return f"({abs(f)*100:.1f}%)" if f < 0 else f"{f*100:.1f}%"


# ── from make_report ──
def _cell_txt(c, v):
    """跟欄名格式化：率／比例→%，數字→千分位（負數括號），其餘原文。"""
    if "比例" in str(c):
        t = _fmt_ratio(v)
        return "-" if t in ("0.0%", "(0.0%)") else t      # 報告：比例 0 出「-」
    if "率" in str(c):
        return fmt_pct(v)
    if _is_num(v):
        return fmt_money(v)
    return "" if v is None else str(v)


# ── from make_report ──
def _df_table(df, first_label=None):
    """DataFrame → (subs, rows, widths, supers)：認 `大組·細名` 做兩層表頭，
    第一欄空 = 範疇 section 行，尾『小計/合計/總計』= shaded。
    第一欄表頭：範疇表用「（萬澳門元）」（跟 scan 角位放單位），其餘用返欄名。"""
    cols = list(df.columns)
    if first_label is None:
        first_label = ("" if cols[0] == "序號" else       # 報告個角位空白，單位喺隔離格
                       ("萬澳門元" if cols[0] == "範疇" else cols[0]))
    grouped = any("·" in c for c in cols)
    widths = [0.4 if c == "序號" else
              (2.0 if c in ("範疇", "項目名稱", "潛在調整事項") else
               (2.8 if c == "主要涉及項目" else 0.95)) for c in cols]
    subs = [first_label] + [(c.split("·")[1] if "·" in c else c) for c in cols[1:]]
    if cols[0] == "序號" and len(cols) > 1 and cols[1] == "範疇":
        subs[1] = "萬澳門元"                          # 單位放範疇欄（跟 scan 角位）
    supers = None
    if grouped:
        groups = [""] + [(c.split("·")[0] if "·" in c else "") for c in cols[1:]]
        supers, ci = [], 1
        while ci < len(cols):
            cj = ci
            while cj + 1 < len(cols) and groups[cj + 1] == groups[ci]:
                cj += 1
            supers.append((groups[ci], ci, cj + 1)); ci = cj + 1
        supers.insert(0, ("", 0, 1))
    rows = []
    # ⚠ 加咗「序號」欄之後，範疇名喺第 2 欄 → 分類（sec/小計/總計）要睇標籤欄，唔可以睇 cols[0]，
    #   否則 小計/總計 全部當 data（冇粗體、冇橫線）。
    li = 1 if (cols[0] == "序號" and len(cols) > 1) else 0
    for _, row in df.iterrows():
        cells = [str(row[cols[0]]).strip()] + [_cell_txt(c, row[c]) for c in cols[1:]]
        if cols[0] in ("範疇", "序號"):        # 範疇名用報告寫法（render 層）
            cells[li] = sub_display(cells[li])
        lab = (cells[li] or cells[0]).strip()
        if all(str(row[c]).strip() == "" for c in cols[li + 1:]):
            rows.append(("sec", cells)); continue
        kind = ("data" if lab == "涉及項目數量" else
                "formula" if lab.endswith("完成率") else      # 1.2 表尾兩條完成率行＝斜體（scan p10）
                "tot" if lab in ("總計", "合計") else
                "subtot" if lab.endswith(("小計", "合計")) or lab.startswith("承諾的") else "data")
        rows.append((kind, cells))
    return subs, rows, widths, supers


# ── from make_report ──
def _draw_table(slide, df, x, y, max_w, font=6.5):
    """喺 slide (x,y) 畫 navy 表（單 chunk，caller 自行分頁）。max_w=可用闊(吋)。"""
    subs, rows, widths, supers = _df_table(df)
    return draw_table(slide, x, y, max_w, subs, rows, widths, supers=supers,
                        font=font, hfont=font - 0.5)


# ── from make_report ──
def _bullets_into(box, bullets, size=8):
    """（保留舊 API）scan 敘述格式：navy 粗體小標題 + body 段落。"""
    prose(box, bullets, head_size=size - 1, body_size=size - 1.5)


# ── from make_report ──
_OV_GROUP = {
    "報告投資金額": "報告投資金額·金額",
    "投資計劃完成率": "報告投資金額·完成率",
    "潛在調整金額": "潛在調整後投資金額·潛在調整金額",
    "潛在調整後投資金額": "潛在調整後投資金額·金額",
    "潛在調整後投資計劃完成率": "潛在調整後投資金額·完成率",
    "設施建設/資本性支出": "潛在調整後投資金額·設施建設/資本性支出",
    "活動舉辦/營運性支出": "潛在調整後投資金額·活動舉辦/營運性支出",
}


# ── from make_report ──
def _hdr_cols(subs, supers):
    """表頭欄組色（項目組 2026-08-17 指定）：預設全部 HDR_FILL #1E49E2，重點欄綠 #098E7E。
      · 1.2／1.3 概覽表 → 「獲批的計劃投資金額」綠
      · 調整表 1.4／2.2／2.4 → 七大類欄組＋潛在調整合計 天藍、比例欄 綠
      · 4.1 金額匯總 → 三個年度欄組全藍，最右「潛在調整後投資金額」欄組 綠
      · 2.5 三年表 → 2023 藍／2024 紫／2025 天藍／三年累計 綠
    ⚠ 4.1 同 2.5 都有 super「2025年度投資計劃」，唔可以齋睇 label 撞色 →
      有「三年累計」個 super 先當係 2.5 嗰套年度配色。"""
    out = {c: HDR_KEY for c, v in enumerate(subs) if str(v).strip() == "獲批的計劃投資金額"}
    sup = list(supers or [])
    tri = any("三年累計" in str(lab) for lab, _, _ in sup)
    # 調整表（1.4／2.2／2.4）：七大類欄組 + 潛在調整合計 = 天藍；比例欄 = 綠
    for lab, c0, c1 in sup:
        if "潛在調整事項" in str(lab):
            for c in range(c0, c1):
                out[c] = HDR_SKY
    for c, v in enumerate(subs):
        t = str(v).strip()
        if t in ("潛在調整合計", "潛在調整金額"):
            out[c] = HDR_SKY
        elif t.endswith("比例"):
            out[c] = HDR_KEY
    if tri:                     # 2.5 三年表
        for lab, c0, c1 in sup:
            t = str(lab)
            col = (HDR_PUR if "2024年度投資計劃" in t else
                   HDR_SKY if "2025年度投資計劃" in t else
                   HDR_KEY if "三年累計" in t else None)
            if col is not None:
                for c in range(c0, c1):
                    out[c] = col
    else:                       # 4.1／期後概覽：最右「潛在調整後投資金額」欄組 = 綠
        for lab, c0, c1 in sup:
            if str(lab).strip() == "潛在調整後投資金額" and c1 == len(subs):
                for c in range(c0, c1):
                    out[c] = HDR_KEY
    return out


# ── from make_report ──
_OV_GROUP_POST = {
    "設施建設/資本性支出": "潛在調整後投資金額·設施建設/資本性支出",
    "活動舉辦/營運性支出": "潛在調整後投資金額·活動舉辦/營運性支出",
}


# ── from make_report ──
def _overview_display(ov):
    """概覽表 → 顯示用：範疇前加【序號】（逐 scope 由 1 數起，section/小計/總計留空）+ 兩層表頭。"""
    if ov is None or ov.empty or "範疇" not in ov.columns:
        return ov
    d = ov.copy()
    seq, n = [], 0
    for _, r in d.iterrows():
        lab = str(r["範疇"]).strip()
        others = [c for c in d.columns if c != "範疇"]
        if all(str(r[c]).strip() == "" for c in others):      # section 行（博彩項目/非博彩項目）
            n = 0; seq.append("")
        elif lab.endswith(("小計", "總計", "合計")) or lab.startswith(("原計劃", "投資執行報告",
                                                                   "承諾的", "2025年")):
            seq.append("")
        else:
            n += 1; seq.append(str(n))
    d.insert(0, "序號", seq)
    d["範疇"] = d["範疇"].map(sub_display)          # 報告寫法（render 層，唔影響算數）
    # 1.2（scan p10）已經冇完成率欄、亦冇設施/活動欄 → 單層表頭，唔使 group map
    gmap = ({} if "獲批的計劃投資金額" in d.columns else _OV_GROUP_POST)
    return d.rename(columns={k: v for k, v in gmap.items() if k in d.columns})


# ── from make_report ──
def render_overview_page(prs, crumb, headline, table_df, bullets, *, sec=0, table_name=None,
                         note=None):
    """報告概述式 2 欄版（對 scan slide 10/15）：crumb + navy 導語，左 表，右 敘述。"""
    slide, W, H, top = _page(prs, sec, crumb, headline)
    left_w = W * 0.60
    tbl_bot = top
    if table_df is not None and not table_df.empty:
        if table_name:
            top = caption_bar(slide, MARGIN, top, left_w, table_name)
        disp = _overview_display(table_df)
        subs, rows, widths, supers = _df_table(disp)
        if "潛在調整金額" in list(table_df.columns):      # 期後表：表頭下面加斜體公式行
            rows = [("formula", overview_formula_row(list(disp.columns)))] + rows
        wid = [w * left_w / sum(widths) for w in widths]
        avail = CONTENT_BOTTOM - top - 0.30          # 留位俾表下面個「註」
        hh = header_h(supers, subs, wid, SZ_TBL_HDR)
        font = SZ_TBL
        # −0.18 安全位：估算同 PowerPoint 實際長高會差少少，唔留位就會 TABLE-GROW
        while font > 4.5 and sum(row_h(c, wid, font) for _, c in rows) > avail - hh - 0.18:
            font -= 0.25
        tbl_bot, _ = draw_table(slide, MARGIN, top, left_w, subs, rows, widths,
                                  supers=supers, font=font, hfont=max(4.5, font - 0.5),
                                  fill_h=avail - 0.18, left_cols=2,   # 序號+範疇 左對齊；−0.18 安全位
                                  hdr_cols=_hdr_cols(subs, supers))
    if note:      # 「註」貼喺表底下，唔可以同底部嘅資料來源疊字（多行註要留夠位）
        nh = 0.16 * (1 + note.count("\n")) + 0.16
        put(slide, MARGIN, min(tbl_bot + 0.06, CONTENT_BOTTOM - nh), left_w, nh,
              note, size=SZ_NOTE - 1, italic=True, color=GREY)
    rx = MARGIN + left_w + 0.22
    prose_box(slide, rx, top - 0.02, W - rx - MARGIN, CONTENT_BOTTOM - top, bullets)
    source_note(slide, W)


# ── from make_report ──
def _draw_adj_table(slide, x, y, w, adjdf, *, font=None):
    """報告 1.4／2.2／2.4 個表：範疇 × 七大類 + 表頭下面嗰行斜體公式（對 scan slide 15）。"""
    subs, rows, widths, supers = _df_table(adjdf, first_label="萬澳門元")
    rows = [("formula", adj_formula_row(list(adjdf.columns)))] + rows
    f = font or (SZ_TBL_WIDE if len(subs) > 11 else SZ_TBL)
    avail = CONTENT_BOTTOM - y - 0.28
    wid = [v * w / sum(widths) for v in widths]
    while f > 4.0:      # 加咗公式行同「涉及項目數量」行之後會高咗 → 自動縮到放得落
        need = (header_h(supers, subs, wid, max(4.5, f - 0.5))
                + sum(row_h(c, wid, f) for _k, c in rows))
        if need <= avail:
            break
        f -= 0.25
    return draw_table(slide, x, y, w, subs, rows, widths, supers=supers,
                        font=f, hfont=max(4.5, f - 0.5), left_cols=1,
                        hdr_cols=_hdr_cols(subs, supers))


# ── from make_report ──
def render_overview_pages(prs, crumb, headline, table_df, bullets, *, sec=0, table_name=None,
                          note=None, grouped=False):
    """同 render_overview_page，但右邊敘述長就自動分版，【左邊同一個表逐版重複】。
    對 scan slide 11-14：1.3 四版全部都係左邊 1.2 嗰個整體概況表 + 右邊唔同段落。
    grouped=True 時 bullets = [(右欄小標題, [(head, body)…])]，每組至少一版（報告 1/4…4/4）。"""
    if not bullets:
        return
    W, _H = size_of(prs)
    left_w = W * 0.60
    rx = MARGIN + left_w + 0.22
    colw = W - rx - MARGIN
    top = HEAD_Y + head_h(f"{headline}（1/9）", W)[0] + 0.10 + (0.20 if table_name else 0)
    avail = CONTENT_BOTTOM - top
    pages = []
    for grp in (bullets if grouped else [(None, bullets)]):
        head, items = grp if grouped else grp
        if not items:
            continue
        chunks = fit_prose(items, colw, avail - (0.24 if head else 0),
                             head_size=SZ_BODY_HEAD, body_size=SZ_BODY)
        for ci, ch in enumerate(chunks):
            pages.append(([(head + ("（續）" if ci else ""), "")] if head else []) + ch)
    for pi, page in enumerate(pages):
        render_overview_page(prs, crumb, headline + _pg(pi + 1, len(pages)), table_df, page,
                             sec=sec, table_name=table_name, note=note)


# ── from make_report ──
def _total_line(df):
    """由表自己嘅『總計』行砌一句機械導語（避免「淨得個表冇文字」）。"""
    cols = list(df.columns)
    tot = df[df[cols[0]].astype(str).str.strip().str.endswith(("總計", "合計"))]
    if tot.empty:
        return ""
    r = tot.iloc[0]
    parts = []
    for c in cols[1:]:
        v = r[c]
        if not _is_num(v) or float(v) == 0:
            continue
        lab = c.replace("·", "－")
        parts.append(f"{lab} {_cell_txt(c, v)}")
    return "；".join(parts[:6]) + "（單位：萬澳門元，除完成率外）。" if parts else ""


# ── from make_report ──
def _table_bullets(df):
    """冇 LLM 時嘅機械表旁敘述（由表自己嘅小計/總計行計，自洽）→ [(小標題, 內容)]。"""
    cols = list(df.columns)
    first = df[cols[0]].astype(str).str.strip()
    num = [c for c in cols[1:] if pd.to_numeric(df[c], errors="coerce").notna().any()]
    if not num:
        return []
    out = []
    tot = df[first.str.endswith(("總計", "合計"))]
    if not tot.empty:
        out.append(("整體情況", "全部範疇合計 " +
                    "；".join(f"{c.replace('·', '－')} {_cell_txt(c, tot.iloc[0][c])}"
                              for c in num[:5]) + "（單位：萬澳門元，除完成率外）。"))
    subs = df[first.str.endswith("小計")]
    if not subs.empty:
        out.append(("博彩／非博彩分佈", "；".join(
            f"{r[cols[0]]} " + "、".join(f"{c.replace('·', '－')} {_cell_txt(c, r[c])}"
                                         for c in num[:3]) for _, r in subs.iterrows()) + "。"))
    val = next((c for c in num if "報告" in c or "合計" in c), num[0])
    data = df[~first.str.endswith(("小計", "總計", "合計")) &
              df[cols[1:]].astype(str).apply(lambda r: any(x.strip() for x in r), axis=1)].copy()
    data["_v"] = pd.to_numeric(data[val], errors="coerce")
    top = data.sort_values("_v", ascending=False).head(4)
    if not top.empty:
        out.append((f"金額最大的範疇（按{val.replace('·', '－')}）", "、".join(
            f"{r[cols[0]]}（{_cell_txt(val, r[val])}）" for _, r in top.iterrows()) + "。"))
    return out


# ── from make_report ──
BUDGET_FILES = ["data/10year_budget.yml", "data/10year_budget.json", "conf/local/10year_budget.yml"]


# ── from make_report ──
def _load_budget(entity):
    for fn in BUDGET_FILES:
        p = Path(fn)
        if not p.exists():
            continue
        try:
            if p.suffix == ".json":
                import json
                d = json.loads(p.read_text(encoding="utf-8"))
            else:
                import yaml
                d = yaml.safe_load(p.read_text(encoding="utf-8"))
            b = (d or {}).get(entity) or (d or {}).get(entity.upper()) or {}
            if b:  # noqa: SIM102
                print(f"    10年投資預算：{fn} → 總計 {b.get('總計', 0):,.0f} 萬澳門元")
                return {k: float(v) for k, v in b.items()}
        except Exception as e:
            print(f"    ⚠ 讀唔到 {fn}: {e}")
    return {}


# ── from make_report ──
def _proj_counts(sdf, plan, ov=None):
    """(n_plan, n_impl, n_zero) —— 表尾數量行同 1.2 敘述共用，保證三個數同一母體、自洽。
    n_plan 優先攞【概況表自己個「總計」項目數量】，咁「總計 = 已實施 + 未實施」一定成立
    （項目組 2026-08-15：總計 84 但 79+10 唔等於 84）。冇表就用清單 2025 計劃金額 > 0 嘅碼
    （0 行唔算「獲批開展」，否則出 256 個）；再冇就退化成 feed 有支出嘅碼（n_zero=0，句子會略去）。"""
    d = sdf[sdf["_bucket"] == BUCKET_ORDER[0]]
    plan25 = {k for k, v in ((plan or {}).get(25, {}) or {}).items()
              if isinstance(v, (int, float)) and v > 0}
    spent = {(str(r["ng_scope"]) == "gaming", _norm(r["dicj code"]))
             for _, r in d[pd.to_numeric(d["調整前_萬"], errors="coerce").fillna(0) != 0]
             .drop_duplicates(["ng_scope", "dicj code"]).iterrows()}
    n_impl = len(plan25 & spent) if plan25 else len(spent)
    n_plan = _tot_projects(ov) or (len(plan25) if plan25 else len(spent))
    n_plan = max(n_plan, n_impl)
    return n_plan, n_impl, n_plan - n_impl


# ── from make_report ──
def _tot_projects(ov):
    """概況表「總計」行嘅項目數量（0 = 攞唔到）。"""
    if ov is None or getattr(ov, "empty", True) or "項目數量" not in getattr(ov, "columns", []):
        return 0
    t = ov[ov["範疇"].astype(str).str.strip() == "總計"]
    if t.empty:
        return 0
    v = pd.to_numeric(pd.Series([t.iloc[0]["項目數量"]]), errors="coerce").iloc[0]
    return int(v) if pd.notna(v) and v > 0 else 0


# ── from make_report ──
_RATE_COLS = ("投資計劃完成率", "潛在調整後投資計劃完成率")


# ── from make_report ──
def _overview_extra(ov, plan, sdf, budget, ent_up):
    """1.2 概況表出表前處理（對 scan p10）：
      ① drop 兩條完成率【欄】—— 報告冇，完成率係表尾三行入面嘅一【行】；
      ② 加尾三行：2025年度投資計劃完成率 ／ 承諾的10年投資預算 ／ 2025年度投資支出佔10年投資預算的完成率。
    10年預算要 config，冇就照出行填「-」（保持報告結構，一眼睇到係待填而唔係漏咗）。"""
    tot = ov[ov["範疇"].astype(str).str.strip() == "總計"]
    ov = ov.drop(columns=[c for c in _RATE_COLS if c in ov.columns])
    cols = list(ov.columns)
    b_all = budget.get("總計") if budget else None

    def _v(c):
        if not len(tot) or c not in tot.columns:
            return None
        x = pd.to_numeric(pd.Series([tot.iloc[0][c]]), errors="coerce").iloc[0]
        return None if pd.isna(x) else float(x)
    plan_tot = _v("獲批的計劃投資金額")
    # ① 完成率行：報告／潛在調整事項／潛在調整後 各自除以「獲批的計劃投資金額」（scan p10：103.9%｜(38.2%)｜65.6%）
    r1 = {cols[0]: "2025年度投資計劃完成率"}
    for c in ("報告投資金額", "投資金額的潛在調整事項", "潛在調整後投資金額"):
        v = _v(c)
        r1[c] = _rate(v, plan_tot) if (v is not None and plan_tot) else "-"
    # ⚠ 冇 budget 都要填「-」，唔可以留空 —— 成行全空會俾 _df_table 當做【範疇 section 行】（→ 變粗體）
    _bcol = "獲批的計劃投資金額" if "獲批的計劃投資金額" in cols else "報告投資金額"
    r2 = {cols[0]: "承諾的10年投資預算", _bcol: b_all if b_all else "-"}
    r3 = {cols[0]: "2025年度投資支出佔10年投資預算的完成率"}
    for c in ("獲批的計劃投資金額", "報告投資金額", "潛在調整後投資金額"):
        v = _v(c)
        r3[c] = _rate(v, b_all) if (v is not None and b_all) else "-"
    return pd.concat([ov, pd.DataFrame([r1, r2, r3])], ignore_index=True)


# ── from make_report ──
def _bucket_adj_table(ov):
    """期後調整事項匯總嘅左表（對 scan p-11）：範疇 × 報告(a)｜潛在調整(b)｜調整後(c=a+b)｜b/a。"""
    keep = ["範疇", "項目數量", "報告投資金額", "潛在調整金額", "潛在調整後投資金額"]
    d = ov[[c for c in keep if c in ov.columns]].copy()
    rep = pd.to_numeric(d.get("報告投資金額"), errors="coerce")
    adj = pd.to_numeric(d.get("潛在調整金額"), errors="coerce")
    ratio = (adj / rep).where(rep.abs() > 0).astype(object)   # object：section 行要填 ""（避 pandas FutureWarning）
    ratio[d["範疇"].astype(str).str.strip().eq("") | rep.isna()] = ""
    # ⚠ 分母 0（例如博彩全 0 嘅小計行）→ 出「-」，唔可以出 nan%
    ratio[ratio.isna()] = "-"
    d["潛在調整金額佔報告投資金額比例"] = ratio
    return d


# ── from make_report ──
def render_bucket_adjustment(prs, ent_up, bk, sdf, ov, narr, llm=None):
    """② 期後【調整事項匯總】（scan p-11 / p-13）：左 = 範疇 × 調整表，
    右 = navy 小標題 + 編號清單『 類型（金額）：說明』，編號跟七大類 canonical 序（會跳號）。"""
    yr = bk[:4]
    d = sdf[sdf["_bucket"] == bk].copy()
    if d.empty:
        return
    d["_adj"] = d["調整一級"].map(CANON).fillna(d["調整一級"])
    llm_bkt = (llm or {}).get("bkt", {})
    items = []
    d.loc[~d["_adj"].isin(ADJ7), "_adj"] = ADJ_POST      # 殘差＝第 8 類（同表一致）
    for t in ADJ_ALL:                          # 編號＝1-8 canonical 序，冇金額嗰類 skip（斷號）
        i = adj_no(t)
        sub = d[d["_adj"] == t]
        amt = pd.to_numeric(sub["調整_萬"], errors="coerce").sum()
        if abs(amt) < 0.5:
            continue
        body = llm_bkt.get(bkt_key(bk, t)) or (
            f"在{yr}年度投資計劃期後投資金額中，{ent_up}同樣申報了{t}。")
        body += (f"與前述的「2025年度投資計劃報告投資金額的潛在調整事項匯總」一致，"
                 f"我們就{yr}年度投資計劃期後投資金額同樣建議了該項調整。詳見後續主要發現「{t}」。")
        items.append((i, f"{t}（{_amt(amt)}）：", body))
    if not items:
        return
    tot = pd.to_numeric(d["調整_萬"], errors="coerce").sum()
    npj = int(d[pd.to_numeric(d["調整_萬"], errors="coerce") != 0]["dicj code"].nunique())
    aft = pd.to_numeric(d["調整後_萬"], errors="coerce").sum()
    head = (f"基於我們的各項審查程序，我們認為{ent_up}報告的{yr}年度投資計劃在2025年繼續執行的支出中，"
            f"存在{_cn(len(items))}大類的調整事項，潛在調減投資金額約{_amt(tot)}"
            f"（涉及{npj}個投資項目），經潛在調減後的{yr}年度計劃投資項目在2025年的投資金額約{_amt(aft)}。")
    S2 = "過往年度投資計劃在2025年繼續執行的審查跟進"
    tname = f"{ent_up} {yr}年度投資計劃於2025年申報的期後投資金額的潛在調整"
    W, H = size_of(prs)
    crumb = f"{S2}  |  {yr}年度投資計劃報告投資金額的潛在調整事項匯總"
    tbl = adjustment_by_sub(sdf, bk)
    if tbl.empty:
        tbl = _bucket_adj_table(ov)
    # ★ 版式跟 scan p21-22（同 1.4 個 p15 唔一樣！）：
    #     第 1 版 = 表【左】+ 逐類說明【右】（右欄裝得落幾多就幾多）
    #     之後   = 全闊兩欄續版（p22），左右欄頂各有 navy 小標題（右邊加「（續）」）
    left_w = W * 0.60
    rx = MARGIN + left_w + 0.22
    rw = W - MARGIN - rx
    # 先【唔起版】計好第 1 版右欄裝得落邊幾項 → 先知總頁數，導語尾寫得出「（1/2）」
    top0 = HEAD_Y + head_h(head, W)[0] + 0.10
    rlim = CONTENT_BOTTOM - top0 - 0.22          # 減右欄頂嗰行小標題
    first, rest, used = [], [], 0.0
    for it in items:
        hh = est_numbered_h([it], rw, size=SZ_BODY)
        if rest or (first and used + hh > rlim):
            rest.append(it); continue
        first.append(it); used += hh
    n_all = 1 + len(_prose_pages(prs, rest, head, tname))
    slide, W, H, top = _page(prs, 1, crumb, head + (f"（1/{n_all}）" if n_all > 1 else ""))
    t2 = caption_bar(slide, MARGIN, top, left_w, tname)
    tbot = (_draw_adj_table(slide, MARGIN, t2, left_w, tbl.fillna("")) or (t2, 0))[0]
    put(slide, MARGIN, min(tbot + 0.06, CONTENT_BOTTOM - 0.26), left_w, 0.3,
          "註：金額單位為萬澳門元；括號表示調減。", size=SZ_NOTE - 1, italic=True, color=GREY)
    put(slide, rx, top, rw, 0.18, tname, size=7, bold=True, color=NAVY)
    prose_numbered(_tb(slide, rx, top + 0.22, rw, CONTENT_BOTTOM - top - 0.22),
                     first, size=SZ_BODY)
    source_note(slide, W, more=(n_all > 1))
    _prose_2col(prs, crumb, rest, sec=1, headline=head, subtitle=tname, pg0=1, pgn=n_all)


# ── from make_report ──
def _cum_table(df, plan, cat=None):
    """三年累計表（scan slide 26）→ DataFrame（`·` = 兩層表頭）。
    每個計劃年 Y：獲批(a)=清單計劃｜2025年前已獲認可(b)=報告年<25 調整後｜2025年期後(c)=報告年25 調整後
    ｜合計(d=b+c)｜完成率(d/a)。2025計劃冇 b。尾段＝三年累計 Σa｜Σd｜Σd/Σa。"""
    d = df[df["dicj code"].astype(str).str.match(r"^項目\s*\d")].copy()
    if d.empty or not plan:
        return pd.DataFrame()
    d["_sub"] = sub_of(d)
    d["_g"] = (d["ng_scope"] == "gaming")
    d["_ry"] = pd.to_numeric(d["報告年"], errors="coerce")
    d["_af"] = pd.to_numeric(d["調整後_萬"], errors="coerce").fillna(0)
    d["_ngn"] = d["ng_code"].map(_ngn)
    d["_go"] = d["_sub"].map(lambda x: GORDER.get(x, 5))
    order = (d.drop_duplicates(["_g", "_sub"]).sort_values(["_g", "_go", "_ngn", "_sub"],
             ascending=[False, True, True, True])[["_g", "_sub"]].values.tolist())
    code_sub = {(bool(r["_g"]), _norm(r["dicj code"])): str(r["_sub"])
                for _, r in d.drop_duplicates(["_g", "dicj code"]).iterrows()}
    # feed 有嘅項目 → 直接知範疇；feed 冇（零投資／往年項目）→ 用清單『項目性質』學返
    d2sub = {}
    for (gm, code), sub in code_sub.items():
        dv = (cat or {}).get((gm, code))
        if dv:
            d2sub.setdefault(str(dv), set()).add(sub)
    d2sub1 = {k: next(iter(v)) for k, v in d2sub.items() if len(v) == 1}
    A, miss = {}, 0.0
    for yr in (23, 24, 25):
        for (gm, code), v in (plan.get(yr, {}) or {}).items():
            sub = code_sub.get((bool(gm), code)) or d2sub1.get(str((cat or {}).get((bool(gm), code), "")))
            if sub:
                A[(yr, bool(gm), sub)] = A.get((yr, bool(gm), sub), 0.0) + float(v or 0)
            else:
                miss += float(v or 0)
    if miss > 0.5:
        print(f"    ⚠ 2.5 三年累計：{miss:,.0f} 萬計劃金額派唔到範疇（清單項目性質對唔到）")

    def _af(yr, gm, sub, pre):
        m = ((d["_plan_year"] == yr) & (d["_g"] == gm) & (d["_sub"] == sub) &
             ((d["_ry"] < 25) if pre else (d["_ry"] >= 25)))
        return round(float(d.loc[m, "_af"].sum()), 1)

    Y3, Y5 = "2023年度投資計劃", "2025年度投資計劃"
    Y4, CUM = "2024年度投資計劃", "截至2025年末三年累計"
    cols = ["範疇"]
    for y in (Y3, Y4):
        cols += [f"{y}·獲批的計劃投資金額", f"{y}·2025年前已獲認可", f"{y}·2025年期後",
                 f"{y}·合計", f"{y}·完成率"]
    cols += [f"{Y5}·獲批的計劃投資金額", f"{Y5}·潛在調整後投資金額", f"{Y5}·完成率"]
    cols += [f"{CUM}·獲批的計劃投資金額", f"{CUM}·獲認可／潛在調整後投資金額", f"{CUM}·完成率"]

    def line(name, cells):
        r = {"範疇": name}
        r.update(dict(zip(cols[1:], cells)))
        return r

    def calc(items):
        """items = [(gm, sub)] → 12 個數（同一套公式，逐個範疇同小計/總計共用）。"""
        out, ta, td = [], 0.0, 0.0
        for yr, y in ((23, Y3), (24, Y4), (25, Y5)):
            a = round(sum(A.get((yr, g, sb), 0.0) for g, sb in items), 1)
            b = round(sum(_af(yr, g, sb, True) for g, sb in items), 1)
            c = round(sum(_af(yr, g, sb, False) for g, sb in items), 1)
            out += ([a, b, c, round(b + c, 1), _rate(b + c, a)] if yr != 25
                    else [a, c, _rate(c, a)])
            ta += a; td += b + c
        return out + [round(ta, 1), round(td, 1), _rate(td, ta)]

    rows, all_items = [], []
    for gm, label in ((True, "博彩項目"), (False, "非博彩項目")):
        items = [(g, sb) for g, sb in order if g == gm]
        if not items:
            continue
        rows.append({"範疇": label})
        for g, sb in items:
            rows.append(line(sb, calc([(g, sb)])))
        rows.append(line(f"{label}小計", calc(items)))
        all_items += items
    rows.append(line("合計", calc(all_items)))
    return pd.DataFrame(rows, columns=cols)


# ── from make_report ──
def _collect_toc(prs, W, H):
    """由【已起好嘅版】反推目錄：深色分隔頁 = 章節；內容頁 crumb「章節 | 子題」= 子項。
    頁碼 = 插入目錄之後嘅位置（+1）。"""
    ent, last = [], None
    for i, sl in enumerate(prs.slides, 1):
        full = any(sh.width / 914400.0 > W - 0.1 and sh.height / 914400.0 > H - 0.1
                   for sh in sl.shapes)
        pg = i + 1                                  # 目錄會插喺第 2 版
        if full:
            texts = [sh.text_frame.text.strip() for sh in sl.shapes
                     if sh.has_text_frame and sh.text_frame.text.strip()]
            no = next((t for t in texts if re.fullmatch(r"\d+\.", t)), "")
            ttl = next((t for t in texts if len(t) >= 2 and t != "KPMG"
                        and not re.fullmatch(r"\d+\.", t)), "")
            if no and ttl:
                ent.append((no, ttl.split("\n")[0], False, pg)); last = None
            continue
        for sh in sl.shapes:
            if sh.has_text_frame and 0.28 < sh.top / 914400.0 < 0.46 and "  |  " in sh.text_frame.text:
                sub = sh.text_frame.text.split("  |  ")[-1].strip()
                sub = re.sub(r"（\d+/\d+）$", "", sub).strip()
                if sub and sub != last:
                    ent.append(("", sub, True, pg)); last = sub
                break
    return ent


# ── from make_report ──
def _renumber_footers(prs, W, H):
    """插咗目錄之後，全份『初稿 N』重編（頁碼喺建版時已寫死）。"""
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if not sh.has_text_frame or sh.top / 914400.0 < H - 0.45:
                continue
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.text.startswith("初稿"):
                        r.text = f"初稿　{i}"


# ── from make_report ──
def _sec_slides(prs, W, H):
    """{章 index(0-based): slide index} —— 深色分隔頁上嘅「」認章號，供 breadcrumb 頁籤跳頁。"""
    out = {}
    for i, sl in enumerate(prs.slides):
        if not any(sh.width / 914400.0 > W - 0.1 and sh.height / 914400.0 > H - 0.1
                   for sh in sl.shapes):
            continue
        for sh in sl.shapes:
            m = sh.has_text_frame and re.fullmatch(r"(\d+)\.", sh.text_frame.text.strip())
            if m:
                out.setdefault(int(m.group(1)) - 1, i)
                break
    return out


# ── from make_report ──
def _move_slide(prs, frm, to):
    """把第 frm 版（0-based）搬去 to。"""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[frm]); lst.insert(to, ids[frm])


# ── from make_report ──
def render_toc(prs, ent_up, entries):
    """報告 slide 7 目錄：六大章節 + 子項 + 頁碼（頁碼喺 build 完先知 → 由 caller 傳）。
    子項多過一版就自動分版。"""
    avail = CONTENT_BOTTOM - (HEAD_Y + 0.06) - 0.55
    pages, cur, used = [], [], 0.0
    for e in entries:
        h = 0.30 if e[2] else 0.34
        if cur and used + h > avail:
            pages.append(cur); cur, used = [], 0.0
        cur.append(e); used += h
    if cur:
        pages.append(cur)
    for pi, page in enumerate(pages):
        _render_toc_page(prs, ent_up, page, _pg(pi + 1, len(pages)))


# ── from make_report ──
def _render_toc_page(prs, ent_up, entries, suffix=""):
    slide, W, H, top = _page(prs, 0, f"{ent_up} 2025年年度投資計劃執行情況審查專項工作報告  |  目錄",
                             None)
    put(slide, MARGIN, top, W - 2 * MARGIN, 0.35, "目錄" + suffix, size=18, bold=True, color=HDR)
    y = top + 0.55
    for no, title, sub, pg in entries:
        put(slide, MARGIN, y, 0.6, 0.26, no, size=SZ_BODY_HEAD, bold=True, color=HDR)
        put(slide, MARGIN + 0.62, y, W - 2 * MARGIN - 1.5, 0.26, title,
              size=SZ_BODY_HEAD, bold=not sub, color=HDR if not sub else INK)
        if pg:
            put(slide, W - MARGIN - 0.7, y, 0.7, 0.26, str(pg), size=SZ_BODY,
                  color=GREY, align=PP_ALIGN.RIGHT)
        y += 0.30 if sub else 0.34
    source_note(slide, W, note="")


# ── from make_report ──
def render_visit_summary(prs, ent_up, df, threshold=2000):
    """報告 slide 71『設施建設項目現場走訪情況匯總』：樣本選取標準 + 樣本量表。
    全部由 feed 計（capex 項目母體 vs 走訪樣本），冇外部資料都做到。"""
    cap = df[df["final_capex_opex"] == "Capex"].copy()
    cap = cap[cap["dicj code"].astype(str).str.match(r"^項目\s*\d")]
    if cap.empty:
        return
    g = cap.groupby(["ng_scope", "dicj code"]).agg(報告=("調整前_萬", "sum")).reset_index()
    n_all, amt_all = len(g), float(g["報告"].sum())
    sel = g[g["報告"] >= threshold]
    n_sel, amt_sel = len(sel), float(sel["報告"].sum())
    rows = pd.DataFrame([
        {"項目": "設施建設（資本性支出）項目母體", "項目數量": n_all, "涉及金額": round(amt_all, 1),
         "佔母體金額比例": 1.0},
        {"項目": f"現場走訪樣本（單一項目資本性支出 ≥ {threshold:,.0f} 萬澳門元）",
         "項目數量": n_sel, "涉及金額": round(amt_sel, 1),
         "佔母體金額比例": (amt_sel / amt_all if amt_all else None)},
        {"項目": "未列入走訪樣本", "項目數量": n_all - n_sel,
         "涉及金額": round(amt_all - amt_sel, 1),
         "佔母體金額比例": ((amt_all - amt_sel) / amt_all if amt_all else None)},
    ])
    head = (f"我們在制定本次審查工作範圍時，計劃就{ent_up}報告投資金額中重大的設施建設項目開展現場走訪。"
            f"我們根據管理層提供的投資項目底層財務明細，就設施建設（資本性支出）項目共{n_all}個"
            f"（涉及{_amt(amt_all)}）作為母體，選取單一項目資本性支出達{threshold:,.0f}萬澳門元或以上之"
            f"{n_sel}個項目進行現場走訪，涉及金額{_amt(amt_sel)}，佔母體金額"
            f"{_pct(amt_sel / amt_all) if amt_all else '—'}。")
    render_generic(prs, f"{ent_up} 設施建設項目現場走訪情況匯總", rows, sec=3,
                   crumb="其他信息  |  本次審查工作執行的程序匯總", headline=head, side=False,
                   note="註：金額單位為萬澳門元；母體為報告投資金額中歸類為設施建設（資本性支出）之項目。")


# ── from make_report ──
def render_artwork(prs, ent_up, biao2_dir="data/表2", entity="mgm"):
    """報告 slide 101『藝術品展出情況清單』：表2 附件『藝術品』sheet 逐件列示。"""
    cols, body = load_artwork(biao2_dir, entity, log=print)
    if not cols or not body:
        return
    KEEP = ["名稱", "類別", "Artist", "購入", "當前位置", "當前狀態", "展出紀錄"]
    idx = [i for i, c in enumerate(cols) if any(k in c for k in KEEP)]
    if not idx:
        idx = list(range(min(7, len(cols))))
    hdr = ["序號"] + [cols[i][:14] for i in idx]
    rows = pd.DataFrame([[str(n)] + [body[n - 1][i][:60] for i in idx]
                         for n in range(1, len(body) + 1)], columns=hdr)
    st_i = next((i for i, c in enumerate(cols) if "當前狀態" in c), None)
    shown = sum(1 for r in body if st_i is not None and "展出" in r[st_i]) if st_i is not None else 0
    head = (f"下表列示{ent_up}已購入之藝術品共{len(body)}件之展出情況"
            + (f"，其中{shown}件現正展出，{len(body) - shown}件未在展出中" if st_i is not None else "")
            + "。藝術品之社會價值須透過面向公眾持續展出方能體現，故其展出情況為本次審查"
              "「未完全實現投資目的的投資支出」之判斷依據。")
    render_generic(prs, f"{ent_up} 藝術品展出情況清單", rows, sec=5,
                   crumb="附件  |  藝術品展出情況清單", headline=head, side=False,
                   note="資料來源：承批公司提供之藝術品清單（審查底稿表2 附件），畢馬威分析")


# ── from make_report ──
def render_cumulative(prs, ent_up, df, plan, cat=None):
    """2.5 截至2025年末投資金額概覽（scan slide 26）。"""
    tbl = _cum_table(df, plan, cat)
    if tbl.empty:
        return
    t = tbl[tbl["範疇"].astype(str).str.strip() == "合計"]
    r = t.iloc[0] if len(t) else None

    def v(c):
        x = r.get(c) if r is not None else None
        return float(x) if isinstance(x, (int, float)) and not pd.isna(x) else 0.0
    CUM = "截至2025年末三年累計"
    a, dd = v(f"{CUM}·獲批的計劃投資金額"), v(f"{CUM}·獲認可／潛在調整後投資金額")
    g = tbl[tbl["範疇"].astype(str).str.strip() == "博彩項目小計"]
    ng = tbl[tbl["範疇"].astype(str).str.strip() == "非博彩項目小計"]

    def sv(x, c):
        return (float(x.iloc[0][c]) if len(x) and isinstance(x.iloc[0][c], (int, float))
                and not pd.isna(x.iloc[0][c]) else 0.0)
    head = (f"截至2025年末，{ent_up}於2023至2025年三年的年度投資計劃累計獲批的計劃投資金額為{_amt(a)}"
            f"（博彩項目{_amt(sv(g, f'{CUM}·獲批的計劃投資金額'))}和"
            f"非博彩項目{_amt(sv(ng, f'{CUM}·獲批的計劃投資金額'))}），"
            f"累計獲認可／潛在調整後投資金額為{_amt(dd)}，"
            f"累計投資計劃金額完成率為{_pct(_rate(dd, a))}"
            f"（博彩項目完成率為{_pct(sv(g, f'{CUM}·完成率'))}，"
            f"非博彩項目完成率為{_pct(sv(ng, f'{CUM}·完成率'))}）。")
    render_generic(prs, f"{ent_up} 截至2025年末投資金額概覽", tbl.fillna(""), sec=1,
                   crumb="過往年度投資計劃在2025年繼續執行的審查跟進  |  截至2025年末投資金額概覽",
                   headline=head, side=False,
                   note="註：金額單位為萬澳門元。「2025年前已獲認可」＝該年度計劃於2025年之前"
                        "（即當年及往年審查）已認可之投資金額；「2025年期後」＝於2025年發生之期後投資。")


# ── from make_report ──
def render_generic(prs, title, df, *, sec=3, crumb=None, headline=None, note=None,
                   llm=None, tbl_id=None, side=None):
    """單張表（範疇/項目 + 數字欄；·=2-row group header）。逐頁：crumb + 導語 + caption bar
    + 表 + 資料來源，按【累積高度】分頁（唔會超出版面）。

    side=True → 報告式 2 欄（表左 + 敘述右，對 scan p-10~p-13）。敘述優先用 LLM 寫嘅
    `llm['tbl'][tbl_id]`，冇就用 _table_bullets 機械生成。欄數少（≤8）而且一版放得落先會用。"""
    if df is None or df.empty:
        return
    raw = (llm or {}).get("tbl", {}).get(tbl_id) or []
    if isinstance(raw, dict):                      # 新 shape：{"導語": …, "段落": [[h,b],…]}
        headline = raw.get("導語") or headline
        raw = raw.get("段落") or []
    bullets = [tuple(x) for x in raw] or _table_bullets(df)
    if side is None:
        side = len(df.columns) <= 8 and bool(bullets)
    if side and bullets:
        W, _H = size_of(prs)
        lw = W * 0.60
        s2, r2, w2, sp2 = _df_table(df)
        wid2 = [w * lw / sum(w2) for w in w2]
        need = header_h(sp2, s2, wid2, 5.0) + sum(row_h(c, wid2, 5.0) for _, c in r2)
        if need <= CONTENT_BOTTOM - 1.9:              # 一版放得落先用 2 欄，否則落返全闊分頁
            render_overview_page(prs, (crumb or title), headline or _total_line(df), df,
                                 bullets, sec=sec, table_name=title, note=note)
            return
    df = _overview_display(df) if "範疇" in df.columns and "項目數量" in df.columns else df
    subs, rows, widths, supers = _df_table(df)
    W, H = size_of(prs)
    tw = W - 2 * MARGIN
    wid = [w * tw / sum(widths) for w in widths]
    head = headline or _total_line(df)
    crumb = crumb or title
    # 先用一版試高度（導語行數會食掉可用高）
    probe_top = HEAD_Y + head_h(head, W)[0] + 0.10
    avail = CONTENT_BOTTOM - probe_top - 0.24
    hh = header_h(supers, subs, wid, 5.5)
    hcols = _hdr_cols(subs, supers)      # 全闊表一樣要派重點欄色（4.1 最右「潛在調整後投資金額」）
    pages = fit_rows(rows, wid, 6.5, avail, hh)
    for pi, chunk in enumerate(pages):
        suffix = f"（{pi+1}/{len(pages)}）" if len(pages) > 1 else ""
        slide, W, H, top = _page(prs, sec, crumb, (head or "") + suffix)
        top = caption_bar(slide, MARGIN, top, tw, title + suffix)
        draw_table(slide, MARGIN, top, tw, subs, chunk, widths, supers=supers,
                     font=SZ_TBL, hfont=SZ_TBL_HDR, fill_h=CONTENT_BOTTOM - top - 0.28,
                     hdr_cols=hcols)
        source_note(slide, W, note=note, more=(pi < len(pages) - 1))


# ── from make_report ──
def _cards(prs, sec, crumb, headline, recs, *, note=None):
    """逐個項目一張 card（navy 標題條 + 敘述段），按【累積高度】排版分頁 → 填滿版面唔留大白位。
    recs = [(bar_text, [(label, body)])]。"""
    W, H = size_of(prs)
    cw = W - 2 * MARGIN
    probe = HEAD_Y + head_h(headline, W)[0] + 0.10
    avail = CONTENT_BOTTOM - probe

    def card_h(items):
        return 0.24 + est_prose_h(items, cw - 0.12, head_size=SZ_BODY, body_size=SZ_BODY, gap=3) + 0.14
    pages, cur, used = [], [], 0.0
    for rec in recs:
        h = card_h(rec[1])
        if cur and used + h > avail:
            pages.append(cur); cur, used = [], 0.0
        cur.append(rec); used += min(h, avail)
    if cur:
        pages.append(cur)
    for pi, page in enumerate(pages):
        suffix = f"（{pi+1}/{len(pages)}）" if len(pages) > 1 else ""
        slide, W, H, y = _page(prs, sec, crumb, (headline or "") + suffix)
        for bar_text, items in page:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(y),
                                         Inches(cw), Inches(0.22))
            bar.fill.solid(); bar.fill.fore_color.rgb = HDR      # #00338D（card 條，唔跟表頭色）
            bar.line.fill.background(); bar.shadow.inherit = False
            btf = bar.text_frame; btf.word_wrap = True
            btf.margin_left = Emu(54000); btf.margin_top = btf.margin_bottom = Emu(0)
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            br = btf.paragraphs[0].add_run(); br.text = bar_text
            setfont(br, 8, bold=True, color=LIGHT)   # Arial latin + Microsoft YaHei ea（同全份一致）
            bh = min(est_prose_h(items, cw - 0.12, head_size=SZ_BODY, body_size=SZ_BODY, gap=3),
                     CONTENT_BOTTOM - y - 0.26)
            prose_box(slide, MARGIN + 0.06, y + 0.26, cw - 0.12, bh, items,
                        head_size=SZ_BODY, body_size=SZ_BODY, gap=3)
            y += 0.24 + bh + 0.14
        source_note(slide, W, note=note, more=(pi < len(pages) - 1))


# ── from make_report ──
def _finding_body(box, find, mgmt, grey=None):
    """（保留舊 API）KPMG分析發現 / 管理層解釋 兩段。"""
    prose(box, [(l + "：", t) for l, t in
                  [("KPMG分析發現", find), ("管理層解釋", mgmt)] if t],
            head_size=SZ_BODY, body_size=SZ_BODY, gap=3)


# ── from make_report ──
def render_findings(prs, ent_up, df, narr, llm=None, b2=None):
    """③ 主要發現：每 canonical 調整類型 → 受影響項目 card = navy 標題條(項目+金額) + body。
    body 優先用 LLM 寫嘅『事項描述』（ground 住表2＋清單），管理層原話照樣保留；
    冇 LLM 就 fallback 清單抄字（KPMG分析發現／管理層解釋）。"""
    llm_proj = (llm or {}).get("proj", {})
    d = df.copy()
    d["_adj"] = d["調整一級"].map(CANON).fillna(d["調整一級"])
    d.loc[~d["_adj"].isin(ADJ7), "_adj"] = ADJ_POST     # 殘差＝第 8 類（同 1.4／2.2 表一致）
    for adj in ADJ_ALL:
        sub = d[(d["_adj"] == adj) & (pd.to_numeric(d["調整_萬"], errors="coerce") != 0)]
        if sub.empty:
            continue
        projs = sub.groupby(["ng_scope", "dicj code"]).agg(名稱=("project", "first"),
                             報告=("調整前_萬", "sum"), 調整=("調整_萬", "sum")).reset_index()
        projs = projs.reindex(projs["調整"].abs().sort_values(ascending=False).index)
        recs = []
        for _, p in projs.iterrows():
            nr = nlook(narr, p["ng_scope"], p["dicj code"])
            desc = llm_proj.get(proj_key(adj, p["ng_scope"], p["dicj code"]), "")
            if desc:      # LLM 寫嘅事項描述（ground 表2＋清單）＋ 管理層原話
                items = [("事項描述：", desc)]
                if nr.get("管理層解釋"):
                    items.append(("管理層解釋：", nr["管理層解釋"]))
            else:
                items = [(l + "：", t) for l, t in
                         [("KPMG分析發現", nr.get("KPMG分析發現", "")),
                          ("管理層解釋", nr.get("管理層解釋", "")),
                          ("跨司工作組／KPMG意見", nr.get("跨司回覆", "") or nr.get("KPMG回覆", ""))] if t]
                if not items and b2:      # 清單冇 → 用表2 抽到嘅原文頂住
                    t2 = b2text(b2, p["ng_scope"], p["dicj code"])
                    if t2:
                        items = [("事項描述：", t2[:600])]
            if not items:
                items = [("", "清單未提供本項目之分析發現，待項目組補充。")]
            recs.append((f"{p['dicj code']}　{str(p['名稱'])[:34]}　│　報告 "
                         f"{fmt_money(p['報告'])}／潛在調整 {fmt_money(p['調整'])} 萬澳門元", items))
        tot = sub["調整_萬"].sum()
        head = (f"{ent_up} 報告的投資金額中，屬「{adj}」之潛在調減金額合計約{abs(tot):,.0f}萬澳門元，"
                f"涉及{len(recs)}個投資項目，逐項說明如下：")
        _cards(prs, 2, f"本年度審查工作的主要發現  |  {adj}", head, recs)


# ── from make_report ──
def render_site_visits(prs, ent_up, df, narr, threshold=2000):
    """附件二 現場走訪（slide 93-100）：capex ≥ 2,000萬 設施項目（報告走訪準則），
    配清單 實施地點 + 實際投資內容 做走訪概述。每頁 2 個項目 card。"""
    cap = df[df["final_capex_opex"] == "Capex"].copy()
    cap = cap[cap["dicj code"].astype(str).str.match(r"^項目\s*\d")]
    g = cap.groupby(["ng_scope", "dicj code"]).agg(
        名稱=("project", "first"), 報告=("調整前_萬", "sum")).reset_index()
    g = g[g["報告"] >= threshold]
    if g.empty:
        return
    g["_s"] = (g["ng_scope"] != "gaming").astype(int)
    g = g.sort_values(["_s", "報告"], ascending=[True, False])
    recs = []
    for _, p in g.iterrows():
        nr = nlook(narr, p["ng_scope"], p["dicj code"])
        items = [(l + "：", t) for l, t in
                 [("實施地點", nr.get("實施地點", "")),
                  ("現場走訪概述", nr.get("實際投資內容", "")),
                  ("現場走訪圖片", "〔待插入〕")] if t]
        recs.append((f"{p['dicj code']}　{str(p['名稱'])[:32]}　│　設施建設（資本性支出）"
                     f"{fmt_money(p['報告'])} 萬澳門元", items))
    head = (f"我們就{ent_up}報告投資金額中設施建設（資本性支出）達{threshold:,.0f}萬澳門元或以上之"
            f"{len(recs)}個投資項目進行了現場走訪，走訪情況如下：")
    _cards(prs, 5, "附件  |  部分項目的現場走訪情況", head, recs,
           note="資料來源：現場走訪記錄、管理層提供之項目資料，畢馬威分析")


# ── from make_report ──
def _prose_slide(prs, title, bullets, headline=None, *, sec=0):
    """一版敘述（crumb + navy 導語 + 段落），按估算高度自動分頁。"""
    W, H = size_of(prs)
    cw = W - 2 * MARGIN
    probe = HEAD_Y + head_h(headline, W)[0] + 0.10
    pages = fit_prose(bullets, cw, CONTENT_BOTTOM - probe, head_size=8, body_size=8)
    for pi, page in enumerate(pages):
        suffix = f"（{pi+1}/{len(pages)}）" if len(pages) > 1 else ""
        slide, W, H, top = _page(prs, sec, title, (headline or "") + suffix)
        prose_box(slide, MARGIN, top, cw, CONTENT_BOTTOM - top, page,
                    head_size=8, body_size=8, gap=7)
        source_note(slide, W, more=(pi < len(pages) - 1))


# ── from make_report ──
def _headline(ent_up, ov, df, plan):
    """slide 10 整體投資支出概況 → 回 (headline 句, bullets)。全 mechanical 由數字生成。"""
    tot = ov[ov["範疇"] == "總計"]
    if not len(tot):
        return "", []
    r = tot.iloc[0]

    def num(x):
        return float(x) if isinstance(x, (int, float)) and not pd.isna(x) else 0.0
    plan_amt = num(r.get("獲批的計劃投資金額"))
    report_amt = num(r.get("報告投資金額"))
    after_amt = num(r.get("潛在調整後投資金額"))
    adj_amt = report_amt - after_amt
    rate = r.get("投資計劃完成率")
    after_rate = r.get("潛在調整後投資計劃完成率")

    d = df.copy()
    d["_adj"] = pd.to_numeric(d["調整_萬"], errors="coerce").fillna(0)
    # ⚠ 三個數要自洽：計劃(母體) ⊇ 已實施；之前 n_impl 數 feed 全部有支出嘅碼，
    #   同 n_plan 唔同母體 → 出現「實施 112 > 計劃 89、0 個未發生」（項目組 2026-08-13 指出）
    n_plan, n_impl, n_zero = _proj_counts(df, plan, ov)
    n_adj = d[d["_adj"] != 0]["dicj code"].nunique()
    headline = (f"{ent_up} 2025年度原獲批計劃開展{n_plan}個投資項目，涉及計劃投資金額約{_amt(plan_amt)}；"
                f"{ent_up}提交的投資執行報告顯示2025年度投資金額約{_amt(report_amt)}"
                f"（投資計劃金額完成率{_pct(rate)}）。本次審查工作識別潛在調減金額約{_amt(adj_amt)}，"
                f"經潛在調整後的2025年度投資金額約{_amt(after_amt)}（經調整後投資計劃金額完成率{_pct(after_rate)}）。")
    bullets = [
        ("2025年度獲批的計劃投資金額與報告的投資金額：",
         f"根據{ent_up}提交的2025年度投資計劃方案與投資執行報告，{ent_up}獲批計劃開展{n_plan}個投資項目，"
         f"計劃投資金額約{_amt(plan_amt)}。投資執行報告顯示實際開展其中{n_impl}個投資項目"
         # 0 個未發生就唔好寫「計劃中有0個項目未發生投資金額」（項目組 2026-08-13 指出係錯）
         + (f"（計劃中有{n_zero}個項目未發生投資金額）" if n_zero else "")
         + f"，報告投資金額約{_amt(report_amt)}，報告投資計劃金額完成率為{_pct(rate)}。"),
        ("2025年度投資支出金額的潛在調整事項：",
         f"我們在本次審查工作中發現，{ent_up}報告投資金額中存在部分投資支出可能不應確認為2025年度計劃的投資支出"
         f"（涉及{n_adj}個投資項目，合計約{_amt(adj_amt)}）。考慮潛在調減事項後，{ent_up} 2025年度投資支出金額"
         f"約{_amt(after_amt)}，投資計劃金額完成率應為{_pct(after_rate)}。"),
    ]
    return headline, bullets


# ── from make_report ──
def _pct(v):
    return f"{v*100:.1f}%" if isinstance(v, (int, float)) and not pd.isna(v) else "—"


# ── from make_report ──
def _trim(s):
    """剝尾標點 —— 清單／表2 原文多數自帶「。」，接落我哋句式會變「。。」。"""
    return str(s or "").strip().rstrip("。．.；;，,、 ")


# ── from make_report ──
_CN_NUM = "零一二三四五六七八九十"


# ── from make_report ──
def _cn(n):
    """1→一、7→七、12→十二（報告寫『存在七大類的調整事項』）。"""
    n = int(n)
    if n <= 10:
        return _CN_NUM[n]
    if n < 20:
        return "十" + (_CN_NUM[n - 10] if n > 10 else "")
    return f"{n}"


# ── from make_report ──
def _amt(wan):
    """萬 → 報告用字。scan：≥1億寫『6.4億澳門元』，<1億寫『5,527萬澳門元』。"""
    try:
        w = abs(float(wan or 0))
    except (TypeError, ValueError):
        return "—"
    return f"{w/10000:.1f}億澳門元" if w >= 10000 else f"{w:,.0f}萬澳門元"


# ── from make_report ──
def _cats_of(ov, col="報告投資金額", n=3, scope=None):
    """由 overview 表攞金額最大嘅幾個範疇名（scan：『主要涉及會議展覽、文化藝術、社區旅遊等…』）。"""
    if ov is None or ov.empty or col not in ov.columns:
        return ""
    d = ov[~ov["範疇"].astype(str).str.endswith(("小計", "總計", "合計", "項目"))].copy()
    d["_v"] = pd.to_numeric(d[col], errors="coerce").fillna(0)
    d = d[d["_v"] > 0].sort_values("_v", ascending=False)
    named = d[~d["範疇"].astype(str).isin(["其他", "其它"])]      # 「其他」唔好排喺點名範疇最前
    d = (named if len(named) >= min(n, 2) else d).head(n)
    return "、".join(str(r["範疇"]) for _, r in d.iterrows())


# ── from make_report ──
def _pg(i, n):
    """scan 導語尾有頁碼標記『（1/2）』。"""
    return f"（{i}/{n}）" if n > 1 else ""


# ── from make_report ──
def _zero_intro(ent_up, zi):
    """報告 1.2 頁尾段：未發生投資項目嘅大致介紹（跨年／內部研究／取消），並指去 1.3。"""
    if not zi:
        return None
    n, tot, groups = zi
    seg = []
    for kind, txt in (("跨年", "為跨年項目，{e}仍在就以前年度計劃項目進行持續投資，於2025年將投資額"
                              "作為2023年度計劃或2024年度計劃期後投資金額進行申報"),
                      ("內部研究", "個項目由於處於內部研究、計劃階段或未收到詳細指引未發生實際支出"),
                      ("取消", "個項目已取消")):
        g = groups.get(kind) or []
        if not g:
            continue
        seg.append((f"其中{len(g)}個" if kind == "跨年" else f"{len(g)}") + txt.format(e=ent_up))
    return ("未發生投資項目的大致介紹",
            f"{ent_up}於2025年度的投資計劃中有{n}個非博彩項目未產生投資金額，"
            + "；".join(seg) + "。請見後續「2025年度投資項目的整體執行概況」。")


# ── from make_report ──
def _bucket_headline(ent_up, bucket, ov):
    """②期後概覽導語（由表自己嘅合計行計，自洽）。"""
    tot = ov[ov["範疇"].astype(str).str.strip().isin(("合計", "總計"))]
    if tot.empty:
        return ""
    r = tot.iloc[0]

    def n(c):
        v = r.get(c)
        return float(v) if isinstance(v, (int, float)) and not pd.isna(v) else 0.0
    rep, a, aft = n("報告投資金額"), n("潛在調整金額"), n("潛在調整後投資金額")
    yr, npj = bucket[:4], int(n("項目數量"))
    cats = _cats_of(ov, "報告投資金額", 3)
    tail = (f"，主要涉及{cats}等非博彩投資範疇的{npj}個項目" if cats else f"，涉及{npj}個投資項目")
    return (f"{ent_up}在2025年度執行報告中申報的「因發生期後事項需作後續調整之{yr}年度博彩／非博彩項目」"
            f"投資金額為{_amt(rep)}。本次審查工作識別潛在調減金額約{_amt(a)}，"
            f"經潛在調減後的{yr}年度計劃投資項目在2025年的投資金額約{_amt(aft)}{tail}。")


# ── from make_report ──
def _rate_of(df, name, col):
    r = df[df["範疇"] == name]
    if not len(r) or col not in df.columns:
        return None
    v = r.iloc[0][col]
    return v if isinstance(v, (int, float)) and not pd.isna(v) else None


# ── from make_report ──
def _prose_paginated(prs, title, bullets, per):
    if not bullets:
        return
    pages = [bullets[i:i + per] for i in range(0, len(bullets), per)]
    for pi, page in enumerate(pages):
        t = title + (f"（{pi+1}/{len(pages)}）" if len(pages) > 1 else "")
        _prose_slide(prs, t, page)


# ── from make_report ──
def _prose_pages(prs, bullets, headline=None, subtitle=None):
    """_prose_2col 嘅分頁計算，抽咗出嚟 —— caller（1.4／2.2／2.4）要【預先知總頁數】
    先寫得出表版嗰句「（1/3）」（scan p15 = 1/3、p16 = 2/3、p17 = 3/3）。"""
    if not bullets:
        return []
    numbered = len(bullets[0]) == 3
    W, _ = size_of(prs)
    colw = (W - 2 * MARGIN - COL_GAP) / 2
    probe = HEAD_Y + head_h(headline, W)[0] + 0.10
    avail = CONTENT_BOTTOM - probe - (0.2 if subtitle else 0)
    if not numbered:
        return fit_prose(bullets, colw, avail * 2, head_size=SZ_BODY_HEAD, body_size=SZ_BODY)
    pages, cur, used = [], [], 0.0
    for b in bullets:
        hh = est_numbered_h([b], colw, size=SZ_BODY)
        if cur and used + hh > avail * 2:
            pages.append(cur); cur, used = [], 0.0
        cur.append(b); used += hh
    if cur:
        pages.append(cur)
    return pages


# ── from make_report ──
def _prose_2col(prs, title, bullets, per=12, subtitle=None, *, sec=0, headline=None,
                pg0=0, pgn=0):
    """報告式 2 欄敘述（對 scan p16-17：左右兩欄，每欄頂有 navy 小標題，右欄加「（續）」）。
    pg0/pgn = 前面已經有幾多版／成節總共幾多版（表版計埋）→ 導語尾寫「（2/3）」。"""
    if not bullets:
        return
    numbered = bool(bullets) and len(bullets[0]) == 3      # (no, head, body) = scan 編號清單
    W, H = size_of(prs)
    colw = (W - 2 * MARGIN - COL_GAP) / 2
    probe = HEAD_Y + head_h(headline, W)[0] + 0.10
    avail = CONTENT_BOTTOM - probe - (0.2 if subtitle else 0)
    half_pages = _prose_pages(prs, bullets, headline, subtitle)
    n_all = pgn or len(half_pages)
    for pi, page in enumerate(half_pages):
        suffix = f"（{pg0 + pi + 1}/{n_all}）" if n_all > 1 else ""
        slide, W, H, top = _page(prs, sec, title, (headline or "") + suffix)
        if subtitle:      # scan p16：小標題喺【每欄】頂，右欄加「（續）」
            put(slide, MARGIN, top, colw, 0.18, subtitle, size=7, bold=True, color=NAVY)
            put(slide, MARGIN + colw + COL_GAP, top, colw, 0.18, subtitle + "（續）",
                  size=7, bold=True, color=NAVY)
            top += 0.22
        # 斷欄：以【總高一半】為目標令左右大致平均（對 scan），但唔可以超過一欄可用高
        lim = CONTENT_BOTTOM - top
        hs = [(est_numbered_h([it], colw, size=SZ_BODY) if numbered
               else est_prose_h([it], colw, head_size=SZ_BODY_HEAD, body_size=SZ_BODY)) for it in page]
        target = sum(hs) / 2.0
        cut, used = len(page), 0.0
        for i, ih in enumerate(hs):
            if i and (used >= target or used + ih > lim):
                cut = i; break
            used += ih
        cut = max(1, cut)
        if numbered:
            prose_numbered(_tb(slide, MARGIN, top, colw, lim), page[:cut], size=SZ_BODY)
            if page[cut:]:
                prose_numbered(_tb(slide, MARGIN + colw + COL_GAP, top, colw, lim),
                                 page[cut:], size=SZ_BODY)
        else:
            prose_box(slide, MARGIN, top, colw, lim, page[:cut], head_size=SZ_BODY_HEAD, body_size=SZ_BODY)
            if page[cut:]:
                prose_box(slide, MARGIN + colw + COL_GAP, top, colw, lim, page[cut:],
                            head_size=SZ_BODY_HEAD, body_size=SZ_BODY)
        source_note(slide, W, more=(pi < len(half_pages) - 1))


# ── from make_report ──
def render_category_overview(prs, ent_up, ov, df, narr, llm=None, ovx=None, note=None):
    """報告 slide 11-14（1.3 整體執行概況）：【左邊照舊擺 1.2 嗰個整體概況表】、右邊逐範疇敘述，
    敘述長就分版（scan 係 1/4 … 4/4，四版嘅表一模一樣）。LLM summary 優先，否則清單抄字。"""
    if not narr:
        return
    llm_cat = (llm or {}).get("cat", {})
    d = df.copy()
    d["_sub"] = sub_of(d)
    proj = d.groupby(["_sub", "dicj code"])["調整前_萬"].sum().reset_index()
    cats = ov[~ov["範疇"].astype(str).str.endswith(("小計", "總計", "項目"))]
    g_bul, n_bul = [], []       # 按範疇概況：博彩 / 非博彩 各自一版（報告 3/4、4/4）
    for _, r in cats.iterrows():
        sub = str(r["範疇"]); rate = r.get("投資計劃完成率")
        if not isinstance(rate, (int, float)) or pd.isna(rate):
            continue
        if sub in llm_cat and llm_cat[sub]:               # LLM 寫嘅摘要優先
            txt = llm_cat[sub]
            txt = txt[len(sub) + 1:] if txt.startswith(sub + "：") else txt
            (g_bul if sub.startswith("博彩") else n_bul).append((f"{sub}：", txt)); continue
        scope = "gaming" if sub.startswith("博彩") else "non_gaming"
        pr = proj[proj["_sub"] == sub].sort_values("調整前_萬", ascending=False)
        content = reason = ""       # content=清單實際投資內容；reason=清單管理層解釋(變更原因)
        for _, pp in pr.iterrows():
            nr = nlook(narr, scope, pp["dicj code"])
            if not content:
                content = nr.get("實際投資內容", "")
            if not reason:
                reason = nr.get("管理層解釋", "")   # 業務原因；唔用 KPMG分析發現（審計腔）
            if content and reason:
                break
        content, reason = _trim(content), _trim(reason)   # 清單原文已有句號 → 唔剝就出「。。」
        # ⚠ 字數上限 = 版面約束：非博彩 11 個範疇要一版放晒（報告 4/4 得一版）→ 每個 ~55 字。
        #   同 build_llm_narrative._cat_prompt 嘅 40-55 字一致，LLM 同 fallback 都唔會撐爆版。
        summ = (content[:34] + "…") if len(content) > 34 else content
        rsn = ("，主要由於" + (reason[:24] + "…" if len(reason) > 24 else reason)) if reason else ""
        body = (f"主要包括{summ}。投資計劃金額完成率為{_pct(rate)}{rsn}。" if summ
                else f"投資計劃金額完成率為{_pct(rate)}{rsn}。")
        (g_bul if sub.startswith("博彩") else n_bul).append((f"{sub}：", body))
    # scan p-06 句式：著重於投入{博彩範疇}等博彩項目，以及{非博彩範疇}等非博彩投資項目 + 前後完成率
    gm = ov[ov["範疇"].astype(str).str.startswith("博彩") &
            ~ov["範疇"].astype(str).str.endswith(("小計", "項目"))]
    ng = ov[~ov["範疇"].astype(str).str.startswith("博彩") &
            ~ov["範疇"].astype(str).str.endswith(("小計", "總計", "項目"))]
    g_cats = _cats_of(gm, "報告投資金額", 2) or "博彩娛樂場場地的優化"
    n_cats = _cats_of(ng, "報告投資金額", 5)
    g_r, ng_r = _rate_of(ov, "博彩項目小計", "投資計劃完成率"), _rate_of(ov, "非博彩項目小計", "投資計劃完成率")
    g_a, ng_a = (_rate_of(ov, "博彩項目小計", "潛在調整後投資計劃完成率"),
                 _rate_of(ov, "非博彩項目小計", "潛在調整後投資計劃完成率"))
    # scan slide 11 尾句：「博彩項目各範疇…均超過100%，非博彩項目…平均完成率為44.2%，未有達到100%完成率的範疇」
    def _sub100(d):
        c = d[d["潛在調整後投資計劃完成率"].apply(
            lambda v: isinstance(v, (int, float)) and not pd.isna(v))] if \
            "潛在調整後投資計劃完成率" in d.columns else d.iloc[0:0]
        return c[c["潛在調整後投資計劃完成率"] < 1.0]
    g_all100 = "博彩項目各範疇的投資計劃金額完成率均超過100%" if len(_sub100(gm)) == 0 else \
               f"博彩項目的投資計劃金額完成率為{_pct(g_a)}"
    ng_lo = _sub100(ng)
    ng_tail = ("未有達到100%完成率的範疇" if len(ng_lo) == len(ng) and len(ng) else
               ("未達到100%完成率的範疇包括" +
                "、".join(f"{r['範疇']}（{_pct(r['潛在調整後投資計劃完成率'])}）"
                          for _, r in ng_lo.sort_values("潛在調整後投資計劃完成率").head(4).iterrows())
                if len(ng_lo) else "各範疇的投資計劃金額完成率均超過100%"))
    head = (f"{ent_up}的2025年度計劃投資項目著重於投入{g_cats}等博彩項目，以及{n_cats}等非博彩投資項目。"
            f"在報告投資金額中，博彩項目的投資計劃金額完成率為{_pct(g_r)}，"
            f"非博彩項目的投資計劃金額完成率為{_pct(ng_r)}。考慮投資金額的潛在調整後，"
            f"{g_all100}，非博彩項目的投資計劃金額平均完成率為{_pct(ng_a)}，{ng_tail}。")
    hi_s = "、".join(f"{r['範疇']}（{_pct(r['投資計劃完成率'])}）" for _, r in
                    cats[cats["投資計劃完成率"].apply(
                        lambda v: isinstance(v, (int, float)) and not pd.isna(v) and v >= 1.0)]
                    .sort_values("投資計劃完成率", ascending=False).iterrows())
    lo_s = "、".join(f"{r['範疇']}（{_pct(r['投資計劃完成率'])}）" for _, r in
                    cats[cats["投資計劃完成率"].apply(
                        lambda v: isinstance(v, (int, float)) and not pd.isna(v) and v < 0.5)]
                    .sort_values("投資計劃完成率").iterrows())
    exec_b = []               # 1/4：整體執行概況
    if hi_s:
        exec_b.append(("報告投資金額完成率較高的範疇", f"{hi_s}。"))
    if lo_s:
        exec_b.append(("報告投資金額完成率相對較低的範疇", f"{lo_s}。"))
    # ★ 報告 1.3 係【4 版、每版右邊唔同主題】（項目組 2026-08-17 指明），
    #   唔係逐範疇順住排落去分頁。左邊 4 版都係同一個 1.2 概況表。
    groups = [("2025年度投資項目的整體執行概況", exec_b),
              ("2025年度投資計劃區分設施建設/活動舉辦的投資金額", _fac_bullets(ent_up, ov)),
              ("按範疇的項目概況 — 博彩項目", g_bul),
              ("按範疇的項目概況 — 非博彩項目", n_bul)]
    render_overview_pages(prs, "2025年度投資計劃執行情況概述  |  2025年度投資項目的整體執行概況",
                          head, ovx if ovx is not None else ov, groups, sec=0, note=note,
                          table_name=f"{ent_up} 2025年度計劃的整體投資支出概況", grouped=True)


# ── from make_report ──
def _fac_bullets(ent_up, ov):
    """1.3 第 2 版：區分設施建設／活動舉辦嘅投資金額（由概況表自己嗰兩欄機械計）。"""
    F, A = "設施建設/資本性支出", "活動舉辦/營運性支出"
    if ov is None or ov.empty or F not in ov.columns:
        return []

    def num(row, c):
        v = row.get(c)
        return float(v) if isinstance(v, (int, float)) and not pd.isna(v) else 0.0

    def line(label):
        r = ov[ov["範疇"].astype(str).str.strip() == label]
        return (num(r.iloc[0], F), num(r.iloc[0], A)) if len(r) else (0.0, 0.0)
    tf, ta = line("總計")
    gf, ga = line("博彩項目小計")
    nf, na = line("非博彩項目小計")
    cat = ov[~ov["範疇"].astype(str).str.endswith(("小計", "總計", "項目"))].copy()
    if cat.empty or (tf + ta) == 0:
        return []
    cat["_f"] = cat.apply(lambda r: num(r, F), axis=1)
    cat["_a"] = cat.apply(lambda r: num(r, A), axis=1)
    top_f = "、".join(f"{r['範疇']}（{_amt(r['_f'])}）" for _, r in
                     cat.sort_values("_f", ascending=False).head(3).iterrows() if r["_f"] > 0)
    top_a = "、".join(f"{r['範疇']}（{_amt(r['_a'])}）" for _, r in
                     cat.sort_values("_a", ascending=False).head(3).iterrows() if r["_a"] > 0)
    return [
        ("設施建設／資本性支出",
         f"考慮潛在調整事項後，設施建設／資本性支出的投資金額約{_amt(tf)}"
         f"（佔{tf / (tf + ta) * 100:.1f}%），其中博彩項目約{_amt(gf)}、非博彩項目約{_amt(nf)}。"
         + (f"金額較大的範疇包括{top_f}。" if top_f else "")),
        ("活動舉辦／營運性支出",
         f"考慮潛在調整事項後，活動舉辦／營運性支出的投資金額約{_amt(ta)}"
         f"（佔{ta / (tf + ta) * 100:.1f}%），其中博彩項目約{_amt(ga)}、非博彩項目約{_amt(na)}。"
         + (f"金額較大的範疇包括{top_a}。" if top_a else "")),
    ]


# ── from make_report ──
def _adj_detail_bullets(ent_up, adj, df, narr, llm=None):
    """slide 16-17 潛在調整事項詳述 → 回 bullets：LLM summary 優先，否則清單分析發現。"""
    if not narr:
        return []
    llm_adj = (llm or {}).get("adj", {})
    pb = BUCKET_ORDER[0]      # 2025計劃 bucket（唔用合計；期後另計）
    d = df.copy()
    d["_adj"] = d["調整一級"].map(CANON).fillna(d["調整一級"])
    bullets = []
    for _, r in adj.iterrows():
        t = str(r["潛在調整事項"])
        if t in ("合計", "跨年及其他調整"):
            continue
        amt = r.get(pb, 0)
        if not isinstance(amt, (int, float)) or abs(amt) < 0.5:
            continue
        no = adj_no(t)
        if t in llm_adj and llm_adj[t]:                   # LLM 寫嘅摘要優先
            bullets.append((no, f"{t}（{_amt(amt)}）：", llm_adj[t])); continue
        sub = d[(d["_adj"] == t) & (d["_bucket"] == pb) & (pd.to_numeric(d["調整_萬"], errors="coerce") != 0)]
        if sub.empty:
            continue
        names = "、".join(str(x) for x in sub.groupby("dicj code")["project"].first().tolist()[:3])
        reason = ruling = ""
        for _, pp in sub.drop_duplicates("dicj code").iterrows():
            nr = nlook(narr, pp["ng_scope"], pp["dicj code"])
            if not reason:
                reason = nr.get("KPMG分析發現", "") or nr.get("調整事項備註", "")
            if not ruling:
                ruling = nr.get("跨司回覆", "") or nr.get("KPMG回覆", "")
            if reason and ruling:
                break
        reason, ruling = _trim(reason), _trim(ruling)
        r2 = (reason[:150] + "…" if len(reason) > 150 else reason + "。") if reason else ""
        rl = ("跨司工作組／KPMG意見：" + (ruling[:90] + "…" if len(ruling) > 90 else ruling + "。")) if ruling else ""
        body = f"主要涉及{names}等項目。{r2}{rl}" if (r2 or rl) else f"主要涉及{names}等項目。"
        bullets.append((no, f"{t}（{_amt(amt)}）：", body))
    return bullets


# ── from make_report ──
def _exec_bullets(ent_up, ov):
    """整體執行概況敘述（報告 slide 11）：完成率高/低範疇，全部由 overview 數字生成 → 回 bullets。"""
    g = _rate_of(ov, "博彩項目小計", "投資計劃完成率")
    ng = _rate_of(ov, "非博彩項目小計", "投資計劃完成率")
    tot = _rate_of(ov, "總計", "投資計劃完成率")
    ga = _rate_of(ov, "博彩項目小計", "潛在調整後投資計劃完成率")
    nga = _rate_of(ov, "非博彩項目小計", "潛在調整後投資計劃完成率")
    cat = ov[~ov["範疇"].astype(str).str.endswith(("小計", "總計", "項目"))].copy()
    cat = cat[cat["投資計劃完成率"].apply(lambda v: isinstance(v, (int, float)) and not pd.isna(v))]
    high = cat[cat["投資計劃完成率"] >= 1.0].sort_values("投資計劃完成率", ascending=False)
    low = cat[cat["投資計劃完成率"] < 0.5].sort_values("投資計劃完成率")
    high_s = "、".join(f"{r['範疇']}（{_pct(r['投資計劃完成率'])}）" for _, r in high.iterrows())
    low_s = "、".join(f"{r['範疇']}（{_pct(r['投資計劃完成率'])}）" for _, r in low.iterrows())
    bullets = [
        ("", f"{ent_up}於2025年度計劃投資項目涵蓋博彩及非博彩範疇。在報告投資金額中，"
             f"博彩項目的投資計劃完成率為{_pct(g)}，非博彩項目為{_pct(ng)}，整體為{_pct(tot)}。"),
        ("", f"考慮投資金額的潛在調整後，博彩項目的投資計劃完成率為{_pct(ga)}，"
             f"非博彩項目的平均完成率為{_pct(nga)}。"),
    ]
    # ⚠ 高/低完成率範疇喺報告係喺【1.3 整體執行概況】嗰版，唔喺 1.2（項目組 2026-08-13）
    return bullets


# ── from make_report ──
def _adj_summary(ent_up, adj, ov=None, sdf=None):
    """潛在調整事項匯總（報告 slide 15）→ 回 (headline, bullets)。逐類型金額。
    ⚠ 用 2025計劃 bucket（唔係合計）：報告調整詳述只計 2025年度計劃，期後另有匯總。"""
    pb = BUCKET_ORDER[0]      # "2025年度投資計劃"
    tot_row = adj[adj["潛在調整事項"] == "合計"]
    total = tot_row.iloc[0][pb] if len(tot_row) else 0
    n_type = sum(1 for _, r in adj.iterrows()
                 if r["潛在調整事項"] not in ("合計", "跨年及其他調整")
                 and isinstance(r.get(pb), (int, float)) and abs(r.get(pb, 0)) >= 0.5)
    n_proj = after = None
    if sdf is not None:
        d = sdf[(sdf["_bucket"] == pb) & (pd.to_numeric(sdf["調整_萬"], errors="coerce") != 0)]
        n_proj = int(d["dicj code"].nunique())
    if ov is not None and not ov.empty:
        t = ov[ov["範疇"].astype(str).str.strip() == "總計"]
        if len(t):
            after = t.iloc[0].get("潛在調整後投資金額")
    # scan p-08 句式：存在七大類的調整事項，潛在調減約6.9億澳門元（涉及22個投資項目），經潛在調減後約11.8億澳門元
    headline = (f"基於我們的各項審查程序，我們認為{ent_up}報告的2025年度投資金額中，"
                f"存在{_cn(n_type)}大類的調整事項，潛在調減投資金額約{_amt(total)}"
                + (f"（涉及{n_proj}個投資項目）" if n_proj else "")
                + (f"，經潛在調減後的2025年度計劃的投資金額約{_amt(after)}。" if after is not None else "。"))
    bullets = []
    for _, r in adj.iterrows():
        name = r["潛在調整事項"]
        if name in ("合計", "跨年及其他調整"):
            continue
        amt = r[pb] if pb in adj.columns else 0
        if not isinstance(amt, (int, float)) or abs(amt) < 0.5:
            continue
        bullets.append((f"{name}：", f"約{abs(amt):,.0f}萬澳門元。"))
    return headline, bullets


def main():
    entity = (sys.argv[1] if len(sys.argv) > 1 and not sys.argv[1].startswith("-") else "mgm").lower()
    feed = Path(FEED)
    if not feed.exists():
        print(f"✗ 揾唔到 feed {feed}（root 應有 tableau_combined_25.csv）"); return
    qingdan = _find("data/投資項目清單", entity, ".xlsx")
    template = _find("data/reports", entity, ".pptx", prefer=["2025"])
    global ENT_UP
    ent_up = ENT_UP = entity.upper()
    print(f"build {BUILD_STAMP}")
    print(f"entity={ent_up}  feed={feed.name}  清單={qingdan.name if qingdan else '(冇)'}  "
          f"template={template.name if template else '(冇→用 13.33x7.5)'}")

    df = pd.read_csv(feed, low_memory=False)
    df = df[df["entity"].astype(str).str.lower() == entity]
    df["報告年"] = pd.to_numeric(df["報告年"], errors="coerce")
    df["_plan_year"] = df["year_bucket"].map(_plan_year)
    plan = load_plan(qingdan) if qingdan else None
    cat = load_category(qingdan) if qingdan else None     # 項目性質(D)→派零投資項目計劃返範疇
    narr = load_narrative(qingdan) if qingdan else {}     # 清單 by-project narrative（抄字）
    if narr:
        print(f"    清單 narrative: {sum(1 for r in narr.values() if r.get('KPMG分析發現'))} 個項目有發現")
    _coverage_probe(df, str(qingdan) if qingdan else None)   # 探 ❓頁（主體/KPI/藝術品）coverage
    # 有 workbench creds 就自動即場生成 LLM 敘述（毋須任何 flag）；冇 creds 靜靜跳過用清單 fallback；--no-llm 強制跳過
    av = sys.argv
    if "--no-llm" not in av:
        model = av[av.index("--model") + 1] if "--model" in av else None
        try:
            has_creds = bool(Workbench(model=model).config_masked().get("key_ok"))
        except Exception:
            has_creds = False
        if has_creds:
            # default 4：8 條並行俾公司網關 WAF 全部擋（2026-08-15，60/60「request is blocked」）
            workers = int(av[av.index("--workers") + 1]) if "--workers" in av else 4
            biao2_dir = av[av.index("--biao2") + 1] if "--biao2" in av else "data/表2"
            print("  由 feed+清單+表2 即場生成 LLM 敘述…")
            try:
                generate_llm_narrative(str(feed), entity, str(qingdan) if qingdan else None,
                                       biao2_dir=biao2_dir, model=model, workers=workers)
            except Exception as e:
                print(f"  ⚠ LLM 生成失敗（{type(e).__name__}: {e}）→ 用現有 json / 清單 fallback")
    llm = _load_llm(entity)     # {entity}_llm_narrative.json 有就用 LLM 文字，否則清單 fallback
    if llm:
        print("    LLM narrative: " + "、".join(f"{k} {len(llm.get(k, {}))}"
                                              for k in ("adj", "cat", "tbl", "proj", "bkt")) + " 段")

    global USE_TEMPLATE
    tmpl = _find_template() if "--use-template" in sys.argv else None   # 預設 fresh 手砌（template 樣式已 hardcode）；--use-template 先開 template
    USE_TEMPLATE = bool(tmpl)
    if tmpl:
        prs = Presentation(str(tmpl))
        _strip_slides(prs)      # drop_rel 正確清走原有 content slides（唔會 duplicate 名 corrupt）
        print(f"    template（master-driven）：{tmpl}  layouts={sum(len(m.slide_layouts) for m in prs.slide_masters)}")
    else:
        prs = Presentation()
        if template:
            prs.slide_width, prs.slide_height = Presentation(str(template)).slide_width, Presentation(str(template)).slide_height
        else:
            prs.slide_width, prs.slide_height = Inches(SLIDE_W), Inches(SLIDE_H)
        print(f"    無 template → fallback 手砌 formatting（{prs.slide_width/914400:.2f}x"
              f"{prs.slide_height/914400:.2f}in）")

    sdf = _load(feed, entity)     # 於2025發生 slice（概述 + 金額匯總 共用）

    render_cover(prs, entity)       # 封面（報告 p1）

    # ① 2025年度投資計劃執行情況概述（報告 slide 8-18）
    S1 = "2025年度投資計劃執行情況概述"
    divider(prs, S1, "1", [
        ("1.1  股權架構簡圖及發生投資支出的主體公司", ""),
        ("1.2  2025年度計劃的整體投資支出概況", ""),
        ("1.3  2025年度投資項目的整體執行概況", ""),
        ("1.4  2025年度投資計劃報告投資金額的潛在調整事項匯總", ""),
    ])
    budget = _load_budget(entity)
    ov = overview_by_bucket(sdf, "2025年度投資計劃", plan, cat)
    adj = adjustment_bridge(sdf)
    # 1.2 表底兩條註（逐字對 scan p10；註釋2 就係「項目數量含零申報」嗰條口徑）
    NOTE_RATE = ("註釋1：上述承諾的10年投資預算包含額外投資部分（澳門全年博彩毛收入達到1,800億澳門元後"
                 "觸發的非博彩範疇額外20%投資）\n"
                 "註釋2：項目數量是指各家承批公司在其2025年度投資執行報告中披露的投資項目數量，"
                 "包含申報的投資支出為零的部分")
    if not ov.empty:      # slide 10-11：表左 + headline/執行敘述右（報告 2 欄式）
        zi = zero_investment_summary(sdf, plan, cat, narr, ent_up)
        hl, hlb = _headline(ent_up, ov, sdf, plan)
        exb = _exec_bullets(ent_up, ov)
        zintro = _zero_intro(ent_up, zi)
        ovx = _overview_extra(ov, plan, sdf, budget, ent_up).fillna("")
        render_overview_page(prs, f"{S1}  |  2025年度計劃的整體投資支出概況",
                             hl, ovx, hlb + exb + ([zintro] if zintro else []), sec=0,
                             table_name=f"{ent_up} 2025年度的整體投資支出概況", note=NOTE_RATE)
        # slide 11-14 逐範疇概況（LLM 優先）；表照 1.2 嗰個逐版重複，同 scan 一致
        render_category_overview(prs, ent_up, ov, sdf, narr, llm, ovx=ovx, note=NOTE_RATE)
        zit = zero_investment_text(zi, ent_up)
        if zit:      # 報告概述尾段：2025計劃申報投資為零嘅項目（跨年/內部研究/取消）
            _prose_slide(prs, f"{S1}  |  2025年度計劃申報投資支出為零的項目",
                         [("", x) for x in zit[1:]], headline=zit[0], sec=0)
    ahl, ab = _adj_summary(ent_up, adj, ov, sdf)   # slide 15：全闊表 + 敘述另起版
    adj2 = adjustment_by_sub(sdf, BUCKET_ORDER[0])
    _c14 = f"{S1}  |  2025年度投資計劃報告投資金額的潛在調整事項匯總"
    _adjb = _adj_detail_bullets(ent_up, adj, sdf, narr, llm)
    _tname14 = f"{ent_up} 2025年度報告投資金額的潛在調整事項"
    # 表版 + 詳述版【一齊數頁】：scan p15 = (1/3)、p16 = (2/3)、p17 = (3/3)
    _n14 = 1 + len(_prose_pages(prs, _adjb, ahl, _tname14))
    _sfx14 = f"（1/{_n14}）" if _n14 > 1 else ""
    if not adj2.empty:
        _sl, _W, _H, _top = _page(prs, 0, _c14, ahl + _sfx14)
        _tw = _W - 2 * MARGIN
        _t2 = caption_bar(_sl, MARGIN, _top, _tw,
                            f"{ent_up} 2025年度投資計劃報告投資金額潛在調整")
        _draw_adj_table(_sl, MARGIN, _t2, _tw, adj2.fillna(""))
        put(_sl, MARGIN, CONTENT_BOTTOM - 0.26, _tw, 0.3,
              "註：金額單位為萬澳門元；括號表示調減。", size=SZ_NOTE - 1, italic=True, color=GREY)
        source_note(_sl, _W)
    else:
        render_overview_page(prs, _c14, ahl + _sfx14, adj.fillna(""), ab, sec=0,
                             table_name=f"{ent_up} 2025年度投資計劃報告投資金額的潛在調整事項匯總",
                             note="註：金額單位為萬澳門元；括號表示調減。")
    _prose_2col(prs, _c14, _adjb, 6, subtitle=_tname14, sec=0,
                headline=ahl, pg0=1, pgn=_n14)   # slide 16-17 詳述（LLM 優先）

    # ② 過往年度投資計劃在2025年繼續執行的審查跟進（報告 slide 19-26）
    S2 = "過往年度投資計劃在2025年繼續執行的審查跟進"
    divider(prs, S2, "2", [      # 子項用返報告字眼（scan slide 18）
        ("2.1  2024年度投資計劃期後投資金額概覽", ""),
        ("2.2  2024年度投資計劃期後報告投資金額的潛在調整事項匯總", ""),
        ("2.3  2023年度投資計劃期後投資金額概覽", ""),
        ("2.4  2023年度投資計劃期後報告投資金額的潛在調整事項匯總", ""),
        ("2.5  截至2025年末投資金額概覽", ""),
    ])
    for bk in ["2024年度計劃期後投資", "2023年度計劃期後投資"]:
        ov = overview_by_bucket(sdf, bk, plan, cat)
        if not ov.empty:
            render_generic(prs, f"{ent_up} {bk}金額概覽", ov.fillna(""), sec=1,
                           crumb=f"{S2}  |  {bk}金額概覽",
                           headline=_bucket_headline(ent_up, bk, ov),
                           note="註：金額單位為萬澳門元；括號表示調減。",
                           llm=llm, tbl_id=tbl_key("期後概覽", bk))
            render_bucket_adjustment(prs, ent_up, bk, sdf, ov, narr, llm)   # 2.2 / 2.4
    render_cumulative(prs, ent_up, df, plan, cat)      # 2.5 截至2025年末投資金額概覽（scan slide 26）

    # ③ 本年度審查工作的主要發現（報告 slide 28-40）
    S3 = "本年度審查工作的主要發現"
    divider(prs, S3, "3")
    fs = finding_summary(sdf)
    if not fs.empty:
        render_generic(prs, f"{ent_up} 本年度審查工作的主要發現摘要", fs.fillna(""), sec=2,
                       crumb=f"{S3}  |  主要發現摘要",
                       headline=(f"本次審查工作就{ent_up}報告的投資金額識別出{len(fs)}類潛在調整事項，"
                                 f"合計潛在調減約{abs(pd.to_numeric(fs['調整額合計'], errors='coerce').sum()):,.0f}"
                                 f"萬澳門元，摘要如下；逐項說明見後頁。"),
                       note="註：金額單位為萬澳門元；括號表示調減。",
                       llm=llm, tbl_id=tbl_key("發現摘要"))
    if narr:      # 逐調整類型 × 項目：金額(feed) + 事項描述(LLM ground 表2＋清單) / 清單抄字
        b2 = {}
        try:            # 表2＝審查底稿，清單冇料時頂住（加密檔，開唔到就靜靜跳過）
            b2 = load_biao2_struct(av[av.index("--biao2") + 1] if "--biao2" in av else "data/表2",
                               entity, log=lambda *a: None)
        except Exception:
            pass
        render_findings(prs, ent_up, sdf, narr, llm=llm, b2=b2)

    # ④ 其他信息（報告 slide 42-63）
    S4 = "其他信息"
    divider(prs, S4, "4")
    render_generic(prs, f"{ent_up} 2025年發生的投資金額匯總",
                   summary_amount(sdf).fillna(""), sec=3,
                   crumb=f"{S4}  |  2025年發生的投資金額匯總",
                   headline=(f"下表匯總{ent_up} 2025年度投資計劃及過往年度計劃期後投資"
                             f"於2025年發生的投資金額（報告投資金額及潛在調整後投資金額）。"),
                   note="註：金額單位為萬澳門元。", llm=llm, tbl_id=tbl_key("金額匯總"))
    for bk in BUCKET_ORDER:
        fa = facility_activity(sdf, bk)
        if not fa.empty:
            render_generic(prs, f"{ent_up} {bk}區分設施建設/活動舉辦的投資金額", fa.fillna(""), sec=3,
                           crumb=f"{S4}  |  2025年發生的投資金額區分設施建設/活動舉辦",
                           headline=(f"下表按範疇列示{ent_up} {bk}於2025年發生的投資金額，"
                                     f"區分設施建設（資本性支出）及活動舉辦（營運性支出）。"),
                           note="註：金額為潛在調整後金額，單位為萬澳門元。",
                           llm=llm, tbl_id=tbl_key("設施活動", bk))
    for yr in (25, 24, 23):     # 單個項目審查匯總（slide 46-63）
        tab, _ = build_year(df, yr, plan.get(yr) if plan else None)
        if tab is not None and not tab.empty:
            render_sheet(prs, f"報告年{yr}", tab.fillna(""), list(tab.columns),
                           ent_up=ent_up, sec=3, crumb=f"{S4}  |  單個項目審查結果匯總")

    render_visit_summary(prs, ent_up, sdf)      # 報告 slide 71 走訪情況匯總（樣本標準+樣本量）

    # ⑥ 附件（slide 93-105）
    divider(prs, "附件", "6")
    if narr:
        render_site_visits(prs, ent_up, sdf, narr)
    render_artwork(prs, ent_up, av[av.index("--biao2") + 1] if "--biao2" in av else "data/表2",
                   entity)                      # 報告 slide 101 藝術品展出情況清單

    # 目錄（報告 slide 7）：起完全部版先知頁碼 → 砌好插去第 2 版，再全份重編頁碼
    _W, _H = size_of(prs)
    toc = _collect_toc(prs, _W, _H)
    if toc:
        n0 = len(prs.slides._sldIdLst)
        render_toc(prs, ent_up, toc)
        for k in range(len(prs.slides._sldIdLst) - n0):        # 目錄可能多過一版
            _move_slide(prs, n0 + k, 1 + k)
        _renumber_footers(prs, _W, _H)
        print(f"    目錄：{sum(1 for e in toc if not e[2])} 章 / {sum(1 for e in toc if e[2])} 子項")
    wire_nav(prs, _sec_slides(prs, _W, _H), home=1 if toc else 0)   # ◀⌂▶ + 頁籤內部跳頁

    if tmpl:      # template mode：重編 slide 高號，徹底避開 template 殘留 orphan part 撞名 corruption
        _renumber_slides(prs)

    apply_theme_fonts(prs)      # deck theme 字體 → KPMG（唔明寫嘅地方唔會跌返 Calibri）
    out = Path(f"{entity}_report_llm.pptx")
    try:
        prs.save(out)
    except PermissionError:      # 舊檔喺 PowerPoint 開住鎖住 → 改名唔 crash
        import time
        out = Path(f"{entity}_report_llm_{time.strftime('%H%M%S')}.pptx")
        prs.save(out)
        print(f"⚠ 原檔開住(鎖住)，改存 → {out.name}（開之前記得閂舊 pptx）")
    print(f"✓ {out.resolve()}  共 {len(list(prs.slides))} 頁（概述 + 主要發現 + 金額匯總 + 設施 + 單項審查）")
    if "--dump" in sys.argv:      # 要 cross-check 先加 --dump（慳空間；ok 咗嘅唔使 dump）
        dump = _dump_pptx_text(prs, entity, with_tables="--dump-tables" in sys.argv)
        print(f"✓ text dump → {dump.name}（逐版文字，cross-check 用）")


if __name__ == "__main__":
    main()
