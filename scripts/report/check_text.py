#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
check_text.py — 文件用字檢查：掃 folder 入面全部 pptx/potx/docx，逐【頁】揾錯別字、簡體字、
繁簡轉換錯誤（后/後、里/裡、干/乾…）同其他用語問題。

點用（乜 flag 都唔使，掉檔入 folder 就得）：
    python scripts\\report\\check_text.py                 # 掃 file_check\\ 全部檔
    python scripts\\report\\check_text.py --dir 其他資料夾
    python scripts\\report\\check_text.py --no-llm        # 只跑機械檢查（唔使 API、即刻出）
    python scripts\\report\\check_text.py --workers 6 --model 5.5
    python scripts\\report\\check_text.py --fresh         # 唔用 cache，全部重新問

出：
    file_check\\_檢查報告\\{檔名}.md   ← 逐頁列問題（原文 → 建議 → 理由）
    console 一個總結表（邊個檔幾多問題、邊類最多）

兩層檢查：
  ① 機械層（唔使 API、100% 穩定）：簡體字表 + 高信心繁簡詞組（「然后」「這里」…）
  ② LLM 層（gpt-5.5）：錯別字、語境相關嘅繁簡（皇后 vs 之後）、標點、數字格式、用語不一致

Cache：逐頁按【文字 hash】記住，改過嘅頁先會再問 LLM → 同一份檔翻查好平。
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

_HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(_HERE))
for _p in (_HERE.parent.parent, _HERE.parent.parent / "src"):      # repo root / src
    sys.path.insert(0, str(_p))

try:
    from tqdm import tqdm
except ImportError:                                    # tqdm 冇裝都跑得（退返逐行 log）
    def tqdm(it=None, **kw):
        return it if it is not None else _Dummy()

    class _Dummy:
        def update(self, *a): pass
        def refresh(self): pass
        def close(self): pass
        @staticmethod
        def write(s): print(s)

class _Ticker:
    """令 tqdm 每秒 refresh 一次 —— LLM 一 call 要幾秒，唔 tick 嘅話 elapsed/rate 好似死咗機。
    update() 同 refresh() 共用一把鎖（tqdm 本身唔係 thread-safe）。"""

    def __init__(self, bar, every=1.0):
        self.bar, self.every = bar, every
        self._lock = threading.Lock()
        self._stop = threading.Event()
        self._t = threading.Thread(target=self._run, daemon=True)
        self._t.start()

    def _run(self):
        while not self._stop.wait(self.every):
            with self._lock:
                try:
                    self.bar.refresh()
                except Exception:
                    return

    def update(self, n=1):
        with self._lock:
            self.bar.update(n)

    def close(self):
        self._stop.set()
        self._t.join(timeout=self.every + 1)
        with self._lock:
            self.bar.close()


DEFAULT_DIR = "file_check"
OUT_SUB = "_檢查報告"
CACHE = ".check_cache.json"

# ── ① 機械層：簡體字（表由 zhconv 生成，唔係手寫 → 唔會誤報）────────────────
#   做法：0x4E00-0x9FFF 逐個字問 zhconv「轉繁體會唔會變」，變咗就係簡體字。
#   ⚠ 已剔走【繁體本身都會用】嘅字（后里干台余范征游云几谷划卷姜岳表采御…）——
#     呢啲要睇語境（皇后 ✓ / 然後 ✗），交畀下面 BAD_PHRASE 同 LLM 判斷。
_S = (
    "专业丛东丝丢两严丧个为丽乔习乡于亏亘亚产亩亲亵亸亿仅仆从仑仓仪们价众优伙会伛伞伟传伡伣伤伥伦伧伪伫体佣佥侠侣侥侦侧侨"
    "侩侪侬侭俣俦俨俩俪俫俭债倾偬偻偾偿傤傥傧储傩儿兑兖党兰关兴兹养兽冁内冈册写军农冯决况冻净凄凉减凑凛凤凫凭凯凶击凿刍刘"
    "则刚创删刬刭刹刽刾刿剀剂剐剑剥剧劝办务劢动励劲劳势勋勚匀匦匮区华协单卖占卢卤卧卫却卺厂厅历厉压厌厍厐厕厘厢厣厦厨厩厮"
    "县叁参叆叇双变叙叠叶号叹叽吁吃吓吕吗吣吨听启吴呐呒呓呕呖呗员呙呛呜咏咙咛咝咤咨响哑哒哓哔哕哗哙哜哝哟唇唛唝唠唡唢唤啧"
    "啬啭啮啯啰啴啸喂喷喽喾嗫嗳嘘嘤嘱噜嚣团园囱围囵国图圆圣圹场坏块坚坛坜坝坞坟坠垄垅垆垒垦垩垫垭垯垱垲垴埘埙埚堑堕塆墙壮"
    "声壳壶壸处备够头夹夺奁奂奋奖奥妆妇妈妩妪妫姗姹娄娅娆娇娈娱娲娴婳婴婵婶媪媭嫒嫔嫱嬷孙学孪宁宝实宠审宪宫宽宾寝对寻导寿"
    "将尔尘尝尧尴尸尽层屃屉届属屡屦屿岁岂岖岗岘岚岛岩岭岽岿峃峄峡峣峤峥峦峰崂崃崄崭嵘嵚嵝巅巩巯币帅师帏帐帜带帧帮帱帻帼幂"
    "并广庄庆床庐庑库应庙庞废庼廪开异弃弑张弥弪弯弹强归当录彟彦彨彻径徕忆忏忧忾怀态怂怃怄怅怆怜总怼怿恋恒恳恶恸恹恺恻恼恽"
    "悦悫悬悭悮悯惊惧惨惩惫惬惭惮惯愠愤愦愿慑慭懑懒懔戆戋戏戗战戬戯户扑托执扩扪扫扬扰抚抛抟抠抡抢护报担拟拢拣拥拦拧拨择挂"
    "挚挛挜挝挞挟挠挡挢挣挤挥挦捝捞损捡换捣据掳掴掷掸掺掼揽揾揿搀搁搂搄搅携摄摅摆摇摈摊撄撑撵撷撸撺擜擞攒敌敚敛敩数斋斓斩"
    "断无旧时旷旸昙昵昼昽显晋晒晓晔晕晖暂暅暧术朴机杀杂权杠条来杨杩杰极构枞枢枣枥枧枨枪枫枭柜柠柽栀栅标栈栉栊栋栌栎栏树栖"
    "栗样栾桠桡桢档桤桥桦桧桨桩桪梦梼梾梿检棁棂椁椝椟椠椢椤椫椭椮楼榄榅榇榈榉榝槚槛槟槠横樯樱橥橱橹橼檩欢欤欧歼殁殇残殒殓"
    "殚殡殴毁毂毕毙毡毵毶氇气氢氩氲汇汉汤汹沟没沣沤沥沦沧沨沩沪泞泪泶泷泸泺泻泼泽泾洁洒洼浃浅浆浇浈浉浊测浍济浏浐浑浒浓浔"
    "浕涂涌涚涛涝涞涟涠涡涢涣涤润涧涨涩渊渌渍渎渐渑渔渖渗温湾湿溁溃溅溆溇滗滚滞滟滠满滢滤滥滦滨滩滪漤潆潇潋潍潜潴澛澜濑濒"
    "灏灭灯灵灾灿炀炉炖炜炝点炼炽烁烂烃烛烟烦烧烨烩烫烬热焕焖焘煴熏爱爷牍牦牵牺犊状犷犸犹狈狝狞独狭狮狯狰狱狲猃猎猕猡猪猫"
    "猬献獭玑玙玚玛玮环现玱玺珐珑珰珲琎琏琐琼瑶瑷瑸璎瓒瓮瓯电画畅畴疖疗疟疠疡疬疭疮疯疱疴痈痉痒痖痨痪痫痳痴瘅瘆瘗瘘瘪瘫瘾"
    "瘿癞癣癫皂皑皱皲盏盐监盖盗盘眍眦眬睁睐睑瞆瞒瞩矫矶矾矿砀码砖砗砚砜砺砻砾础硁硕硖硗硙硚确硵硷碍碛碜碱礼祃祎祢祯祷祸禀"
    "禄禅离秃秆秘积称秽秾稆税稣稳穑穞穷窃窍窎窑窜窝窥窦窭竖竞笃笋笔笕笺笼笾筚筛筜筝筹筼签筿简箓箦箧箨箩箪箫篑篓篮篯篱簖籁"
    "籴类籼粜粝粤粪粮粽糁糇糍紧絷緼縆纟纠纡红纣纥约级纨纩纪纫纬纭纮纯纰纱纲纳纴纵纶纷纸纹纺纻纼纽纾线绀绁绂练组绅细织终绉"
    "绊绋绌绍绎经绐绑绒结绔绕绖绗绘给绚绛络绝绞统绠绡绢绣绤绥绦继绨绩绪绫绬续绮绯绰绱绲绳维绵绶绷绸绹绺绻综绽绾绿缀缁缂缃"
    "缄缅缆缇缈缉缊缋缌缍缎缏缐缑缒缓缔缕编缗缘缙缚缛缜缝缞缟缠缡缢缣缤缥缦缧缨缩缪缫缬缭缮缯缰缱缲缳缴缵罂网罗罚罢罴羁羟"
    "羡群翘翙翚耢耧耸耻聂聋职聍联聩聪肃肠肤肮肴肾肿胀胁胆胧胨胪胫胶脉脍脐脑脓脔脚脱脶脸腊腌腘腭腻腼腽腾膑膻臜舆舣舰舱舻艰"
    "艳艺节芈芗芜芦苁苇苈苋苌苍苎苏苧茎茏茑茔茕茧荆荐荙荚荛荜荝荞荟荠荡荣荤荥荦荧荨荩荪荫荬荭荮药莅莱莲莳莴莶获莸莹莺莼萚"
    "萝萤营萦萧萨葱蒀蒇蒉蒋蒌蒏蓝蓟蓠蓣蓥蓦蔂蔷蔹蔺蔼蕰蕲蕴薮藓藴蘖虏虑虚虫虬虮虱虽虾虿蚀蚁蚂蚃蚕蚝蚬蛊蛎蛏蛮蛰蛱蛲蛳蛴蜕"
    "蜗蜡蝇蝈蝉蝎蝼蝾螀螨蟏衅衔补衬衮袄袅袆袜袭袯装裆裈裢裣裤裥褛褴襕见观觃规觅视觇览觉觊觋觌觍觎觏觐觑觞触觯訚詟誉誊讠计"
    "订讣认讥讦讧讨让讪讫讬训议讯记讱讲讳讴讵讶讷许讹论讻讼讽设访诀证诂诃评诅识诇诈诉诊诋诌词诎诏诐译诒诓诔试诖诗诘诙诚诛"
    "诜话诞诟诠诡询诣诤该详诧诨诩诪诫诬语诮误诰诱诲诳说诵诶请诸诹诺读诼诽课诿谀谁谂调谄谅谆谇谈谉谊谋谌谍谎谏谐谑谒谓谔谕"
    "谖谗谘谙谚谛谜谝谞谟谠谡谢谣谤谥谦谧谨谩谪谫谬谭谮谯谰谱谲谳谴谵谶豮贝贞负贠贡财责贤败账货质贩贪贫贬购贮贯贰贱贲贳贴"
    "贵贶贷贸费贺贻贼贽贾贿赀赁赂赃资赅赆赇赈赉赊赋赌赍赎赏赐赑赒赓赔赕赖赗赘赙赚赛赜赝赞赟赠赡赢赣赪赵赶趋趱趸跃跄跖跞践"
    "跶跷跸跹跻踊踌踪踬踯蹑蹒蹰蹿躏躜躯輼车轧轨轩轪轫转轭轮软轰轱轲轳轴轵轶轷轸轹轺轻轼载轾轿辀辁辂较辄辅辆辇辈辉辊辋辌辍"
    "辎辏辐辑辒输辔辕辖辗辘辙辚辞辟辩辫边辽达迁过迈运还这进远违连迟迩迳迹选逊递逦逻遗遥邓邝邬邮邹邺邻郏郐郑郓郦郧郸酂酝酦"
    "酱酽酾酿醖释鉴銮錾钅钆钇针钉钊钋钌钍钎钏钐钑钒钓钔钕钖钗钘钙钚钛钜钝钞钠钡钢钣钤钥钦钧钨钩钪钫钬钭钮钯钰钱钲钳钴钵钶"
    "钷钸钹钺钻钼钽钾钿铀铁铂铃铄铅铆铇铈铉铊铋铌铍铎铏铐铑铒铓铔铕铖铗铘铙铚铛铜铝铞铟铠铡铢铣铤铥铦铧铨铩铪铫铬铭铮铯铰"
    "铱铲铳铴铵银铷铸铹铺铻铼铽链铿销锁锂锃锄锅锆锇锈锉锊锋锌锍锎锏锐锑锒锓锔锕锖锗锘错锚锛锜锝锞锟锠锡锢锣锤锥锦锧锨锩锪"
    "锫锬锭键锯锰锱锲锳锴锵锶锷锸锹锺锻锼锽锾锿镀镁镂镃镄镅镆镇镈镉镊镋镌镍镎镏镐镑镒镓镔镕镖镗镘镙镚镛镜镝镞镟镠镡镢镣镤"
    "镥镦镧镨镩镪镫镬镭镮镯镰镱镲镳镴镵镶闩闪闫闬闭问闯闰闱闲闳间闵闶闷闸闹闺闻闼闽闾闿阀阁阂阃阄阅阆阇阈阉阊阋阌阍阎阏阐"
    "阑阒阓阔阕阖阗阘阙阚阛队阳阴阵阶际陆陇陈陉陕陦陧陨险随隐隶隽难雇雏雠雳雾霁霉霡霭靓靔静靥鞑鞒鞯鞲韦韧韨韩韪韫韬韵页顶"
    "顷顸项顺顼顽顾顿颀颁颂颃预颅领颇颈颉颊颋颌颍颎颏颐频颒颓颔颕颖颗题颙颚颛颜额颞颟颠颡颢颣颤颥颦颧风飏飐飑飒飓飔飕飖飗"
    "飘飙飚飞飨餍饣饤饥饦饧饨饩饪饫饬饭饮饯饰饱饲饳饴饵饶饷饸饹饺饻饼饽饾饿馀馁馂馃馄馅馆馇馈馉馊馋馌馍馎馏馐馑馒馓馔馕马"
    "驭驮驯驰驱驲驳驴驵驶驷驸驹驺驻驼驽驾驿骀骁骂骃骄骅骆骇骈骉骊骋验骍骎骏骐骑骒骓骔骕骖骗骘骙骚骛骜骝骞骟骠骡骢骣骤骥骦"
    "骧髅髋髌鬓鬶魇魉鱼鱽鱾鱿鲀鲁鲂鲃鲄鲅鲆鲇鲈鲉鲊鲋鲌鲍鲎鲏鲐鲑鲒鲓鲔鲕鲖鲗鲘鲙鲚鲛鲜鲝鲞鲟鲠鲡鲢鲣鲤鲥鲦鲧鲨鲩鲪鲫鲬鲭"
    "鲮鲯鲰鲱鲲鲳鲴鲵鲶鲷鲸鲹鲺鲻鲼鲽鲾鲿鳀鳁鳂鳃鳄鳅鳆鳇鳈鳉鳊鳋鳌鳍鳎鳏鳐鳑鳒鳓鳔鳕鳖鳗鳘鳙鳚鳛鳜鳝鳞鳟鳠鳡鳢鳣鳤鸟鸠鸡"
    "鸢鸣鸤鸥鸦鸧鸨鸩鸪鸫鸬鸭鸮鸯鸰鸱鸲鸳鸴鸵鸶鸷鸸鸹鸺鸻鸼鸽鸾鸿鹀鹁鹂鹃鹄鹅鹆鹇鹈鹉鹊鹋鹌鹍鹎鹏鹐鹑鹒鹓鹔鹕鹖鹗鹘鹙鹚鹛"
    "鹜鹝鹞鹟鹠鹡鹢鹣鹤鹥鹦鹧鹨鹩鹪鹫鹬鹭鹮鹯鹰鹱鹲鹳鹴鹾麦麸麹麺麽黄黉黡黩黪黾鼋鼌鼍鼗鼹齐齑齿龀龁龂龃龄龅龆龇龈龉龊龋龌"
    "龙龚龛龟鿎鿏鿒鿔鿭"
)
_T = (
    "專業叢東絲丟兩嚴喪個爲麗喬習鄉於虧亙亞產畝親褻嚲億僅僕從侖倉儀們價衆優夥會傴傘偉傳俥俔傷倀倫傖僞佇體傭僉俠侶僥偵側僑"
    "儈儕儂儘俁儔儼倆儷倈儉債傾傯僂僨償儎儻儐儲儺兒兌兗黨蘭關興茲養獸囅內岡冊寫軍農馮決況凍淨悽涼減湊凜鳳鳧憑凱兇擊鑿芻劉"
    "則剛創刪剗剄剎劊㓨劌剴劑剮劍剝劇勸辦務勱動勵勁勞勢勳勩勻匭匱區華協單賣佔盧滷臥衛卻巹廠廳歷厲壓厭厙龎廁釐廂厴廈廚廄廝"
    "縣叄參靉靆雙變敘疊葉號嘆嘰籲喫嚇呂嗎唚噸聽啓吳吶嘸囈嘔嚦唄員咼嗆嗚詠嚨嚀噝吒諮響啞噠嘵嗶噦譁噲嚌噥喲脣嘜嗊嘮啢嗩喚嘖"
    "嗇囀齧嘓囉嘽嘯餵噴嘍嚳囁噯噓嚶囑嚕囂團園囪圍圇國圖圓聖壙場壞塊堅壇壢壩塢墳墜壟壠壚壘墾堊墊埡墶壋塏堖塒壎堝塹墮壪牆壯"
    "聲殼壺壼處備夠頭夾奪奩奐奮獎奧妝婦媽嫵嫗嬀姍奼婁婭嬈嬌孌娛媧嫺嫿嬰嬋嬸媼嬃嬡嬪嬙嬤孫學孿寧寶實寵審憲宮寬賓寢對尋導壽"
    "將爾塵嘗堯尷屍盡層屓屜屆屬屢屨嶼歲豈嶇崗峴嵐島巖嶺崬巋嶨嶧峽嶢嶠崢巒峯嶗崍嶮嶄嶸嶔嶁巔鞏巰幣帥師幃帳幟帶幀幫幬幘幗冪"
    "並廣莊慶牀廬廡庫應廟龐廢廎廩開異棄弒張彌弳彎彈強歸當錄彠彥彲徹徑徠憶懺憂愾懷態慫憮慪悵愴憐總懟懌戀恆懇惡慟懨愷惻惱惲"
    "悅愨懸慳悞憫驚懼慘懲憊愜慚憚慣慍憤憒願懾憖懣懶懍戇戔戲戧戰戩戱戶撲託執擴捫掃揚擾撫拋摶摳掄搶護報擔擬攏揀擁攔擰撥擇掛"
    "摯攣掗撾撻挾撓擋撟掙擠揮撏挩撈損撿換搗據擄摑擲撣摻摜攬搵撳攙擱摟揯攪攜攝攄擺搖擯攤攖撐攆擷擼攛㩵擻攢敵敓斂斆數齋斕斬"
    "斷無舊時曠暘曇暱晝曨顯晉曬曉曄暈暉暫𣈶曖術樸機殺雜權槓條來楊榪傑極構樅樞棗櫪梘棖槍楓梟櫃檸檉梔柵標棧櫛櫳棟櫨櫟欄樹棲"
    "慄樣欒椏橈楨檔榿橋樺檜槳樁樳夢檮棶槤檢梲欞槨槼櫝槧槶欏樿橢槮樓欖榲櫬櫚櫸樧檟檻檳櫧橫檣櫻櫫櫥櫓櫞檁歡歟歐殲歿殤殘殞殮"
    "殫殯毆毀轂畢斃氈毿𣯶氌氣氫氬氳匯漢湯洶溝沒灃漚瀝淪滄渢潙滬濘淚澩瀧瀘濼瀉潑澤涇潔灑窪浹淺漿澆湞溮濁測澮濟瀏滻渾滸濃潯"
    "濜塗湧涗濤澇淶漣潿渦溳渙滌潤澗漲澀淵淥漬瀆漸澠漁瀋滲溫灣溼濚潰濺漵漊潷滾滯灩灄滿瀅濾濫灤濱灘澦灠瀠瀟瀲濰潛瀦瀂瀾瀨瀕"
    "灝滅燈靈災燦煬爐燉煒熗點煉熾爍爛烴燭煙煩燒燁燴燙燼熱煥燜燾熅燻愛爺牘犛牽犧犢狀獷獁猶狽獮獰獨狹獅獪猙獄猻獫獵獼玀豬貓"
    "蝟獻獺璣璵瑒瑪瑋環現瑲璽琺瓏璫琿璡璉瑣瓊瑤璦璸瓔瓚甕甌電畫暢疇癤療瘧癘瘍癧瘲瘡瘋皰痾癰痙癢瘂癆瘓癇痲癡癉瘮瘞瘻癟癱癮"
    "癭癩癬癲皁皚皺皸盞鹽監蓋盜盤瞘眥矓睜睞瞼瞶瞞矚矯磯礬礦碭碼磚硨硯碸礪礱礫礎硜碩硤磽磑礄確磠礆礙磧磣鹼禮禡禕禰禎禱禍稟"
    "祿禪離禿稈祕積稱穢穠穭稅穌穩穡穭窮竊竅窵窯竄窩窺竇窶豎競篤筍筆筧箋籠籩篳篩簹箏籌篔籤篠簡籙簀篋籜籮簞簫簣簍籃籛籬籪籟"
    "糴類秈糶糲粵糞糧糉糝餱餈緊縶縕緪糹糾紆紅紂紇約級紈纊紀紉緯紜紘純紕紗綱納紝縱綸紛紙紋紡紵紖紐紓線紺紲紱練組紳細織終縐"
    "絆紼絀紹繹經紿綁絨結絝繞絰絎繪給絢絳絡絕絞統綆綃絹繡綌綏絛繼綈績緒綾緓續綺緋綽鞝緄繩維綿綬繃綢綯綹綣綜綻綰綠綴緇緙緗"
    "緘緬纜緹緲緝縕繢緦綞緞緶線緱縋緩締縷編緡緣縉縛縟縝縫縗縞纏縭縊縑繽縹縵縲纓縮繆繅纈繚繕繒繮繾繰繯繳纘罌網羅罰罷羆羈羥"
    "羨羣翹翽翬耮耬聳恥聶聾職聹聯聵聰肅腸膚骯餚腎腫脹脅膽朧腖臚脛膠脈膾臍腦膿臠腳脫腡臉臘醃膕齶膩靦膃騰臏羶臢輿艤艦艙艫艱"
    "豔藝節羋薌蕪蘆蓯葦藶莧萇蒼薴蘇薴莖蘢蔦塋煢繭荊薦薘莢蕘蓽萴蕎薈薺蕩榮葷滎犖熒蕁藎蓀蔭蕒葒葤藥蒞萊蓮蒔萵薟獲蕕瑩鶯蓴蘀"
    "蘿螢營縈蕭薩蔥蒕蕆蕢蔣蔞醟藍薊蘺蕷鎣驀虆薔蘞藺藹薀蘄蘊藪蘚蘊櫱虜慮虛蟲虯蟣蝨雖蝦蠆蝕蟻螞蠁蠶蠔蜆蠱蠣蟶蠻蟄蛺蟯螄蠐蛻"
    "蝸蠟蠅蟈蟬蠍螻蠑螿蟎蠨釁銜補襯袞襖嫋褘襪襲襏裝襠褌褳襝褲襉褸襤襴見觀覎規覓視覘覽覺覬覡覿覥覦覯覲覷觴觸觶誾讋譽謄訁計"
    "訂訃認譏訐訌討讓訕訖託訓議訊記訒講諱謳詎訝訥許訛論訩訟諷設訪訣證詁訶評詛識詗詐訴診詆謅詞詘詔詖譯詒誆誄試詿詩詰詼誠誅"
    "詵話誕詬詮詭詢詣諍該詳詫諢詡譸誡誣語誚誤誥誘誨誑說誦誒請諸諏諾讀諑誹課諉諛誰諗調諂諒諄誶談讅誼謀諶諜謊諫諧謔謁謂諤諭"
    "諼讒諮諳諺諦謎諞諝謨讜謖謝謠謗諡謙謐謹謾謫譾謬譚譖譙讕譜譎讞譴譫讖豶貝貞負貟貢財責賢敗賬貨質販貪貧貶購貯貫貳賤賁貰貼"
    "貴貺貸貿費賀貽賊贄賈賄貲賃賂贓資賅贐賕賑賚賒賦賭齎贖賞賜贔賙賡賠賧賴賵贅賻賺賽賾贗贊贇贈贍贏贛赬趙趕趨趲躉躍蹌蹠躒踐"
    "躂蹺蹕躚躋踴躊蹤躓躑躡蹣躕躥躪躦軀轀車軋軌軒軑軔轉軛輪軟轟軲軻轤軸軹軼軤軫轢軺輕軾載輊轎輈輇輅較輒輔輛輦輩輝輥輞輬輟"
    "輜輳輻輯轀輸轡轅轄輾轆轍轔辭闢辯辮邊遼達遷過邁運還這進遠違連遲邇逕跡選遜遞邐邏遺遙鄧鄺鄔郵鄒鄴鄰郟鄶鄭鄆酈鄖鄲酇醞醱"
    "醬釅釃釀醞釋鑑鑾鏨釒釓釔針釘釗釙釕釷釺釧釤鈒釩釣鍆釹鍚釵鈃鈣鈈鈦鉅鈍鈔鈉鋇鋼鈑鈐鑰欽鈞鎢鉤鈧鈁鈥鈄鈕鈀鈺錢鉦鉗鈷鉢鈳"
    "鉕鈽鈸鉞鑽鉬鉭鉀鈿鈾鐵鉑鈴鑠鉛鉚鉋鈰鉉鉈鉍鈮鈹鐸鉶銬銠鉺鋩錏銪鋮鋏鋣鐃銍鐺銅鋁銱銦鎧鍘銖銑鋌銩銛鏵銓鎩鉿銚鉻銘錚銫鉸"
    "銥鏟銃鐋銨銀銣鑄鐒鋪鋙錸鋱鏈鏗銷鎖鋰鋥鋤鍋鋯鋨鏽銼鋝鋒鋅鋶鐦鐧銳銻鋃鋟鋦錒錆鍺鍩錯錨錛錡鍀錁錕錩錫錮鑼錘錐錦鑕鍁錈鍃"
    "錇錟錠鍵鋸錳錙鍥鍈鍇鏘鍶鍔鍤鍬鍾鍛鎪鍠鍰鎄鍍鎂鏤鎡鐨鎇鏌鎮鎛鎘鑷钂鐫鎳鎿鎦鎬鎊鎰鎵鑌鎔鏢鏜鏝鏍鏰鏞鏡鏑鏃鏇鏐鐔钁鐐鏷"
    "鑥鐓鑭鐠鑹鏹鐙鑊鐳鐶鐲鐮鐿鑔鑣鑞鑱鑲閂閃閆閈閉問闖閏闈閒閎間閔閌悶閘鬧閨聞闥閩閭闓閥閣閡閫鬮閱閬闍閾閹閶鬩閿閽閻閼闡"
    "闌闃闠闊闋闔闐闒闕闞闤隊陽陰陣階際陸隴陳陘陝隯隉隕險隨隱隸雋難僱雛讎靂霧霽黴霢靄靚靝靜靨韃鞽韉韝韋韌韍韓韙韞韜韻頁頂"
    "頃頇項順頊頑顧頓頎頒頌頏預顱領頗頸頡頰頲頜潁熲頦頤頻頮頹頷頴穎顆題顒顎顓顏額顳顢顛顙顥纇顫顬顰顴風颺颭颮颯颶颸颼颻飀"
    "飄飆飈飛饗饜飠飣飢飥餳飩餼飪飫飭飯飲餞飾飽飼飿飴餌饒餉餄餎餃餏餅餑餖餓餘餒餕餜餛餡館餷饋餶餿饞饁饃餺餾饈饉饅饊饌饢馬"
    "馭馱馴馳驅馹駁驢駔駛駟駙駒騶駐駝駑駕驛駘驍罵駰驕驊駱駭駢驫驪騁驗騂駸駿騏騎騍騅騌驌驂騙騭騤騷騖驁騮騫騸驃騾驄驏驟驥驦"
    "驤髏髖髕鬢鬹魘魎魚魛魢魷魨魯魴䰾魺鮁鮃鮎鱸鮋鮓鮒鮊鮑鱟鮍鮐鮭鮚鮳鮪鮞鮦鰂鮜鱠鱭鮫鮮鮺鯗鱘鯁鱺鰱鰹鯉鰣鰷鯀鯊鯇鮶鯽鯒鯖"
    "鯪鯕鯫鯡鯤鯧鯝鯢鯰鯛鯨鰺鯴鯔鱝鰈鰏鱨鯷鰮鰃鰓鱷鰍鰒鰉鰁鱂鯿鰠鰲鰭鰨鰥鰩鰟鰜鰳鰾鱈鱉鰻鰵鱅䲁鰼鱖鱔鱗鱒鱯鱤鱧鱣䲘鳥鳩雞"
    "鳶鳴鳲鷗鴉鶬鴇鴆鴣鶇鸕鴨鴞鴦鴒鴟鴝鴛鷽鴕鷥鷙鴯鴰鵂鴴鵃鴿鸞鴻鵐鵓鸝鵑鵠鵝鵒鷳鵜鵡鵲鶓鵪鵾鵯鵬鵮鶉鶊鵷鷫鶘鶡鶚鶻鶖鷀鶥"
    "鶩鷊鷂鶲鶹鶺鷁鶼鶴鷖鸚鷓鷚鷯鷦鷲鷸鷺䴉鸇鷹鸌鸏鸛鸘鹺麥麩麴麪麼黃黌黶黷黲黽黿鼂鼉鞀鼴齊齏齒齔齕齗齟齡齙齠齜齦齬齪齲齷"
    "龍龔龕龜䃮䥑鿓鎶鉨"
)

SIMP = set(_S)
S2T = dict(zip(_S, _T))          # 簡 → 繁（出建議用）

# 高信心繁簡／別字詞組（機械層直接標，唔使問 LLM）。key = 錯，value = (啱, 理由)
BAD_PHRASE = {
    "然后": ("然後", "時間先後用「後」，「后」係皇后/太后"),
    "以后": ("以後", "時間先後用「後」"),
    "之后": ("之後", "時間先後用「後」"),
    "最后": ("最後", "時間先後用「後」"),
    "后来": ("後來", "時間先後用「後」"),
    "后續": ("後續", "時間先後用「後」"),
    "后期": ("後期", "時間先後用「後」"),
    "期后": ("期後", "時間先後用「後」"),
    "調整后": ("調整後", "時間先後用「後」"),
    "落后": ("落後", "時間先後用「後」"),
    "背后": ("背後", "方位用「後」"),
    "后面": ("後面", "方位用「後」"),
    "這里": ("這裡", "處所用「裡」"),
    "那里": ("那裡", "處所用「裡」"),
    "哪里": ("哪裡", "處所用「裡」"),
    "里面": ("裡面", "處所用「裡」（「里」係長度單位／鄉里）"),
    "干凈": ("乾淨", "清潔用「乾」"),
    "干燥": ("乾燥", "清潔／無水用「乾」"),
    "干預": ("干預", ""),                       # 正確，放住做對照唔會標
    "松動": ("鬆動", "緊鬆用「鬆」，「松」係松樹"),
    "放松": ("放鬆", "緊鬆用「鬆」"),
    "頭發": ("頭髮", "毛髮用「髮」"),
    "發型": ("髮型", "毛髮用「髮」"),
    "面條": ("麵條", "食物用「麵」"),
    "范圍": ("範圍", "「范」係姓氏，範圍用「範」"),
    "規范": ("規範", "「范」係姓氏"),
    "布置": ("佈置", "安排用「佈」"),
    "公布": ("公佈", "宣告用「佈」"),
    "分布": ("分佈", "散開用「佈」"),
    "帳號": ("賬號", ""),                        # 兩者皆可，唔標
    "並且": ("並且", ""),
    "台灣": ("台灣", ""),
    "沖擊": ("衝擊", "碰撞用「衝」"),
    "沖突": ("衝突", "碰撞用「衝」"),
    "余額": ("餘額", "剩餘用「餘」，「余」係第一人稱"),
    "其余": ("其餘", "剩餘用「餘」"),
    "征收": ("徵收", "徵集用「徵」"),
    "象征": ("象徵", "徵集用「徵」"),
    "系統性": ("系統性", ""),
    "聯系": ("聯繫", "聯絡用「繫」"),
    "關系": ("關係", "「系」→「係」"),
    "沒有": ("沒有", ""),
    "了解": ("了解", ""),
    "制定": ("制定", ""),
    "制造": ("製造", "生產用「製」"),
    "复制": ("複製", "重複／生產"),
    "复核": ("覆核", "審查用「覆」"),
    "回复": ("回覆", "答覆用「覆」"),
    "答复": ("答覆", "答覆用「覆」"),
    "覆蓋": ("覆蓋", ""),
    "重复": ("重複", "重疊用「複」"),
    "游客": ("旅客", "報告一律用「旅客」（如原文係遊覽義則用「遊客」）"),
    "旅游": ("旅遊", "遊覽用「遊」"),
    "上游": ("上游", ""),
    "只": ("只", ""),
    "面積": ("面積", ""),
}
_BAD = {k: v for k, v in BAD_PHRASE.items() if v[1]}          # 有理由先算「錯」

SYS = (
    "你係繁體中文（澳門／香港用字）校對員，專門校對畢馬威（KPMG）投資計劃執行情況審查報告。"
    "任務：喺俾你嘅頁面文字入面揾【真係錯】嘅嘢，逐個報返。要揾嘅類型：\n"
    "1 簡體字 —— 任何簡體字都要標。\n"
    "2 繁簡轉換錯誤（同音／一簡對多繁揀錯字）：后↔後、里↔裡、干↔乾/幹、發↔髮、"
    "松↔鬆、范↔範、余↔餘、征↔徵、系↔係/繫、复↔複/覆、么↔麼、只↔隻、"
    "布↔佈、沖↔衝、台↔臺、制↔製、志↔誌、注↔註、采↔採、划↔劃、"
    "游↔遊、卷↔捲、朴↔樸、面↔麵、板↔闆、尽↔盡/儘。\n"
    "3 錯別字／同音別字（例如「帳目」寫成「賬目」以外嘅真錯字、人名地名寫錯）。\n"
    "4 標點：全形半形混用、中文用咗英文逗號句號、括號唔配對、重覆句號。\n"
    "5 數字／單位格式唔一致（例如同一頁又「萬澳門元」又「萬元」、千分位有時有有時冇）。\n"
    "6 用語前後唔一致（同一個概念頁內用咗兩個名）。\n\n"
    "★ 極重要：寧可漏報，都唔好誤報。以下【唔算錯】，一律唔好報：\n"
    "  · 台灣／臺灣、著／着、裡／裏、為／爲 呢類兩地皆可嘅異體；\n"
    "  · 專有名詞、品牌名、英文、數字本身；\n"
    "  · 你覺得「可以更好」但原文冇錯嘅寫法（唔好做潤色建議）；\n"
    "  · 排版斷行造成嘅斷字。\n"
    "★ 「后」要睇語境：皇后／太后／后冠 正確；時間先後（然後／之後／期後）先算錯。\n"
    "★ 只可以引用頁面入面真係出現過嘅字串做 original，唔可以自己改寫。"
)

USER_TMPL = (
    "以下係《{doc}》第 {page} 頁嘅全部文字（逐段編號）。請校對。\n\n{body}\n\n"
    "回覆【純 JSON】，格式：\n"
    '{{"issues":[{{"seg":<段號int>,"type":"簡體字|繁簡|錯別字|標點|格式|用語","original":"<原文片段>",'
    '"suggested":"<建議>","reason":"<10-30字理由>"}}]}}\n'
    "冇問題就回 {{\"issues\":[]}}。"
)


# ── 抽文字：pptx 逐 slide、docx 逐頁（用分頁符／每 N 段）──────────────────
def _shape_texts(shapes, out):
    for sh in shapes:
        if getattr(sh, "shape_type", None) is not None and getattr(sh, "has_table", False):
            for ri, row in enumerate(sh.table.rows):
                for ci, cell in enumerate(row.cells):
                    t = (cell.text or "").strip()
                    if t:
                        out.append((f"表格 r{ri}c{ci}", t))
            continue
        if getattr(sh, "shape_type", None) == 6 or hasattr(sh, "shapes"):      # GROUP
            try:
                _shape_texts(sh.shapes, out); continue
            except Exception:
                pass
        if getattr(sh, "has_text_frame", False):
            t = (sh.text_frame.text or "").strip()
            if t:
                out.append(("文字框", t))


def pages_of_pptx(path):
    from pptx import Presentation
    prs = Presentation(str(path))
    for i, sl in enumerate(prs.slides, 1):
        segs = []
        _shape_texts(sl.shapes, segs)
        if getattr(sl, "has_notes_slide", False) and sl.notes_slide.notes_text_frame.text.strip():
            segs.append(("備註", sl.notes_slide.notes_text_frame.text.strip()))
        yield i, segs


def pages_of_docx(path, per=40):
    import docx
    d = docx.Document(str(path))
    segs, page = [], 1
    for p in d.paragraphs:
        t = (p.text or "").strip()
        if t:
            segs.append(("段落", t))
        if len(segs) >= per:
            yield page, segs; page += 1; segs = []
    for ti, tb in enumerate(d.tables):
        for ri, row in enumerate(tb.rows):
            for ci, cell in enumerate(row.cells):
                t = (cell.text or "").strip()
                if t:
                    segs.append((f"表{ti} r{ri}c{ci}", t))
    if segs:
        yield page, segs


# .potx = PowerPoint 範本，同 pptx 一樣係 OOXML package，python-pptx 開得（我哋主要用 potx）
READERS = {".pptx": pages_of_pptx, ".potx": pages_of_pptx, ".pptm": pages_of_pptx,
           ".docx": pages_of_docx, ".dotx": pages_of_docx}


# ── ① 機械層 ────────────────────────────────────────────────────────
def mech_check(segs):
    out = []
    for si, (loc, t) in enumerate(segs, 1):
        bad = sorted({c for c in t if c in SIMP})
        if bad:
            out.append({"seg": si, "loc": loc, "type": "簡體字", "original": "".join(bad),
                        "suggested": "".join(S2T.get(c, c) for c in bad),
                        "reason": "、".join(f"{c}→{S2T.get(c, '?')}" for c in bad), "src": "機械"})
        for k, (good, why) in _BAD.items():
            if k in t and k != good:
                out.append({"seg": si, "loc": loc, "type": "繁簡", "original": k,
                            "suggested": good, "reason": why, "src": "機械"})
        if re.search(r"[，。；：？！][,.;:?!]|[,.;:?!][，。；：？！]", t):
            out.append({"seg": si, "loc": loc, "type": "標點", "original": t[:40],
                        "suggested": "統一用全形中文標點", "reason": "全形半形標點混用", "src": "機械"})
        if "。。" in t or "，，" in t:
            out.append({"seg": si, "loc": loc, "type": "標點", "original": "。。／，，",
                        "suggested": "刪重覆", "reason": "重覆標點", "src": "機械"})
    return out


# ── ② LLM 層 ───────────────────────────────────────────────────────
def llm_check(wb, doc, page, segs, model=None):
    body = "\n".join(f"[{i}] ({loc}) {t}" for i, (loc, t) in enumerate(segs, 1))
    if len(body) > 12000:
        body = body[:12000] + "\n…（過長已截）"
    r = wb.chat_json(USER_TMPL.format(doc=doc, page=page, body=body), system=SYS,
                     model=model, reasoning_effort="high")
    out = []
    for it in (r or {}).get("issues", []) or []:
        si = it.get("seg")
        try:
            loc = segs[int(si) - 1][0]
        except Exception:
            loc, si = "?", si
        out.append({"seg": si, "loc": loc, "type": str(it.get("type", "?")),
                    "original": str(it.get("original", ""))[:120],
                    "suggested": str(it.get("suggested", ""))[:120],
                    "reason": str(it.get("reason", ""))[:120], "src": "LLM"})
    return out


def _key(doc, page, segs):
    h = hashlib.sha1("\u0001".join(t for _l, t in segs).encode("utf-8")).hexdigest()[:16]
    return f"{doc}|{page}|{h}"


def main():
    ap = argparse.ArgumentParser(description="掃 folder 入面嘅 pptx/docx，逐頁校對繁體中文用字")
    ap.add_argument("--dir", default=DEFAULT_DIR, help=f"要掃嘅資料夾（預設 {DEFAULT_DIR}）")
    ap.add_argument("--no-llm", action="store_true", help="只跑機械檢查，唔叫 API")
    ap.add_argument("--workers", type=int, default=4,
                    help="LLM 並行數（預設 4；8 條試過俾公司網關 WAF 擋，同 build_report 一致）")
    ap.add_argument("--model", default=None, help="workbench model alias（預設 5.5）")
    ap.add_argument("--fresh", action="store_true", help="唔用 cache")
    a = ap.parse_args()

    root = Path(a.dir)
    if not root.exists():
        root.mkdir(parents=True, exist_ok=True)
        print(f"✓ 開咗 {root.resolve()} —— 掉 pptx／docx 入去再跑一次就得"); return
    files = sorted(p for p in root.rglob("*")
                   if p.suffix.lower() in READERS and not p.name.startswith("~$")
                   and OUT_SUB not in p.parts)
    if not files:
        print(f"✗ {root.resolve()} 入面搵唔到 pptx／docx（掉檔入去再跑）"); return
    outdir = root / OUT_SUB
    outdir.mkdir(exist_ok=True)
    cachef = outdir / CACHE
    cache = {}
    if not a.fresh and cachef.exists():
        try:
            cache = json.loads(cachef.read_text(encoding="utf-8"))
        except Exception:
            cache = {}

    wb = None
    if not a.no_llm:
        try:
            from kpi.lib.workbench import Workbench
            wb = Workbench(model=a.model)
            if not wb.config_masked().get("key_ok"):
                print("  ⚠ 冇 workbench api key → 只跑機械檢查（--no-llm 一樣）"); wb = None
            else:
                wb.ping()      # preflight：403/blocked 即刻知，唔好逐頁白等
        except Exception as e:
            print(f"  ⚠ LLM 連唔到（{str(e)[:120]}）→ 只跑機械檢查"); wb = None

    print(f"── 掃 {root.resolve()}：{len(files)} 個檔"
          f"｜LLM {'開' if wb else '關'}｜cache {len(cache)} 頁")
    # 先【全部檔一次過】抽晒文字 + 跑機械層，再把所有未 cache 嘅頁排成【一條 task list】：
    #   一個 thread pool、一條 tqdm 橫跨成個 folder（同 build_report.py 個做法一致）。
    docs, todo = [], []
    for f in files:
        doc = f.stem
        try:
            pages = [(i, s) for i, s in READERS[f.suffix.lower()](f) if s]
        except Exception as e:
            print(f"  ✗ {f.name} 讀唔到：{str(e)[:140]}"); continue
        found = []
        for pi, segs in pages:
            found += [{**x, "page": pi} for x in mech_check(segs)]
            if wb is None:
                continue
            k = _key(doc, pi, segs)
            if k in cache:
                found += [{**x, "page": pi} for x in cache[k]]
            else:
                todo.append((len(docs), doc, pi, segs, k))
        docs.append({"file": f, "pages": pages, "found": found})
    if todo:
        bar = _Ticker(tqdm(total=len(todo), desc=f"校對 {len(files)} 檔", unit="頁",
                           ncols=80, mininterval=0.5, smoothing=0.1))
        with ThreadPoolExecutor(max_workers=a.workers) as ex:
            futs = {ex.submit(llm_check, wb, doc, pi, segs, a.model): (di, doc, pi, k)
                    for di, doc, pi, segs, k in todo}
            for fu in as_completed(futs):
                di, doc, pi, k = futs[fu]
                try:
                    res = fu.result()
                    cache[k] = res
                    docs[di]["found"] += [{**x, "page": pi} for x in res]
                except Exception as e:
                    tqdm.write(f"    ⚠ {doc} p{pi} LLM 失敗：{str(e)[:100]}")
                bar.update(1)
        bar.close()
        cachef.write_text(json.dumps(cache, ensure_ascii=False), encoding="utf-8")
    summary = []
    for d in docs:
        d["found"].sort(key=lambda x: (x["page"], str(x.get("seg"))))
        write_report(outdir / f"{d['file'].stem}.md", d["file"], d["pages"], d["found"])
        summary.append((d["file"].name, len(d["pages"]), d["found"]))

    print(f"\n{'檔名':<44}{'頁':>4}{'問題':>6}   分類")
    print("-" * 96)
    for name, npg, found in summary:
        from collections import Counter
        c = Counter(x["type"] for x in found)
        print(f"{name[:42]:<44}{npg:>4}{len(found):>6}   "
              + "、".join(f"{k}{v}" for k, v in c.most_common()))
    print(f"\n✓ 逐頁報告寫咗落 {outdir.resolve()}")


def write_report(path, src, pages, found):
    from collections import Counter
    L = [f"# {src.name} 用字檢查", "",
         f"- 來源：`{src.resolve()}`", f"- 頁數：{len(pages)}", f"- 問題：**{len(found)}** 個",
         "- 分類：" + ("、".join(f"{k} {v}" for k, v in
                               Counter(x['type'] for x in found).most_common()) or "冇"), ""]
    if not found:
        L += ["✅ 冇揾到問題。", ""]
    cur = None
    for x in found:
        if x["page"] != cur:
            cur = x["page"]; L += ["", f"## 第 {cur} 頁", "",
                                   "| 位置 | 類型 | 原文 | 建議 | 理由 | 來源 |",
                                   "|---|---|---|---|---|---|"]
        def esc(s):
            return str(s).replace("|", "\\|").replace("\n", " ")
        L.append(f"| {esc(x['loc'])} | {esc(x['type'])} | `{esc(x['original'])}` | "
                 f"`{esc(x['suggested'])}` | {esc(x['reason'])} | {x['src']} |")
    path.write_text("\n".join(L) + "\n", encoding="utf-8")


if __name__ == "__main__":
    main()
