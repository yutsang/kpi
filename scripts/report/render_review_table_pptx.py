#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
render_review_table_pptx.py — 把 build_project_review_table 出嘅 xlsx render 落 pptx
native table（報告 slide 46-63「單個項目審查結果匯總表」，取代 Tableau 截圖）。

策略（user 2026-07-21）：**data 行先**（行/欄/範疇/小計/數字擺正 + 分頁），
formatting 顏色後補；**尺寸跟原報告**（--template 讀 slide 尺寸）。

用法（Windows，kpi-main 底下）：
    pip install python-pptx openpyxl
    python scripts\\report\\render_review_table_pptx.py mgm_項目審查匯總.xlsx ^
        --template "data\\reports\\MGM.2025年年度投資計劃執行情況審查專項工作報告.初稿.pptx"
  出 mgm_項目審查匯總.native.pptx（開嚟同報告 slide 46-63 對）+ console 印每 sheet 分幾頁、欄。
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
try:
    import pandas as pd
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("✗ 需要 pandas + python-pptx → pip install pandas python-pptx openpyxl"); sys.exit(1)

import layout as L                       # 版式引擎（bundler 會 inline）

RATE_COLS = {"投資計劃完成率", "潛在調整後投資計劃完成率"}
TEXT_COLS = {"項目序號", "項目名稱"}
# 3 欄組（畫 group header 用）
G1 = ["項目序號", "項目名稱", "計劃投資金額", "報告投資金額", "投資計劃完成率"]
G3 = ["調整後投資金額", "潛在調整後投資計劃完成率", "設施建設/資本性支出", "活動舉辦/營運性支出"]
GROUP_LABEL = {"G1": "項目基本信息", "G2": "投資金額的潛在調整事項", "G3": "潛在調整後投資金額"}

YEAR_TITLE = {
    "報告年25": "{e} 2025年度投資計劃單個項目審查結果匯總表",
    "報告年24": "{e} 2024年度投資計劃單個項目截至2025年末的審查結果匯總表",
    "報告年23": "{e} 2023年度投資計劃單個項目截至2025年末的審查結果匯總表",
}
# 長欄名 → 表頭短名（18 欄要迫入 10.83in，跟 scan 用兩行短標）
SHORT = {
    "項目序號": "項目\n序號", "項目名稱": "項目名稱",
    "計劃投資金額": "計劃\n投資金額", "報告投資金額": "報告\n投資金額", "投資計劃完成率": "投資計劃\n完成率",
    "一般支持性部門的人工成本": "一般支持\n人工成本", "其他日常營運支出調整": "其他日常\n營運調整",
    "超出可計入範圍的內部資源支出": "超出可計入\n內部資源", "酒店客房改造支出": "酒店客房\n改造",
    "不符合“吸引外國客源”定義的相關投資支出": "不符吸引\n外國客源",
    "未完全實現投資目的的投資支出": "未完全\n實現目的",
    "投資計劃獲批前發生且未被認可的投資支出": "獲批前\n未認可",
    "潛在調整合計": "潛在調整\n合計", "調整後投資金額": "調整後\n投資金額",
    "潛在調整後投資計劃完成率": "潛在調整後\n完成率",
    "設施建設/資本性支出": "設施建設/\n資本性", "活動舉辦/營運性支出": "活動舉辦/\n營運性",
}


def fmt_money(v):
    if v is None or v == "" or (isinstance(v, float) and pd.isna(v)):
        return ""
    try:
        f = float(v)
    except (TypeError, ValueError):
        return str(v)
    n = int(round(f))
    if n == 0:
        return "-"
    return f"({abs(n):,})" if n < 0 else f"{n:,}"


def fmt_pct(v):
    if v is None or v == "" or (isinstance(v, float) and pd.isna(v)):
        return ""
    try:
        return f"{float(v)*100:.1f}%"
    except (TypeError, ValueError):
        return str(v)


def row_kind(seq: str) -> str:
    s = str(seq)
    if s.endswith("小計") or s.endswith("合計"):
        return "subtotal"
    if "—" in s or "－" in s:
        return "section"
    return "data"


def col_group(c):
    if c in G1:
        return "G1"
    if c in G3:
        return "G3"
    return "G2"


def _blocks(df, cols):
    """df → 逐【範疇】block：[(kind, cells)]。kind = sec / data / subtot / tot。
    scan 每個範疇有自己嘅小計、全份報告冇「續」→ 分頁時整個 block 唔拆。
    ⚠ 範疇／小計標籤要擺【項目名稱】欄（闊），唔可以留喺「項目序號」窄欄 —— 否則 wrap 5 行撐爆表。"""
    lab_c = cols.index("項目名稱") if "項目名稱" in cols else 0
    out, cur = [], []
    for _, row in df.iterrows():
        seq = str(row["項目序號"])
        kind = row_kind(seq)
        if kind == "section":
            if cur:
                out.append(cur)
            cells = [""] * len(cols); cells[lab_c] = seq
            cur = [("sec", cells)]
            continue
        cells = [("" if row.get(c, "") is None else str(row.get(c, ""))) if c in TEXT_COLS
                 else (fmt_pct(row.get(c, "")) if c in RATE_COLS else fmt_money(row.get(c, "")))
                 for c in cols]
        k = "tot" if seq.endswith("合計") else ("subtot" if kind == "subtotal" else "data")
        if k != "data":                       # 小計／合計：標籤搬去項目名稱欄
            cells[lab_c] = seq; cells[0] = ""
        cur.append((k, cells))
    if cur:
        out.append(cur)
    return out


def render_sheet(prs, sheet_name, df, cols, *, ent_up="MGM", sec=3, crumb=None, page_cb=None):
    """單個項目審查結果匯總表（對 scan slide 46-63）：navy 2 層表頭、逐範疇 block 唔拆頁、
    按【累積高度】分頁（唔會超出版面）、表頂 caption bar、表底 資料來源／（下頁待續）。"""
    ncol = len(cols)
    W, H = L.size_of(prs)
    tw = W - 2 * L.MARGIN
    title = YEAR_TITLE.get(sheet_name, sheet_name).format(e=ent_up)
    subs = [SHORT.get(c, c) for c in cols]
    # 欄寬（相對）：序號窄、名闊、其餘等闊
    widths = [0.45 if c == "項目序號" else (2.0 if c == "項目名稱" else 0.62) for c in cols]
    wid = [w * tw / sum(widths) for w in widths]
    # super header（3 大組）
    supers, ci = [], 0
    while ci < ncol:
        g = col_group(cols[ci]); cj = ci
        while cj + 1 < ncol and col_group(cols[cj + 1]) == g:
            cj += 1
        supers.append((GROUP_LABEL[g], ci, cj + 1)); ci = cj + 1
    font = 5.8 if ncol > 16 else 6.3
    yr = "20" + (sheet_name[-2:] if sheet_name[-2:].isdigit() else "25")
    head = (f"下表匯總了我們在審查{ent_up} {yr}年度投資計劃各項目投資執行情況時，識別出的各項目"
            f"投資支出涉及的潛在調整事項，以及相關的影響金額。")
    probe = L.HEAD_Y + L.head_h(head, W)[0] + 0.10
    avail = L.CONTENT_BOTTOM - probe - 0.24 - 0.17          # 減 caption bar
    hh = L.header_h(supers, subs, wid, font - 0.5)
    pages = L.fit_blocks(_blocks(df, cols), wid, font, avail, hh)
    for pi, chunk in enumerate(pages):
        suffix = f"（{pi+1}/{len(pages)}）"
        slide = L.blank(prs)
        L.breadcrumb(slide, W, sec, ent_up)
        L.footer(slide, W, H, len(prs.slides._sldIdLst))
        top = L.page_head(slide, W, (crumb or "其他信息  |  單個項目審查結果匯總") + suffix, head)
        top = L.caption_bar(slide, L.MARGIN, top, tw, title + suffix)
        L.draw_table(slide, L.MARGIN, top, tw, subs, chunk, widths, supers=supers,
                     font=font, hfont=font - 0.5, left_cols=2,
                     fill_h=L.CONTENT_BOTTOM - top - 0.24)
        L.source_note(slide, W, note="註：金額單位為萬澳門元；括號表示調減。",
                      more=(pi < len(pages) - 1))
        print(f"    {sheet_name} 第 {pi+1}/{len(pages)} 頁：{len(chunk)} 行 × {ncol} 欄")


def main():
    args = sys.argv[1:]
    template = None
    if "--template" in args:
        i = args.index("--template"); template = args[i + 1]; del args[i:i + 2]
    if not args:
        print("俾 build_project_review_table 出嘅 xlsx（--template <報告.pptx> 跟 slide 尺寸）"); return
    xlsx = Path(args[0])
    sheets = pd.read_excel(xlsx, sheet_name=None)

    # fresh 包（唔可以 clone template 再刪 slide：sldId 刪咗但 slide part 仲喺 → 撞名 corrupt）
    prs = Presentation()
    if template and Path(template).exists():
        ref = Presentation(template)
        prs.slide_width = ref.slide_width
        prs.slide_height = ref.slide_height
        print(f"── 跟 template 尺寸: {prs.slide_width/914400:.2f}x{prs.slide_height/914400:.2f}in（fresh 包避免撞名）")
    else:
        prs.slide_width = Inches(L.SLIDE_W); prs.slide_height = Inches(L.SLIDE_H)
        print(f"── 冇 template，用 {L.SLIDE_W}x{L.SLIDE_H}in（想跟報告尺寸請 --template 報告.pptx）")

    for sn, df in sheets.items():
        df = df.fillna("")
        cols = list(df.columns)
        print(f"\n# sheet {sn}：{len(df)} 行，{len(cols)} 欄")
        print("  欄:", cols)
        render_sheet(prs, sn, df, cols)

    out = xlsx.with_suffix(".native.pptx")
    prs.save(out)
    print(f"\n✓ 寫入 {out.resolve()}（開嚟同報告 slide 46-63 對；顏色/精細格式後補）")


if __name__ == "__main__":
    main()
